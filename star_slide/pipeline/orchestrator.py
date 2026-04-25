"""파이프라인 오케스트레이터 — Phase 1 단일 프로세스 실행.

CLI `star-slide convert input.pptx -o output.pptx --report report.json` end-to-end.

단계 (PRD §7.1):
  1. 파일 검증
  2. 슬라이드 추출 (PPTX/PDF/이미지)
  3. (옵션) SAM 마스크 생성
  4. OCR (텍스트 객체 추출)
  5. PPTX 조립 (현재는 텍스트 + 원본 배경 합성)
  6. 품질 리포트 JSON 출력

Phase 1 MVP 범위:
- 텍스트 객체 → PowerPoint 텍스트박스 (편집 가능)
- 슬라이드 배경 = 원본 PNG (인페인팅은 P1-T06에서 추가)
- 도형/아이콘 → vtracer + custGeom (P1-T08에서 추가, 현재는 스킵)
- 표/차트 → 이미지 fallback
"""

from __future__ import annotations

import contextlib
import shutil
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from xml.sax.saxutils import escape as xml_escape

import numpy as np
from numpy.typing import NDArray
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from star_slide.composer.inject import replace_geometry_with_custgeom
from star_slide.composer.vectorize import VectorizedShape, vectorize_shape
from star_slide.inpaint.lama import (
    _bboxes_to_mask,
    inpaint_with_mask,
)
from star_slide.input.pptx_extractor import inspect_pptx
from star_slide.input.validator import validate
from star_slide.ocr.metrics import cer
from star_slide.ocr.paddleocr_worker import OcrLine, run_ocr
from star_slide.ocr.text_attributes import (
    detect_korean,
    estimate_font_color,
    estimate_font_size_pt,
)
from star_slide.rasterize.coords import CoordTransform
from star_slide.rasterize.libreoffice import (
    LibreOfficeNotFoundError,
    fallback_render_from_embedded,
    render_pptx_to_pngs,
)
from star_slide.schema import (
    EditableLevel,
    JobState,
    Object,
    ObjectType,
    Project,
    Qa,
    QaReport,
    ShapePayload,
    Slide,
    SlideSize,
    SourceAssets,
    TextPayload,
)
from star_slide.segmentation.classify import (
    ClassificationResult,
    MaskClass,
    classify_masks,
)
from star_slide.segmentation.iou import bbox_iou
from star_slide.segmentation.sam2_auto import Sam2Result, run_sam2_auto
from star_slide.segmentation.sam3 import (
    Sam3Mask,
    Sam3Result,
    run_sam3,
    run_sam3_box_prompts,
    run_sam3_text_prompt,
)
from star_slide.vision_llm import VisionExtractor


@dataclass(frozen=True)
class ConvertOptions:
    """convert 명령 옵션."""

    use_libreoffice: bool = True  # False면 임베드 이미지 직접 추출
    ocr_min_confidence: float = 0.7  # textbox 생성 임계 (보수적)
    inpaint: bool = True  # P1-T06: LaMa로 텍스트 자리 배경 복원
    inpaint_padding_px: int = 12
    inpaint_dilate_px: int = 7
    inpaint_min_confidence: float = 0.3  # 마스크에 포함할 OCR 임계 (관대)
    vectorize_shapes: bool = True  # P1-T08: vtracer로 native shape (custGeom) 변환
    use_sam: bool = True  # P1-T03/T04: SAM 객체 분리 + SAM 마스크 기반 인페인팅
    sam_text_dilate_px: int = 3  # SAM TEXT 마스크용 dilate (정밀해서 작게)
    sam_pred_iou_thresh: float = 0.86
    sam_stability_thresh: float = 0.92
    sam_points_per_side: int = 32
    sam_max_masks: int = 200
    vectorize_shape_min_area_ratio: float = 0.001  # vtracer 입력 최소 면적 (0.1%)
    vectorize_shape_max_area_ratio: float = 0.15  # 너무 큰 SHAPE은 카드 배경 → vector 부적합
    # SAM 3 (PCS) — 텍스트 정밀 마스크 (잔재 해소)
    use_sam3_text_masks: bool = True  # OCR bbox → SAM 3 box-prompt → 정밀 글자 마스크
    sam3_supplement_text_prompt: bool = True  # OCR 미검출 텍스트 보충 (prompt='text')
    sam3_text_threshold: float = 0.3  # 'text' prompt 임계 (관대 — 모든 텍스트 capture)
    sam3_text_dilate_px: int = 5  # SAM 3 마스크 union 후 dilate (안티앨리어싱 가장자리 흡수)
    extract_text_attributes: bool = True  # SAM 3 마스크로 글자 색상/크기 추정
    # SAM 3 요소 분리 → SVG/PPTX — NotebookLM 단일 이미지 슬라이드 복원 경로
    use_sam3_elements: bool = False
    sam3_element_prompts: tuple[str, ...] = (
        "icon",
        "logo",
        "shape",
        "graphic",
        "illustration",
        "chart",
        "diagram",
        "arrow",
        "line",
        "circle",
        "rectangle",
        "photo",
        "table",
        "text",
    )
    sam3_element_threshold: float = 0.5
    sam3_element_mask_threshold: float = 0.5
    sam3_element_max_masks_per_concept: int = 30
    sam3_element_text_overlap_threshold: float = 0.45
    emit_svg: bool = True
    render_max_width_px: int | None = 1600
    visible_text_overlay: bool = False
    preserve_original_background: bool = True
    text_size_scale: float = 0.92
    # Vision LLM (cliproxy) — 슬라이드 PNG를 vision LLM에 직접 줘서 구조 JSON 추출
    use_vision_llm: bool = False  # 켜면 OCR+SAM 경로 우회, vision JSON으로 직접 조립
    vision_base_url: str = "http://localhost:8300/v1"
    vision_api_key: str = ""  # 비어있으면 VISION_PROXY_API_KEY/LOCAL_CLAUDE_API_KEY 환경변수
    vision_model: str = "claude-opus-4-6"
    vision_timeout_sec: float = 240.0
    vision_inpaint_text_pad_px: int = 8
    vision_inpaint_shape_pad_px: int = 0
    vision_inpaint_dilate_px: int = 5


def convert(
    input_path: Path,
    output_path: Path,
    *,
    workdir: Path | None = None,
    options: ConvertOptions | None = None,
    progress: object = None,
) -> tuple[Project, QaReport]:
    """input PPTX/PDF/이미지 → 편집 가능 PPTX + QaReport.

    Args:
        input_path: 입력 파일
        output_path: 출력 PPTX 경로
        workdir: 중간 산출물(렌더 PNG, 마스크 등) 저장 위치. None이면 output 옆에 자동.
        options: 변환 옵션
        progress: rich.progress 또는 None (단순 print)
    """
    options = options or ConvertOptions()
    workdir = workdir or output_path.parent / f"_workdir_{output_path.stem}"
    workdir.mkdir(parents=True, exist_ok=True)

    project = Project(
        id=f"prj_{uuid.uuid4().hex[:8]}",
        source_file=input_path,
        source_kind=validate(input_path),
        state=JobState.QUEUED,
    )

    # PDF/ZIP은 후속. NotebookLM locked PPTX와 단일 이미지부터 지원한다.
    if project.source_kind not in {"pptx", "image"}:
        raise NotImplementedError(
            f"현재 변환은 PPTX/이미지만 지원. PDF/ZIP은 후속 phase. (kind={project.source_kind})"
        )

    project.state = JobState.RASTERIZING

    # 1. 슬라이드 추출 + 렌더링
    render_dir = workdir / "renders"
    render_dir.mkdir(parents=True, exist_ok=True)

    if project.source_kind == "pptx":
        slide_w_emu, slide_h_emu, infos = inspect_pptx(input_path)

        if options.use_libreoffice:
            try:
                png_paths = render_pptx_to_pngs(input_path, render_dir)
            except LibreOfficeNotFoundError:
                png_paths = fallback_render_from_embedded(input_path, render_dir)
        else:
            png_paths = fallback_render_from_embedded(input_path, render_dir)

        if len(png_paths) != len(infos):
            # 렌더 슬라이드 수와 PPTX 슬라이드 수 다름 → 짧은 쪽 기준
            n = min(len(png_paths), len(infos))
            png_paths = png_paths[:n]
            infos = infos[:n]
    else:
        slide_png = render_dir / "slide_001.png"
        with Image.open(input_path) as im:
            rgb = im.convert("RGB")
            rgb.save(slide_png)
            im_w, im_h = rgb.size
        slide_h_emu = 6_858_000
        slide_w_emu = round(slide_h_emu * im_w / max(1, im_h))
        png_paths = [slide_png]
        infos = [object()]

    png_paths = _prepare_analysis_renders(
        png_paths,
        workdir=workdir,
        max_width_px=options.render_max_width_px,
    )

    # 2. 슬라이드별 OCR + 객체 생성
    project.state = JobState.DETECTING

    qa_warnings: list[str] = []
    n_text_total = 0
    n_text_editable = 0
    failed_slide_ids: list[str] = []

    # Vision LLM 사용 시 임시 HTTP 서버 띄움
    vision_extractor: VisionExtractor | None = None
    if options.use_vision_llm:
        vision_extractor = VisionExtractor(
            base_url=options.vision_base_url,
            api_key=options.vision_api_key,
            model=options.vision_model,
            timeout_sec=options.vision_timeout_sec,
        )
        vision_extractor.start_server(render_dir)

    for idx, (_info, png_path) in enumerate(zip(infos, png_paths, strict=False), start=1):
        try:
            with Image.open(png_path) as im:
                im_w, im_h = im.size
        except Exception as exc:
            failed_slide_ids.append(f"sld_{idx:03d}")
            qa_warnings.append(f"slide {idx} 이미지 로드 실패: {exc}")
            continue

        coord = CoordTransform(
            slide_width_emu=slide_w_emu,
            slide_height_emu=slide_h_emu,
            image_width_px=im_w,
            image_height_px=im_h,
        )

        slide_size = SlideSize(
            width_emu=slide_w_emu,
            height_emu=slide_h_emu,
            width_px=im_w,
            height_px=im_h,
            ratio=_ratio_str(slide_w_emu, slide_h_emu),
        )
        render_dpi = _render_dpi(
            image_size=(im_w, im_h),
            slide_size_emu=(slide_w_emu, slide_h_emu),
        )

        # === Vision LLM 경로 (use_vision_llm) ===
        if options.use_vision_llm and vision_extractor is not None:
            try:
                slide_v, n_text_v = _process_slide_with_vision(
                    png_path=png_path,
                    image_size=(im_w, im_h),
                    slide_size=slide_size,
                    coord=coord,
                    extractor=vision_extractor,
                    idx=idx,
                    workdir=workdir,
                    options=options,
                    qa_warnings=qa_warnings,
                )
                project.slides.append(slide_v)
                n_text_total += n_text_v
                n_text_editable += n_text_v
            except Exception as exc:
                failed_slide_ids.append(f"sld_{idx:03d}")
                qa_warnings.append(f"slide {idx} vision LLM 실패: {exc}")
            continue

        # OCR
        try:
            ocr_lines = run_ocr(png_path, lang="korean")
        except Exception as exc:
            qa_warnings.append(f"slide {idx} OCR 실패: {exc}")
            ocr_lines = []

        accepted_lines = [ln for ln in ocr_lines if ln.confidence >= options.ocr_min_confidence]
        # 인페인팅용은 더 관대한 임계 (작은 영문 라벨까지 지움)
        mask_lines = [
            ln for ln in ocr_lines if ln.confidence >= options.inpaint_min_confidence
        ]

        sam3_element_result: Sam3Result | None = None
        if options.use_sam3_elements:
            try:
                sam3_element_result = run_sam3(
                    png_path,
                    prompts=options.sam3_element_prompts,
                    threshold=options.sam3_element_threshold,
                    mask_threshold=options.sam3_element_mask_threshold,
                    max_masks_per_concept=options.sam3_element_max_masks_per_concept,
                )
                by_concept: dict[str, int] = {}
                for sm in sam3_element_result.masks:
                    by_concept[sm.concept] = by_concept.get(sm.concept, 0) + 1
                qa_warnings.append(f"slide {idx} SAM3 elements: {by_concept}")
            except Exception as exc:
                qa_warnings.append(f"slide {idx} SAM3 elements 실패: {exc}")
                sam3_element_result = None

        # SAM 2.1 객체 분리 (도형 검출 + vectorize 입력) — SAM 3는 grid auto-mask 미지원이라 SAM 2.1 유지
        sam_result: Sam2Result | None = None
        classification: ClassificationResult | None = None
        if options.use_sam and not options.use_sam3_elements:
            try:
                sam_result = run_sam2_auto(
                    png_path,
                    pred_iou_thresh=options.sam_pred_iou_thresh,
                    stability_score_thresh=options.sam_stability_thresh,
                    points_per_side=options.sam_points_per_side,
                    max_masks=options.sam_max_masks,
                )
                classification = classify_masks(
                    sam_result.masks,
                    mask_lines,
                    image_size=(im_w, im_h),
                )
            except Exception as exc:
                qa_warnings.append(f"slide {idx} SAM2 실패 (OCR bbox fallback): {exc}")
                sam_result = None
                classification = None

        # SAM 3 box-prompt — OCR bbox 별 정밀 텍스트 마스크 (잔재 해소)
        sam3_box_result: Sam3Result | None = None
        sam3_text_result: Sam3Result | None = None
        if options.use_sam3_text_masks and mask_lines:
            try:
                sam3_box_result = run_sam3_box_prompts(
                    png_path,
                    [ln.bbox for ln in mask_lines],
                )
            except Exception as exc:
                qa_warnings.append(f"slide {idx} SAM3 box-prompt 실패: {exc}")
                sam3_box_result = None

        # SAM 3 text-prompt 'text' — OCR 미검출 텍스트 영역 보충 (인페인팅에만 사용)
        if options.use_sam3_text_masks and options.sam3_supplement_text_prompt:
            if sam3_element_result is not None:
                sam3_text_masks = [
                    sm for sm in sam3_element_result.masks if sm.concept == "text"
                ]
                sam3_text_result = Sam3Result(
                    image_size=sam3_element_result.image_size,
                    masks=sam3_text_masks,
                    device=sam3_element_result.device,
                )
            else:
                try:
                    sam3_text_result = run_sam3_text_prompt(
                        png_path,
                        "text",
                        threshold=options.sam3_text_threshold,
                    )
                except Exception as exc:
                    qa_warnings.append(f"slide {idx} SAM3 text-prompt 실패: {exc}")
                    sam3_text_result = None

        # 인페인팅
        background_path = png_path
        has_sam3_text_masks = bool(sam3_text_result and sam3_text_result.masks)
        if options.inpaint and (mask_lines or has_sam3_text_masks):
            try:
                inpainted = _inpaint_slide(
                    png_path=png_path,
                    image_size=(im_w, im_h),
                    mask_lines=mask_lines,
                    classification=classification,
                    sam3_box_result=sam3_box_result,
                    sam3_text_result=sam3_text_result,
                    options=options,
                )
                inpaint_dir = workdir / "inpainted"
                inpaint_dir.mkdir(parents=True, exist_ok=True)
                background_path = inpaint_dir / png_path.name
                inpainted.save(background_path)
            except Exception as exc:
                qa_warnings.append(f"slide {idx} 인페인팅 실패 (원본 사용): {exc}")
                background_path = png_path

        slide = Slide(
            id=f"sld_{idx:03d}",
            page_no=idx,
            size=slide_size,
            render_path=png_path,
            background_path=background_path,
        )

        # SAM 3 box-prompt 결과를 mask_lines 인덱스로 매핑 (글자 속성 추출용)
        sam3_mask_by_box: dict[int, np.ndarray] = {}
        if sam3_box_result is not None:
            for sm in sam3_box_result.masks:
                if sm.source_box_idx is not None:
                    sam3_mask_by_box[sm.source_box_idx] = sm.segmentation
        # mask_lines → accepted_lines 인덱스 매핑 (mask_lines가 더 큰 집합)
        accepted_to_mask_idx: dict[int, int] = {}
        for ai, al in enumerate(accepted_lines):
            for mli, ml in enumerate(mask_lines):
                if ml is al:
                    accepted_to_mask_idx[ai] = mli
                    break

        slide_pil_for_attr: Image.Image | None = None
        for k, line in enumerate(accepted_lines):
            n_text_total += 1
            n_text_editable += 1

            bbox_px = line.bbox
            bbox_emu = coord.px_bbox_to_emu(bbox_px)

            # 글자 색상/크기 추출 (SAM 3 정밀 마스크가 있을 때만)
            color = "#000000"
            font_size_pt: float | None = None
            mask_h_px = bbox_px[3]
            if options.extract_text_attributes:
                mi: int | None = accepted_to_mask_idx.get(k)
                seg: np.ndarray | None = (
                    sam3_mask_by_box.get(mi) if mi is not None else None
                )
                if seg is not None:
                    if slide_pil_for_attr is None:
                        slide_pil_for_attr = Image.open(png_path).convert("RGB")
                    est_color = estimate_font_color(slide_pil_for_attr, seg)
                    if est_color is not None:
                        color = est_color
                    # 마스크 height (실제 글자 영역) — bbox보다 더 정확
                    ys_seg = np.nonzero(seg)[0]
                    if ys_seg.size > 0:
                        mask_h_px = float(ys_seg.max() - ys_seg.min() + 1)
                font_size_pt = estimate_font_size_pt(
                    mask_h_px,
                    is_korean=detect_korean(line.text),
                    dpi=render_dpi,
                )
                font_size_pt *= options.text_size_scale

            obj = Object(
                id=f"{slide.id}_obj_{k:03d}",
                type=ObjectType.TEXT,
                bbox_px=bbox_px,
                bbox_emu=bbox_emu,
                confidence=line.confidence,
                editable_level=EditableLevel.NATIVE,
                source=SourceAssets(crop_path=None, detector="paddleocr_ppocrv5"),
                qa=Qa(),
                text=TextPayload(
                    content=line.text,
                    confidence=line.confidence,
                    color=color,
                    font_size_pt=font_size_pt,
                ),
            )
            slide.objects.append(obj)

        if sam3_element_result is not None:
            try:
                _append_sam3_element_objects(
                    slide=slide,
                    png_path=png_path,
                    sam3_result=sam3_element_result,
                    ocr_lines=accepted_lines,
                    coord=coord,
                    workdir=workdir,
                    options=options,
                    qa_warnings=qa_warnings,
                )
                if options.emit_svg:
                    svg_dir = workdir / "sam3_svg"
                    svg_dir.mkdir(parents=True, exist_ok=True)
                    svg_path = svg_dir / f"{slide.id}.svg"
                    _write_slide_svg(slide, svg_path)
                    qa_warnings.append(f"slide {idx} SVG: {svg_path}")
            except Exception as exc:
                qa_warnings.append(f"slide {idx} SAM3 객체 조립 실패: {exc}")

        # SAM SHAPE → vtracer → custGeom (P1-T08)
        if classification is not None and not options.use_sam3_elements:
            shape_offset = len(slide.objects)
            slide_pil_for_vec: Image.Image | None = None
            shape_min_px = options.vectorize_shape_min_area_ratio * im_w * im_h
            shape_max_px = options.vectorize_shape_max_area_ratio * im_w * im_h
            for j, cm in enumerate(classification.classified):
                if cm.cls != MaskClass.SHAPE:
                    continue

                vec: VectorizedShape | None = None
                if (
                    options.vectorize_shapes
                    and shape_min_px <= cm.mask.area <= shape_max_px
                ):
                    if slide_pil_for_vec is None:
                        slide_pil_for_vec = Image.open(png_path).convert("RGB")
                    try:
                        vec = vectorize_shape(slide_pil_for_vec, cm.mask.segmentation)
                    except Exception as exc:
                        qa_warnings.append(
                            f"slide {idx} shape {j} vectorize 실패: {exc}"
                        )
                        vec = None

                bbox_px = vec.bbox if vec is not None else cm.mask.bbox
                bbox_emu = coord.px_bbox_to_emu(bbox_px)
                shape_obj = Object(
                    id=f"{slide.id}_shape_{shape_offset + j:03d}",
                    type=ObjectType.SHAPE,
                    bbox_px=bbox_px,
                    bbox_emu=bbox_emu,
                    confidence=cm.mask.score,
                    editable_level=(
                        EditableLevel.VECTOR if vec is not None else EditableLevel.RASTER
                    ),
                    source=SourceAssets(detector="sam2.1_hiera_large"),
                    qa=Qa(),
                    shape=vec.payload if vec is not None else None,
                )
                slide.objects.append(shape_obj)

        project.slides.append(slide)

    # vision LLM HTTP 서버 종료
    if vision_extractor is not None:
        vision_extractor.stop_server()

    # 3. PPTX 조립
    project.state = JobState.RECONSTRUCTING
    _compose_pptx(project, output_path, options=options)

    # 4. 품질 리포트 — NATIVE(텍스트박스) + VECTOR(custGeom 도형) 모두 편집 가능으로 카운트
    project.state = JobState.READY
    n_objects = sum(len(s.objects) for s in project.slides)
    editable_levels = {EditableLevel.NATIVE, EditableLevel.VECTOR}
    avg_editable = sum(
        1 for s in project.slides for o in s.objects if o.editable_level in editable_levels
    ) / max(1, n_objects)
    report = QaReport(
        project_id=project.id,
        n_slides=len(project.slides),
        n_objects=n_objects,
        avg_editable_ratio=avg_editable,
        text_objects_editable=n_text_editable,
        text_objects_total=n_text_total,
        failed_slide_ids=failed_slide_ids,
        warnings=qa_warnings,
    )

    return project, report


def _prepare_analysis_renders(
    png_paths: list[Path],
    *,
    workdir: Path,
    max_width_px: int | None,
) -> list[Path]:
    """OCR/SAM 입력 PNG를 적정 크기로 정규화한다.

    LibreOffice는 retina급 큰 PNG를 만들 수 있는데, SAM3는 픽셀 수에 민감하다.
    좌표 변환은 실제 분석 PNG 크기를 기준으로 하므로 비율만 유지하면 PPTX 배치는
    그대로 보존된다.
    """
    if max_width_px is None or max_width_px <= 0:
        return png_paths

    out_dir = workdir / "analysis_renders"
    out_dir.mkdir(parents=True, exist_ok=True)
    normalized: list[Path] = []

    for src in png_paths:
        try:
            with Image.open(src) as im:
                im = im.convert("RGB")
                if im.width <= max_width_px:
                    normalized.append(src)
                    continue
                scale = max_width_px / im.width
                new_size = (max_width_px, max(1, round(im.height * scale)))
                resized = im.resize(new_size, Image.Resampling.LANCZOS)
                dst = out_dir / src.name
                resized.save(dst)
                normalized.append(dst)
        except Exception:
            normalized.append(src)

    return normalized


def _bbox_overlap_ratio(
    inner: tuple[float, float, float, float],
    outer: tuple[float, float, float, float],
) -> float:
    """inner bbox 면적 중 outer bbox와 겹치는 비율."""
    ix, iy, iw, ih = inner
    ox, oy, ow, oh = outer
    if iw <= 0 or ih <= 0:
        return 0.0
    x1 = max(ix, ox)
    y1 = max(iy, oy)
    x2 = min(ix + iw, ox + ow)
    y2 = min(iy + ih, oy + oh)
    if x2 <= x1 or y2 <= y1:
        return 0.0
    return float(((x2 - x1) * (y2 - y1)) / (iw * ih))


def _is_text_like_sam3_mask(
    mask: Sam3Mask,
    ocr_lines: list[OcrLine],
    *,
    overlap_threshold: float,
) -> bool:
    if mask.concept == "text":
        return True
    return any(
        _bbox_overlap_ratio(mask.bbox, line.bbox) >= overlap_threshold
        for line in ocr_lines
    )


def _dedupe_sam3_masks(
    masks: list[Sam3Mask],
    *,
    iou_threshold: float = 0.88,
) -> list[Sam3Mask]:
    kept: list[Sam3Mask] = []
    for mask in sorted(masks, key=lambda m: m.score, reverse=True):
        if any(bbox_iou(mask.bbox, other.bbox) >= iou_threshold for other in kept):
            continue
        kept.append(mask)
    return kept


def _append_sam3_element_objects(
    *,
    slide: Slide,
    png_path: Path,
    sam3_result: Sam3Result,
    ocr_lines: list[OcrLine],
    coord: CoordTransform,
    workdir: Path,
    options: ConvertOptions,
    qa_warnings: list[str],
) -> None:
    """SAM3 concept masks를 PPTX/SVG 조립용 객체로 변환한다.

    텍스트와 배경성 대형 마스크는 제외한다. 단순 도형 계열은 custGeom으로
    벡터화하고, 사진/차트/표/로고처럼 시각 충실도가 중요한 요소는 선택 가능한
    crop picture로 보존한다.
    """
    pil = Image.open(png_path).convert("RGB")
    im_w, im_h = pil.size
    slide_area = float(im_w * im_h) if im_w > 0 and im_h > 0 else 1.0

    candidates = [
        m
        for m in sam3_result.masks
        if not _is_text_like_sam3_mask(
            m,
            ocr_lines,
            overlap_threshold=options.sam3_element_text_overlap_threshold,
        )
    ]
    candidates = _dedupe_sam3_masks(candidates)

    crop_dir = workdir / "sam3_crops" / slide.id
    crop_dir.mkdir(parents=True, exist_ok=True)

    vector_concepts = {
        "shape",
        "icon",
        "arrow",
        "diagram",
        "graphic",
        "line",
        "circle",
        "rectangle",
    }
    raster_first_concepts = {"photo", "chart", "table", "logo"}
    min_area = options.vectorize_shape_min_area_ratio * slide_area
    max_vector_area = options.vectorize_shape_max_area_ratio * slide_area
    start = len(slide.objects)
    n_vector = 0
    n_raster = 0

    for j, mask in enumerate(candidates):
        mx, my, mw, mh = mask.bbox
        bbox_area = mw * mh
        if bbox_area < slide_area * 0.0002:
            continue
        if bbox_area > slide_area * 0.72:
            continue

        vec: VectorizedShape | None = None
        should_vectorize = (
            options.vectorize_shapes
            and mask.concept in vector_concepts
            and min_area <= bbox_area <= max_vector_area
        )
        if should_vectorize:
            try:
                vec = vectorize_shape(pil, mask.segmentation)
            except Exception as exc:
                qa_warnings.append(
                    f"slide {slide.page_no} SAM3 {mask.concept} vectorize 실패: {exc}"
                )
                vec = None

        if vec is not None:
            bbox_px = vec.bbox
            obj = Object(
                id=f"{slide.id}_sam3_{start + j:03d}",
                type=ObjectType.SHAPE,
                subtype=mask.concept,
                bbox_px=bbox_px,
                bbox_emu=coord.px_bbox_to_emu(bbox_px),
                confidence=mask.score,
                editable_level=EditableLevel.VECTOR,
                source=SourceAssets(detector=f"sam3:{mask.concept}"),
                qa=Qa(),
                shape=vec.payload,
            )
            slide.objects.append(obj)
            n_vector += 1
            continue

        if mask.concept in raster_first_concepts or bbox_area >= min_area:
            x1 = max(0, int(mx))
            y1 = max(0, int(my))
            x2 = min(im_w, int(mx + mw))
            y2 = min(im_h, int(my + mh))
            if x2 <= x1 or y2 <= y1:
                continue
            crop_path = crop_dir / f"{mask.concept}_{j:03d}.png"
            pil.crop((x1, y1, x2, y2)).save(crop_path)
            bbox_px = (float(x1), float(y1), float(x2 - x1), float(y2 - y1))
            obj = Object(
                id=f"{slide.id}_sam3_{start + j:03d}",
                type=ObjectType.PHOTO,
                subtype=mask.concept,
                bbox_px=bbox_px,
                bbox_emu=coord.px_bbox_to_emu(bbox_px),
                confidence=mask.score,
                editable_level=EditableLevel.RASTER,
                source=SourceAssets(crop_path=crop_path, detector=f"sam3:{mask.concept}"),
                qa=Qa(),
            )
            slide.objects.append(obj)
            n_raster += 1

    qa_warnings.append(
        f"slide {slide.page_no} SAM3 objects: vector={n_vector}, selectable_raster={n_raster}"
    )


def _write_slide_svg(slide: Slide, svg_path: Path) -> None:
    """Layer 객체를 디버깅/후처리용 SVG로 저장."""
    svg_path.parent.mkdir(parents=True, exist_ok=True)
    assets_dir = svg_path.parent / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    bg_src = slide.background_path or slide.render_path
    bg_dst = assets_dir / f"{slide.id}_background{bg_src.suffix.lower() or '.png'}"
    if bg_src.exists():
        shutil.copy2(bg_src, bg_dst)
    bg_href = f"assets/{bg_dst.name}"

    w = slide.size.width_px
    h = slide.size.height_px
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {w} {h}" width="{w}" height="{h}">',
        f'<image href="{xml_escape(bg_href)}" x="0" y="0" width="{w}" height="{h}" opacity="1"/>',
    ]

    for obj in slide.objects:
        if obj.type == ObjectType.SHAPE and obj.shape and obj.shape.svg_path_d:
            x, y, _bw, _bh = obj.bbox_px
            fill = obj.shape.fill or "#808080"
            d = xml_escape(obj.shape.svg_path_d)
            parts.append(
                f'<g id="{xml_escape(obj.id)}" data-concept="{xml_escape(obj.subtype or "")}" '
                f'transform="translate({x:.3f},{y:.3f})">'
                f'<path d="{d}" fill="{xml_escape(fill)}"/>'
                "</g>"
            )
        elif obj.type == ObjectType.PHOTO and obj.source.crop_path:
            crop_src = Path(obj.source.crop_path)
            if not crop_src.exists():
                continue
            crop_dst = assets_dir / f"{obj.id}{crop_src.suffix.lower() or '.png'}"
            shutil.copy2(crop_src, crop_dst)
            x, y, bw, bh = obj.bbox_px
            parts.append(
                f'<image id="{xml_escape(obj.id)}" data-concept="{xml_escape(obj.subtype or "")}" '
                f'href="assets/{xml_escape(crop_dst.name)}" x="{x:.3f}" y="{y:.3f}" '
                f'width="{bw:.3f}" height="{bh:.3f}"/>'
            )
        elif obj.type == ObjectType.TEXT and obj.text:
            x, y, _bw, bh = obj.bbox_px
            font_size = obj.text.font_size_pt or max(8.0, bh * 0.7)
            color = obj.text.color or "#000000"
            baseline = y + bh * 0.78
            parts.append(
                f'<text id="{xml_escape(obj.id)}" x="{x:.3f}" y="{baseline:.3f}" '
                f'font-size="{font_size:.2f}" fill="{xml_escape(color)}" '
                'font-family="Noto Sans KR, Arial, sans-serif">'
                f"{xml_escape(obj.text.content)}</text>"
            )

    parts.append("</svg>")
    svg_path.write_text("\n".join(parts), encoding="utf-8")


def _preset_for_subtype(subtype: str | None) -> Any:
    """vision LLM subtype 문자열 → python-pptx MSO_SHAPE."""
    s = (subtype or "rectangle").lower()
    return {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
        "oval": MSO_SHAPE.OVAL,
        "circle": MSO_SHAPE.OVAL,
        "arrow": MSO_SHAPE.RIGHT_ARROW,
        "right_arrow": MSO_SHAPE.RIGHT_ARROW,
        "line": MSO_SHAPE.RECTANGLE,
        "polygon": MSO_SHAPE.RECTANGLE,
    }.get(s, MSO_SHAPE.RECTANGLE)


def _ratio_str(w_emu: int, h_emu: int) -> str:
    if h_emu == 0:
        return "?"
    r = w_emu / h_emu
    if abs(r - 16 / 9) < 0.01:
        return "16:9"
    if abs(r - 4 / 3) < 0.01:
        return "4:3"
    return f"{r:.2f}"


def _render_dpi(
    *,
    image_size: tuple[int, int],
    slide_size_emu: tuple[int, int],
) -> float:
    """분석 렌더의 실제 DPI 추정."""
    im_w, im_h = image_size
    slide_w, slide_h = slide_size_emu
    w_in = slide_w / 914_400 if slide_w > 0 else 0
    h_in = slide_h / 914_400 if slide_h > 0 else 0
    values = []
    if w_in > 0:
        values.append(im_w / w_in)
    if h_in > 0:
        values.append(im_h / h_in)
    return float(sum(values) / len(values)) if values else 96.0


def _compose_pptx(
    project: Project,
    output_path: Path,
    *,
    options: ConvertOptions,
) -> None:
    """python-pptx로 결과 PPTX 조립.

    각 슬라이드에 대해:
    - 슬라이드 배경 = 원본 PNG (전체 슬라이드 영역)
    - 텍스트 객체 → 투명 텍스트박스 (편집 가능)

    Phase 1 단순화: 모든 텍스트박스는 흰 배경 텍스트 위에 덮어 그리지 않고,
    원본 PNG를 배경으로 두고 그 위에 OCR 텍스트박스를 투명하게 배치.
    Phase 2에서 인페인팅 + 텍스트박스로 교체.
    """
    if not project.slides:
        raise ValueError("프로젝트에 슬라이드가 없습니다")

    first = project.slides[0]
    prs = Presentation()
    prs.slide_width = Emu(first.size.width_emu)
    prs.slide_height = Emu(first.size.height_emu)

    blank_layout = prs.slide_layouts[6]

    for slide_data in project.slides:
        slide = prs.slides.add_slide(blank_layout)

        # fidelity 기본값: 원본 렌더 배경을 보존하고 편집 레이어를 투명하게 얹는다.
        bg_path = (
            slide_data.render_path
            if options.preserve_original_background
            else slide_data.background_path or slide_data.render_path
        )
        slide.shapes.add_picture(
            str(bg_path),
            Emu(0),
            Emu(0),
            width=Emu(slide_data.size.width_emu),
            height=Emu(slide_data.size.height_emu),
        )

        # 도형 (SHAPE) — preset geom (vision) 또는 vtracer custGeom (SAM 2.1)
        for obj in slide_data.objects:
            if (
                obj.type != ObjectType.SHAPE
                or obj.shape is None
                or obj.bbox_emu is None
            ):
                continue
            x, y, w, h = obj.bbox_emu
            if w <= 0 or h <= 0:
                continue
            geom = obj.shape.geom_type
            if geom == "preset":
                # vision LLM이 알려준 subtype을 PowerPoint preset으로 매핑
                preset = _preset_for_subtype(obj.shape.preset_name or obj.subtype)
                sh = slide.shapes.add_shape(preset, Emu(x), Emu(y), Emu(w), Emu(h))
            elif geom == "custGeom" and obj.shape.custgeom_xml:
                sh = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h)
                )
                try:
                    replace_geometry_with_custgeom(sh, obj.shape.custgeom_xml)
                except Exception:
                    continue
            elif geom == "path":
                # vision LLM이 path로 표시 (게이지 등) — 일단 OVAL placeholder
                sh = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, Emu(x), Emu(y), Emu(w), Emu(h)
                )
            else:
                continue

            if obj.shape.fill:
                with contextlib.suppress(Exception):
                    rgb = obj.shape.fill.lstrip("#")
                    sh.fill.solid()
                    sh.fill.fore_color.rgb = RGBColor.from_string(rgb)  # type: ignore[no-untyped-call]
            else:
                with contextlib.suppress(Exception):
                    sh.fill.background()
            if obj.shape.stroke:
                with contextlib.suppress(Exception):
                    rgb = obj.shape.stroke.lstrip("#")
                    sh.line.color.rgb = RGBColor.from_string(rgb)  # type: ignore[no-untyped-call]
                    if obj.shape.stroke_width_pt > 0:
                        sh.line.width = Pt(obj.shape.stroke_width_pt)
            else:
                with contextlib.suppress(Exception):
                    sh.line.fill.background()

        # PHOTO (vision LLM image type) — 원본 crop 임베드
        for obj in slide_data.objects:
            if obj.type != ObjectType.PHOTO or obj.bbox_emu is None:
                continue
            crop_path = obj.source.crop_path if obj.source else None
            if crop_path is None or not Path(crop_path).exists():
                continue
            x, y, w, h = obj.bbox_emu
            slide.shapes.add_picture(
                str(crop_path), Emu(x), Emu(y), width=Emu(w), height=Emu(h)
            )

        # 텍스트 객체
        for obj in slide_data.objects:
            if obj.type != ObjectType.TEXT or obj.text is None or obj.bbox_emu is None:
                continue
            x, y, w, h = obj.bbox_emu
            tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
            tf = tb.text_frame
            tf.word_wrap = False
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = obj.text.content
            run.font.name = "Malgun Gothic"
            if obj.text.font_size_pt:
                run.font.size = Pt(obj.text.font_size_pt)
            if obj.text.color and obj.text.color != "#000000":
                with contextlib.suppress(Exception):
                    rgb = obj.text.color.lstrip("#")
                    run.font.color.rgb = RGBColor.from_string(rgb)  # type: ignore[no-untyped-call]
            if obj.rotation:
                with contextlib.suppress(Exception):
                    tb.rotation = float(obj.rotation)
            align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
            with contextlib.suppress(Exception):
                p.alignment = align_map.get(obj.text.align, PP_ALIGN.LEFT)
            if not options.visible_text_overlay:
                _make_text_run_transparent(run)
                with contextlib.suppress(Exception):
                    tb.name = f"editable_text_hidden_{obj.id}"
            # 배경 PNG와 겹치지 않도록 텍스트박스 자체는 채움 없음 (기본)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def _make_text_run_transparent(run: Any) -> None:
    """PowerPoint에서 선택/편집은 가능하지만 화면에는 보이지 않게 한다."""
    from pptx.oxml import parse_xml

    # python-pptx 공개 API에는 글자 alpha가 없어 DrawingML을 직접 보강한다.
    run.font.color.rgb = RGBColor(0, 0, 0)
    r_pr = run._r.get_or_add_rPr()
    r_pr.set("noProof", "1")

    for child in list(r_pr):
        if child.tag.endswith("solidFill"):
            r_pr.remove(child)

    solid = parse_xml(
        '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:srgbClr val="000000"><a:alpha val="0"/></a:srgbClr>'
        "</a:solidFill>"
    )
    r_pr.append(solid)


def _inpaint_slide(
    *,
    png_path: Path,
    image_size: tuple[int, int],
    mask_lines: list[OcrLine],
    classification: ClassificationResult | None,
    sam3_box_result: Sam3Result | None = None,
    sam3_text_result: Sam3Result | None = None,
    options: ConvertOptions,
) -> Image.Image:
    """슬라이드 1장 인페인팅.

    전략 (품질 우선):
      - **SAM 3 box-prompt 마스크** (OCR bbox별 정밀 글자 ink): 가장 정확,
        안티앨리어싱까지 fit → 잔재 거의 없음.
      - **SAM 3 text-prompt 'text' 마스크**: OCR이 못 잡은 텍스트 영역 보충.
      - **OCR 사각 bbox fallback**: SAM 3가 박스에 매칭 못한 OCR 라인.
      - **SHAPE 픽셀 보호**: SAM 2.1 SHAPE 영역은 인페인팅 제외 (게이지 등).
    """
    import cv2

    pil = Image.open(png_path).convert("RGB")
    w, h = image_size

    # 1. 인페인팅 마스크 빌드 — 하이브리드 (각 영역의 장점 결합):
    #   (a) OCR rect + padding: 큰 한글 외곽선까지 안전하게 덮음 (proven)
    #   (b) SAM 3 box-prompt: OCR bbox 안 정밀 글자 ink (잔재 fade 더 줄임)
    #   (c) SAM 3 text-prompt: OCR 미검출 텍스트 영역 (영문 라벨 등)
    #   (d) 모두 union → 큰 dilate(7)로 가장자리 흡수
    inpaint_arr = np.zeros((h, w), dtype=np.uint8)

    # (a) OCR rect mask (기본, 큰 잔재 방지)
    if mask_lines:
        ocr_rect = _bboxes_to_mask(
            [ln.bbox for ln in mask_lines],
            image_size=image_size,
            padding=options.inpaint_padding_px,
            dilate_kernel=options.inpaint_dilate_px,
        )
        inpaint_arr = np.maximum(
            inpaint_arr, np.asarray(ocr_rect.convert("L"), dtype=np.uint8)
        )

    # (b) SAM 3 box-prompt 정밀 마스크 추가 union
    if sam3_box_result is not None and options.use_sam3_text_masks:
        for sm in sam3_box_result.masks:
            if sm.segmentation.shape == (h, w):
                inpaint_arr[sm.segmentation] = 255

    # (c) SAM 3 text-prompt 'text' 보충 (OCR 미검출 영역)
    if sam3_text_result is not None and options.sam3_supplement_text_prompt:
        for sm in sam3_text_result.masks:
            if sm.segmentation.shape == (h, w):
                inpaint_arr[sm.segmentation] = 255

    # (d) 추가 dilate — SAM 3 정밀 마스크의 안티앨리어싱 흡수
    if options.sam3_text_dilate_px > 0:
        kernel = cv2.getStructuringElement(
            cv2.MORPH_ELLIPSE, (options.sam3_text_dilate_px, options.sam3_text_dilate_px)
        )
        inpaint_arr = cv2.dilate(inpaint_arr, kernel, iterations=1).astype(np.uint8)

    if classification is None:
        # SHAPE 보호 없이 그대로
        return inpaint_with_mask(pil, Image.fromarray(inpaint_arr, mode="L"))

    # legacy ocr_mask_arr (SHAPE 보호 로직 호환용 — 현재는 SAM 3 마스크가 정확하므로 보호 마스크 영향 적음)
    ocr_mask_arr = inpaint_arr

    # 2. SHAPE 보호 마스크: SAM이 "도형"으로 분류한 픽셀 중,
    #    OCR과 (실질적으로) 겹치지 않는 것만 보호한다.
    #    → 게이지/아이콘은 살리지만, OCR bbox 내부의 글자 획 잔재는 보호 X.
    shape_segs: list[NDArray[np.bool_]] = []
    w, h = image_size
    slide_area = float(w * h) if w > 0 and h > 0 else 1.0
    shape_max_area_ratio = 0.15
    shape_overlap_with_ocr_max = 0.05  # OCR과의 픽셀 겹침이 5% 이상이면 텍스트 잔재로 의심 → 보호 X

    for c in classification.classified:
        if c.cls != MaskClass.SHAPE:
            continue
        if c.mask.area / slide_area > shape_max_area_ratio:
            continue
        seg = c.mask.segmentation
        if seg.shape != (h, w):
            continue
        overlap = int(np.logical_and(seg, ocr_mask_arr > 0).sum())
        if c.mask.area > 0 and overlap / c.mask.area > shape_overlap_with_ocr_max:
            continue
        shape_segs.append(seg)

    if shape_segs:
        shape_arr = np.zeros((h, w), dtype=np.uint8)
        for seg in shape_segs:
            shape_arr[seg] = 255
        final_arr = ocr_mask_arr.copy()
        final_arr[shape_arr > 0] = 0
    else:
        final_arr = ocr_mask_arr

    final_mask = Image.fromarray(final_arr.astype(np.uint8), mode="L")
    return inpaint_with_mask(pil, final_mask)


def _process_slide_with_vision(
    *,
    png_path: Path,
    image_size: tuple[int, int],
    slide_size: SlideSize,
    coord: CoordTransform,
    extractor: VisionExtractor,
    idx: int,
    workdir: Path,
    options: ConvertOptions,
    qa_warnings: list[str],
) -> tuple[Slide, int]:
    """Vision LLM JSON → Layer Schema 슬라이드 1장.

    흐름:
      1. vision LLM 호출 → VisionSlide
      2. 인페인팅 마스크 = 모든 element bbox union (image type 제외)
      3. LaMa 인페인팅
      4. element → Object (text/shape/path/table/image)
    """
    import cv2

    vision = extractor.extract(png_path)
    qa_warnings.append(
        f"slide {idx} vision: {len(vision.elements)} elements, "
        f"{extractor.last_usage.elapsed_sec:.1f}s, "
        f"{extractor.last_usage.total_tokens} tokens"
    )

    w, h = image_size
    pil = Image.open(png_path).convert("RGB")

    # 1. 인페인팅 마스크 빌드
    inpaint_arr = np.zeros((h, w), dtype=np.uint8)
    for el in vision.elements:
        if el.type == "image":
            continue  # 원본 raster 보존
        bx, by, bw, bh = el.bbox.as_tuple()
        pad = (
            options.vision_inpaint_text_pad_px
            if el.type == "text"
            else options.vision_inpaint_shape_pad_px
        )
        x1 = max(0, int(bx) - pad)
        y1 = max(0, int(by) - pad)
        x2 = min(w, int(bx + bw) + pad)
        y2 = min(h, int(by + bh) + pad)
        if x2 > x1 and y2 > y1:
            inpaint_arr[y1:y2, x1:x2] = 255

    if options.vision_inpaint_dilate_px > 0:
        kernel = cv2.getStructuringElement(
            cv2.MORPH_ELLIPSE,
            (options.vision_inpaint_dilate_px, options.vision_inpaint_dilate_px),
        )
        inpaint_arr = cv2.dilate(inpaint_arr, kernel, iterations=1).astype(np.uint8)

    background_path = png_path
    if options.inpaint and inpaint_arr.any():
        try:
            inpainted = inpaint_with_mask(pil, Image.fromarray(inpaint_arr, mode="L"))
            inpaint_dir = workdir / "inpainted"
            inpaint_dir.mkdir(parents=True, exist_ok=True)
            background_path = inpaint_dir / png_path.name
            inpainted.save(background_path)
        except Exception as exc:
            qa_warnings.append(f"slide {idx} 인페인팅 실패: {exc}")
            background_path = png_path

    slide = Slide(
        id=f"sld_{idx:03d}",
        page_no=idx,
        size=slide_size,
        render_path=png_path,
        background_path=background_path,
    )

    # 2. element → Object
    n_text_added = 0
    image_crops_dir = workdir / "image_crops" / f"sld_{idx:03d}"
    for k, el in enumerate(vision.elements):
        bx, by, bw, bh = el.bbox.as_tuple()
        bbox_px = (float(bx), float(by), float(bw), float(bh))
        bbox_emu = coord.px_bbox_to_emu(bbox_px)

        if el.type == "text" and el.content:
            obj = Object(
                id=f"{slide.id}_obj_{k:03d}",
                type=ObjectType.TEXT,
                bbox_px=bbox_px,
                bbox_emu=bbox_emu,
                rotation=el.rotation_deg,
                confidence=0.95,  # vision LLM 가정 신뢰도
                editable_level=EditableLevel.NATIVE,
                source=SourceAssets(crop_path=None, detector="vision_llm"),
                qa=Qa(),
                text=TextPayload(
                    content=el.content,
                    language=el.language,
                    color=el.color or "#000000",
                    font_size_pt=el.font_size_pt,
                    align=el.align,
                ),
            )
            slide.objects.append(obj)
            n_text_added += 1

        elif el.type in ("shape", "path"):
            payload = ShapePayload(
                geom_type="preset" if el.type == "shape" else "path",
                preset_name=el.subtype if el.type == "shape" else None,
                fill=el.fill,
                stroke=el.stroke,
                stroke_width_pt=el.stroke_width_pt,
            )
            obj = Object(
                id=f"{slide.id}_obj_{k:03d}",
                type=ObjectType.SHAPE,
                subtype=el.subtype,
                bbox_px=bbox_px,
                bbox_emu=bbox_emu,
                confidence=0.9,
                editable_level=EditableLevel.VECTOR,
                source=SourceAssets(detector="vision_llm"),
                qa=Qa(),
                shape=payload,
            )
            slide.objects.append(obj)

        elif el.type == "image":
            # 원본에서 crop → picture 임베드용 파일 저장
            image_crops_dir.mkdir(parents=True, exist_ok=True)
            crop_path = image_crops_dir / f"img_{k:03d}.png"
            crop = pil.crop((int(bx), int(by), int(bx + bw), int(by + bh)))
            crop.save(crop_path)
            obj = Object(
                id=f"{slide.id}_obj_{k:03d}",
                type=ObjectType.PHOTO,
                bbox_px=bbox_px,
                bbox_emu=bbox_emu,
                confidence=0.9,
                editable_level=EditableLevel.RASTER,
                source=SourceAssets(crop_path=crop_path, detector="vision_llm"),
                qa=Qa(),
            )
            slide.objects.append(obj)

        # 'table' 은 후속 — 현재 cell 단위 textbox + grid shape로 분해 가능하지만
        # 단순화 위해 일단 보류

    return slide, n_text_added


def quick_self_check_cer(predicted: str, ground_truth: str) -> float:
    """간단한 CER 측정 헬퍼 (CLI report 출력용)."""
    return cer(predicted, ground_truth)
