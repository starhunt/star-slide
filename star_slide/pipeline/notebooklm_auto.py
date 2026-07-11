"""NotebookLM image-locked deck auto reconstruction pipeline.

This module intentionally orchestrates the current JSON-driven tools instead of
folding every experiment into the legacy OCR/SAM pipeline. The product target is
one command/web request:

    input.pptx -> editable-ish PPTX + QA report

For each slide it builds both a vector reconstruction and a hybrid
raster-group reconstruction, renders both, then picks the safer layout.
"""

from __future__ import annotations

import contextlib
import json
import re
import shutil
import subprocess
import sys
from collections.abc import Callable
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pydantic import TypeAdapter

from star_slide.config import get_settings
from star_slide.input.pptx_extractor import extract_embedded_images, extract_pdf_pages


@dataclass(frozen=True)
class NotebookLmAutoOptions:
    # 설정 (config.py Settings — .env / STAR_SLIDE_* 환경변수) 의 vision_* 값을
    # default 로 사용. 호출자가 명시하면 override.
    base_url: str = field(default_factory=lambda: get_settings().vision_base_url)
    model: str = field(default_factory=lambda: get_settings().vision_model)
    api_key: str = field(default_factory=lambda: get_settings().vision_api_key)
    timeout_sec: float = 600.0
    retries: int = 2
    use_sam3: bool = False
    hybrid_allowed_delta: float = 0.0
    min_objects: int = 3
    llm_parallel: int = 5
    editable_embedded_text: bool = True
    font_scale: float = 0.93
    keep_intermediates: bool = False
    layout_failure_mode: str = "image_fallback"
    # 워터마크-only 모드:
    #   "off"    — 사용 안 함 (기본 전체 변환)
    #   "fast"   — 우측 하단 영역을 가장자리 평균색으로 단순 페인트 (수초)
    #   "detail" — LaMa 인페인팅으로 자연스럽게 복원 (슬라이드당 1~3초 추가,
    #              첫 호출 시 ~196MB 모델 다운로드)
    watermark_mode: str = "off"
    # 재구성 모드:
    #   "auto"        — 기존 vector + hybrid 자동 선택 (기본)
    #   "image_split" — Codex Vision + image_gen + SAM2 누끼 (단일 이미지 + deck 모두)
    reconstruction_mode: str = "auto"
    # image_split 모드에서 텍스트 제거 방식:
    #   "codex_imagegen" — codex CLI subprocess (builtin image_gen 2.0, 그라데이션 보존)
    #   "solid"          — 주변색 ring sample fill (~1s, 단색 배경 적합)
    text_erase_mode: str = "codex_imagegen"
    # image_split 텍스트 제거 (codex_imagegen 모드) 에 사용할 모델.
    # 빈 문자열이면 codex CLI 기본 모델 사용 (권장). cliproxy 는 input-image edit
    # endpoint 가 없어 image_gen 단계만 codex CLI subprocess 로 직접 호출한다.
    # default 는 STAR_SLIDE_IMAGE_GEN_MODEL 환경변수에서.
    image_gen_model: str = field(default_factory=lambda: get_settings().image_gen_model)
    # image_split 모드의 슬라이드 배경:
    #   "white"       — 흰 캔버스 + alpha 객체 + textbox (default, 깔끔)
    #   "transparent" — 배경 picture 없음 (PowerPoint 기본 슬라이드 배경)
    #   "clean"       — Codex 가 만든 텍스트 제거 이미지 통째 깔기 (시각 충실, 객체 이중 합성)
    background_mode: str = "white"
    # image_split 객체 추출 방식:
    #   True  — Codex Vision 으로 도형 분석 → PPT native shape (rect/oval/arrow) +
    #           픽토그램 bbox alpha crop. 진짜 편집 가능한 layered 객체. (default)
    #   False — SAM2 auto-mask 로 alpha PNG 만 (이전 동작)
    use_native_shapes: bool = True
    # image_split 모드의 입력 이미지에서 우측 하단 NotebookLM 워터마크 자동 제거 (default ON).
    remove_notebooklm_watermark: bool = True
    # hierarchical_overlay 모드에서 raster parent 안의 작은 editable child shape/image를
    # 보존하고 parent crop에서는 해당 영역을 punchout 처리한다.
    child_object_max_area_ratio: float = 0.25
    quiet_subprocesses: bool = False


@dataclass(frozen=True)
class NotebookLmAutoResult:
    output: Path
    workdir: Path
    selected_layout_dir: Path
    report: Path
    vector_pptx: Path
    hybrid_pptx: Path
    artifact_dir: Path
    montage: Path | None = None


class JobCancelledError(RuntimeError):
    """Raised by convert_notebooklm_auto when the cancel callback returns True."""


def repo_root() -> Path:
    return Path(__file__).resolve().parents[2]


def script_path(name: str) -> Path:
    return repo_root() / "scripts" / name


def safe_arg_token(value: str) -> str:
    token = re.sub(r"[^A-Za-z0-9_.-]+", "_", value).strip("._-")
    return token or "deck"


def redact_cmd(cmd: list[str]) -> list[str]:
    redacted: list[str] = []
    hide_next = False
    for part in cmd:
        if hide_next:
            redacted.append("********")
            hide_next = False
            continue
        redacted.append(part)
        if part in {"--api-key", "--vision-api-key"}:
            hide_next = True
    return redacted


def run_cmd(cmd: list[str], *, cwd: Path | None = None, quiet: bool = False) -> None:
    completed = subprocess.run(
        cmd,
        cwd=str(cwd or repo_root()),
        check=False,
        capture_output=quiet,
        text=quiet,
    )
    if completed.returncode != 0:
        printable = " ".join(redact_cmd(cmd))
        raise RuntimeError(f"command failed with exit code {completed.returncode}: {printable}")


def copy_layouts(src_dir: Path, dst_dir: Path) -> None:
    dst_dir.mkdir(parents=True, exist_ok=True)
    for path in sorted(src_dir.glob("*.layout.json")):
        shutil.copy2(path, dst_dir / path.name)


def read_qa(path: Path) -> list[dict[str, Any]]:
    return TypeAdapter(list[dict[str, Any]]).validate_json(path.read_bytes())


def dir_size_bytes(path: Path) -> int:
    if path.is_file():
        return path.stat().st_size
    if not path.exists():
        return 0
    return sum(item.stat().st_size for item in path.rglob("*") if item.is_file())


def report_by_slide(path: Path) -> dict[int, dict[str, Any]]:
    return {int(item["slide_no"]): item for item in read_qa(path)}


def has_raster_groups(groups_dir: Path, slide_stem: str) -> bool:
    path = groups_dir / f"{slide_stem}_raster_groups.json"
    if not path.exists():
        return False
    payload = json.loads(path.read_text(encoding="utf-8"))
    groups = payload.get("raster_groups")
    return isinstance(groups, list) and len(groups) > 0


def build_layout_qa(
    *,
    layout_dir: Path,
    image_root: Path,
    output: Path,
    render_dir: Path,
    allow_images: bool,
    font_scale: float,
    quiet: bool = False,
) -> None:
    layout_args = []
    for path in sorted(layout_dir.glob("*.layout.json")):
        layout_args.extend(["--layout", str(path)])
    cmd = [
        sys.executable,
        str(script_path("layout_qa.py")),
        *layout_args,
        "--image-root",
        str(image_root),
        "-o",
        str(output),
        "--render-dir",
        str(render_dir),
        "--font-scale",
        str(font_scale),
    ]
    if allow_images:
        cmd.append("--allow-images")
    run_cmd(cmd, quiet=quiet)


def generate_layouts(
    *,
    images: list[Path],
    out_dir: Path,
    options: NotebookLmAutoOptions,
    cache_buster: str,
) -> None:
    if options.layout_failure_mode not in {"image_fallback", "fail"}:
        raise ValueError("layout_failure_mode must be 'image_fallback' or 'fail'")
    cmd = [
        sys.executable,
        str(script_path("generate_layout_batch.py")),
        "--images",
        *(str(path) for path in images),
        "--out-dir",
        str(out_dir),
        "--base-url",
        options.base_url,
        "--model",
        options.model,
        "--timeout",
        str(options.timeout_sec),
        "--min-objects",
        str(options.min_objects),
        "--cache-buster",
        cache_buster,
        "--retries",
        str(options.retries),
        "--parallel",
        str(options.llm_parallel),
    ]
    if options.layout_failure_mode == "image_fallback":
        cmd.extend(["--continue-on-error", "--fallback-on-error"])
    if options.api_key:
        cmd.extend(["--api-key", options.api_key])
    run_cmd(cmd, quiet=options.quiet_subprocesses)


def detect_groups(
    *,
    images: list[Path],
    out_dir: Path,
    options: NotebookLmAutoOptions,
    nonce: str,
) -> None:
    cmd = [
        sys.executable,
        str(script_path("detect_raster_groups.py")),
        "--images",
        *(str(path) for path in images),
        "-o",
        str(out_dir),
        "--base-url",
        options.base_url,
        "--model",
        options.model,
        "--timeout",
        str(options.timeout_sec),
        "--nonce",
        nonce,
        "--parallel",
        str(options.llm_parallel),
    ]
    if options.api_key:
        cmd.extend(["--api-key", options.api_key])
    run_cmd(cmd, quiet=options.quiet_subprocesses)


def refine_groups_sam3(
    *, images: list[Path], groups_dir: Path, out_dir: Path, quiet: bool = False
) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    for image in images:
        groups_json = groups_dir / f"{image.stem}_raster_groups.json"
        if not groups_json.exists() or not has_raster_groups(groups_dir, image.stem):
            continue
        run_cmd(
            [
                sys.executable,
                str(script_path("apply_sam3_box_to_raster_groups.py")),
                "--image",
                str(image),
                "--groups-json",
                str(groups_json),
                "-o",
                str(out_dir),
                "--device",
                "auto",
            ],
            quiet=quiet,
        )


def build_hybrid_layouts(
    *,
    layout_dir: Path,
    image_root: Path,
    groups_dir: Path,
    sam_dir: Path | None,
    out_dir: Path,
    slide_count: int,
    editable_embedded_text: bool,
    hierarchical_overlay: bool = False,
    child_object_max_area_ratio: float = 0.25,
    quiet: bool = False,
) -> None:
    cmd = [
        sys.executable,
        str(script_path("apply_raster_groups_to_layout.py")),
        "--layout-dir",
        str(layout_dir),
        "--layout-template",
        "slide_{slide_no:03d}.layout.json",
        "--image-root",
        str(image_root),
        "--groups-dir",
        str(groups_dir),
        "-o",
        str(out_dir),
        "--slides",
        *(str(i) for i in range(1, slide_count + 1)),
        "--erase-mode",
        "inpaint",
        "--drop-full-slide-grid",
    ]
    if not editable_embedded_text:
        cmd.append("--rasterize-embedded-labels")
    if hierarchical_overlay:
        cmd.extend(
            [
                "--peel-child-objects",
                "--child-object-max-area-ratio",
                f"{child_object_max_area_ratio:g}",
            ]
        )
    if sam_dir is not None:
        cmd.extend(["--sam-dir", str(sam_dir)])
    run_cmd(cmd, quiet=quiet)


def select_layouts(
    *,
    vector_layout_dir: Path,
    hybrid_layout_dir: Path,
    groups_dir: Path,
    vector_report: Path,
    hybrid_report: Path,
    out_dir: Path,
    allowed_delta: float,
) -> Path:
    out_dir.mkdir(parents=True, exist_ok=True)
    vector = report_by_slide(vector_report)
    hybrid = report_by_slide(hybrid_report)
    decisions: list[dict[str, Any]] = []

    for layout in sorted(vector_layout_dir.glob("*.layout.json")):
        stem = layout.stem.removesuffix(".layout")
        slide_no = int(stem.split("_")[-1])
        h_layout = hybrid_layout_dir / layout.name
        v_diff = vector.get(slide_no, {}).get("mean_abs_diff")
        h_diff = hybrid.get(slide_no, {}).get("mean_abs_diff")
        has_groups = has_raster_groups(groups_dir, stem)
        choose_hybrid = bool(
            has_groups
            and h_layout.exists()
            and isinstance(v_diff, int | float)
            and isinstance(h_diff, int | float)
            and h_diff <= v_diff + allowed_delta
        )
        chosen = h_layout if choose_hybrid else layout
        reason = "hybrid_diff_better" if choose_hybrid else "vector_default"
        if not has_groups:
            reason = "vector_no_raster_groups"
        elif not h_layout.exists():
            reason = "vector_missing_hybrid_layout"
        elif not (isinstance(v_diff, int | float) and isinstance(h_diff, int | float)):
            reason = "vector_missing_diff"
        shutil.copy2(chosen, out_dir / layout.name)
        decisions.append(
            {
                "slide_no": slide_no,
                "layout": layout.name,
                "chosen": "hybrid" if choose_hybrid else "vector",
                "has_raster_groups": has_groups,
                "vector_mean_abs_diff": v_diff,
                "hybrid_mean_abs_diff": h_diff,
                "reason": reason,
            }
        )

    (out_dir / "selection_report.json").write_text(
        json.dumps(decisions, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return out_dir / "selection_report.json"


def copy_tree_files(src: Path, dst: Path, pattern: str = "*") -> None:
    if not src.exists():
        return
    dst.mkdir(parents=True, exist_ok=True)
    for path in sorted(src.glob(pattern)):
        if path.is_file():
            shutil.copy2(path, dst / path.name)


def make_zip(source_dir: Path, zip_base: Path) -> Path | None:
    if not source_dir.exists():
        return None
    archive = shutil.make_archive(str(zip_base), "zip", root_dir=source_dir)
    return Path(archive)


def remove_generated_sidecars(output_path: Path) -> None:
    sidecar = output_path.parent / f"_{output_path.stem}_assets"
    if sidecar.exists():
        shutil.rmtree(sidecar)


def collect_artifacts(
    *,
    input_path: Path,
    output_path: Path,
    workdir: Path,
    report_path: Path,
    vector_pptx: Path,
    hybrid_pptx: Path,
    vector_layout_dir: Path,
    hybrid_layout_dir: Path,
    selected_layout_dir: Path,
    selected_qa_dir: Path,
    keep_intermediates: bool,
    pre_cleanup_hook: Callable[[Path], None] | None = None,
) -> dict[str, Path | None]:
    artifact_dir = workdir.parent / "artifacts"
    artifact_dir.mkdir(parents=True, exist_ok=True)

    final_output = output_path
    vector_artifact = artifact_dir / "candidate_vector.pptx"
    hybrid_artifact = artifact_dir / "candidate_hybrid.pptx"
    report_artifact = artifact_dir / "report.json"
    montage_artifact = artifact_dir / "montage.png"
    layout_dir = artifact_dir / "layouts"

    if vector_pptx.exists():
        shutil.copy2(vector_pptx, vector_artifact)
    if hybrid_pptx.exists():
        shutil.copy2(hybrid_pptx, hybrid_artifact)
    if report_path.exists():
        shutil.copy2(report_path, report_artifact)
    if (selected_qa_dir / "montage.png").exists():
        shutil.copy2(selected_qa_dir / "montage.png", montage_artifact)

    copy_tree_files(vector_layout_dir, layout_dir / "vector", "*.layout.json")
    copy_tree_files(vector_layout_dir, layout_dir / "vector", "*.usage.json")
    copy_tree_files(hybrid_layout_dir, layout_dir / "hybrid", "*.layout.json")
    copy_tree_files(selected_layout_dir, layout_dir / "selected", "*.layout.json")
    if (selected_layout_dir / "selection_report.json").exists():
        shutil.copy2(
            selected_layout_dir / "selection_report.json",
            layout_dir / "selected" / "selection_report.json",
        )
    layout_zip = make_zip(layout_dir, artifact_dir / "layout_json")

    summary = {
        "input": str(input_path),
        "output": str(final_output),
        "artifacts": {
            "candidate_vector": str(vector_artifact) if vector_artifact.exists() else None,
            "candidate_hybrid": str(hybrid_artifact) if hybrid_artifact.exists() else None,
            "report": str(report_artifact) if report_artifact.exists() else None,
            "montage": str(montage_artifact) if montage_artifact.exists() else None,
            "layout_json_zip": str(layout_zip) if layout_zip else None,
        },
        "sizes": {
            "input_bytes": dir_size_bytes(input_path),
            "output_bytes": dir_size_bytes(final_output),
            "artifact_dir_bytes": dir_size_bytes(artifact_dir),
            "workdir_before_cleanup_bytes": dir_size_bytes(workdir),
        },
        "keep_intermediates": keep_intermediates,
    }
    (artifact_dir / "artifact_manifest.json").write_text(
        json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    remove_generated_sidecars(output_path)
    remove_generated_sidecars(vector_pptx)
    remove_generated_sidecars(hybrid_pptx)
    if pre_cleanup_hook is not None and workdir.exists():
        with contextlib.suppress(Exception):
            pre_cleanup_hook(workdir)
    if not keep_intermediates and workdir.exists():
        shutil.rmtree(workdir)

    return {
        "artifact_dir": artifact_dir,
        "report": report_artifact if report_artifact.exists() else report_path,
        "montage": montage_artifact
        if montage_artifact.exists()
        else selected_qa_dir / "montage.png",
        "vector_pptx": vector_artifact if vector_artifact.exists() else vector_pptx,
        "hybrid_pptx": hybrid_artifact if hybrid_artifact.exists() else hybrid_pptx,
        "selected_layout_dir": layout_dir / "selected"
        if (layout_dir / "selected").exists()
        else selected_layout_dir,
    }


def _build_image_only_pptx(images: list[Path], output_path: Path) -> None:
    """슬라이드당 단일 PICTURE shape 로 16:9 PPTX 를 빌드한다."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    # 16:9 — PowerPoint Widescreen 기본 (12192000 x 6858000 EMU)
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank_layout = prs.slide_layouts[6]
    for img_path in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
    prs.save(str(output_path))


def _convert_watermark_only(
    *,
    input_path: Path,
    output_path: Path,
    workdir: Path,
    images: list[Path],
    emit: Callable[[str, float], None],
    check_cancel: Callable[[], None],
    pre_cleanup_hook: Callable[[Path], None] | None,
    keep_intermediates: bool,
    detail: bool = False,
) -> NotebookLmAutoResult:
    """LLM 호출 없이 우측 하단 워터마크 영역만 제거해 PPTX 를 빌드한다.

    detail=True 면 LaMa 인페인팅으로 배경을 자연스럽게 복원하고,
    detail=False 면 가장자리 평균색으로 단순 페인트한다 (수초).
    """
    from star_slide.input.watermark_remover import (
        remove_watermarks,
        remove_watermarks_inpaint,
    )

    check_cancel()
    if detail:
        emit("NotebookLM 워터마크 제거 중 (디테일 모드, LaMa 인페인팅)", 30)
        remove_watermarks_inpaint(images)
    else:
        emit("NotebookLM 워터마크 제거 중 (빠른 모드)", 30)
        remove_watermarks(images)

    # 미리보기/썸네일 호환을 위해 마스킹된 이미지를 qa_selected 에도 복사한다.
    qa_selected_dir = workdir / "qa_selected"
    qa_selected_dir.mkdir(parents=True, exist_ok=True)
    for src in images:
        dst = qa_selected_dir / src.name
        with contextlib.suppress(OSError):
            shutil.copy2(src, dst)

    check_cancel()
    emit("PPTX 빌드 중", 70)
    _build_image_only_pptx(images, output_path)

    # 간단 report 작성 (web 의 report-summary 가 빈 객체라도 파싱할 수 있도록 최소 구조)
    artifact_dir = workdir.parent / "artifacts"
    artifact_dir.mkdir(parents=True, exist_ok=True)
    report_path = artifact_dir / "report.json"
    report_path.write_text(
        json.dumps(
            {
                "input": str(input_path),
                "output": str(output_path),
                "mode": "watermark_only",
                "slides": len(images),
                "selected_qa": [],
                "vector_qa": [],
                "hybrid_qa": [],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    # cleanup 직전 hook (web_app 의 썸네일/캐시 prefetch)
    if pre_cleanup_hook is not None and workdir.exists():
        with contextlib.suppress(Exception):
            pre_cleanup_hook(workdir)
    if not keep_intermediates and workdir.exists():
        shutil.rmtree(workdir)

    emit("완료", 100)
    return NotebookLmAutoResult(
        output=output_path,
        workdir=workdir,
        selected_layout_dir=workdir,  # 빠른 모드는 layout 산출물 없음 — 더미 경로
        report=report_path,
        vector_pptx=output_path,
        hybrid_pptx=output_path,
        artifact_dir=artifact_dir,
        montage=None,
    )


def _convert_image_split(
    *,
    input_images: list[Path],
    input_path: Path,
    output_path: Path,
    workdir: Path,
    options: NotebookLmAutoOptions,
    emit: Callable[[str, float], None],
    check_cancel: Callable[[], None],
    cancel: Callable[[], bool] | None,
    pre_cleanup_hook: Callable[[Path], None] | None,
    keep_intermediates: bool,
) -> NotebookLmAutoResult:
    """N장 이미지 → 슬라이드별 Codex Vision + image_gen + SAM2 누끼 → editable PPTX.

    convert_image_split_multi 모듈을 호출하고 결과물을 web 의 기존 artifact/preview
    구조 (qa_selected/montage.png, artifacts/report.json) 와 호환되도록 후처리.
    """
    from star_slide.pipeline.codex_image_split import (
        ImageSplitOptions,
        convert_image_split_multi,
    )

    n = len(input_images)
    check_cancel()
    emit(f"image_split 파이프라인 시작 ({n} 슬라이드, Codex + SAM2)", 3)

    split_options = ImageSplitOptions(
        text_erase_mode=options.text_erase_mode,
        background_mode=options.background_mode,
        use_native_shapes=options.use_native_shapes,
        remove_notebooklm_watermark=options.remove_notebooklm_watermark,
        slide_parallel=max(1, int(options.llm_parallel)),
        vision_base_url=options.base_url,
        vision_model=options.model,
        image_gen_model=options.image_gen_model,
        vision_api_key=options.api_key,
        vision_timeout_sec=options.timeout_sec,
    )
    split_workdir = workdir / "image_split"
    result = convert_image_split_multi(
        input_images=input_images,
        output_pptx=output_path,
        workdir=split_workdir,
        options=split_options,
        progress=emit,
        cancel=cancel,
    )

    # 미리보기/썸네일 호환: clean bg 를 qa_selected/slide-1.png 로 복사
    qa_selected_dir = workdir / "qa_selected"
    qa_selected_dir.mkdir(parents=True, exist_ok=True)
    if result.clean_bg_png.exists():
        with contextlib.suppress(OSError):
            shutil.copy2(result.clean_bg_png, qa_selected_dir / "slide-1.png")
    # web preview 가 montage.png 를 찾으므로 같은 이미지를 montage 로도 사용
    if result.clean_bg_png.exists():
        with contextlib.suppress(OSError):
            shutil.copy2(result.clean_bg_png, qa_selected_dir / "montage.png")

    artifact_dir = workdir.parent / "artifacts"
    artifact_dir.mkdir(parents=True, exist_ok=True)
    report_path = artifact_dir / "report.json"
    report_path.write_text(
        json.dumps(
            {
                "input": str(input_path),
                "output": str(output_path),
                "mode": "image_split",
                "text_erase_mode": options.text_erase_mode,
                "elapsed_sec": result.elapsed_sec,
                "text_layout_json": str(result.text_layout_json),
                "object_layers_json": str(result.object_layers_json),
                "clean_bg_png": str(result.clean_bg_png),
                "selected_qa": [],
                "vector_qa": [],
                "hybrid_qa": [],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    montage_artifact = artifact_dir / "montage.png"
    if result.clean_bg_png.exists() and not montage_artifact.exists():
        with contextlib.suppress(OSError):
            shutil.copy2(result.clean_bg_png, montage_artifact)

    if pre_cleanup_hook is not None and workdir.exists():
        with contextlib.suppress(Exception):
            pre_cleanup_hook(workdir)
    if not keep_intermediates and workdir.exists():
        shutil.rmtree(workdir)

    emit("완료", 100)
    return NotebookLmAutoResult(
        output=output_path,
        workdir=workdir,
        selected_layout_dir=split_workdir,
        report=report_path,
        vector_pptx=output_path,
        hybrid_pptx=output_path,
        artifact_dir=artifact_dir,
        montage=montage_artifact if montage_artifact.exists() else None,
    )


def convert_notebooklm_auto(
    *,
    input_path: Path,
    output_path: Path,
    workdir: Path,
    options: NotebookLmAutoOptions,
    progress: Callable[..., None] | None = None,
    cancel: Callable[[], bool] | None = None,
    pre_cleanup_hook: Callable[[Path], None] | None = None,
) -> NotebookLmAutoResult:
    reconstruction_modes = {"auto", "hierarchical_overlay", "image_split"}
    if options.reconstruction_mode not in reconstruction_modes:
        raise ValueError(
            "reconstruction_mode must be one of "
            f"{sorted(reconstruction_modes)}: {options.reconstruction_mode!r}"
        )

    def emit(message: str, percent: float) -> None:
        if progress is not None:
            progress(message, percent)

    def check_cancel() -> None:
        if cancel is not None and cancel():
            raise JobCancelledError("convert_notebooklm_auto cancelled")

    workdir.mkdir(parents=True, exist_ok=True)

    check_cancel()
    emit("슬라이드 이미지 추출 중", 3)
    image_dir = workdir / "images"
    suffix = input_path.suffix.lower()
    if suffix == ".pdf":
        images = extract_pdf_pages(input_path, image_dir)
    elif suffix == ".pptx":
        images = extract_embedded_images(input_path, image_dir)
    else:
        raise RuntimeError(f"unsupported input file type: {input_path.suffix}")
    if not images:
        raise RuntimeError(f"no slide images extracted from {input_path}")
    images = sorted(images)

    # 워터마크-only 모드 — LLM/SAM 호출 없이 곧장 PPTX 빌드.
    if options.watermark_mode in {"fast", "detail"}:
        return _convert_watermark_only(
            input_path=input_path,
            output_path=output_path,
            workdir=workdir,
            images=images,
            emit=emit,
            check_cancel=check_cancel,
            pre_cleanup_hook=pre_cleanup_hook,
            keep_intermediates=options.keep_intermediates,
            detail=(options.watermark_mode == "detail"),
        )

    # image_split 모드 — 단일 이미지 / deck 모두 지원.
    # 각 슬라이드마다 Codex Vision + image_gen + SAM2 누끼 (~80s/슬라이드).
    if options.reconstruction_mode == "image_split":
        return _convert_image_split(
            input_images=images,
            input_path=input_path,
            output_path=output_path,
            workdir=workdir,
            options=options,
            emit=emit,
            check_cancel=check_cancel,
            cancel=cancel,
            pre_cleanup_hook=pre_cleanup_hook,
            keep_intermediates=options.keep_intermediates,
        )

    check_cancel()
    emit("Vision LLM layout JSON 생성 중", 10)
    vector_layout_dir = workdir / "layouts_vector"
    generate_layouts(
        images=images,
        out_dir=vector_layout_dir,
        options=options,
        cache_buster=f"job-{safe_arg_token(input_path.stem)}-layout",
    )

    check_cancel()
    emit("vector PPTX 렌더 QA 중", 38)
    vector_pptx = workdir / "vector.pptx"
    vector_qa_dir = workdir / "qa_vector"
    build_layout_qa(
        layout_dir=vector_layout_dir,
        image_root=image_dir,
        output=vector_pptx,
        render_dir=vector_qa_dir,
        allow_images=options.layout_failure_mode == "image_fallback",
        font_scale=options.font_scale,
        quiet=options.quiet_subprocesses,
    )

    check_cancel()
    emit("큰 이미지 그룹 탐지 중", 48)
    groups_dir = workdir / "raster_groups"
    detect_groups(
        images=images,
        out_dir=groups_dir,
        options=options,
        nonce=f"job-{safe_arg_token(input_path.stem)}-raster",
    )

    check_cancel()
    emit("SAM3 이미지 그룹 보정 중", 60)
    sam_dir = workdir / "raster_groups_sam3" if options.use_sam3 else None
    if sam_dir is not None:
        refine_groups_sam3(
            images=images,
            groups_dir=groups_dir,
            out_dir=sam_dir,
            quiet=options.quiet_subprocesses,
        )

    check_cancel()
    emit("hybrid layout 생성 중", 70)
    hybrid_layout_dir = workdir / "layouts_hybrid"
    build_hybrid_layouts(
        layout_dir=vector_layout_dir,
        image_root=image_dir,
        groups_dir=groups_dir,
        sam_dir=sam_dir,
        out_dir=hybrid_layout_dir,
        slide_count=len(images),
        editable_embedded_text=options.editable_embedded_text,
        hierarchical_overlay=options.reconstruction_mode == "hierarchical_overlay",
        child_object_max_area_ratio=options.child_object_max_area_ratio,
        quiet=options.quiet_subprocesses,
    )

    check_cancel()
    emit("hybrid PPTX 렌더 QA 중", 78)
    hybrid_pptx = workdir / "hybrid.pptx"
    hybrid_qa_dir = workdir / "qa_hybrid"
    build_layout_qa(
        layout_dir=hybrid_layout_dir,
        image_root=image_dir,
        output=hybrid_pptx,
        render_dir=hybrid_qa_dir,
        allow_images=True,
        font_scale=options.font_scale,
        quiet=options.quiet_subprocesses,
    )

    check_cancel()
    emit("최종 layout 자동 선택 중", 88)
    selected_layout_dir = workdir / "layouts_selected"
    selection_report = select_layouts(
        vector_layout_dir=vector_layout_dir,
        hybrid_layout_dir=hybrid_layout_dir,
        groups_dir=groups_dir,
        vector_report=vector_qa_dir / "qa_report.json",
        hybrid_report=hybrid_qa_dir / "qa_report.json",
        out_dir=selected_layout_dir,
        allowed_delta=options.hybrid_allowed_delta,
    )

    check_cancel()
    emit("최종 PPTX 생성 및 미리보기 렌더링 중", 93)
    selected_qa_dir = workdir / "qa_selected"
    build_layout_qa(
        layout_dir=selected_layout_dir,
        image_root=image_dir,
        output=output_path,
        render_dir=selected_qa_dir,
        allow_images=True,
        font_scale=options.font_scale,
        quiet=options.quiet_subprocesses,
    )

    combined_report = {
        "input": str(input_path),
        "output": str(output_path),
        "workdir": str(workdir),
        "selection_report": json.loads(selection_report.read_text(encoding="utf-8")),
        "selected_qa": read_qa(selected_qa_dir / "qa_report.json"),
        "vector_qa": read_qa(vector_qa_dir / "qa_report.json"),
        "hybrid_qa": read_qa(hybrid_qa_dir / "qa_report.json"),
    }
    report_path = workdir / "notebooklm_auto_report.json"
    report_path.write_text(
        json.dumps(combined_report, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    artifacts = collect_artifacts(
        input_path=input_path,
        output_path=output_path,
        workdir=workdir,
        report_path=report_path,
        vector_pptx=vector_pptx,
        hybrid_pptx=hybrid_pptx,
        vector_layout_dir=vector_layout_dir,
        hybrid_layout_dir=hybrid_layout_dir,
        selected_layout_dir=selected_layout_dir,
        selected_qa_dir=selected_qa_dir,
        keep_intermediates=options.keep_intermediates,
        pre_cleanup_hook=pre_cleanup_hook,
    )
    emit("완료", 100)

    return NotebookLmAutoResult(
        output=output_path,
        workdir=workdir,
        selected_layout_dir=artifacts["selected_layout_dir"] or selected_layout_dir,
        report=artifacts["report"] or report_path,
        vector_pptx=artifacts["vector_pptx"] or vector_pptx,
        hybrid_pptx=artifacts["hybrid_pptx"] or hybrid_pptx,
        artifact_dir=artifacts["artifact_dir"] or workdir.parent / "artifacts",
        montage=artifacts["montage"],
    )
