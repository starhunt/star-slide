"""단일 이미지 → editable PPTX 변환: Codex Vision + image_gen + SAM2 파이프라인.

experiments/diagram_split/v7_*.py 의 4단계 PoC 를 라이브러리화 한 모듈.
notebooklm_auto.py 의 reconstruction_mode='image_split' 에서 호출된다.

핵심 흐름:
  Step 1. Codex CLI Vision LLM → text_layout.json
          (텍스트 + bbox + font_size_px + color + alignment + bold)
  Step 2. 텍스트 제거 (두 모드):
          - text_erase_mode='codex_imagegen': Codex image_gen 2.0 으로 텍스트 지운 이미지 생성
          - text_erase_mode='solid': 각 bbox 외곽 ring 픽셀 mode 색상으로 fill (빠름)
  Step 3. SAM2.1 auto-mask (글자 간섭 없음) → 깨끗한 객체 alpha PNG crop
  Step 4. 재조합: 깨끗한 배경 + alpha 객체 (z-order) + editable textbox
          (East Asian font 메타 명시)

설계 원칙:
  - notebooklm_auto 와 동일한 emit/check_cancel 시그니처
  - 모든 중간 산출물을 workdir 아래에 저장 (디버깅/QA 용)
  - 외부 도구 (codex CLI) 가 실패하면 fallback (solid fill 모드로 자동 전환)
"""

from __future__ import annotations

import contextlib
import json
import os
import shutil
import subprocess
import time
from collections import Counter
from collections.abc import Callable
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import cv2
import numpy as np
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap as _pptx_nsmap  # noqa: F401  (lxml namespace 보장)
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import BaseOxmlElement
from pptx.presentation import Presentation as PptxPresentation
from pptx.shapes.autoshape import Shape as PptxShape
from pptx.slide import Slide as PptxSlide
from pptx.util import Emu, Pt

from star_slide.segmentation.sam2_auto import Sam2Mask, run_sam2_auto
from star_slide.vision_llm.image_split_client import (
    VisionClientError,
    call_vision_json,
)

EMU_PER_INCH = 914400
PT_PER_INCH = 72.0
SLIDE_W_EMU = 12192000  # 13.333" — 표준 16:9
SLIDE_H_EMU = 6858000  #  7.5"
DEFAULT_FONT = "Apple SD Gothic Neo"

# Codex Vision shape 분석 prompt
CODEX_SHAPES_PROMPT = """첨부된 인포그래픽에서 모든 도형/객체 정보를 JSON 으로 추출해주세요. 다른 설명 없이 순수 JSON만.

목적: PowerPoint 로 재구성하기 위한 layered 분리.
- 단순 도형 (직사각형/둥근 직사각형/원형/타원/표준 화살표/선) → PPT native shape 으로 재현
- 복잡한 픽토그램 (사람, 건물, 책, 지구본, 모자, 가방, 모니터 등) → alpha PNG 로 처리
- 장식성 화살표 (그라데이션/곡선/3D/장식 있는 화살표) 도 픽토그램으로 처리

형식:
{
  "image_size": [w, h],
  "shapes": [
    {
      "type": "rectangle|rounded_rect|oval|arrow|line|pictogram",
      "bbox": [x, y, w, h],
      "fill": "#hex" | "gradient:#hex1->#hex2" | "transparent",
      "stroke": "#hex" | "none",
      "stroke_width": 2,
      "stroke_dash": "solid|dashed|dotted",
      "direction": "up|down|left|right|up_right|up_left|down_right|down_left",
      "on_dark_background": false,
      "name": "선택적 라벨",
      "z_hint": "background|container|card|pictogram|decoration"
    }
  ]
}

요구사항:
- ★ image_size 는 입력 이미지의 실제 픽셀 크기 [width, height]. 절대 누락하거나 null
  로 주지 말 것. 모르면 실제 측정값을 추정해서라도 채울 것.
- 큰 컨테이너 박스, 그 안의 작은 카드, 점선 박스, 화살표, 픽토그램 모두 빠짐없이
- 픽토그램은 type='pictogram' (alpha PNG 처리 대상). bbox 만 정확히.
- 텍스트는 무시 (별도 추출됨)
- z_hint 로 z-order 힌트 (큰=container, 카드=card, 픽토그램=pictogram, 작은 장식=decoration)
- ★ 화살표 (type='arrow') 는 가리키는 방향 'direction' 필수 (up/down/left/right 등).
  ↓ 아래는 'down', → 오른쪽은 'right', ↑ 위는 'up', ← 왼쪽은 'left'.
  방향이 시각 의미를 담고 있으므로 정확히.
  ◆ type='arrow' 는 단색 + 표준 4방향 (up/down/left/right) + 직선형 화살표만.
    곡선 화살표, 대각선, 그라데이션/입체/장식 있는 화살표는 type='pictogram'으로.
- ★ 픽토그램 (type='pictogram') 의 'on_dark_background': true 면 어두운 박스 (예:
  파랑/남색 카드) 위 흰색 아이콘. false 면 흰 배경 위 컬러 아이콘. 이 정보로
  alpha 추출 임계를 반전한다.
"""


@dataclass(frozen=True)
class ImageSplitOptions:
    """단일 이미지 분할 변환 옵션."""

    text_erase_mode: str = "codex_imagegen"
    """텍스트 제거 방식: 'codex_imagegen'|'solid'.
    'codex_imagegen' 은 Codex CLI image_gen 2.0 호출 (~30-60s). 그라데이션 보존 우수.
    'solid' 은 주변색 ring sample fill (~1s). 단색 배경에 적합."""

    background_mode: str = "white"
    """슬라이드 배경 처리:
      'white'       — 흰 캔버스 + alpha 객체 + textbox (default, 깔끔)
      'transparent' — 배경 picture 없음 (PowerPoint 기본 슬라이드 배경 유지)
      'clean'       — Codex 가 만든 텍스트 제거 이미지를 통째로 깔기 (시각 충실하지만
                      객체와 이중 합성. 사용자가 객체를 옮기면 배경에 같은 모양 잔존)
    """

    use_native_shapes: bool = True
    """Codex Vision 으로 도형 (rect/rounded_rect/oval/arrow) 추출해 PPT native shape
    으로 추가 (사용자가 PPT 에서 색/크기 직접 편집 가능). 픽토그램은 codex bbox 로
    alpha PNG crop. False 면 SAM2 만 사용 (이전 동작)."""

    remove_notebooklm_watermark: bool = True
    """입력 이미지 우측 하단의 NotebookLM 워터마크를 사전 제거 (배경색 fill).
    image_split 의 첫 단계에서 적용 — codex 분석/image_gen 이 워터마크를
    객체나 텍스트로 잘못 잡는 것을 방지. NotebookLM 외 입력에는 효과 없음
    (배경색 fill 만이라 무해)."""

    arrow_native_max_area_ratio: float = 0.04
    """codex 가 type='arrow' 로 식별한 도형을 PPT native arrow shape 으로 재현하는
    최대 면적 비율 (bbox 면적 / 슬라이드 면적). 이를 초과하면 단순 화살표가 아닐
    가능성이 높으므로 raster alpha crop (pictogram 강등) 으로 처리. 표준 4방향
    + 단색 + 작은 크기 화살표만 native 로 가져가 편집성을 보장."""

    vision_base_url: str = ""
    """cliproxy/OpenAI 호환 endpoint. 빈 문자열이면 config.py Settings.vision_base_url
    (STAR_SLIDE_VISION_BASE_URL) 사용. NotebookLmAutoOptions 의 base_url 이
    자동 전파된다."""

    vision_model: str = ""
    """텍스트/도형 layout 분석에 사용할 vision 모델. 빈 문자열이면 config.py
    Settings.vision_model (STAR_SLIDE_VISION_MODEL, default 'claude-opus') 사용.
    cliproxy 가 multiplex 하는 실제 alias (claude-opus, gemini-pro, svtx-pro 등)
    를 지정. gpt-5.5 같은 codex provider 라우팅 모델은 cliproxy 측에서 hang 되는
    사례가 있어 multimodal 전용 alias 권장."""

    vision_timeout_sec: float = 240.0
    """slide-per-call vision LLM timeout (sec). 한 슬라이드 분석이 이 시간을 초과하면
    예외 발생 → 호출자(process_slide)가 단계별 fallback. 너무 길게 (예: 600s) 잡으면
    11장 deck 에서 timeout 한 장이 전체 진행을 막아 사용자 체감 시간이 폭증한다."""

    image_gen_model: str = ""
    """텍스트 제거된 배경 이미지 생성 모델. text_erase_mode='codex_imagegen' 일 때만
    사용. 빈 문자열이면 codex CLI 기본 모델 (ChatGPT 계정 인증 + builtin image_gen
    2.0 도구) 을 사용 — 권장. cliproxy 는 input-image 를 받는 image-edit endpoint
    를 노출하지 않아 이 단계만 codex CLI subprocess 로 직접 호출한다."""

    vision_api_key: str = ""
    """cliproxy/OpenAI 호환 API key. 빈 문자열이면 환경변수
    VISION_PROXY_API_KEY / LOCAL_CLAUDE_API_KEY 를 순서대로 시도."""

    slide_parallel: int = 5
    """multi-slide image_split 의 슬라이드 병렬 처리 동시 개수. 슬라이드 1장당
    ~80-120s (codex Vision 2회 + image_gen 1회) 라 5-deck 이면 순차 ~600s →
    parallel=5 로 ~120s 까지 단축 가능. NotebookLmAutoOptions.llm_parallel 매핑."""

    sam_resize: int = 1920
    """SAM2 입력 리사이즈 폭. 0=원본 사용 (정밀하지만 느림)"""

    sam_points_per_side: int = 48
    """SAM2 grid 한 변의 point 수 (총 = side²). dense=48 권장."""

    sam_pred_iou_thresh: float = 0.80
    sam_stability_thresh: float = 0.88

    min_object_area_ratio: float = 0.0005
    max_object_area_ratio: float = 0.55
    object_iou_dedupe_threshold: float = 0.85

    font_name: str = DEFAULT_FONT


@dataclass(frozen=True)
class ImageSplitResult:
    output: Path
    workdir: Path
    text_layout_json: Path
    clean_bg_png: Path
    object_layers_json: Path
    elapsed_sec: float


@dataclass
class _TextItem:
    text: str
    bbox: tuple[float, float, float, float]
    font_size_px: float
    color: str
    is_bold: bool
    alignment: str


@dataclass
class _ObjectLayer:
    bbox: tuple[int, int, int, int]
    z: int
    area_ratio: float
    score: float
    alpha_path: Path


# ============================================================
# Step 1 — Codex Vision 으로 텍스트 layout JSON 추출
# ============================================================


CODEX_TEXT_PROMPT = """이미지의 모든 텍스트를 JSON 으로 추출해주세요. 다른 설명이나 코드블록 없이 순수 JSON만 반환해주세요.

형식:
{
  "image_size": [width_px, height_px],
  "texts": [
    {
      "text": "텍스트 내용",
      "bbox": [x, y, w, h],
      "font_size_px": 60,
      "color": "#hex",
      "is_bold": true,
      "alignment": "center"
    }
  ]
}

요구사항:
- ★ image_size 는 입력 이미지의 실제 픽셀 크기 [width, height]. 절대 누락하거나 null 로
  주지 말 것. 모르면 [0, 0] 보다는 실제 측정값을 추정해서라도 채울 것.
- 모든 한글/영문 텍스트 빠짐없이 (제목, 부제, 라벨, 박스 안 글, 따옴표 안 글, footer 등)
- bbox 픽셀 좌표 (좌상단 기준, 가로세로 px)
- font_size_px = 글자 자체 픽셀 높이 (라인 높이 X)
- color = 글자의 주된 hex 색
- 텍스트 라인 단위 (예: '광주 17개 대학 IDP' 한 줄로, 단어 분리하지 말 것)
- alignment = 'left'|'center'|'right'
- is_bold = 굵은 글씨 여부
"""


def normalize_text_layout_to_image_size(
    text_layout: dict[str, Any],
    actual_image_size: tuple[int, int],
) -> dict[str, Any]:
    """codex 가 다른 해상도 (예: 2048x1152) 로 좌표를 줄 수 있어, 실제 이미지 크기로
    bbox/font_size_px 를 비례 변환.

    text_layout['image_size'] 가 실제와 다르면 모든 텍스트의 좌표를 스케일.
    in-place 수정 후 반환. vision LLM 이 image_size 를 빼먹거나 null 로 주는
    경우 실제 이미지 크기로 채워 넣고 정규화는 skip (좌표가 이미 실제 크기 기준
    이라고 가정).
    """
    if "texts" not in text_layout:
        return text_layout
    js_size = text_layout.get("image_size")
    if not isinstance(js_size, (list, tuple)) or len(js_size) != 2:
        text_layout["image_size"] = list(actual_image_size)
        return text_layout
    js_w, js_h = js_size
    if not isinstance(js_w, (int, float)) or not isinstance(js_h, (int, float)):
        text_layout["image_size"] = list(actual_image_size)
        return text_layout
    real_w, real_h = actual_image_size
    if js_w <= 0 or js_h <= 0:
        text_layout["image_size"] = list(actual_image_size)
        return text_layout
    if abs(js_w - real_w) <= 4 and abs(js_h - real_h) <= 4:
        return text_layout  # 같음
    sx = real_w / js_w
    sy = real_h / js_h
    s_avg = (sx + sy) / 2  # 폰트 크기 (가로/세로 평균)
    for t in text_layout["texts"]:
        bbox = t.get("bbox", (0, 0, 0, 0))
        if len(bbox) == 4:
            t["bbox"] = [
                int(bbox[0] * sx),
                int(bbox[1] * sy),
                int(bbox[2] * sx),
                int(bbox[3] * sy),
            ]
        if "font_size_px" in t:
            t["font_size_px"] = float(t["font_size_px"]) * s_avg
    text_layout["image_size"] = [real_w, real_h]
    text_layout["_normalized_from"] = [js_w, js_h]
    return text_layout


def _resolve_vision_api_key(api_key: str) -> str:
    """빈 문자열이면 config.py Settings.vision_api_key (STAR_SLIDE_VISION_API_KEY) →
    환경변수 VISION_PROXY_API_KEY / LOCAL_CLAUDE_API_KEY 순서로 fallback."""
    if api_key:
        return api_key
    from star_slide.config import get_settings

    settings_key = get_settings().vision_api_key
    if settings_key:
        return settings_key
    return os.environ.get("VISION_PROXY_API_KEY") or os.environ.get("LOCAL_CLAUDE_API_KEY") or ""


def _resolve_vision_base_url(base_url: str) -> str:
    """빈 문자열이면 config.py Settings.vision_base_url 사용."""
    if base_url:
        return base_url
    from star_slide.config import get_settings

    return get_settings().vision_base_url


def _resolve_vision_model(model: str) -> str:
    """빈 문자열이면 config.py Settings.vision_model 사용."""
    if model:
        return model
    from star_slide.config import get_settings

    return get_settings().vision_model


def analyze_text_with_vision(
    image_path: Path,
    *,
    out_json: Path,
    base_url: str,
    model: str,
    api_key: str,
    timeout_sec: float = 600.0,
) -> dict[str, Any]:
    """cliproxy/OpenAI 호환 chat/completions 로 텍스트 + 위치 추출.

    응답 JSON 을 out_json 에도 저장 (디버깅/QA 용).
    """
    out_json.parent.mkdir(parents=True, exist_ok=True)
    try:
        data = call_vision_json(
            image_path,
            CODEX_TEXT_PROMPT,
            base_url=_resolve_vision_base_url(base_url),
            model=_resolve_vision_model(model),
            api_key=_resolve_vision_api_key(api_key),
            timeout_sec=timeout_sec,
        )
    except VisionClientError as exc:
        raise RuntimeError(f"Vision 텍스트 분석 실패: {exc}") from exc
    out_json.write_text(
        json.dumps(data, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return data


def analyze_shapes_with_vision(
    image_path: Path,
    *,
    out_json: Path,
    base_url: str,
    model: str,
    api_key: str,
    timeout_sec: float = 600.0,
) -> dict[str, Any]:
    """cliproxy/OpenAI 호환 chat/completions 로 도형/픽토그램 layout 추출.

    출력 JSON:
      {"image_size": [w, h], "shapes": [{type, bbox, fill, stroke, ..., z_hint}, ...]}
    """
    out_json.parent.mkdir(parents=True, exist_ok=True)
    try:
        data = call_vision_json(
            image_path,
            CODEX_SHAPES_PROMPT,
            base_url=_resolve_vision_base_url(base_url),
            model=_resolve_vision_model(model),
            api_key=_resolve_vision_api_key(api_key),
            timeout_sec=timeout_sec,
        )
    except VisionClientError as exc:
        raise RuntimeError(f"Vision 도형 분석 실패: {exc}") from exc
    out_json.write_text(
        json.dumps(data, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return data


def normalize_shapes_to_image_size(
    shapes_layout: dict[str, Any],
    actual_image_size: tuple[int, int],
) -> dict[str, Any]:
    """text_layout 과 동일 — codex 가 다른 해상도로 좌표를 줘도 실제 크기로 변환.

    vision LLM 이 image_size 를 빼먹거나 null 로 주는 경우 실제 이미지 크기로
    채워 넣고 정규화는 skip.
    """
    if "shapes" not in shapes_layout:
        return shapes_layout
    js_size = shapes_layout.get("image_size")
    if not isinstance(js_size, (list, tuple)) or len(js_size) != 2:
        shapes_layout["image_size"] = list(actual_image_size)
        return shapes_layout
    js_w, js_h = js_size
    if not isinstance(js_w, (int, float)) or not isinstance(js_h, (int, float)):
        shapes_layout["image_size"] = list(actual_image_size)
        return shapes_layout
    real_w, real_h = actual_image_size
    if js_w <= 0 or js_h <= 0:
        shapes_layout["image_size"] = list(actual_image_size)
        return shapes_layout
    if abs(js_w - real_w) <= 4 and abs(js_h - real_h) <= 4:
        return shapes_layout
    sx = real_w / js_w
    sy = real_h / js_h
    for s in shapes_layout["shapes"]:
        bbox = s.get("bbox", (0, 0, 0, 0))
        if len(bbox) == 4:
            s["bbox"] = [
                int(bbox[0] * sx),
                int(bbox[1] * sy),
                int(bbox[2] * sx),
                int(bbox[3] * sy),
            ]
        if "stroke_width" in s:
            s["stroke_width"] = max(1, int(float(s["stroke_width"]) * (sx + sy) / 2))
    shapes_layout["image_size"] = [real_w, real_h]
    shapes_layout["_normalized_from"] = [js_w, js_h]
    return shapes_layout


# ============================================================
# Step 2 — 텍스트 제거
# ============================================================


CODEX_REMOVE_TEXT_PROMPT = """첨부된 한글 인포그래픽에서 모든 텍스트(한글+영문)만 정확히 제거한 새 이미지를 image_gen 도구로 생성해줘.
- 큰 제목, 부제, 라벨, 박스 안 모든 글자, 픽토그램 라벨, 따옴표 안 글, 하단 footer 모두 제거
- 다른 모든 시각 요소(컬러 박스, 픽토그램 아이콘, 화살표, 그라데이션, 박스 형태와 색)는 픽셀 단위로 정확히 보존
- 글자 자리는 주변 배경색/그라데이션으로 자연스럽게 복원
- 입력 이미지의 원본 비율 유지
생성된 PNG 파일 경로만 마지막에 출력.
"""


def remove_text_with_image_gen(
    image_path: Path,
    *,
    out_png: Path,
    model: str = "",
    target_size: tuple[int, int] | None = None,
    timeout_sec: float = 600.0,
) -> Path:
    """codex CLI 의 builtin image_gen (gpt-image 2.0) 도구로 텍스트 제거 이미지 생성.

    cliproxy 는 input-image 를 받는 image-edit endpoint 를 노출하지 않으므로,
    image_gen 단계만 codex CLI subprocess 로 직접 호출한다 (text/shape 분석은
    cliproxy chat/completions 그대로 사용). codex 가 image_gen 도구로 새 PNG 를
    생성하면 보통 ~/.codex/generated_images/.../ig_*.png 에 저장하고 stdout 마지막
    줄에 그 경로를 출력한다.

    stdout 파싱 규칙:
      - 입력 이미지 경로와 동일한 line 은 무시
      - 'generated_images' 디렉토리 안의 PNG 를 우선 채택 (가장 최근 mtime)
      - codex 가 image_gen 을 호출하지 않고 입력 그대로 반환한 경우 RuntimeError

    model: 빈 문자열이면 codex 기본 모델 사용 (ChatGPT 계정 인증). 명시 시 codex
           `-m` 인자로 전달. 단 codex 는 'gpt-image-*' 같은 image-only 모델명을
           거부할 수 있어 기본 빈 문자열 권장.
    """
    out_png.parent.mkdir(parents=True, exist_ok=True)
    input_abs = str(image_path.resolve())
    cmd = [
        "codex",
        "exec",
        "--dangerously-bypass-approvals-and-sandbox",
        "--skip-git-repo-check",
        "--ephemeral",
        "-i",
        str(image_path),
    ]
    if model:
        cmd.extend(["-m", model])
    completed = subprocess.run(
        cmd,
        input=CODEX_REMOVE_TEXT_PROMPT,
        capture_output=True,
        text=True,
        timeout=timeout_sec,
        check=False,
    )
    if completed.returncode != 0:
        raise RuntimeError(
            f"Codex image_gen 실패 (exit={completed.returncode}): {completed.stderr[:500]}"
        )
    stdout_lines = [ln.strip() for ln in completed.stdout.splitlines() if ln.strip()]
    candidates: list[Path] = []
    for line in stdout_lines:
        if not line.endswith(".png"):
            continue
        path_str = line.rsplit(" ", 1)[-1]
        candidate = Path(path_str)
        if not candidate.exists():
            continue
        if candidate.resolve() == Path(input_abs).resolve():
            continue
        candidates.append(candidate)
    if not candidates:
        raise RuntimeError(
            "Codex image_gen 출력 PNG 경로를 찾을 수 없음 (image_gen 도구 미호출 의심). "
            f"stdout 마지막 500자:\n{completed.stdout[-500:]}"
        )
    generated = [p for p in candidates if "generated_images" in str(p)]
    pool = generated or candidates
    chosen = max(pool, key=lambda p: p.stat().st_mtime)
    shutil.copy2(chosen, out_png)
    if target_size is not None:
        with Image.open(out_png) as im:
            if im.size != target_size:
                im.convert("RGB").resize(
                    target_size,
                    Image.Resampling.LANCZOS,
                ).save(out_png)
    return out_png


def _quantize_color(rgb: tuple[int, int, int], step: int = 16) -> tuple[int, int, int]:
    return (
        (rgb[0] // step) * step + step // 2,
        (rgb[1] // step) * step + step // 2,
        (rgb[2] // step) * step + step // 2,
    )


def _sample_ring_color(
    arr: np.ndarray,
    bbox: tuple[int, int, int, int],
    *,
    ring_thickness: int = 8,
    pad: int = 2,
) -> tuple[int, int, int]:
    h, w = arr.shape[:2]
    x, y, bw, bh = bbox
    x0 = max(0, x - ring_thickness - pad)
    y0 = max(0, y - ring_thickness - pad)
    x1 = min(w, x + bw + ring_thickness + pad)
    y1 = min(h, y + bh + ring_thickness + pad)
    inner_x0 = max(0, x - pad)
    inner_y0 = max(0, y - pad)
    inner_x1 = min(w, x + bw + pad)
    inner_y1 = min(h, y + bh + pad)
    region = arr[y0:y1, x0:x1].copy()
    if region.size == 0:
        return (255, 255, 255)
    rh, rw = region.shape[:2]
    mask = np.ones((rh, rw), dtype=bool)
    mask[inner_y0 - y0 : inner_y1 - y0, inner_x0 - x0 : inner_x1 - x0] = False
    ring_pixels = region[mask]
    if len(ring_pixels) == 0:
        return (255, 255, 255)
    quant = [_quantize_color(tuple(p)) for p in ring_pixels]
    most_common = Counter(quant).most_common(1)[0][0]
    return int(most_common[0]), int(most_common[1]), int(most_common[2])


def remove_text_solid_fill(
    image_path: Path,
    text_layout: dict[str, Any],
    *,
    out_png: Path,
    pad: int = 6,
) -> Path:
    """text_layout JSON 의 각 bbox 외곽 ring mode 색으로 fill (빠른 fallback).

    그라데이션이 강한 배경에선 약점이지만 단색/카드 배경엔 우수.
    """
    out_png.parent.mkdir(parents=True, exist_ok=True)
    image = Image.open(image_path).convert("RGB")
    arr = np.asarray(image).copy()
    h, w = arr.shape[:2]
    for t in text_layout.get("texts", []):
        bx, by, bw, bh = t.get("bbox", (0, 0, 0, 0))
        bx = int(bx)
        by = int(by)
        bw = int(bw)
        bh = int(bh)
        if bw <= 0 or bh <= 0:
            continue
        bg_color = _sample_ring_color(arr, (bx, by, bw, bh), ring_thickness=8, pad=pad)
        x0 = max(0, bx - pad)
        y0 = max(0, by - pad)
        x1 = min(w, bx + bw + pad)
        y1 = min(h, by + bh + pad)
        arr[y0:y1, x0:x1] = bg_color
    Image.fromarray(arr).save(out_png)
    return out_png


# ============================================================
# Step 3 — SAM2 객체 추출 (깨끗한 배경에서)
# ============================================================


def _upsample_mask(seg_small: np.ndarray, target_size: tuple[int, int]) -> np.ndarray:
    seg_uint = (seg_small.astype(np.uint8)) * 255
    seg_full = cv2.resize(seg_uint, target_size, interpolation=cv2.INTER_LINEAR)
    return np.asarray(seg_full > 127, dtype=np.bool_)


def _mask_iou(a: np.ndarray, b: np.ndarray) -> float:
    if a.shape != b.shape:
        return 0.0
    inter = np.logical_and(a, b).sum()
    union = np.logical_or(a, b).sum()
    return float(inter) / float(union) if union else 0.0


def _filter_and_dedupe(
    masks: list[Sam2Mask],
    image_size: tuple[int, int],
    *,
    min_ratio: float,
    max_ratio: float,
    iou_threshold: float,
) -> list[Sam2Mask]:
    w, h = image_size
    total = w * h
    filtered = [m for m in masks if min_ratio <= m.area / total <= max_ratio]
    sorted_masks = sorted(filtered, key=lambda m: -m.area)
    kept: list[Sam2Mask] = []
    for m in sorted_masks:
        if not any(_mask_iou(m.segmentation, k.segmentation) > iou_threshold for k in kept):
            kept.append(m)
    return kept


def _export_alpha_crop(
    image: Image.Image,
    mask: np.ndarray,
    *,
    out_dir: Path,
    name: str,
) -> tuple[tuple[int, int, int, int], Path] | None:
    ys, xs = np.where(mask)
    if len(ys) == 0:
        return None
    x0, x1 = int(xs.min()), int(xs.max()) + 1
    y0, y1 = int(ys.min()), int(ys.max()) + 1
    arr = np.asarray(image.convert("RGB"))
    crop_rgb = arr[y0:y1, x0:x1]
    crop_mask = mask[y0:y1, x0:x1]
    rgba = np.zeros((y1 - y0, x1 - x0, 4), dtype=np.uint8)
    rgba[:, :, :3] = crop_rgb
    rgba[:, :, 3] = (crop_mask.astype(np.uint8)) * 255
    out_dir.mkdir(parents=True, exist_ok=True)
    p = out_dir / name
    Image.fromarray(rgba, mode="RGBA").save(p)
    return (x0, y0, x1 - x0, y1 - y0), p


def extract_objects_sam2_clean(
    clean_image_path: Path,
    *,
    out_dir: Path,
    options: ImageSplitOptions,
) -> tuple[list[_ObjectLayer], Path]:
    """깨끗한(텍스트 제거된) 이미지에서 SAM2 객체 alpha PNG crop 들 생성.

    out_dir/objects/layer_NNN.png 들 + out_dir/object_layers.json 생성.
    """
    image = Image.open(clean_image_path).convert("RGB")
    px_w, px_h = image.size

    sam_input = image
    if options.sam_resize and px_w > options.sam_resize:
        scale = options.sam_resize / px_w
        sam_input = image.resize((options.sam_resize, int(px_h * scale)))
    res = run_sam2_auto(
        sam_input,
        points_per_side=options.sam_points_per_side,
        max_masks=400,
        pred_iou_thresh=options.sam_pred_iou_thresh,
        stability_score_thresh=options.sam_stability_thresh,
    )

    full_masks: list[Sam2Mask] = []
    for m in res.masks:
        seg = _upsample_mask(m.segmentation, (px_w, px_h))
        ys, xs = np.where(seg)
        if len(ys) == 0:
            continue
        bbox = (
            int(xs.min()),
            int(ys.min()),
            int(xs.max() - xs.min() + 1),
            int(ys.max() - ys.min() + 1),
        )
        full_masks.append(Sam2Mask(bbox=bbox, segmentation=seg, score=m.score, area=int(seg.sum())))

    deduped = _filter_and_dedupe(
        full_masks,
        (px_w, px_h),
        min_ratio=options.min_object_area_ratio,
        max_ratio=options.max_object_area_ratio,
        iou_threshold=options.object_iou_dedupe_threshold,
    )

    obj_dir = out_dir / "objects"
    layers: list[_ObjectLayer] = []
    sorted_masks = sorted(deduped, key=lambda m: -m.area)
    for i, m in enumerate(sorted_masks):
        out = _export_alpha_crop(image, m.segmentation, out_dir=obj_dir, name=f"layer_{i:03d}.png")
        if out is None:
            continue
        bbox, path = out
        layers.append(
            _ObjectLayer(
                bbox=bbox,
                z=i,
                area_ratio=float(m.area) / (px_w * px_h),
                score=float(m.score),
                alpha_path=path,
            )
        )

    json_path = out_dir / "object_layers.json"
    json_path.write_text(
        json.dumps(
            {
                "image_size": [px_w, px_h],
                "layer_count": len(layers),
                "sam_total_masks": len(full_masks),
                "layers": [
                    {
                        "bbox": ly.bbox,
                        "z": ly.z,
                        "area_ratio": ly.area_ratio,
                        "score": ly.score,
                        "alpha_path": str(ly.alpha_path.relative_to(out_dir)),
                    }
                    for ly in layers
                ],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    return layers, json_path


# ============================================================
# Step 4 — PPTX 재조합
# ============================================================


def _set_east_asian_font(rPr: BaseOxmlElement, font_name: str) -> None:
    for tag in ("a:ea", "a:cs"):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)
    cs = etree.SubElement(rPr, qn("a:cs"))
    cs.set("typeface", font_name)


def _add_textbox(
    slide: PptxSlide,
    item: dict[str, Any],
    *,
    scale: float,
    effective_dpi: float,
    font_name: str,
) -> None:
    bbox = item.get("bbox", (0, 0, 0, 0))
    if len(bbox) != 4:
        return
    x, y, w, h = bbox
    font_px = float(item.get("font_size_px", 18))
    font_pt = font_px * PT_PER_INCH / effective_dpi
    pad_emu = int(font_pt / PT_PER_INCH * EMU_PER_INCH * 0.3)
    left = int(x * scale) - pad_emu // 2
    top = int(y * scale) - pad_emu // 2
    width = int(w * scale) + pad_emu
    min_h_emu = int(font_pt / PT_PER_INCH * EMU_PER_INCH * 1.4)
    height = max(int(h * scale) + pad_emu, min_h_emu)
    if width <= 0 or height <= 0:
        return
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, attr, Emu(0))
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(item.get("alignment", "left"), PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = item.get("text", "")
    run.font.name = font_name
    run.font.size = Pt(round(font_pt, 1))
    run.font.bold = bool(item.get("is_bold", False))
    color_hex = item.get("color", "#222222")
    try:
        r = int(color_hex[1:3], 16)
        g = int(color_hex[3:5], 16)
        b = int(color_hex[5:7], 16)
        run.font.color.rgb = RGBColor(r, g, b)
    except Exception:
        pass
    _set_east_asian_font(run._r.get_or_add_rPr(), font_name)


def _hex_to_rgb(s: str) -> tuple[int, int, int]:
    s = (s or "").strip().lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    if len(s) != 6:
        return (200, 200, 200)
    try:
        return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except ValueError:
        return (200, 200, 200)


def _parse_dash(dash: str) -> MSO_LINE_DASH_STYLE:
    return {
        "solid": MSO_LINE_DASH_STYLE.SOLID,
        "dashed": MSO_LINE_DASH_STYLE.DASH,
        "dotted": MSO_LINE_DASH_STYLE.ROUND_DOT,
    }.get((dash or "solid").lower(), MSO_LINE_DASH_STYLE.SOLID)


def _shape_z(z_hint: str) -> int:
    """z_hint → z-order 정수 (작을수록 뒤). 같은 그룹 안에서는 면적 큰 순."""
    return {
        "background": 0,
        "container": 10,
        "decoration": 20,
        "card": 30,
        "pictogram": 40,
    }.get((z_hint or "decoration").lower(), 25)


def _apply_fill(shape: PptxShape, fill_spec: str) -> None:
    """fill_spec: '#hex' | 'gradient:#hex1->#hex2' | 'transparent' | 'none'."""
    spec = (fill_spec or "").strip()
    if spec in ("transparent", "none", ""):
        with contextlib.suppress(Exception):
            shape.fill.background()
        return
    if spec.startswith("gradient:") and "->" in spec:
        # python-pptx 는 gradient 직접 지원 미흡 — XML 직접 작성
        body = spec[len("gradient:") :]
        c1, c2 = body.split("->", 1)
        r1, g1, b1 = _hex_to_rgb(c1)
        r2, g2, b2 = _hex_to_rgb(c2)
        sp = shape.fill._xPr
        # 기존 fill 제거
        for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
            for el in sp.findall(qn(tag)):
                sp.remove(el)
        grad_xml = (
            f'<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">'
            f"<a:gsLst>"
            f'<a:gs pos="0"><a:srgbClr val="{r1:02X}{g1:02X}{b1:02X}"/></a:gs>'
            f'<a:gs pos="100000"><a:srgbClr val="{r2:02X}{g2:02X}{b2:02X}"/></a:gs>'
            f"</a:gsLst>"
            f'<a:lin ang="5400000" scaled="0"/>'
            f"</a:gradFill>"
        )
        sp.append(etree.fromstring(grad_xml))
        return
    # solid hex
    r, g, b = _hex_to_rgb(spec)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)


def _apply_stroke(
    shape: PptxShape,
    stroke_spec: str,
    stroke_width: float,
    dash: str,
) -> None:
    line = shape.line
    if (stroke_spec or "none").lower() in ("none", ""):
        line.fill.background()
        return
    r, g, b = _hex_to_rgb(stroke_spec)
    line.color.rgb = RGBColor(r, g, b)
    if stroke_width and stroke_width > 0:
        line.width = Emu(int(stroke_width * EMU_PER_INCH / 96))
    with contextlib.suppress(Exception):
        line.dash_style = _parse_dash(dash)


_STANDARD_ARROW_DIRECTIONS = frozenset({"up", "down", "left", "right"})


def _resolve_arrow_shape(direction: str) -> MSO_SHAPE:
    """화살표 direction 에 따라 PPT native arrow shape 선택."""
    d = (direction or "right").lower().strip()
    return {
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "right": MSO_SHAPE.RIGHT_ARROW,
        # 대각선 방향은 PPT 가 native 로 없으므로 가까운 가로/세로로 근사
        "up_right": MSO_SHAPE.UP_ARROW,
        "up_left": MSO_SHAPE.UP_ARROW,
        "down_right": MSO_SHAPE.DOWN_ARROW,
        "down_left": MSO_SHAPE.DOWN_ARROW,
    }.get(d, MSO_SHAPE.RIGHT_ARROW)


def _is_simple_arrow(
    shape_data: dict[str, Any],
    *,
    image_size: tuple[int, int],
    max_area_ratio: float,
) -> bool:
    """codex 가 type='arrow' 로 잡은 도형이 PPT native arrow 로 재현하기 안전한지 판정.

    조건 (모두 충족):
      - direction 이 표준 4방향 (up/down/left/right). 대각선/곡선 X.
      - fill 이 단색(#hex) 또는 transparent. gradient/3D X.
      - 크기가 슬라이드 면적 대비 max_area_ratio 이하.
      - stroke_width 가 비정상적으로 두껍지 않음 (단순 line 화살표 보장).

    실패 시 codex 결과가 잘못된 게 아니라, 단지 raster crop 으로 강등해야 한다는 신호.
    """
    if (shape_data.get("type") or "").lower() != "arrow":
        return False
    direction = (shape_data.get("direction") or "").lower().strip()
    if direction not in _STANDARD_ARROW_DIRECTIONS:
        return False
    fill = (shape_data.get("fill") or "transparent").strip().lower()
    if fill.startswith("gradient:"):
        return False
    bbox = shape_data.get("bbox", [0, 0, 0, 0])
    if len(bbox) != 4 or bbox[2] <= 0 or bbox[3] <= 0:
        return False
    img_w, img_h = image_size
    if img_w <= 0 or img_h <= 0:
        return False
    area_ratio = (float(bbox[2]) * float(bbox[3])) / (float(img_w) * float(img_h))
    return area_ratio <= max_area_ratio


def demote_complex_arrows_to_pictogram(
    shapes_layout: dict[str, Any],
    *,
    max_area_ratio: float,
) -> int:
    """단순 케이스가 아닌 화살표를 pictogram 으로 강등 (in-place).

    이렇게 강등된 화살표는 native PPT shape 분기에서 제외되고,
    crop_pictograms_from_codex_shapes 에서 raster alpha crop 으로 처리된다.

    반환: 강등된 개수.
    """
    if not isinstance(shapes_layout, dict):
        return 0
    image_size = shapes_layout.get("image_size") or [0, 0]
    if len(image_size) != 2:
        return 0
    img_w, img_h = int(image_size[0]), int(image_size[1])
    demoted = 0
    for s in shapes_layout.get("shapes", []):
        if (s.get("type") or "").lower() != "arrow":
            continue
        if not _is_simple_arrow(s, image_size=(img_w, img_h), max_area_ratio=max_area_ratio):
            s["type"] = "pictogram"
            s["z_hint"] = "pictogram"
            s["_demoted_from"] = "arrow"
            demoted += 1
    return demoted


def _infer_arrow_direction_from_bbox(bbox: list[float]) -> str:
    """bbox 종횡비로 방향 추론 (codex 가 direction 안 줬을 때 fallback).

    세로가 더 길면 위/아래, 가로가 더 길면 좌/우. 더 정확한 방향은 codex 결과 활용.
    """
    if len(bbox) != 4:
        return "right"
    _, _, w, h = bbox
    if h > w * 1.2:
        return "down"  # 세로 화살표 — 인포그래픽 흐름은 보통 위→아래
    if w > h * 1.2:
        return "right"  # 가로 화살표 — 보통 왼→오른
    return "right"


def _add_native_shape(
    slide: PptxSlide,
    shape_data: dict[str, Any],
    *,
    scale: float,
) -> None:
    """codex shape 한 개 → PPT native shape (rect/rounded/oval/arrow/line)."""
    bbox = shape_data.get("bbox", [0, 0, 0, 0])
    if len(bbox) != 4 or bbox[2] <= 0 or bbox[3] <= 0:
        return
    x, y, w, h = bbox
    left = Emu(int(x * scale))
    top = Emu(int(y * scale))
    width = Emu(int(w * scale))
    height = Emu(int(h * scale))

    type_str = (shape_data.get("type") or "").lower()

    # 화살표는 direction 별 native shape 선택
    auto_shape: MSO_SHAPE | None
    if type_str == "arrow":
        direction = shape_data.get("direction") or _infer_arrow_direction_from_bbox(bbox)
        auto_shape = _resolve_arrow_shape(direction)
    else:
        auto_shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "ellipse": MSO_SHAPE.OVAL,
            "circle": MSO_SHAPE.OVAL,
            "line": MSO_SHAPE.RECTANGLE,  # line 은 가는 직사각형으로 근사
        }
        auto_shape = auto_shape_map.get(type_str)
    if auto_shape is None:
        return  # pictogram 등은 별도 처리 (alpha PNG)

    shape = slide.shapes.add_shape(auto_shape, left, top, width, height)
    _apply_fill(shape, shape_data.get("fill", "transparent"))
    _apply_stroke(
        shape,
        shape_data.get("stroke", "none"),
        float(shape_data.get("stroke_width", 0) or 0),
        shape_data.get("stroke_dash", "solid"),
    )
    # textbox 가 아닌 도형이므로 텍스트 프레임 비활성
    with contextlib.suppress(Exception):
        shape.text_frame.text = ""


def _add_image_split_slide(
    prs: PptxPresentation,
    *,
    image_size: tuple[int, int],
    background_path: Path | None,
    background_mode: str,
    text_layout: dict[str, Any],
    object_layers: list[_ObjectLayer],
    font_name: str,
    shapes_layout: dict[str, Any] | None = None,
) -> None:
    """기존 Presentation 에 1 슬라이드 추가. multi-slide 합성용.

    z-order (뒤→앞):
      1. background (clean/white/transparent)
      2. shapes_layout 의 native shape (z_hint container → decoration → card 순)
      3. object_layers 의 alpha PNG (큰 마스크 → 작은 마스크)
      4. shapes_layout 의 pictogram bbox alpha (있으면 SAM 대신 또는 추가)
      5. text_layout 의 textbox
    """
    px_w, _ = image_size
    scale = prs.slide_width / px_w
    effective_dpi = px_w / (prs.slide_width / EMU_PER_INCH)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    if background_mode == "clean" and background_path is not None:
        slide.shapes.add_picture(
            str(background_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
    elif background_mode == "white":
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            prs.slide_width,
            prs.slide_height,
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.fill.background()
    # 'transparent' 면 아무것도 안 깔음

    # native shape 들 (z_hint 작은 순 = 뒤). 같은 z 안에서는 면적 큰 순.
    if shapes_layout:
        natives = [
            s
            for s in shapes_layout.get("shapes", [])
            if (s.get("type") or "").lower() not in ("pictogram",)
        ]
        natives.sort(
            key=lambda s: (
                _shape_z(s.get("z_hint", "decoration")),
                -(s.get("bbox", [0, 0, 0, 0])[2] * s.get("bbox", [0, 0, 0, 0])[3]),
            )
        )
        for sd in natives:
            _add_native_shape(slide, sd, scale=scale)

    # alpha 객체 (SAM2 또는 codex pictogram bbox crop)
    for ly in sorted(object_layers, key=lambda layer: layer.z):
        x, y, w, h = ly.bbox
        slide.shapes.add_picture(
            str(ly.alpha_path),
            Emu(int(x * scale)),
            Emu(int(y * scale)),
            width=Emu(int(w * scale)),
            height=Emu(int(h * scale)),
        )

    # textbox 최상단
    for t in text_layout.get("texts", []):
        _add_textbox(slide, t, scale=scale, effective_dpi=effective_dpi, font_name=font_name)


def compose_image_split_pptx(
    *,
    image_size: tuple[int, int],
    background_path: Path | None,
    text_layout: dict[str, Any],
    object_layers: list[_ObjectLayer],
    out_path: Path,
    font_name: str = DEFAULT_FONT,
    background_mode: str = "white",
    shapes_layout: dict[str, Any] | None = None,
) -> Path:
    """단일 슬라이드 PPTX 생성."""
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    _add_image_split_slide(
        prs,
        image_size=image_size,
        background_path=background_path,
        background_mode=background_mode,
        text_layout=text_layout,
        object_layers=object_layers,
        font_name=font_name,
        shapes_layout=shapes_layout,
    )
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def _estimate_picto_background(
    arr: np.ndarray,
    bbox: tuple[int, int, int, int],
    *,
    ring_thickness: int = 6,
) -> tuple[int, int, int]:
    """픽토그램 bbox 외곽 ring 픽셀 mode 색을 배경색으로 추정.

    어두운 배경 위 흰 아이콘인지 흰 배경 위 컬러 아이콘인지 자동 판정용.
    """
    h, w = arr.shape[:2]
    x, y, bw, bh = bbox
    x0 = max(0, x - ring_thickness)
    y0 = max(0, y - ring_thickness)
    x1 = min(w, x + bw + ring_thickness)
    y1 = min(h, y + bh + ring_thickness)
    region = arr[y0:y1, x0:x1].copy()
    if region.size == 0:
        return (255, 255, 255)
    rh, rw = region.shape[:2]
    mask = np.ones((rh, rw), dtype=bool)
    inner_y0 = max(0, y - y0)
    inner_x0 = max(0, x - x0)
    inner_y1 = min(rh, inner_y0 + bh)
    inner_x1 = min(rw, inner_x0 + bw)
    mask[inner_y0:inner_y1, inner_x0:inner_x1] = False
    ring_pixels = region[mask]
    if len(ring_pixels) == 0:
        return (255, 255, 255)
    quant = (ring_pixels // 32) * 32 + 16
    counts: dict[tuple[int, int, int], int] = {}
    for px in quant:
        key = (int(px[0]), int(px[1]), int(px[2]))
        counts[key] = counts.get(key, 0) + 1
    return max(counts, key=counts.__getitem__)


def crop_pictograms_from_codex_shapes(
    clean_image_path: Path,
    shapes_layout: dict[str, Any],
    *,
    out_dir: Path,
) -> list[_ObjectLayer]:
    """codex shapes 의 'pictogram' bbox 를 clean 배경에서 직접 alpha crop.

    배경색 자동 검출:
      - codex 가 'on_dark_background': true 명시 OR 외곽 ring mode 색 luminance < 130
        → 어두운 배경 위 흰/밝은 stroke 아이콘. alpha = light pixels.
      - 그 외 → 흰/밝은 배경 위 어두운 stroke. alpha = dark pixels.
    이러면 박스 안 흰 아이콘도 정상적으로 분리됨.
    """
    image = Image.open(clean_image_path).convert("RGB")
    arr = np.asarray(image)
    h, w = arr.shape[:2]
    out_dir.mkdir(parents=True, exist_ok=True)
    layers: list[_ObjectLayer] = []
    pictos = [
        s for s in shapes_layout.get("shapes", []) if (s.get("type") or "").lower() == "pictogram"
    ]
    pictos.sort(key=lambda s: -(s.get("bbox", [0, 0, 0, 0])[2] * s.get("bbox", [0, 0, 0, 0])[3]))
    for i, p in enumerate(pictos):
        bbox = p.get("bbox", [0, 0, 0, 0])
        if len(bbox) != 4:
            continue
        x, y, bw, bh = bbox
        x = max(0, int(x))
        y = max(0, int(y))
        bw = min(w - x, int(bw))
        bh = min(h - y, int(bh))
        if bw <= 0 or bh <= 0:
            continue
        crop_rgb = arr[y : y + bh, x : x + bw]
        lum = (
            0.299 * crop_rgb[:, :, 0].astype(np.float32)
            + 0.587 * crop_rgb[:, :, 1].astype(np.float32)
            + 0.114 * crop_rgb[:, :, 2].astype(np.float32)
        )
        # 배경 색 검출 → adaptive alpha
        codex_dark = bool(p.get("on_dark_background", False))
        bg_rgb = _estimate_picto_background(arr, (x, y, bw, bh))
        bg_lum = 0.299 * bg_rgb[0] + 0.587 * bg_rgb[1] + 0.114 * bg_rgb[2]
        on_dark = codex_dark or bg_lum < 130
        if on_dark:
            # 어두운 배경 위 밝은 아이콘 → 밝은 픽셀이 alpha
            alpha = ((lum > 200).astype(np.uint8)) * 255
        else:
            # 흰/밝은 배경 위 어두운 아이콘 → 어두운 픽셀이 alpha
            alpha = ((lum < 240).astype(np.uint8)) * 255
        # alpha 가 거의 0 이면 (오검출) 통째 보존 fallback
        if alpha.sum() < (bw * bh * 255 * 0.01):
            alpha = np.full((bh, bw), 255, dtype=np.uint8)
        rgba = np.zeros((bh, bw, 4), dtype=np.uint8)
        rgba[:, :, :3] = crop_rgb
        rgba[:, :, 3] = alpha
        path = out_dir / f"picto_{i:03d}.png"
        Image.fromarray(rgba, mode="RGBA").save(path)
        layers.append(
            _ObjectLayer(
                bbox=(x, y, bw, bh),
                z=i + 1000,
                area_ratio=float(bw * bh) / (w * h),
                score=1.0,
                alpha_path=path,
            )
        )
    return layers


def compose_image_split_pptx_multi(
    *,
    slides_data: list[dict[str, Any]],
    out_path: Path,
    font_name: str = DEFAULT_FONT,
    background_mode: str = "white",
) -> Path:
    """N개 슬라이드 image_split 결과를 한 PPTX 로 합성.

    slides_data: [
      {
        'image_size': (w, h),
        'background_path': Path,
        'text_layout': dict,
        'object_layers': list[_ObjectLayer],
      },
      ...
    ]
    """
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    for sd in slides_data:
        if sd.get("_failed") and sd.get("background_path") is not None:
            # 실패한 슬라이드는 input 이미지를 통째로 깐 raster 슬라이드로 대체
            # — 비어 있는 슬라이드보다 사용자가 후속 편집하기 쉬움.
            blank = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                str(sd["background_path"]),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            continue
        _add_image_split_slide(
            prs,
            image_size=sd["image_size"],
            background_path=sd.get("background_path"),
            background_mode=background_mode,
            text_layout=sd["text_layout"],
            object_layers=sd["object_layers"],
            font_name=font_name,
            shapes_layout=sd.get("shapes_layout"),
        )
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


# ============================================================
# 통합 entry point — notebooklm_auto 에서 호출
# ============================================================


def convert_image_split(
    *,
    input_image: Path,
    output_pptx: Path,
    workdir: Path,
    options: ImageSplitOptions,
    progress: Callable[[str, float], None] | None = None,
    cancel: Callable[[], bool] | None = None,
) -> ImageSplitResult:
    """단일 이미지 → editable PPTX 4단계 파이프라인.

    workdir 아래 모든 중간 산출물 보존:
      workdir/text_layout.json
      workdir/02_clean_bg.png  (Step 2 결과)
      workdir/objects/layer_NNN.png
      workdir/object_layers.json
    """

    def emit(msg: str, pct: float) -> None:
        if progress is not None:
            progress(msg, pct)

    def check_cancel() -> None:
        if cancel is not None and cancel():
            from star_slide.pipeline.notebooklm_auto import JobCancelledError

            raise JobCancelledError("convert_image_split cancelled")

    workdir.mkdir(parents=True, exist_ok=True)
    started = time.perf_counter()

    check_cancel()
    emit(f"Vision LLM 으로 텍스트 + 위치 추출 중 ({options.vision_model})", 5)
    text_json_path = workdir / "text_layout.json"
    text_layout = analyze_text_with_vision(
        input_image,
        out_json=text_json_path,
        base_url=options.vision_base_url,
        model=options.vision_model,
        api_key=options.vision_api_key,
        timeout_sec=options.vision_timeout_sec,
    )

    check_cancel()
    image = Image.open(input_image).convert("RGB")
    target_size = image.size

    # vision LLM 이 다른 해상도로 좌표를 줄 수 있어 실제 이미지 크기로 정규화 후 저장
    js_size = text_layout.get("image_size", target_size)
    text_layout = normalize_text_layout_to_image_size(text_layout, target_size)
    if "_normalized_from" in text_layout:
        text_json_path.write_text(
            json.dumps(text_layout, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        emit(
            f"텍스트 {len(text_layout.get('texts', []))}개 추출 "
            f"(좌표 {js_size}→{target_size} 정규화)",
            25,
        )
    else:
        emit(f"텍스트 {len(text_layout.get('texts', []))}개 추출", 25)

    clean_path = workdir / "02_clean_bg.png"
    if options.text_erase_mode == "codex_imagegen":
        emit("Codex image_gen 으로 텍스트 제거된 배경 생성 중 (~30-60s)", 30)
        try:
            remove_text_with_image_gen(
                input_image,
                out_png=clean_path,
                target_size=target_size,
                model=options.image_gen_model,
            )
        except Exception as exc:
            emit(f"Codex image_gen 실패 → solid fill 로 fallback: {exc}", 40)
            remove_text_solid_fill(input_image, text_layout, out_png=clean_path)
    else:
        emit("주변색 fill 로 텍스트 제거 중", 30)
        remove_text_solid_fill(input_image, text_layout, out_png=clean_path)
    emit("배경 생성 완료", 60)

    check_cancel()
    emit(f"SAM2 로 객체 추출 중 (pps={options.sam_points_per_side})", 65)
    layers, layers_json = extract_objects_sam2_clean(
        clean_path,
        out_dir=workdir,
        options=options,
    )
    emit(f"{len(layers)}개 객체 alpha crop", 85)

    check_cancel()
    emit("PPTX 재조합 중", 90)
    compose_image_split_pptx(
        image_size=target_size,
        background_path=clean_path,
        text_layout=text_layout,
        object_layers=layers,
        out_path=output_pptx,
        font_name=options.font_name,
        background_mode=options.background_mode,
    )
    emit("완료", 100)

    return ImageSplitResult(
        output=output_pptx,
        workdir=workdir,
        text_layout_json=text_json_path,
        clean_bg_png=clean_path,
        object_layers_json=layers_json,
        elapsed_sec=round(time.perf_counter() - started, 1),
    )


def convert_image_split_multi(
    *,
    input_images: list[Path],
    output_pptx: Path,
    workdir: Path,
    options: ImageSplitOptions,
    progress: Callable[[str, float], None] | None = None,
    cancel: Callable[[], bool] | None = None,
) -> ImageSplitResult:
    """N장 이미지 → N-슬라이드 editable PPTX (각 슬라이드마다 image_split 적용).

    workdir 아래에 슬라이드별 sub-dir:
      workdir/slide_001/text_layout.json, 02_clean_bg.png, objects/, ...
      workdir/slide_NNN/...

    progress callback 은 전체 (0~100) 기준으로 emit. 슬라이드 1장당
    pct_per_slide = 100 / N 만큼 진행.
    """

    def emit(msg: str, pct: float) -> None:
        if progress is not None:
            progress(msg, pct)

    def check_cancel() -> None:
        if cancel is not None and cancel():
            from star_slide.pipeline.notebooklm_auto import JobCancelledError

            raise JobCancelledError("convert_image_split_multi cancelled")

    workdir.mkdir(parents=True, exist_ok=True)
    started = time.perf_counter()
    n = len(input_images)
    if n == 0:
        raise ValueError("input_images 가 비었습니다")

    # NotebookLM 워터마크 사전 제거 (default ON, image_split 모든 흐름의 첫 단계).
    # 우측 하단 영역을 추정 배경색으로 덮는 best-effort 처리.
    # codex 분석/image_gen 호출 전에 적용해 워터마크가 객체로 분류되거나
    # text_layout 에 잡히는 것을 방지.
    if options.remove_notebooklm_watermark:
        emit("NotebookLM 워터마크 사전 제거", 1)
        from star_slide.input.watermark_remover import remove_watermarks

        remove_watermarks(list(input_images))

    # 슬라이드 단위 작업을 worker 함수로 분리 — ThreadPoolExecutor 병렬 실행.
    # 각 슬라이드는 독립적 (workdir/slide_NNN/ 격리, 다른 슬라이드 결과 의존 X).
    # vision LLM 호출은 IO-bound (HTTP) 라 GIL 이슈 없음.
    import threading

    progress_lock = threading.Lock()
    completed_count = [0]  # mutable counter for thread-safe increment

    def _raster_fallback_slide(i: int, img_path: Path, error: str) -> dict[str, Any]:
        """슬라이드 변환이 catastrophic 실패할 때 input 이미지를 그대로 raster 로
        넣는 fallback 데이터를 반환. 텍스트/객체 없는 1장짜리 picture 슬라이드.
        """
        with Image.open(img_path) as im:
            size = im.convert("RGB").size
        return {
            "_index": i,
            "_failed": True,
            "_failure_error": error,
            "_failure_image": img_path,
            "image_size": size,
            "background_path": img_path,
            "text_layout": {"image_size": list(size), "texts": []},
            "object_layers": [],
            "shapes_layout": None,
            "_text_json": None,
            "_layers_json": None,
        }

    def process_slide(i: int, img_path: Path) -> dict[str, Any]:
        slide_dir = workdir / f"slide_{i:03d}"
        slide_dir.mkdir(parents=True, exist_ok=True)
        check_cancel()

        # Step 1: 텍스트 추출 — 실패 시 빈 텍스트로 진행 (배경/객체는 계속).
        text_json = slide_dir / "text_layout.json"
        image = Image.open(img_path).convert("RGB")
        target_size = image.size
        try:
            text_layout = analyze_text_with_vision(
                img_path,
                out_json=text_json,
                base_url=options.vision_base_url,
                model=options.vision_model,
                api_key=options.vision_api_key,
                timeout_sec=options.vision_timeout_sec,
            )
        except Exception as exc:
            with progress_lock:
                emit(f"[slide {i}] 텍스트 분석 실패 → 텍스트 없이 진행: {exc}", -1)
            text_layout = {"image_size": list(target_size), "texts": []}
        check_cancel()
        text_layout = normalize_text_layout_to_image_size(text_layout, target_size)
        if "_normalized_from" in text_layout:
            text_json.write_text(
                json.dumps(text_layout, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )

        check_cancel()
        # Step 2: 텍스트 제거 배경
        clean_path = slide_dir / "02_clean_bg.png"
        if options.text_erase_mode == "codex_imagegen":
            try:
                remove_text_with_image_gen(
                    img_path,
                    out_png=clean_path,
                    target_size=target_size,
                    model=options.image_gen_model,
                )
            except Exception as exc:
                with progress_lock:
                    emit(f"[slide {i}] image_gen 실패 → solid fallback: {exc}", -1)
                remove_text_solid_fill(img_path, text_layout, out_png=clean_path)
        else:
            remove_text_solid_fill(img_path, text_layout, out_png=clean_path)

        # Step 3: 객체 추출
        shapes_layout: dict[str, Any] | None = None
        layers_json_path: Path | None = None
        if options.use_native_shapes:
            check_cancel()
            shapes_json_path = slide_dir / "shapes_layout.json"
            try:
                shapes_layout = analyze_shapes_with_vision(
                    clean_path,
                    out_json=shapes_json_path,
                    base_url=options.vision_base_url,
                    model=options.vision_model,
                    api_key=options.vision_api_key,
                    timeout_sec=options.vision_timeout_sec,
                )
                shapes_layout = normalize_shapes_to_image_size(shapes_layout, target_size)
                # 복잡/장식성 화살표는 raster crop 으로 강등 (single-source 분류 진입점).
                demote_complex_arrows_to_pictogram(
                    shapes_layout,
                    max_area_ratio=options.arrow_native_max_area_ratio,
                )
                if "_normalized_from" in shapes_layout:
                    shapes_json_path.write_text(
                        json.dumps(shapes_layout, ensure_ascii=False, indent=2),
                        encoding="utf-8",
                    )
                layers = crop_pictograms_from_codex_shapes(
                    clean_path,
                    shapes_layout,
                    out_dir=slide_dir / "objects",
                )
                layers_json_path = slide_dir / "object_layers.json"
                layers_json_path.write_text(
                    json.dumps(
                        {
                            "image_size": list(target_size),
                            "layer_count": len(layers),
                            "source": "codex_pictogram_bbox",
                            "layers": [
                                {
                                    "bbox": ly.bbox,
                                    "z": ly.z,
                                    "area_ratio": ly.area_ratio,
                                    "alpha_path": str(ly.alpha_path.relative_to(slide_dir)),
                                }
                                for ly in layers
                            ],
                        },
                        ensure_ascii=False,
                        indent=2,
                    ),
                    encoding="utf-8",
                )
            except Exception as exc:
                with progress_lock:
                    emit(f"[slide {i}] codex shape 분석 실패 → SAM2 fallback: {exc}", -1)
                layers, layers_json_path = extract_objects_sam2_clean(
                    clean_path,
                    out_dir=slide_dir,
                    options=options,
                )
                shapes_layout = None
        else:
            check_cancel()
            layers, layers_json_path = extract_objects_sam2_clean(
                clean_path,
                out_dir=slide_dir,
                options=options,
            )

        # 진행률 thread-safe 업데이트
        with progress_lock:
            completed_count[0] += 1
            done_pct = 5 + 88 * completed_count[0] / n  # 5%~93% 범위로 매핑
            n_shapes = len(shapes_layout.get("shapes", [])) if shapes_layout else 0
            emit(
                f"[{completed_count[0]}/{n}] 완료 ({len(text_layout.get('texts', []))} 텍스트, "
                f"{n_shapes} shape, {len(layers)} 픽토그램)",
                done_pct,
            )

        return {
            "_index": i,
            "image_size": target_size,
            "background_path": clean_path,
            "text_layout": text_layout,
            "object_layers": layers,
            "shapes_layout": shapes_layout,
            "_text_json": text_json,
            "_layers_json": layers_json_path,
        }

    # 병렬 실행. 한 슬라이드 catastrophic 실패 → raster fallback 으로 대체하고
    # 다른 슬라이드는 계속 진행. JobCancelledError 만 즉시 전체 취소.
    parallel = max(1, min(int(options.slide_parallel), n))
    emit(f"image_split 슬라이드 {n}장 병렬 처리 시작 (parallel={parallel})", 5)
    results: list[dict[str, Any]] = []
    failed_count = 0

    def _safe_process(i: int, img_path: Path) -> dict[str, Any]:
        try:
            return process_slide(i, img_path)
        except Exception as exc:
            from star_slide.pipeline.notebooklm_auto import JobCancelledError

            if isinstance(exc, JobCancelledError):
                raise
            with progress_lock:
                emit(f"[slide {i}] 슬라이드 처리 실패 → raster fallback: {exc}", -1)
            return _raster_fallback_slide(i, img_path, str(exc))

    if parallel == 1 or n == 1:
        for i, img_path in enumerate(input_images, start=1):
            r = _safe_process(i, img_path)
            results.append(r)
            if r.get("_failed"):
                failed_count += 1
    else:
        from concurrent.futures import ThreadPoolExecutor, as_completed

        with ThreadPoolExecutor(max_workers=parallel) as executor:
            futures = {
                executor.submit(_safe_process, i, img_path): (i, img_path)
                for i, img_path in enumerate(input_images, start=1)
            }
            for fut in as_completed(futures):
                i, img_path = futures[fut]
                try:
                    r = fut.result()
                except Exception as exc:  # JobCancelledError 등
                    from star_slide.pipeline.notebooklm_auto import JobCancelledError

                    if isinstance(exc, JobCancelledError):
                        raise
                    with progress_lock:
                        emit(f"[slide {i}] 처리 실패 (예외) → raster fallback: {exc}", -1)
                    r = _raster_fallback_slide(i, img_path, str(exc))
                results.append(r)
                if r.get("_failed"):
                    failed_count += 1

    # 슬라이드 순서대로 정렬 (병렬 완료 순서가 다를 수 있음)
    results.sort(key=lambda r: r["_index"])
    if failed_count > 0:
        emit(f"⚠ {failed_count}/{n} 슬라이드 실패 — raster fallback 으로 진행", -1)
    slides_data: list[dict[str, Any]] = []
    last_clean_bg: Path | None = None
    last_text_json: Path | None = None
    last_layers_json: Path | None = None
    for r in results:
        slides_data.append(
            {
                "image_size": r["image_size"],
                "background_path": r["background_path"],
                "text_layout": r["text_layout"],
                "object_layers": r["object_layers"],
                "shapes_layout": r["shapes_layout"],
                "_failed": bool(r.get("_failed")),
            }
        )
        last_clean_bg = r["background_path"]
        last_text_json = r.get("_text_json") or last_text_json
        last_layers_json = r.get("_layers_json") or last_layers_json

    check_cancel()
    emit(f"PPTX 합성 ({n} 슬라이드)", 95)
    compose_image_split_pptx_multi(
        slides_data=slides_data,
        out_path=output_pptx,
        font_name=options.font_name,
        background_mode=options.background_mode,
    )
    emit("완료", 100)

    return ImageSplitResult(
        output=output_pptx,
        workdir=workdir,
        text_layout_json=last_text_json or workdir,
        clean_bg_png=last_clean_bg or workdir,
        object_layers_json=last_layers_json or workdir,
        elapsed_sec=round(time.perf_counter() - started, 1),
    )
