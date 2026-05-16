"""Local FastAPI web app for NotebookLM PPTX conversion."""

from __future__ import annotations

import asyncio
import contextlib
import importlib.util
import json
import re
import shutil
import time
import traceback
import uuid
import zipfile
from collections.abc import AsyncIterator
from concurrent.futures import Future, ThreadPoolExecutor
from dataclasses import dataclass, field
from pathlib import Path
from threading import Event, Lock
from typing import Any

from star_slide.api.events import EventBus, format_sse
from star_slide.api.preview_assets import (
    PREVIEW_KINDS,
    generate_previews,
    index_previews,
)
from star_slide.pipeline.notebooklm_auto import (
    JobCancelledError,
    NotebookLmAutoOptions,
    convert_notebooklm_auto,
)

try:
    from fastapi import Body, FastAPI, HTTPException, Request
    from fastapi.responses import (
        FileResponse,
        HTMLResponse,
        JSONResponse,
        StreamingResponse,
    )
except ImportError as exc:  # pragma: no cover - exercised by runtime environment
    raise RuntimeError(
        "FastAPI 의존성이 필요합니다. `uv run --extra api star-slide web run`으로 실행하세요."
    ) from exc


WEB_ROOT = Path("output/web_jobs")
ALLOWED_SUFFIXES = {".pptx", ".pdf"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg"}
UPLOAD_SUFFIXES = ALLOWED_SUFFIXES | IMAGE_SUFFIXES
MAX_UPLOAD_BYTES = 300 * 1024 * 1024
SSE_HEARTBEAT_SECONDS = 15.0
TERMINAL_STATUSES = frozenset({"done", "failed", "cancelled"})


@dataclass
class JobState:
    id: str
    filename: str
    status: str = "queued"
    phase: str = "대기 중"
    progress: float = 0.0
    created_at: float = field(default_factory=time.time)
    updated_at: float = field(default_factory=time.time)
    error: str | None = None
    workdir: str | None = None
    output: str | None = None
    report: str | None = None
    montage: str | None = None
    input_path: str | None = None
    options: dict[str, Any] = field(default_factory=dict)
    raw_options: dict[str, Any] = field(default_factory=dict)
    cancel_event: Event = field(default_factory=Event)
    bus: EventBus = field(default_factory=EventBus)


jobs: dict[str, JobState] = {}
futures: dict[str, Future[None]] = {}
jobs_lock = Lock()
executor = ThreadPoolExecutor(max_workers=2)


def create_app() -> FastAPI:
    app = FastAPI(title="Star-Slide", version="0.1.0")

    @app.get("/", response_class=HTMLResponse)
    def index() -> str:
        return INDEX_HTML

    @app.get("/api/presets")
    def presets() -> dict[str, Any]:
        return {
            "providers": [
                {
                    "id": "local",
                    "label": "Local Proxy",
                    "baseUrl": "http://localhost:8300/v1",
                    "model": "gpt-5.5",
                    "hint": "star-cliproxy 같은 로컬 OpenAI 호환 프록시",
                },
                {
                    "id": "openai",
                    "label": "OpenAI",
                    "baseUrl": "https://api.openai.com/v1",
                    "model": "",
                    "hint": "OpenAI API key와 vision 지원 모델명을 입력",
                },
                {
                    "id": "gemini",
                    "label": "Gemini",
                    "baseUrl": "https://generativelanguage.googleapis.com/v1beta/openai",
                    "model": "gemini-3.1-pro-preview",
                    "hint": "Gemini OpenAI-compatible endpoint 또는 로컬 프록시 모델명 사용",
                },
                {
                    "id": "ollama",
                    "label": "Ollama",
                    "baseUrl": "http://localhost:11434/v1",
                    "model": "llama3.2",
                    "hint": "로컬 Ollama. Base URL에 /v1 prefix 필수. API key는 비워두세요.",
                },
                {
                    "id": "custom",
                    "label": "Custom",
                    "baseUrl": "",
                    "model": "",
                    "hint": "OpenAI 호환 /v1/chat/completions endpoint",
                },
            ],
            "defaults": {
                "timeout": 600,
                "retries": 2,
                "llmParallel": 5,
                "fontScale": 0.93,
                "keepIntermediates": False,
                "sam3": False,
                "hybridAllowedDelta": 0.0,
                "editableEmbeddedText": True,
                "layoutFailureMode": "image_fallback",
                "watermarkMode": "off",
                "reconstructionMode": "auto",
                "textEraseMode": "codex_imagegen",
                "backgroundMode": "white",
                "useNativeShapes": True,
            },
        }

    @app.post("/api/test-llm")
    async def test_llm(request: Request) -> JSONResponse:
        try:
            payload = await request.json()
        except Exception as exc:
            raise HTTPException(status_code=400, detail="요청 본문이 JSON 형식이어야 합니다.") from exc
        result = await asyncio.to_thread(probe_llm_endpoint, payload)
        return JSONResponse(result)

    @app.post("/api/list-models")
    async def list_models_endpoint(request: Request) -> JSONResponse:
        try:
            payload = await request.json()
        except Exception as exc:
            raise HTTPException(status_code=400, detail="요청 본문이 JSON 형식이어야 합니다.") from exc
        base_url = str(payload.get("baseUrl") or "").strip().rstrip("/")
        api_key = str(payload.get("apiKey") or "").strip()
        timeout = float(payload.get("timeout") or 10)
        timeout = max(2.0, min(30.0, timeout))
        if not base_url:
            return JSONResponse({"models": [], "error": "Base URL이 없습니다."}, status_code=400)
        models = await asyncio.to_thread(_list_models, base_url, api_key, timeout)
        return JSONResponse({"models": models})

    @app.get("/api/system-check")
    def system_check() -> dict[str, Any]:
        libreoffice = first_executable(("soffice", "libreoffice"))
        poppler = first_executable(("pdftoppm", "pdfinfo"))
        sam3_ready = python_module_available("torch") and python_module_available("transformers")
        return {
            "items": [
                {
                    "id": "libreoffice",
                    "label": "LibreOffice",
                    "ok": libreoffice is not None,
                    "path": libreoffice,
                    "required": True,
                    "message": "PPTX 렌더 QA와 자동 선택에 필요합니다.",
                },
                {
                    "id": "poppler",
                    "label": "Poppler",
                    "ok": poppler is not None,
                    "path": poppler,
                    "required": False,
                    "message": "PDF 입력 또는 LibreOffice PDF 렌더 경유에 필요할 수 있습니다.",
                },
                {
                    "id": "sam3",
                    "label": "SAM3",
                    "ok": sam3_ready,
                    "path": "torch + transformers" if sam3_ready else None,
                    "required": False,
                    "message": "고품질 bbox 보정 옵션입니다. 사용하려면 gpu-segmentation extra와 SAM3 모델 접근 권한이 필요합니다.",
                },
            ]
        }

    @app.post("/api/jobs")
    async def submit_job(request: Request) -> JSONResponse:
        filename = request.query_params.get("filename", "upload.pptx")
        safe_name = Path(filename).name
        if not safe_name or safe_name in {".", ".."}:
            raise HTTPException(status_code=400, detail="파일명이 올바르지 않습니다.")
        upload_suffix = Path(safe_name).suffix.lower()
        if upload_suffix not in UPLOAD_SUFFIXES:
            raise HTTPException(
                status_code=400,
                detail="PPTX/PDF 또는 이미지(PNG/JPG) 파일만 업로드할 수 있습니다.",
            )

        raw_options = request.headers.get("x-star-slide-options", "{}")
        try:
            options_payload = json.loads(raw_options)
        except json.JSONDecodeError as exc:
            raise HTTPException(status_code=400, detail="옵션 JSON을 파싱할 수 없습니다.") from exc

        job_id = uuid.uuid4().hex
        job_dir = WEB_ROOT / job_id
        job_dir.mkdir(parents=True, exist_ok=True)
        input_path = job_dir / safe_name

        # 스트리밍 업로드: 청크 단위로 디스크에 직접 기록하고 누적 크기를 검사한다.
        # 한도 초과 시 부분 파일을 즉시 정리해 디스크 누수를 막는다.
        total = 0
        with input_path.open("wb") as fh:
            async for chunk in request.stream():
                if not chunk:
                    continue
                total += len(chunk)
                if total > MAX_UPLOAD_BYTES:
                    fh.close()
                    with contextlib.suppress(OSError):
                        input_path.unlink()
                    raise HTTPException(
                        status_code=413,
                        detail=f"업로드 파일이 너무 큽니다 (한도 {MAX_UPLOAD_BYTES // (1024 * 1024)}MB).",
                    )
                fh.write(chunk)
        if total == 0:
            with contextlib.suppress(OSError):
                input_path.unlink()
            raise HTTPException(status_code=400, detail="업로드 파일이 비어 있습니다.")

        # 단일 이미지 업로드는 1-슬라이드 PPTX 로 래핑해서 PPTX 파이프라인에 그대로 흘려보낸다.
        # 원본 이미지는 그대로 두고, input_path/filename 만 wrapper PPTX 로 교체한다.
        if upload_suffix in IMAGE_SUFFIXES:
            try:
                wrapper_path = wrap_image_as_pptx(input_path)
            except Exception as exc:
                with contextlib.suppress(OSError):
                    input_path.unlink()
                raise HTTPException(
                    status_code=400,
                    detail=f"이미지 파일을 PPTX 로 변환하지 못했습니다: {sanitize_error(str(exc))}",
                ) from exc
            input_path = wrapper_path
            safe_name = wrapper_path.name

        state = JobState(
            id=job_id,
            filename=safe_name,
            workdir=str(job_dir / "work"),
            output=str(job_dir / "result.pptx"),
            report=str(job_dir / "artifacts" / "report.json"),
            montage=str(job_dir / "artifacts" / "montage.png"),
            input_path=str(input_path),
            options=public_options(options_payload),
            raw_options=dict(options_payload),
        )
        with jobs_lock:
            jobs[job_id] = state
            futures[job_id] = executor.submit(run_job, job_id, input_path, job_dir, options_payload)
        return JSONResponse(job_snapshot(state))

    @app.get("/api/jobs")
    def list_jobs() -> list[dict[str, Any]]:
        sync_saved_jobs()
        with jobs_lock:
            current_jobs = sorted(jobs.values(), key=lambda item: item.created_at, reverse=True)
        return [job_snapshot(job) for job in current_jobs]

    @app.get("/api/jobs/{job_id}")
    def get_job(job_id: str) -> dict[str, Any]:
        return job_snapshot(require_job(job_id))

    @app.get("/api/jobs/{job_id}/download")
    def download(job_id: str) -> FileResponse:
        job = require_job(job_id)
        if job.status != "done" or not job.output or not Path(job.output).exists():
            raise HTTPException(status_code=404, detail="완료된 PPTX가 없습니다.")
        return FileResponse(
            job.output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"{Path(job.filename).stem}_editable.pptx",
        )

    @app.get("/api/jobs/{job_id}/montage")
    def montage(job_id: str) -> FileResponse:
        job = require_job(job_id)
        if not job.montage or not Path(job.montage).exists():
            raise HTTPException(status_code=404, detail="미리보기 이미지가 없습니다.")
        return FileResponse(job.montage, media_type="image/png")

    @app.get("/api/jobs/{job_id}/thumbnail")
    def thumbnail(job_id: str) -> FileResponse:
        """Return the first-slide thumbnail for the job list UI.

        Tries cheap pre-rendered sources first (per-slide preview JPEGs,
        cached page PNGs, then workdir slide PNGs). Returns 404 when no
        rendered slide image exists yet so the client can hide the slot.
        """
        job = require_job(job_id)
        path = first_slide_thumbnail_path(job)
        if path is None:
            raise HTTPException(status_code=404, detail="썸네일이 아직 없습니다.")
        media_type = "image/jpeg" if path.suffix.lower() in {".jpg", ".jpeg"} else "image/png"
        return FileResponse(
            path,
            media_type=media_type,
            headers={"Cache-Control": "max-age=300"},
        )

    @app.get("/api/jobs/{job_id}/report")
    def report(job_id: str) -> FileResponse:
        job = require_job(job_id)
        if not job.report or not Path(job.report).exists():
            raise HTTPException(status_code=404, detail="리포트가 없습니다.")
        return FileResponse(job.report, media_type="application/json", filename="notebooklm_auto_report.json")

    @app.get("/api/jobs/{job_id}/report-summary")
    def report_summary(job_id: str) -> dict[str, Any]:
        job = require_job(job_id)
        if not job.report or not Path(job.report).exists():
            raise HTTPException(status_code=404, detail="리포트가 없습니다.")
        return build_report_summary(job, json.loads(Path(job.report).read_text(encoding="utf-8")))

    @app.get("/api/jobs/{job_id}/artifact/{name}")
    def artifact(job_id: str, name: str) -> FileResponse:
        job = require_job(job_id)
        path = resolve_artifact_path(job, name)
        if path is None or not path.exists():
            raise HTTPException(status_code=404, detail="산출물을 찾을 수 없습니다.")
        media_type = "application/zip" if path.suffix == ".zip" else "application/octet-stream"
        if path.suffix == ".pptx":
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        return FileResponse(path, media_type=media_type, filename=path.name)

    @app.get("/api/jobs/{job_id}/layout-summary")
    def layout_summary(job_id: str) -> dict[str, Any]:
        job = require_job(job_id)
        path = resolve_artifact_path(job, "layout-json")
        if path is None or not path.exists():
            raise HTTPException(status_code=404, detail="Layout JSON 산출물이 없습니다.")
        return build_layout_summary(job, path)

    @app.get("/api/jobs/{job_id}/events")
    async def events(job_id: str, request: Request) -> StreamingResponse:
        job = require_job(job_id)
        bus = job.bus
        loop = asyncio.get_running_loop()
        queue = bus.subscribe(loop)

        async def stream() -> AsyncIterator[bytes]:
            try:
                yield format_sse(job_snapshot(job), event="snapshot").encode("utf-8")
                while True:
                    if await request.is_disconnected():
                        return
                    try:
                        item = await asyncio.wait_for(queue.get(), timeout=SSE_HEARTBEAT_SECONDS)
                    except TimeoutError:
                        yield b": heartbeat\n\n"
                        continue
                    if item is None:
                        return
                    yield format_sse(item, event="snapshot").encode("utf-8")
                    if item.get("status") in TERMINAL_STATUSES:
                        return
            finally:
                bus.unsubscribe(loop, queue)

        return StreamingResponse(
            stream(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "X-Accel-Buffering": "no",
                "Connection": "keep-alive",
            },
        )

    @app.post("/api/jobs/{job_id}/cancel")
    def cancel(request: Request, job_id: str) -> JSONResponse:
        require_json_content_type(request)
        job = require_job(job_id)
        if job.status in TERMINAL_STATUSES:
            return JSONResponse(job_snapshot(job))
        job.cancel_event.set()
        future = futures.get(job_id)
        if future is not None and not future.running():
            future.cancel()
        if job.status == "queued":
            update_job(job_id, status="cancelled", phase="취소됨", progress=100)
        else:
            update_job(job_id, phase="취소 중...", status="cancelling")
        return JSONResponse(job_snapshot(require_job(job_id)))

    @app.post("/api/jobs/delete")
    def delete_jobs_bulk(
        request: Request,
        payload: dict[str, Any] = Body(default_factory=dict),  # noqa: B008
    ) -> JSONResponse:
        require_json_content_type(request)
        ids = payload.get("ids") or []
        mode = (str(payload.get("mode") or "trash")).strip().lower()
        if mode not in {"permanent", "trash"}:
            raise HTTPException(status_code=400, detail="mode 는 permanent 또는 trash 여야 합니다.")
        if not isinstance(ids, list) or not ids:
            raise HTTPException(status_code=400, detail="ids 가 비어 있습니다.")
        deleted: list[str] = []
        skipped: list[dict[str, str]] = []
        for raw_id in ids:
            job_id = str(raw_id)
            try:
                _delete_one_job(job_id, mode=mode)
                deleted.append(job_id)
            except HTTPException as exc:
                skipped.append({"id": job_id, "reason": str(exc.detail)})
            except Exception as exc:  # pragma: no cover - 방어적
                skipped.append({"id": job_id, "reason": sanitize_error(str(exc))})
        return JSONResponse({"deleted": deleted, "skipped": skipped, "mode": mode})

    @app.delete("/api/jobs/{job_id}")
    def delete_job(job_id: str, mode: str = "trash") -> JSONResponse:
        normalized = (mode or "trash").strip().lower()
        if normalized not in {"permanent", "trash"}:
            raise HTTPException(status_code=400, detail="mode 는 permanent 또는 trash 여야 합니다.")
        _delete_one_job(job_id, mode=normalized)
        return JSONResponse({"id": job_id, "mode": normalized})

    @app.post("/api/jobs/{job_id}/rerun")
    def rerun(request: Request, job_id: str) -> JSONResponse:
        require_json_content_type(request)
        job = require_job(job_id)
        if job.status not in TERMINAL_STATUSES:
            raise HTTPException(status_code=409, detail="진행 중인 작업은 다시 실행할 수 없습니다.")
        if not job.input_path or not Path(job.input_path).exists():
            raise HTTPException(status_code=410, detail="원본 입력 파일을 찾을 수 없어 재실행이 불가능합니다.")
        new_job_id = uuid.uuid4().hex
        job_dir = WEB_ROOT / new_job_id
        job_dir.mkdir(parents=True, exist_ok=True)
        new_input = job_dir / Path(job.input_path).name
        shutil.copy2(job.input_path, new_input)
        new_state = JobState(
            id=new_job_id,
            filename=job.filename,
            workdir=str(job_dir / "work"),
            output=str(job_dir / "result.pptx"),
            report=str(job_dir / "artifacts" / "report.json"),
            montage=str(job_dir / "artifacts" / "montage.png"),
            input_path=str(new_input),
            options=public_options(job.raw_options),
            raw_options=dict(job.raw_options),
        )
        with jobs_lock:
            jobs[new_job_id] = new_state
            futures[new_job_id] = executor.submit(run_job, new_job_id, new_input, job_dir, dict(job.raw_options))
        return JSONResponse(job_snapshot(new_state))

    @app.get("/api/jobs/{job_id}/previews")
    def previews_index(job_id: str) -> dict[str, Any]:
        job = require_job(job_id)
        out_dir = previews_dir_for(job)
        entries = index_previews(out_dir)
        return {
            "job": {"id": job.id, "filename": job.filename},
            "kinds": list(PREVIEW_KINDS),
            "slides": [{"slide_no": entry.slide_no, "kinds": list(entry.kinds)} for entry in entries],
        }

    @app.get("/api/jobs/{job_id}/previews/{slide_no}/{kind}")
    def preview_image(job_id: str, slide_no: int, kind: str) -> FileResponse:
        if kind not in PREVIEW_KINDS:
            raise HTTPException(status_code=404, detail="지원하지 않는 preview kind 입니다.")
        job = require_job(job_id)
        out_dir = previews_dir_for(job)
        path = out_dir / f"{slide_no:03d}_{kind}.jpg"
        if not path.exists():
            raise HTTPException(status_code=404, detail="해당 슬라이드 미리보기가 없습니다.")
        return FileResponse(
            path,
            media_type="image/jpeg",
            headers={"Cache-Control": "max-age=3600"},
        )

    @app.get("/api/jobs/{job_id}/pptx-file")
    def pptx_file(job_id: str, which: str = "result") -> FileResponse:
        """Stream the raw PPTX so the browser can render it client-side.

        `which=result` returns the converted PPTX, `which=original` returns
        the user's input. Used by the slide viewer that loads pptx-preview
        from a CDN, eliminating the LibreOffice round-trip.
        """
        job = require_job(job_id)
        try:
            source = _resolve_source_for_pages(job, which)
        except FileNotFoundError as exc:
            raise HTTPException(status_code=404, detail=str(exc)) from exc
        if source.suffix.lower() == ".pdf":
            media_type = "application/pdf"
        else:
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        return FileResponse(
            source,
            media_type=media_type,
            headers={"Cache-Control": "max-age=300"},
        )

    @app.get("/api/jobs/{job_id}/pptx-pages")
    async def pptx_pages_index(job_id: str, which: str = "result") -> JSONResponse:
        job = require_job(job_id)
        try:
            pages = await asyncio.to_thread(render_pptx_pages, job, which)
        except FileNotFoundError as exc:
            raise HTTPException(status_code=404, detail=str(exc)) from exc
        except RuntimeError as exc:
            raise HTTPException(status_code=503, detail=str(exc)) from exc
        return JSONResponse(
            {
                "job": {"id": job.id, "filename": job.filename},
                "which": which,
                "pages": len(pages),
            }
        )

    @app.get("/api/jobs/{job_id}/pptx-pages/{page_no}")
    def pptx_page_image(job_id: str, page_no: int, which: str = "result") -> FileResponse:
        job = require_job(job_id)
        cache_dir = pptx_pages_cache_dir(job, which)
        path = cache_dir / f"page_{page_no:03d}.png"
        if not path.exists():
            raise HTTPException(status_code=404, detail="해당 페이지가 없습니다. /pptx-pages를 먼저 호출하세요.")
        return FileResponse(path, media_type="image/png", headers={"Cache-Control": "max-age=3600"})

    return app


def previews_dir_for(job: JobState) -> Path:
    if job.report:
        return Path(job.report).parent / "previews"
    return WEB_ROOT / job.id / "artifacts" / "previews"


def first_slide_thumbnail_path(job: JobState) -> Path | None:
    """Locate an existing first-slide image for the job-card thumbnail.

    Walks cheap pre-rendered sources from highest-fidelity to fallback.
    Returns None when no rendered slide image is available yet.
    """
    previews = previews_dir_for(job)
    for kind in ("selected", "hybrid", "vector", "original"):
        candidate = previews / f"001_{kind}.jpg"
        if candidate.exists():
            return candidate
    for which in ("result", "original"):
        candidate = pptx_pages_cache_dir(job, which) / "page_001.png"
        if candidate.exists():
            return candidate
    if job.workdir:
        workdir = Path(job.workdir)
        for sub in ("qa_selected", "images"):
            directory = workdir / sub
            if not directory.is_dir():
                continue
            pngs = sorted(directory.glob("*.png"))
            if pngs:
                return pngs[0]
    return None


def _resolve_source_for_pages(job: JobState, which: str) -> Path:
    if which == "result":
        candidate = Path(job.output) if job.output else None
    elif which == "original":
        candidate = Path(job.input_path) if job.input_path else None
    elif which in ("vector", "hybrid"):
        candidate = resolve_artifact_path(job, which)
    else:
        raise FileNotFoundError(f"지원하지 않는 미리보기 종류 '{which}' (result|original|vector|hybrid)")
    if candidate is None or not candidate.exists():
        raise FileNotFoundError(f"미리보기 원본 파일을 찾을 수 없습니다 ({which}).")
    return candidate


def pptx_pages_cache_dir(job: JobState, which: str) -> Path:
    return WEB_ROOT / job.id / "preview_cache" / which


def prefetch_pptx_pages_from_workdir(
    job_id: str,
    job_dir: Path,
    workdir: Path,
    result_pptx: Path,
    input_pptx: Path,
) -> None:
    """Reuse workdir's per-page PNGs as the slide viewer cache.

    convert_notebooklm_auto already renders every slide to PNG during QA
    (workdir/qa_selected for the converted result, workdir/images for the
    original input). Copying those into preview_cache/{result,original}/
    means the user's first click on "미리보기/원본 보기/비교 보기" hits the
    cache instead of triggering a fresh LibreOffice run.
    """
    sources = {
        "result": (workdir / "qa_selected", result_pptx),
        "original": (workdir / "images", input_pptx),
    }
    for which, (src_dir, source_file) in sources.items():
        if not src_dir.exists() or not source_file.exists():
            continue
        pngs = sorted(src_dir.glob("*.png"))
        if not pngs:
            continue
        cache_dir = WEB_ROOT / job_id / "preview_cache" / which
        if cache_dir.exists():
            for stale in cache_dir.glob("page_*.png"):
                with contextlib.suppress(OSError):
                    stale.unlink()
        cache_dir.mkdir(parents=True, exist_ok=True)
        for i, src in enumerate(pngs, start=1):
            dst = cache_dir / f"page_{i:03d}.png"
            with contextlib.suppress(OSError):
                shutil.copy2(src, dst)
        stat = source_file.stat()
        fingerprint = f"{int(stat.st_mtime)}-{stat.st_size}"
        (cache_dir / "fingerprint").write_text(fingerprint, encoding="utf-8")


def render_pptx_pages(job: JobState, which: str) -> list[Path]:
    """Render the job's PPTX/PDF (result or original) into per-page PNGs (cached).

    Cache keyed on source file mtime+size. Returns sorted PNG paths.
    Raises FileNotFoundError when the source file is missing,
    RuntimeError when LibreOffice/pdftoppm dependencies are unavailable.
    """
    source = _resolve_source_for_pages(job, which)
    cache_dir = pptx_pages_cache_dir(job, which)
    stat = source.stat()
    fingerprint = f"{int(stat.st_mtime)}-{stat.st_size}"
    fingerprint_file = cache_dir / "fingerprint"

    if (
        cache_dir.exists()
        and fingerprint_file.exists()
        and fingerprint_file.read_text(encoding="utf-8").strip() == fingerprint
    ):
        existing = sorted(cache_dir.glob("page_*.png"))
        if existing:
            return existing

    if cache_dir.exists():
        for stale in cache_dir.glob("page_*.png"):
            with contextlib.suppress(OSError):
                stale.unlink()
    cache_dir.mkdir(parents=True, exist_ok=True)

    suffix = source.suffix.lower()
    try:
        if suffix == ".pptx":
            from star_slide.rasterize.libreoffice import render_pptx_to_pngs

            rendered = render_pptx_to_pngs(source, cache_dir)
        elif suffix == ".pdf":
            from pdf2image import convert_from_path

            images = convert_from_path(str(source), dpi=160)
            rendered = []
            for i, img in enumerate(images, start=1):
                out_path = cache_dir / f"page_{i:03d}.png"
                img.save(out_path, "PNG")
                rendered.append(out_path)
        else:
            raise FileNotFoundError(f"지원하지 않는 입력 형식 '{suffix}'")
    except Exception as exc:
        raise RuntimeError(f"미리보기 렌더링 실패: {sanitize_error(str(exc))}") from exc

    # 파일명을 page_NNN.png 규약으로 정렬해서 다시 저장 (libreoffice 출력은 slide_NNN)
    final_pages: list[Path] = []
    for i, src in enumerate(sorted(rendered), start=1):
        target = cache_dir / f"page_{i:03d}.png"
        if src != target:
            with contextlib.suppress(OSError):
                if target.exists():
                    target.unlink()
                src.rename(target)
        final_pages.append(target)

    # libreoffice 헬퍼가 만든 임시 디렉토리 정리
    pdf_dir = cache_dir / "_pdf"
    if pdf_dir.exists():
        with contextlib.suppress(OSError):
            shutil.rmtree(pdf_dir)

    fingerprint_file.write_text(fingerprint, encoding="utf-8")
    return sorted(cache_dir.glob("page_*.png"))


def sync_saved_jobs() -> None:
    if not WEB_ROOT.exists():
        return
    with jobs_lock:
        known = set(jobs)
    for job_dir in sorted(WEB_ROOT.iterdir()):
        if not job_dir.is_dir() or job_dir.name in known:
            continue
        # _trash, _backup 등 underscore prefix 디렉토리는 잡 디렉토리가 아님
        if job_dir.name.startswith("_"):
            continue
        result = job_dir / "result.pptx"
        report = job_dir / "artifacts" / "report.json"
        montage = job_dir / "artifacts" / "montage.png"
        if not report.exists():
            report = job_dir / "work" / "notebooklm_auto_report.json"
        if not montage.exists():
            montage = job_dir / "work" / "qa_selected" / "montage.png"
        source_files = [
            path
            for suffix in sorted(ALLOWED_SUFFIXES)
            for path in job_dir.glob(f"*{suffix}")
            if path.name != "result.pptx" and not path.name.startswith("result_")
        ]
        source = source_files[0] if source_files else None
        if not result.exists() and not report.exists():
            continue
        created_at = source.stat().st_mtime if source and source.exists() else job_dir.stat().st_mtime
        state = JobState(
            id=job_dir.name,
            filename=source.name if source else job_dir.name,
            status="done" if result.exists() else "failed",
            phase="완료" if result.exists() else "결과 파일 없음",
            progress=100,
            created_at=created_at,
            updated_at=max(result.stat().st_mtime if result.exists() else created_at, report.stat().st_mtime if report.exists() else created_at),
            workdir=str(job_dir / "work"),
            output=str(result),
            report=str(report) if report.exists() else None,
            montage=str(montage) if montage.exists() else None,
            input_path=str(source) if source else None,
        )
        with jobs_lock:
            jobs.setdefault(state.id, state)


def resolve_artifact_path(job: JobState, name: str) -> Path | None:
    allowed = {
        "vector": "candidate_vector.pptx",
        "hybrid": "candidate_hybrid.pptx",
        "layout-json": "layout_json.zip",
        "manifest": "artifact_manifest.json",
    }
    filename = allowed.get(name)
    if filename is None:
        return None
    candidates: list[Path] = []
    if job.report:
        candidates.append(Path(job.report).parent / filename)
    candidates.append(WEB_ROOT / job.id / "artifacts" / filename)
    for path in candidates:
        if path.exists():
            return path
    return candidates[0] if candidates else None


def build_layout_summary(job: JobState, layout_zip: Path) -> dict[str, Any]:
    slides: list[dict[str, Any]] = []
    with zipfile.ZipFile(layout_zip) as archive:
        names = sorted(name for name in archive.namelist() if name.startswith("selected/") and name.endswith(".layout.json"))
        for name in names:
            payload = json.loads(archive.read(name).decode("utf-8"))
            objects = payload.get("objects") if isinstance(payload.get("objects"), list) else []
            counts: dict[str, int] = {}
            for obj in objects:
                kind = str(obj.get("type", "unknown"))
                counts[kind] = counts.get(kind, 0) + 1
            title = next(
                (
                    object_text(obj)
                    for obj in objects
                    if obj.get("type") == "text" and "title" in str(obj.get("name", "")).lower()
                ),
                "",
            )
            raster_meta = payload.get("metadata", {}).get("raster_group_replacement", {})
            groups = raster_meta.get("groups") if isinstance(raster_meta, dict) else []
            slides.append(
                {
                    "slide_no": len(slides) + 1,
                    "id": payload.get("id"),
                    "title": title,
                    "object_count": len(objects),
                    "type_counts": counts,
                    "raster_group_count": len(groups) if isinstance(groups, list) else 0,
                    "punched_text_regions": raster_meta.get("punched_text_regions") if isinstance(raster_meta, dict) else None,
                }
            )
    return {
        "job": {"id": job.id, "filename": job.filename},
        "layout_zip": str(layout_zip),
        "slide_count": len(slides),
        "slides": slides,
    }


def require_job(job_id: str) -> JobState:
    with jobs_lock:
        job = jobs.get(job_id)
    if job is None:
        raise HTTPException(status_code=404, detail="작업을 찾을 수 없습니다.")
    return job


TRASH_DIR_NAME = "_trash"


def _job_dir_for(job: JobState) -> Path:
    """잡의 실제 디렉토리. report/output/workdir/input_path 중 가장 신뢰할 수 있는 경로에서 추정."""
    candidates: list[Path] = []
    for src in (job.output, job.report, job.workdir, job.input_path):
        if src:
            candidates.append(Path(src))
    for path in candidates:
        # 잡 디렉토리는 WEB_ROOT/<id>
        for ancestor in path.parents:
            if ancestor.parent == WEB_ROOT and ancestor.name == job.id:
                return ancestor
    return WEB_ROOT / job.id


def _delete_one_job(job_id: str, *, mode: str) -> None:
    job = require_job(job_id)
    if job.status in {"running", "queued", "cancelling"}:
        raise HTTPException(
            status_code=409,
            detail="진행 중인 작업은 삭제할 수 없습니다. 먼저 취소하세요.",
        )
    job_dir = _job_dir_for(job)
    if mode == "trash":
        trash_root = WEB_ROOT / TRASH_DIR_NAME
        trash_root.mkdir(parents=True, exist_ok=True)
        if job_dir.exists():
            target = trash_root / job.id
            if target.exists():
                # 동일 id가 trash에 이미 있으면 timestamp 접미사로 회피
                target = trash_root / f"{job.id}.{int(time.time())}"
            shutil.move(str(job_dir), str(target))
    else:  # permanent
        if job_dir.exists():
            shutil.rmtree(job_dir, ignore_errors=True)
    # 메모리 dict + future 정리
    with jobs_lock:
        jobs.pop(job_id, None)
        futures.pop(job_id, None)


SERIALIZABLE_JOB_FIELDS = (
    "id",
    "filename",
    "status",
    "phase",
    "progress",
    "created_at",
    "updated_at",
    "error",
    "workdir",
    "output",
    "report",
    "montage",
    "input_path",
    "options",
)


def job_snapshot(job: JobState) -> dict[str, Any]:
    data: dict[str, Any] = {key: getattr(job, key) for key in SERIALIZABLE_JOB_FIELDS}
    if job.status == "done":
        data["artifacts"] = {
            "layout_json": bool((path := resolve_artifact_path(job, "layout-json")) and path.exists()),
            "vector": bool((path := resolve_artifact_path(job, "vector")) and path.exists()),
            "hybrid": bool((path := resolve_artifact_path(job, "hybrid")) and path.exists()),
        }
    if job.status != "running" or not job.workdir:
        return data

    # image_split 모드는 vector/hybrid layout 디렉토리를 만들지 않으므로
    # 자동 phase override (vector_layouts/raster_groups 카운트 기반) 를 건너뜀.
    # 실제 progress/phase 는 _convert_image_split 내부 emit() 가 그대로 남김.
    if (job.options or {}).get("reconstructionMode") == "image_split":
        return data

    workdir = Path(job.workdir)
    images = sorted((workdir / "images").glob("*.png"))
    total = len(images)
    if total <= 0:
        return data

    vector_layouts = len(list((workdir / "layouts_vector").glob("*.layout.json")))
    raster_groups = len(list((workdir / "raster_groups").glob("*_raster_groups.json")))
    sam_reports = len(list((workdir / "raster_groups_sam3").glob("*_sam3_report.json")))
    hybrid_layouts = len(list((workdir / "layouts_hybrid").glob("*.layout.json")))
    selected_layouts = len(list((workdir / "layouts_selected").glob("*.layout.json")))

    progress = float(data["progress"])
    if progress < 38 and vector_layouts < total:
        data["phase"] = f"Vision LLM layout JSON 생성 중 ({vector_layouts}/{total})"
        data["progress"] = max(progress, 10 + 28 * vector_layouts / total)
    elif progress < 60 and raster_groups < total:
        data["phase"] = f"큰 이미지 그룹 탐지 중 ({raster_groups}/{total})"
        data["progress"] = max(progress, 48 + 12 * raster_groups / total)
    elif progress < 70 and sam_reports < total and (workdir / "raster_groups_sam3").exists():
        data["phase"] = f"SAM3 이미지 그룹 보정 중 ({sam_reports}/{total})"
        data["progress"] = max(progress, 60 + 10 * sam_reports / total)
    elif progress < 78 and hybrid_layouts < total:
        data["phase"] = f"hybrid layout 생성 중 ({hybrid_layouts}/{total})"
        data["progress"] = max(progress, 70 + 8 * hybrid_layouts / total)
    elif progress < 93 and selected_layouts < total:
        data["phase"] = f"최종 layout 자동 선택 중 ({selected_layouts}/{total})"
        data["progress"] = max(progress, 88 + 5 * selected_layouts / total)
    return data


def public_options(options_payload: dict[str, Any]) -> dict[str, Any]:
    public = dict(options_payload)
    if public.get("apiKey"):
        public["apiKey"] = "********"
    return public


def first_executable(names: tuple[str, ...]) -> str | None:
    for name in names:
        path = shutil.which(name)
        if path:
            return path
    return None


def python_module_available(name: str) -> bool:
    return importlib.util.find_spec(name) is not None


def wrap_image_as_pptx(image_path: Path) -> Path:
    """단일 이미지(PNG/JPG)를 1-슬라이드 PPTX 로 래핑한다.

    파이프라인은 PPTX/PDF 만 직접 처리하므로, 이미지 입력은 동일 디렉토리에
    `<stem>.pptx` (이미 존재하면 `<stem>__input.pptx`) 로 감싸서 흘려보낸다.
    이미지 비율을 유지하기 위해 슬라이드 크기는 이미지 픽셀 크기를 기반으로 잡고,
    PICTURE shape 는 슬라이드 전체를 채운다.
    """
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Emu

    target = image_path.with_suffix(".pptx")
    if target.exists():
        target = image_path.with_name(f"{image_path.stem}__input.pptx")

    with Image.open(image_path) as im:
        width_px, height_px = im.size
    if width_px <= 0 or height_px <= 0:
        raise ValueError(f"잘못된 이미지 크기입니다: {image_path}")

    # PowerPoint 슬라이드 크기 한도(약 51,206,400 EMU ≈ 56인치)를 넘기지 않도록 96 DPI 가정.
    emu_per_inch = 914400
    width_emu = int(width_px / 96 * emu_per_inch)
    height_emu = int(height_px / 96 * emu_per_inch)
    max_emu = 50_000_000
    if width_emu > max_emu or height_emu > max_emu:
        scale = max_emu / max(width_emu, height_emu)
        width_emu = int(width_emu * scale)
        height_emu = int(height_emu * scale)

    prs = Presentation()
    prs.slide_width = Emu(width_emu)
    prs.slide_height = Emu(height_emu)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(
        str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height,
    )
    prs.save(str(target))
    return target


def update_job(job_id: str, *, force: bool = False, **changes: Any) -> None:
    with jobs_lock:
        job = jobs.get(job_id)
        if job is None:
            return
        # terminal-state guard: cancel/done/failed 후 도착하는 late progress가
        # status를 비-terminal로 되돌리지 못하도록 차단한다. force=True (run_job 자체
        # 종료 분기)일 때는 통과.
        if not force and job.status in TERMINAL_STATUSES:
            new_status = changes.get("status", job.status)
            if new_status not in TERMINAL_STATUSES:
                return
        for key, value in changes.items():
            setattr(job, key, value)
        job.updated_at = time.time()
        snapshot = job_snapshot(job)
        bus = job.bus
        terminal = job.status in TERMINAL_STATUSES
    bus.publish(snapshot)
    if terminal:
        bus.close()


def run_job(job_id: str, input_path: Path, job_dir: Path, payload: dict[str, Any]) -> None:
    def progress(message: str, percent: float, *_: Any, **__: Any) -> None:
        update_job(job_id, status="running", phase=message, progress=percent)

    with jobs_lock:
        job_for_cancel = jobs.get(job_id)
    cancel_event = job_for_cancel.cancel_event if job_for_cancel else None

    def cancel_check() -> bool:
        return cancel_event is not None and cancel_event.is_set()

    try:
        update_job(job_id, status="running", phase="작업 시작", progress=1, error=None)
        output_path = job_dir / "result.pptx"
        workdir = job_dir / "work"
        options = NotebookLmAutoOptions(
            base_url=str(payload.get("baseUrl") or "http://localhost:8300/v1"),
            model=str(payload.get("model") or "gpt-5.5"),
            api_key=str(payload.get("apiKey") or ""),
            timeout_sec=float(payload.get("timeout") or 600),
            retries=int(payload.get("retries") or 2),
            llm_parallel=int(payload.get("llmParallel") or 5),
            font_scale=float(payload.get("fontScale") or 0.93),
            keep_intermediates=bool(payload.get("keepIntermediates", False)),
            use_sam3=bool(payload.get("sam3", False)),
            hybrid_allowed_delta=float(payload.get("hybridAllowedDelta") or 0.0),
            editable_embedded_text=bool(payload.get("editableEmbeddedText", True)),
            layout_failure_mode=str(payload.get("layoutFailureMode") or "image_fallback"),
            watermark_mode=str(payload.get("watermarkMode") or "off"),
            reconstruction_mode=str(payload.get("reconstructionMode") or "auto"),
            text_erase_mode=str(payload.get("textEraseMode") or "codex_imagegen"),
            background_mode=str(payload.get("backgroundMode") or "white"),
            use_native_shapes=bool(payload.get("useNativeShapes", True)),
        )
        previews_dir = job_dir / "artifacts" / "previews"

        def _pre_cleanup(workdir_at_cleanup: Path) -> None:
            # workdir이 정리되기 직전에 썸네일/캐시를 만든다. 이 시점에는
            # workdir/{images,qa_*}/ PNG 들이 아직 살아있다.
            with contextlib.suppress(Exception):
                generate_previews(workdir=workdir_at_cleanup, out_dir=previews_dir)
            with contextlib.suppress(Exception):
                prefetch_pptx_pages_from_workdir(
                    job_id, job_dir, workdir_at_cleanup, output_path, input_path
                )

        result = convert_notebooklm_auto(
            input_path=input_path,
            output_path=output_path,
            workdir=workdir,
            options=options,
            progress=progress,
            cancel=cancel_check,
            pre_cleanup_hook=_pre_cleanup,
        )
        update_job(
            job_id,
            force=True,
            status="done",
            phase="완료",
            progress=100,
            output=str(result.output),
            report=str(result.report),
            montage=str(result.montage) if result.montage else None,
        )
    except JobCancelledError:
        update_job(job_id, force=True, status="cancelled", phase="취소됨", progress=100, error=None)
    except Exception as exc:  # pragma: no cover - depends on external tools/model
        error_path = job_dir / "error.log"
        error_path.write_text(sanitize_error(traceback.format_exc()), encoding="utf-8")
        update_job(job_id, force=True, status="failed", phase="실패", error=sanitize_error(str(exc)), progress=100)


_LOOPBACK_HOSTNAMES = frozenset({"localhost", "127.0.0.1", "::1"})

# 사설 네트워크(RFC1918) LLM 호출 허용 여부.
# Default True — 우리 서버가 loopback에만 바인딩됐다는 가정 하에 사내 GPU 서버 등
# 정당한 LAN LLM 사용을 막지 않는다. 우리 서버가 비-loopback host에 바인딩되면
# (cli/web.py가 자동으로) False로 전환해 SSRF 위험을 차단한다.
_allow_private_networks: bool = True


def set_allow_private_networks(allow: bool) -> None:
    """Toggle SSRF guard policy for RFC1918 private IPs.

    Loopback / link-local (169.254.x → cloud IMDS) / multicast / unspecified
    addresses are always rejected regardless of this flag.
    """
    global _allow_private_networks
    _allow_private_networks = bool(allow)


def _validate_outbound_url(url: str) -> tuple[bool, str]:
    """SSRF guard for user-supplied OpenAI-compatible base URLs.

    Always rejects: non-http(s) schemes, link-local (cloud IMDS), multicast,
    unspecified addresses, and resolved loopback unless the hostname literal
    is explicitly localhost/127.0.0.1/::1.
    Conditionally rejects RFC1918 private IPs based on _allow_private_networks.
    """
    import ipaddress
    import socket
    from urllib.parse import urlparse

    parsed = urlparse(url)
    if parsed.scheme not in {"http", "https"}:
        return False, "http 또는 https URL만 허용됩니다."
    hostname = parsed.hostname
    if not hostname:
        return False, "URL에 호스트가 없습니다."
    if hostname in _LOOPBACK_HOSTNAMES:
        return True, ""
    try:
        infos = socket.getaddrinfo(hostname, parsed.port or 80, type=socket.SOCK_STREAM)
    except socket.gaierror:
        return False, f"호스트 '{hostname}'를 해석할 수 없습니다."
    # is_reserved는 NAT64 prefix(64:ff9b::/96) 등 정상 외부 트래픽도 포함하므로
    # 제외한다. link-local(169.254.0.0/16)에는 cloud IMDS가 포함되어 정책과 무관하게
    # 항상 차단한다.
    for _family, _socktype, _proto, _canon, sockaddr in infos:
        try:
            ip = ipaddress.ip_address(sockaddr[0])
        except (ValueError, IndexError):
            continue
        if ip.is_loopback or ip.is_link_local or ip.is_multicast or ip.is_unspecified:
            return False, f"내부 네트워크 주소({ip})는 허용되지 않습니다."
        if ip.is_private and not _allow_private_networks:
            return False, (
                f"사설 네트워크 주소({ip})는 현재 정책에서 차단됩니다 "
                "(서버가 비-loopback host에 바인딩되어 SSRF 방지가 활성화됨)."
            )
    return True, ""


def require_json_content_type(request: Request) -> None:
    """Reject requests without an application/json Content-Type.

    This forces a CORS preflight on cross-origin browser callers, which the
    default same-origin policy then blocks. Trivial CSRF defense for the
    state-mutating POST endpoints (cancel/rerun) that otherwise accept empty
    bodies.
    """
    raw = request.headers.get("content-type", "")
    media_type = raw.split(";", 1)[0].strip().lower()
    if media_type != "application/json":
        raise HTTPException(
            status_code=415,
            detail="Content-Type: application/json 필요 (CSRF 방지).",
        )


def probe_llm_endpoint(payload: dict[str, Any]) -> dict[str, Any]:
    """Issue a tiny chat completion to verify the OpenAI-compatible endpoint.

    Empty apiKey is allowed (Ollama/local proxies). On 404 with a base URL
    that does not end with /v1, retries once with /v1 appended and reports
    the auto-correction in the result note.
    """
    base_url = str(payload.get("baseUrl") or "").strip().rstrip("/")
    model = str(payload.get("model") or "").strip()
    api_key = str(payload.get("apiKey") or "").strip()
    timeout = float(payload.get("timeout") or 15)
    timeout = max(3.0, min(60.0, timeout))

    if not base_url:
        return {"ok": False, "error": "Base URL을 입력하세요."}
    if not model:
        return {"ok": False, "error": "Model 이름을 입력하세요."}

    ok, err = _validate_outbound_url(base_url)
    if not ok:
        return {"ok": False, "error": err}

    first = _probe_chat_completions(base_url, model, api_key, timeout)
    if first.get("ok"):
        return first
    if first.get("status_code") == 404 and not base_url.endswith("/v1"):
        retry_base = f"{base_url}/v1"
        retry = _probe_chat_completions(retry_base, model, api_key, timeout)
        if retry.get("ok"):
            retry["note"] = (
                f"Base URL에 '/v1'이 빠져있어 자동 보정했습니다. "
                f"Base URL을 '{retry_base}'로 갱신하는 것을 권장합니다."
            )
            return retry
        first = dict(first)
        first["error"] = (
            f"{first.get('error', '')} · Base URL이 /v1으로 끝나야 할 수 있습니다 "
            f"(예: {base_url}/v1)"
        )
        if _looks_like_model_missing(retry):
            first["available_models"] = _list_models(retry_base, api_key, timeout)
            first["error"] += " · 모델명이 잘못된 것 같습니다. 사용 가능한 모델은 아래 참고."
        return first
    if _looks_like_model_missing(first):
        first = dict(first)
        first["available_models"] = _list_models(base_url, api_key, timeout)
        suffix = " · `ollama pull <model>`로 모델을 다운로드하거나 사용 가능한 모델 중에서 선택하세요."
        first["error"] = f"{first.get('error', '')}{suffix}"
    return first


def _looks_like_model_missing(result: dict[str, Any]) -> bool:
    if result.get("ok"):
        return False
    return bool(result.get("model_missing"))


def _list_models(base_url: str, api_key: str, timeout: float) -> list[str]:
    import httpx

    ok, _ = _validate_outbound_url(base_url)
    if not ok:
        return []

    headers: dict[str, str] = {"Accept": "application/json"}
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"
    try:
        with httpx.Client(timeout=min(timeout, 8.0), follow_redirects=False) as client:
            response = client.get(f"{base_url}/models", headers=headers)
        if response.status_code >= 400:
            return []
        data = response.json()
    except (httpx.HTTPError, ValueError):
        return []
    items: list[str] = []
    candidates = data.get("data") if isinstance(data, dict) else None
    if isinstance(candidates, list):
        for entry in candidates:
            if isinstance(entry, dict):
                value = entry.get("id") or entry.get("name")
                if isinstance(value, str) and value:
                    items.append(value)
    return items[:12]


def _probe_chat_completions(
    base_url: str,
    model: str,
    api_key: str,
    timeout: float,
) -> dict[str, Any]:
    import httpx

    url = f"{base_url}/chat/completions"
    headers: dict[str, str] = {"Content-Type": "application/json"}
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"
    body = {
        "model": model,
        "messages": [{"role": "user", "content": "Reply with the single word: ok"}],
        "max_tokens": 5,
        "temperature": 0,
    }

    started = time.perf_counter()
    try:
        with httpx.Client(timeout=timeout, follow_redirects=False) as client:
            response = client.post(url, json=body, headers=headers)
    except httpx.TimeoutException:
        return {
            "ok": False,
            "error": f"{timeout:.0f}초 안에 응답이 없습니다.",
            "latency_ms": int((time.perf_counter() - started) * 1000),
        }
    except httpx.HTTPError as exc:
        return {
            "ok": False,
            "error": f"네트워크 오류: {sanitize_error(str(exc))}",
            "latency_ms": int((time.perf_counter() - started) * 1000),
        }
    latency_ms = int((time.perf_counter() - started) * 1000)

    if 300 <= response.status_code < 400:
        # SSRF 회피: 리다이렉트는 허용하지 않음 (내부 메타데이터/사설망 우회 차단).
        return {
            "ok": False,
            "status_code": response.status_code,
            "latency_ms": latency_ms,
            "error": f"HTTP {response.status_code}: 리다이렉트는 보안상 허용되지 않습니다.",
        }

    if response.status_code >= 400:
        # 응답 본문은 반사하지 않음 (외부 endpoint가 내부 정보/시크릿을 에코할 수 있음).
        # status code + 일반 hint만 클라이언트에 전달.
        hint = ""
        if response.status_code in (401, 403):
            hint = " (API key 또는 권한을 확인하세요)"
        elif response.status_code == 404:
            hint = " (Base URL 또는 모델명을 확인하세요)"
        elif response.status_code == 429:
            hint = " (rate limit 초과 또는 quota 부족)"
        # 내부 판별: model not found 안내를 자동으로 추가하기 위해 본문 일부를 검사.
        # 본문은 응답에 포함하지 않고 model_missing 플래그만 기록한다.
        result: dict[str, Any] = {
            "ok": False,
            "status_code": response.status_code,
            "latency_ms": latency_ms,
            "error": f"HTTP {response.status_code}{hint}",
        }
        if response.status_code == 404:
            body_lower = response.text[:500].lower()
            if "not found" in body_lower or "no such model" in body_lower or "does not exist" in body_lower:
                result["model_missing"] = True
        return result

    try:
        data = response.json()
    except ValueError:
        return {
            "ok": False,
            "status_code": response.status_code,
            "latency_ms": latency_ms,
            "error": "응답이 JSON 형식이 아닙니다.",
        }

    sample = ""
    if isinstance(data, dict):
        choices = data.get("choices")
        if isinstance(choices, list) and choices:
            message = choices[0].get("message") if isinstance(choices[0], dict) else None
            if isinstance(message, dict):
                content = message.get("content")
                if isinstance(content, str):
                    sample = content.strip()[:120]
    return {
        "ok": True,
        "status_code": response.status_code,
        "latency_ms": latency_ms,
        "sample": sample,
        "model": model,
    }


def sanitize_error(message: str) -> str:
    message = re.sub(r"sk-[A-Za-z0-9_-]+", "sk-********", message)
    message = re.sub(r"(--api-key['\", ]+)([^'\",\\]]+)", r"\1********", message)
    message = re.sub(r"(apiKey['\": ]+)([^'\",\\]]+)", r"\1********", message)
    # Bearer / Google API key / URL-embedded basic auth credentials
    message = re.sub(r"Bearer\s+\S+", "Bearer ********", message, flags=re.IGNORECASE)
    message = re.sub(r"AIza[0-9A-Za-z_\-]{35}", "AIza********", message)
    message = re.sub(r"(https?://)([^/@\s]+):([^/@\s]+)@", r"\1********:********@", message)
    return message


def _avg_diff(items: list[dict[str, Any]]) -> float | None:
    values = [item.get("mean_abs_diff") for item in items]
    numeric = [float(value) for value in values if isinstance(value, int | float)]
    if not numeric:
        return None
    return round(sum(numeric) / len(numeric), 2)


def _file_size(path: str | None) -> int | None:
    if not path:
        return None
    p = Path(path)
    return p.stat().st_size if p.exists() else None


def object_text(obj: dict[str, Any]) -> str:
    text = obj.get("text")
    if isinstance(text, str):
        return text
    lines = obj.get("lines")
    if isinstance(lines, list):
        return " ".join(str(line) for line in lines)
    return ""


def build_report_summary(job: JobState, report: dict[str, Any]) -> dict[str, Any]:
    selected_qa = report.get("selected_qa") if isinstance(report.get("selected_qa"), list) else []
    vector_qa = report.get("vector_qa") if isinstance(report.get("vector_qa"), list) else []
    hybrid_qa = report.get("hybrid_qa") if isinstance(report.get("hybrid_qa"), list) else []
    decisions = report.get("selection_report") if isinstance(report.get("selection_report"), list) else []

    chosen_counts = {"vector": 0, "hybrid": 0}
    for item in decisions:
        chosen = item.get("chosen")
        if chosen in chosen_counts:
            chosen_counts[chosen] += 1

    worst = sorted(
        (
            item
            for item in selected_qa
            if isinstance(item.get("mean_abs_diff"), int | float)
        ),
        key=lambda item: float(item["mean_abs_diff"]),
        reverse=True,
    )[:5]
    vector_artifact = resolve_artifact_path(job, "vector")
    hybrid_artifact = resolve_artifact_path(job, "hybrid")
    layout_artifact = resolve_artifact_path(job, "layout-json")

    return {
        "job": {
            "id": job.id,
            "filename": job.filename,
            "status": job.status,
            "created_at": job.created_at,
            "updated_at": job.updated_at,
            "options": job.options,
        },
        "summary": {
            "slide_count": len(selected_qa),
            "selected_avg_diff": _avg_diff(selected_qa),
            "vector_avg_diff": _avg_diff(vector_qa),
            "hybrid_avg_diff": _avg_diff(hybrid_qa),
            "chosen_vector_count": chosen_counts["vector"],
            "chosen_hybrid_count": chosen_counts["hybrid"],
        },
        "worst_slides": [
            {
                "slide_no": item.get("slide_no"),
                "layout_id": item.get("layout_id"),
                "object_count": item.get("object_count"),
                "picture_count": item.get("picture_count"),
                "mean_abs_diff": item.get("mean_abs_diff"),
            }
            for item in worst
        ],
        "decisions": decisions,
        "files": {
            "output": job.output,
            "report": job.report,
            "montage": job.montage,
            "candidate_vector": str(vector_artifact) if vector_artifact else None,
            "candidate_hybrid": str(hybrid_artifact) if hybrid_artifact else None,
            "layout_json": str(layout_artifact) if layout_artifact else None,
            "output_bytes": _file_size(job.output),
            "report_bytes": _file_size(job.report),
            "montage_bytes": _file_size(job.montage),
            "candidate_vector_bytes": _file_size(str(vector_artifact) if vector_artifact else None),
            "candidate_hybrid_bytes": _file_size(str(hybrid_artifact) if hybrid_artifact else None),
            "layout_json_bytes": _file_size(str(layout_artifact) if layout_artifact else None),
        },
    }


app = create_app()


INDEX_HTML = r"""<!doctype html>
<html lang="ko">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>Star-Slide</title>
  <link rel="icon" type="image/svg+xml" href="data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 64 64'%3E%3Crect width='64' height='64' rx='12' fill='%230b111b'/%3E%3Crect x='13' y='14' width='38' height='28' rx='4' fill='%233b82f6'/%3E%3Cpath d='M20 24h24M20 31h16' stroke='white' stroke-width='4' stroke-linecap='round'/%3E%3Cpath d='M18 49h28' stroke='%232dd4bf' stroke-width='5' stroke-linecap='round'/%3E%3Ccircle cx='49' cy='42' r='6' fill='%232dd4bf'/%3E%3C/svg%3E" />
  <style>
    :root {
      color-scheme: dark;
      --bg: #0b111b;
      --surface: #101827;
      --surface-2: #131f30;
      --surface-3: #172337;
      --ink: #eef4ff;
      --muted: #94a3b8;
      --line: #2a3a50;
      --field: #0d1624;
      --field-border: #34465f;
      --blue: #3b82f6;
      --teal: #2dd4bf;
      --orange: #fb923c;
      --danger: #f87171;
      --shadow: rgba(0, 0, 0, .42);
    }
    body[data-theme="light"] {
      color-scheme: light;
      --bg: #ffffff;
      --surface: #fbfcfe;
      --surface-2: #ffffff;
      --surface-3: #f8fafc;
      --ink: #111827;
      --muted: #667085;
      --line: #d8dee8;
      --field: #ffffff;
      --field-border: #cfd6e4;
      --blue: #2563eb;
      --teal: #0f766e;
      --orange: #ea580c;
      --danger: #b42318;
      --shadow: rgba(15, 23, 42, .20);
    }
    * { box-sizing: border-box; }
    body {
      margin: 0;
      color: var(--ink);
      font-family: ui-sans-serif, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
      background: var(--bg);
    }
    header {
      height: 68px;
      display: grid;
      grid-template-columns: minmax(220px, auto) minmax(0, 1fr) minmax(280px, auto);
      align-items: center;
      padding: 0 28px;
      border-bottom: 1px solid var(--line);
      background: linear-gradient(180deg, var(--surface) 0%, color-mix(in srgb, var(--surface) 92%, var(--blue)) 100%);
      box-shadow: 0 1px 0 color-mix(in srgb, var(--blue) 18%, transparent);
    }
    .brand {
      display: inline-flex;
      align-items: center;
      gap: 10px;
      font-size: 22px;
      font-weight: 800;
      letter-spacing: -0.01em;
      background: linear-gradient(90deg, var(--ink) 0%, var(--teal) 110%);
      -webkit-background-clip: text;
      background-clip: text;
      color: transparent;
    }
    .brand-mark {
      width: 28px;
      height: 28px;
      border-radius: 7px;
      background: linear-gradient(135deg, var(--blue), var(--teal));
      display: inline-grid;
      place-items: center;
      color: #fff;
      font-size: 14px;
      font-weight: 800;
      box-shadow: 0 4px 14px color-mix(in srgb, var(--blue) 35%, transparent);
    }
    .tagline {
      min-width: 0;
      text-align: center;
      font-size: 15px;
      font-weight: 700;
      letter-spacing: -0.005em;
      color: color-mix(in srgb, var(--ink) 88%, var(--muted));
      overflow: hidden;
      text-overflow: ellipsis;
      white-space: nowrap;
    }
    .tagline strong {
      color: var(--teal);
      font-weight: 800;
    }
    .header-actions {
      display: flex;
      align-items: center;
      justify-content: flex-end;
      gap: 12px;
    }
    .theme-toggle {
      min-width: 104px;
      color: var(--ink);
      border-color: var(--field-border);
      background: var(--surface-2);
    }
    .shell {
      display: grid;
      grid-template-columns: minmax(420px, 520px) minmax(0, 1fr);
      min-height: calc(100vh - 68px);
    }
    aside {
      border-right: 1px solid var(--line);
      padding: 22px;
      background: var(--surface);
      overflow-y: auto;
    }
    main { padding: 26px 28px; overflow: auto; }
    h1, h2, h3 { margin: 0; letter-spacing: 0; }
    h1 {
      font-size: 22px;
      font-weight: 800;
      letter-spacing: -0.01em;
      display: inline-flex;
      align-items: center;
      gap: 10px;
      margin-bottom: 18px;
    }
    h1::before {
      content: "";
      width: 4px;
      height: 22px;
      border-radius: 4px;
      background: linear-gradient(180deg, var(--blue), var(--teal));
    }
    h2 {
      font-size: 11px;
      font-weight: 800;
      text-transform: uppercase;
      letter-spacing: 0.08em;
      color: var(--muted);
      margin-bottom: 10px;
      padding-bottom: 8px;
      border-bottom: 1px solid var(--line);
      display: flex;
      align-items: center;
      gap: 8px;
    }
    h2::before {
      content: "";
      width: 6px;
      height: 6px;
      border-radius: 50%;
      background: var(--teal);
    }
    label { display: block; font-size: 12px; font-weight: 700; color: var(--ink); margin: 12px 0 6px; }
    .label-row {
      display: flex;
      align-items: center;
      gap: 6px;
      margin: 12px 0 6px;
    }
    .label-row label { margin: 0; }
    .help {
      width: 18px;
      height: 18px;
      display: inline-flex;
      align-items: center;
      justify-content: center;
      border: 1px solid var(--field-border);
      border-radius: 50%;
      color: var(--muted);
      background: var(--surface-2);
      font-size: 12px;
      font-weight: 800;
      cursor: help;
    }
    input, select {
      width: 100%;
      height: 38px;
      border: 1px solid var(--field-border);
      border-radius: 7px;
      padding: 0 10px;
      font-size: 14px;
      color: var(--ink);
      background: var(--field);
    }
    input[type="checkbox"] { width: 16px; height: 16px; }
    .row { display: grid; grid-template-columns: 1fr 1fr; gap: 10px; }
    .provider-controls {
      display: grid;
      grid-template-columns: minmax(0, 1fr) auto;
      align-items: end;
      gap: 10px;
    }
    .provider-controls button { white-space: nowrap; }
    .info-box {
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 11px 12px;
      color: var(--muted);
      background: var(--surface-2);
      font-size: 12px;
      line-height: 1.5;
    }
    .system-box {
      display: grid;
      gap: 8px;
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 12px;
      background: var(--surface-2);
      margin-top: 16px;
    }
    .system-row {
      display: grid;
      grid-template-columns: minmax(92px, auto) 74px minmax(0, 1fr);
      align-items: center;
      gap: 8px;
      font-size: 12px;
    }
    .system-name { font-weight: 750; color: var(--ink); }
    .status-pill {
      display: inline-flex;
      justify-content: center;
      border-radius: 999px;
      padding: 3px 8px;
      font-weight: 800;
      background: color-mix(in srgb, var(--danger) 18%, var(--surface-3));
      color: var(--danger);
    }
    .status-pill.ok {
      background: color-mix(in srgb, var(--teal) 18%, var(--surface-3));
      color: var(--teal);
    }
    .custom-provider-fields[hidden] { display: none; }
    .provider-field-actions {
      margin-top: 10px;
      justify-content: flex-start;
    }
    .drop {
      display: grid;
      place-items: center;
      min-height: 168px;
      border: 2px dashed var(--field-border);
      border-radius: 8px;
      background: var(--surface-2);
      text-align: center;
      padding: 20px;
      cursor: pointer;
    }
    .drop.drag { border-color: var(--blue); background: color-mix(in srgb, var(--blue) 13%, var(--surface-2)); }
    .drop strong { display: block; font-size: 17px; margin-bottom: 7px; }
    .hint { color: var(--muted); font-size: 12px; line-height: 1.45; }
    .actions { display: flex; gap: 10px; margin-top: 16px; }
    .upload-actions { margin-top: 12px; }
    button, .button {
      height: 38px;
      display: inline-flex;
      align-items: center;
      justify-content: center;
      gap: 7px;
      border: 1px solid #1d4ed8;
      border-radius: 7px;
      padding: 0 14px;
      color: #fff;
      background: var(--blue);
      font-weight: 750;
      text-decoration: none;
      cursor: pointer;
    }
    button.secondary, .button.secondary {
      color: var(--ink);
      border-color: var(--field-border);
      background: var(--surface-2);
    }
    button:disabled { opacity: .55; cursor: not-allowed; }
    .checkline { display: flex; align-items: center; gap: 8px; margin-top: 12px; font-size: 13px; }
    .image-split-group {
      margin-top: 14px;
      padding: 12px 12px 12px 16px;
      border-left: 3px solid var(--blue);
      border-radius: 6px;
      background: color-mix(in srgb, var(--blue) 6%, var(--surface-2));
      display: flex; flex-direction: column; gap: 10px;
    }
    .image-split-group.disabled { opacity: 0.45; pointer-events: none; border-left-color: var(--line); background: var(--surface-2); }
    .image-split-header { font-size: 12px; font-weight: 600; color: var(--blue); margin-bottom: 4px; }
    .image-split-group.disabled .image-split-header { color: var(--muted); }
    .jobs {
      display: grid;
      gap: 12px;
    }
    .job {
      display: flex;
      gap: 14px;
      align-items: flex-start;
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 14px;
      background: var(--surface-2);
    }
    .job-thumb {
      flex: 0 0 auto;
      width: 220px;
      aspect-ratio: 16 / 9;
      object-fit: cover;
      border-radius: 6px;
      border: 1px solid var(--line);
      background: var(--surface-3);
      display: block;
    }
    .job-thumb.is-missing { display: none; }
    .job-check {
      flex: 0 0 auto;
      width: 18px;
      height: 18px;
      margin-top: 4px;
      accent-color: var(--blue);
      cursor: pointer;
    }
    .jobs-toolbar {
      display: flex;
      align-items: center;
      gap: 12px;
      flex-wrap: wrap;
      padding: 10px 12px;
      margin-bottom: 12px;
      border: 1px solid var(--line);
      border-radius: 8px;
      background: var(--surface-2);
    }
    .jobs-toolbar-check {
      display: inline-flex;
      align-items: center;
      gap: 8px;
      font-size: 13px;
      color: var(--ink);
      cursor: pointer;
    }
    .jobs-toolbar-check input { accent-color: var(--blue); }
    .pager {
      display: flex;
      align-items: center;
      gap: 6px;
      flex-wrap: wrap;
      margin-top: 14px;
      padding: 10px 0;
    }
    .pager-arrow,
    .pager-page {
      min-width: 34px;
      height: 32px;
      padding: 0 10px;
      border: 1px solid var(--field-border);
      border-radius: 6px;
      background: var(--surface-2);
      color: var(--ink);
      font-size: 13px;
      cursor: pointer;
    }
    .pager-arrow:disabled,
    .pager-page:disabled { opacity: 0.4; cursor: default; }
    .pager-page.is-active {
      background: var(--blue);
      color: #fff;
      border-color: var(--blue);
    }
    .pager-ellipsis {
      padding: 0 6px;
      color: var(--muted);
    }
    .job-body {
      flex: 1 1 auto;
      min-width: 0;
    }
    @media (max-width: 720px) {
      .job { flex-direction: column; }
      .job-thumb { width: 100%; max-width: 320px; }
    }
    .job-head {
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 12px;
      margin-bottom: 10px;
    }
    .job-title { font-weight: 760; overflow-wrap: anywhere; }
    .job-meta {
      display: flex;
      flex-wrap: wrap;
      gap: 8px;
      margin-top: 8px;
    }
    .metric {
      border: 1px solid var(--line);
      border-radius: 999px;
      padding: 4px 8px;
      color: var(--muted);
      background: var(--surface-3);
      font-size: 12px;
    }
    .badge {
      border-radius: 999px;
      padding: 4px 9px;
      font-size: 12px;
      font-weight: 750;
      background: color-mix(in srgb, var(--blue) 18%, var(--surface-3));
      color: var(--blue);
      white-space: nowrap;
    }
    .badge.done { background: color-mix(in srgb, var(--teal) 18%, var(--surface-3)); color: var(--teal); }
    .badge.failed { background: color-mix(in srgb, var(--danger) 18%, var(--surface-3)); color: var(--danger); }
    .badge.cancelled, .badge.cancelling { background: color-mix(in srgb, var(--muted) 22%, var(--surface-3)); color: var(--muted); }
    .keystore-banner {
      padding: 10px 28px;
      font-size: 13px;
      font-weight: 600;
      background: color-mix(in srgb, var(--danger) 18%, var(--surface));
      color: var(--danger);
      border-bottom: 1px solid color-mix(in srgb, var(--danger) 40%, var(--line));
    }
    .api-key-row {
      display: grid;
      grid-template-columns: minmax(0, 1fr) auto;
      gap: 8px;
      align-items: center;
    }
    .api-key-row button { white-space: nowrap; }
    .test-result {
      margin-top: 8px;
      padding: 8px 10px;
      border-radius: 6px;
      border: 1px solid var(--field-border);
      background: var(--surface-3);
      font-size: 12px;
      line-height: 1.5;
    }
    .test-result.ok { border-color: color-mix(in srgb, var(--teal) 60%, var(--field-border)); color: var(--teal); }
    .test-result.fail { border-color: color-mix(in srgb, var(--danger) 50%, var(--field-border)); color: var(--danger); }
    .test-result.busy { color: var(--muted); }
    .model-chips { display: flex; flex-wrap: wrap; gap: 6px; margin-top: 6px; }
    .model-chip {
      height: auto;
      padding: 4px 10px;
      font-size: 12px;
      font-weight: 600;
      border: 1px solid var(--field-border);
      border-radius: 999px;
      background: var(--surface-3);
      color: var(--ink);
      cursor: pointer;
    }
    .model-chip:hover { border-color: var(--blue); color: var(--blue); }
    .model-chip.active {
      border-color: var(--blue);
      background: color-mix(in srgb, var(--blue) 22%, var(--surface-3));
      color: var(--blue);
    }
    .upload-progress { margin-top: 12px; display: grid; gap: 6px; }
    .upload-progress-bar { height: 6px; border-radius: 999px; background: var(--surface-3); overflow: hidden; }
    .upload-progress-bar > div { height: 100%; width: 0%; background: linear-gradient(90deg, var(--blue), var(--teal)); transition: width 120ms linear; }
    .modal.viewer-modal { width: min(1400px, 98vw); max-height: 94vh; padding: 0; }
    .modal.viewer-modal.fullscreen { width: 99vw; max-height: 98vh; }
    .viewer {
      display: grid;
      grid-template-rows: auto 1fr auto;
      height: min(90vh, 880px);
      background: var(--surface);
    }
    .viewer.fullscreen { height: 92vh; }
    .viewer-head {
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 12px;
      padding: 10px 14px;
      border-bottom: 1px solid var(--line);
      background: var(--surface-2);
    }
    .viewer-title {
      display: flex;
      align-items: center;
      gap: 10px;
      min-width: 0;
      font-size: 13px;
      font-weight: 700;
      color: var(--ink);
    }
    .viewer-title .kind-pill {
      padding: 2px 8px;
      border-radius: 999px;
      font-size: 11px;
      font-weight: 700;
      letter-spacing: 0.04em;
      text-transform: uppercase;
      background: color-mix(in srgb, var(--blue) 22%, var(--surface-3));
      color: var(--blue);
      white-space: nowrap;
    }
    .viewer-title .kind-pill.original { background: color-mix(in srgb, var(--muted) 26%, var(--surface-3)); color: var(--muted); }
    .viewer-title .filename {
      overflow: hidden;
      text-overflow: ellipsis;
      white-space: nowrap;
    }
    .viewer-title .counter {
      color: var(--muted);
      font-weight: 600;
      white-space: nowrap;
    }
    .viewer-actions { display: flex; align-items: center; gap: 4px; }
    .viewer-icon {
      height: 32px;
      width: 32px;
      padding: 0;
      border: 1px solid transparent;
      background: transparent;
      color: var(--muted);
      border-radius: 7px;
      display: inline-grid;
      place-items: center;
      cursor: pointer;
    }
    .viewer-icon:hover { color: var(--ink); background: var(--surface-3); border-color: var(--field-border); }
    .viewer-stage {
      position: relative;
      display: grid;
      place-items: center;
      /* 다크모드: 모달 본문 톤과 통일해 letterbox 시인성 거슬림 제거 */
      background: var(--surface);
      overflow: hidden;
    }
    body[data-theme="light"] .viewer-stage { background: #e5e7eb; }
    .viewer-stage img {
      max-width: 100%;
      max-height: 100%;
      object-fit: contain;
      box-shadow: 0 6px 28px rgba(0, 0, 0, .35);
      border-radius: 6px;
      background: #fff;
    }
    .viewer-stage .pptx-host {
      width: 100%;
      height: 100%;
      overflow: hidden;
      background: var(--surface);
      display: grid;
      place-items: center;
    }
    body[data-theme="light"] .viewer-stage .pptx-host { background: #e5e7eb; }
    .viewer-stage .pptx-host .pptx-preview-wrapper {
      background: #fff;
      box-shadow: 0 6px 28px rgba(0, 0, 0, .35);
      border-radius: 6px;
      overflow: hidden;
    }
    .viewer-stage.compare {
      display: grid;
      grid-template-columns: 1fr 1fr;
      gap: 6px;
      /* 좌우 패딩을 줄여 슬라이드를 키운다. 화살표(.viewer-nav)는 stage 가장자리에 overlay. */
      padding: 4px 8px;
      align-items: stretch;
      justify-items: stretch;
    }
    .viewer-stage.compare .compare-pane {
      display: block;
      width: 100%;
      height: 100%;
      min-height: 0;
    }
    .viewer-stage.compare .compare-pane .pane-label { display: none; }
    .viewer-stage.compare .compare-pane .pane-img-wrap {
      width: 100%;
      height: 100%;
      min-height: 0;
      overflow: hidden;
      display: grid;
      place-items: center;
    }
    .viewer-stage.compare .compare-pane .pptx-host {
      width: 100%;
      height: 100%;
    }
    @media (max-width: 900px) {
      .viewer-stage.compare {
        grid-template-columns: 1fr;
        grid-template-rows: 1fr 1fr;
        padding: 4px 8px;
        gap: 6px;
      }
    }
    .viewer-empty { color: var(--muted); font-size: 13px; }
    .viewer-nav {
      position: absolute;
      top: 50%;
      transform: translateY(-50%);
      width: 44px;
      height: 44px;
      border-radius: 50%;
      border: 0;
      display: grid;
      place-items: center;
      cursor: pointer;
      background: rgba(0, 0, 0, .45);
      color: #fff;
    }
    .viewer-nav:hover:not(:disabled) { background: rgba(0, 0, 0, .65); }
    .viewer-nav:disabled { opacity: .25; cursor: not-allowed; }
    .viewer-nav.prev { left: 6px; }
    .viewer-nav.next { right: 6px; }
    .viewer-foot {
      display: flex;
      align-items: center;
      justify-content: center;
      gap: 14px;
      padding: 10px 14px;
      border-top: 1px solid var(--line);
      background: var(--surface-2);
      color: var(--muted);
      font-size: 13px;
    }
    .viewer-foot button {
      height: 32px;
      width: 32px;
      padding: 0;
      border-radius: 50%;
      border: 1px solid var(--field-border);
      background: var(--surface-3);
      color: var(--ink);
      display: inline-grid;
      place-items: center;
      cursor: pointer;
    }
    .viewer-foot button:disabled { opacity: .4; cursor: not-allowed; }
    .preview-grid {
      display: grid;
      grid-template-columns: repeat(auto-fill, minmax(180px, 1fr));
      gap: 12px;
      margin-top: 16px;
    }
    .preview-card {
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 8px;
      background: var(--surface-3);
      cursor: pointer;
      display: grid;
      gap: 6px;
    }
    .preview-card:hover { border-color: var(--blue); }
    .preview-card img { width: 100%; height: auto; border-radius: 6px; display: block; background: var(--surface); }
    .preview-card .meta { display: flex; justify-content: space-between; font-size: 12px; color: var(--muted); }
    .compare {
      display: grid;
      gap: 12px;
    }
    .compare-toolbar {
      display: flex;
      flex-wrap: wrap;
      align-items: center;
      gap: 8px;
    }
    .compare-toolbar .spacer { flex: 1; }
    .compare-toolbar .kind-toggle {
      display: inline-flex;
      gap: 4px;
      padding: 3px;
      border: 1px solid var(--field-border);
      border-radius: 8px;
      background: var(--surface-3);
    }
    .compare-toolbar .kind-toggle button {
      height: 28px;
      padding: 0 10px;
      border-radius: 5px;
      border: 0;
      background: transparent;
      color: var(--ink);
    }
    .compare-toolbar .kind-toggle button.active { background: var(--blue); color: #fff; }
    .compare-pair {
      display: grid;
      grid-template-columns: 1fr 1fr;
      gap: 12px;
    }
    .compare-pane {
      display: grid;
      gap: 6px;
    }
    .compare-pane img {
      width: 100%;
      border: 1px solid var(--line);
      border-radius: 8px;
      background: var(--surface);
      display: block;
    }
    .compare-pane .label { font-size: 12px; color: var(--muted); }
    @media (max-width: 720px) {
      .compare-pair { grid-template-columns: 1fr; }
    }
    .bar {
      height: 8px;
      border-radius: 999px;
      background: var(--surface-3);
      overflow: hidden;
      margin: 8px 0;
    }
    .bar > div { height: 100%; background: linear-gradient(90deg, var(--blue), var(--teal)); }
    .job-actions { display: flex; flex-wrap: wrap; gap: 8px; margin-top: 12px; }
    .pager {
      display: flex;
      align-items: center;
      justify-content: flex-end;
      gap: 8px;
      margin-bottom: 12px;
    }
    .pager button {
      height: 32px;
      padding: 0 10px;
    }
    .modal-backdrop {
      position: fixed;
      inset: 0;
      display: none;
      align-items: center;
      justify-content: center;
      padding: 24px;
      /* 다크모드 기본 — 다크 페이지 위에서 모달이 명확히 떠보이도록 진한 backdrop */
      background: rgba(2, 6, 14, .68);
      backdrop-filter: blur(2px);
      z-index: 20;
    }
    body[data-theme="light"] .modal-backdrop {
      background: rgba(15, 23, 42, .40);
    }
    .modal-backdrop.open { display: flex; }
    .modal {
      width: min(980px, 96vw);
      max-height: 88vh;
      overflow: auto;
      border-radius: 10px;
      background: var(--surface-3);
      border: 1px solid color-mix(in srgb, var(--line) 80%, var(--ink) 20%);
      box-shadow: 0 32px 90px rgba(0, 0, 0, .55), 0 0 0 1px color-mix(in srgb, var(--blue) 14%, transparent);
    }
    body[data-theme="light"] .modal {
      background: var(--surface-2);
      box-shadow: 0 24px 80px var(--shadow);
    }
    .modal.settings-modal {
      width: min(640px, 96vw);
    }
    .settings-tabs {
      display: flex;
      gap: 2px;
      border-bottom: 1px solid var(--line);
      margin: -4px -4px 14px;
    }
    .settings-tab {
      flex: 1 1 auto;
      padding: 10px 14px;
      border: 0;
      background: transparent;
      color: var(--muted);
      font-size: 13px;
      font-weight: 700;
      cursor: pointer;
      border-bottom: 2px solid transparent;
      margin-bottom: -1px;
      text-align: center;
    }
    .settings-tab:hover { color: var(--ink); }
    .settings-tab.is-active {
      color: var(--ink);
      border-bottom-color: var(--blue);
    }
    .settings-tab-pane.is-hidden { display: none; }
    .modal-head {
      position: sticky;
      top: 0;
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 12px;
      padding: 16px 18px;
      border-bottom: 1px solid var(--line);
      background: var(--surface-2);
    }
    .modal-body { padding: 18px; }
    .modal-head h2,
    .modal-body h2 {
      font-size: 15px;
      font-weight: 750;
      text-transform: none;
      letter-spacing: 0;
      color: var(--ink);
      border-bottom: 0;
      padding-bottom: 0;
      margin-bottom: 12px;
      display: block;
    }
    .modal-head h2::before,
    .modal-body h2::before { content: none; }
    .preview {
      display: block;
      width: 100%;
      border: 1px solid var(--line);
      border-radius: 8px;
      background: var(--surface);
    }
    .report-grid {
      display: grid;
      grid-template-columns: repeat(4, minmax(0, 1fr));
      gap: 10px;
      margin: 14px 0 18px;
    }
    .report-card {
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 12px;
      background: var(--surface-3);
    }
    .report-card strong { display: block; font-size: 20px; margin-top: 4px; }
    table { width: 100%; border-collapse: collapse; font-size: 13px; }
    th, td { border-bottom: 1px solid var(--line); padding: 8px; text-align: left; }
    th { color: var(--muted); background: var(--surface-3); }
    .empty {
      border: 1px solid var(--line);
      border-radius: 8px;
      padding: 26px;
      color: var(--muted);
      background: var(--surface-2);
    }
    .setting-box {
      border-top: 1px solid var(--line);
      margin-top: 18px;
      padding-top: 16px;
    }
    @media (max-width: 980px) {
      header { grid-template-columns: 1fr auto; height: auto; min-height: 64px; gap: 10px; padding: 14px 18px; }
      .tagline { grid-column: 1 / -1; grid-row: 2; text-align: left; white-space: normal; }
      .shell { grid-template-columns: 1fr; }
      aside { border-right: 0; border-bottom: 1px solid var(--line); }
    }
  </style>
</head>
<body>
  <header>
    <div class="brand"><span class="brand-mark" aria-hidden="true">S</span><span>Star-Slide</span></div>
    <div class="tagline" data-i18n-html="tagline">NotebookLM에서 생성한 이미지 슬라이드를 <strong>편집 가능한 PPTX</strong>로 변환</div>
    <div class="header-actions">
      <button id="languageToggle" class="secondary theme-toggle" type="button" title="Language">EN</button>
      <button id="settingsButton" class="secondary theme-toggle" type="button" title="설정" data-i18n="settings.button" data-i18n-title="settings.title">⚙ 설정</button>
      <button id="themeToggle" class="secondary theme-toggle" type="button">라이트 모드</button>
    </div>
  </header>
  <div id="keystoreBanner" class="keystore-banner" hidden role="alert"></div>
  <div class="shell">
    <aside>
      <h2 data-i18n="upload.title">파일 업로드</h2>
      <div id="drop" class="drop">
        <div>
          <strong data-i18n="upload.drop">PPTX/PDF/이미지를 드래그하거나 클릭</strong>
          <span id="fileLabel" class="hint" data-i18n="upload.fileHint">.pptx · .pdf · .png · .jpg (이미지는 1-슬라이드 PPTX 로 변환)</span>
        </div>
      </div>
      <input id="file" type="file" accept=".pptx,.pdf,.png,.jpg,.jpeg,image/png,image/jpeg" hidden />
      <div class="actions upload-actions">
        <button id="start" data-i18n="upload.start">변환 시작</button>
      </div>
      <div id="uploadProgress" class="upload-progress" hidden>
        <div class="upload-progress-bar"><div></div></div>
        <div class="upload-progress-label hint">업로드 0%</div>
      </div>
      <div class="setting-box">
        <h2>LLM Provider</h2>
        <div class="info-box" data-i18n-html="provider.info">
          LLM은 슬라이드 이미지를 layout JSON으로 해석하고 큰 이미지 영역을 판단하는 데 사용됩니다.
          기본 흐름은 슬라이드당 약 2회 호출이며, retry 설정에 따라 추가 호출될 수 있습니다.
        </div>
        <div class="provider-controls" style="grid-template-columns: minmax(0, 1fr) minmax(0, 1fr);">
          <div>
            <label for="s_provider">Provider</label>
            <select id="s_provider"></select>
          </div>
          <div>
            <label for="s_model">Model</label>
            <input id="s_model" list="modelOptions" placeholder="gpt-5.5 (로컬은 자동 조회)" data-i18n-placeholder="model.placeholder" autocomplete="off" />
          </div>
        </div>
        <div class="hint" style="margin-top:6px;" data-i18n-html="provider.sidebarHint">
          Base URL · API Key · Custom Provider 관리 · 시스템 상태는 우측 상단 <strong>⚙ 설정</strong>에서.
        </div>
      </div>

      <div class="setting-box">
        <h2 data-i18n="options.title">변환 옵션</h2>
        <div class="row">
          <div>
            <label for="s_timeout">Timeout</label>
            <input id="s_timeout" type="number" min="60" step="30" />
          </div>
          <div>
            <div class="label-row">
              <label for="s_retries">Retries</label>
              <span class="help" title="layout JSON 생성이 실패했을 때 같은 슬라이드를 다시 호출하는 횟수입니다. 기본 2회이며, 네트워크/LLM 일시 실패를 흡수합니다." data-i18n-title="help.retries">?</span>
            </div>
            <input id="s_retries" type="number" min="0" step="1" />
          </div>
        </div>
        <div class="row">
          <div>
            <div class="label-row">
              <label for="s_llmParallel" data-i18n="options.llmParallel">LLM 병렬 수</label>
              <span class="help" title="여러 슬라이드의 LLM 분석 요청을 동시에 몇 개까지 실행할지 정합니다." data-i18n-title="help.llmParallel">?</span>
            </div>
            <input id="s_llmParallel" type="number" min="1" max="10" step="1" />
          </div>
          <div>
            <div class="label-row">
              <label for="s_fontScale" data-i18n="options.fontScale">폰트 배율</label>
              <span class="help" title="PPTX로 다시 렌더링할 때 모든 편집 가능 텍스트 크기에 곱하는 값입니다." data-i18n-title="help.fontScale">?</span>
            </div>
            <input id="s_fontScale" type="number" min="0.5" max="1.5" step="0.01" />
          </div>
        </div>
        <div class="row">
          <div>
            <div class="label-row">
              <label for="s_hybridAllowedDelta" data-i18n="options.hybridAllowedDelta">Hybrid 허용 diff</label>
              <span class="help" title="원본 대비 픽셀 차이(diff)가 vector보다 이 값만큼 더 나빠도 hybrid를 선택합니다." data-i18n-title="help.hybridAllowedDelta">?</span>
            </div>
            <input id="s_hybridAllowedDelta" type="number" step="0.1" />
          </div>
          <div></div>
        </div>
        <div>
          <div class="label-row">
            <label for="s_layoutFailureMode" data-i18n="options.layoutFailureMode">Layout 실패 처리</label>
            <span class="help" title="재시도 후에도 layout JSON이 생성되지 않은 슬라이드 처리 방식." data-i18n-title="help.layoutFailureMode">?</span>
          </div>
          <select id="s_layoutFailureMode">
            <option value="image_fallback" data-i18n="options.imageFallback">이미지 폴백으로 계속</option>
            <option value="fail" data-i18n="options.failStop">작업 실패로 중단</option>
          </select>
        </div>
        <label class="checkline"><input id="s_sam3" type="checkbox" /> SAM3 bbox refinement</label>
        <label class="checkline"><input id="s_editableEmbeddedText" type="checkbox" /> <span data-i18n="options.editableEmbeddedText">큰 이미지 내부 텍스트 편집 가능 유지</span></label>
        <label class="checkline"><input id="s_keepIntermediates" type="checkbox" /> <span data-i18n="options.keepIntermediates">큰 중간 산출물 보존</span></label>
        <div>
          <div class="label-row">
            <label for="s_watermarkMode" data-i18n="options.watermarkMode">워터마크만 제거 모드</label>
            <span class="help" title="LLM/SAM 변환을 생략하고 입력 슬라이드 이미지 우측 하단의 NotebookLM 워터마크만 제거합니다. 텍스트 편집은 불가합니다.&#10;빠름: 가장자리 평균색으로 단순 페인트 (수초).&#10;디테일: LaMa 인페인팅으로 배경을 자연스럽게 복원 (슬라이드당 1~3초, 첫 호출 시 모델 다운로드)." data-i18n-title="help.watermarkMode">?</span>
          </div>
          <select id="s_watermarkMode">
            <option value="off" data-i18n="options.watermarkOff">사용 안 함 (워터마크제거포함 전체변환)</option>
            <option value="fast" data-i18n="options.watermarkFast">빠름 — 단순 페인트</option>
            <option value="detail" data-i18n="options.watermarkDetail">디테일 — LaMa 인페인팅</option>
          </select>
        </div>
        <div>
          <div class="label-row">
            <label for="s_reconstructionMode" data-i18n="options.reconstructionMode">재구성 모드</label>
            <span class="help" title="기본: vector + hybrid 자동 선택 (모든 입력). image_split: Codex Vision + image_gen + SAM2 누끼 (단일 이미지/deck, 슬라이드당 60-120초)." data-i18n-title="help.reconstructionMode">?</span>
          </div>
          <select id="s_reconstructionMode" onchange="toggleImageSplitOptions('s_')">
            <option value="auto" data-i18n="options.reconstructionAuto">기본 (vector + hybrid 자동 선택)</option>
            <option value="image_split" data-i18n="options.reconstructionImageSplit">image_split — Codex + SAM2 누끼</option>
          </select>
        </div>
        <div id="s_imageSplitGroup" class="image-split-group">
          <div class="image-split-header" data-i18n="options.imageSplitGroup">↳ image_split 세부 옵션</div>
          <div>
            <div class="label-row">
              <label for="s_textEraseMode" data-i18n="options.textEraseMode">텍스트 제거 방식</label>
              <span class="help" title="codex_imagegen: Codex CLI image_gen 으로 텍스트 제거 (~30-60초, 그라데이션 보존). solid: 주변색 ring fill (~1초, 단색 배경)." data-i18n-title="help.textEraseMode">?</span>
            </div>
            <select id="s_textEraseMode">
              <option value="codex_imagegen" data-i18n="options.eraseCodex">Codex image_gen (느림, 고품질)</option>
              <option value="solid" data-i18n="options.eraseSolid">주변색 fill (빠름)</option>
            </select>
          </div>
          <div>
            <div class="label-row">
              <label for="s_backgroundMode" data-i18n="options.backgroundMode">슬라이드 배경</label>
              <span class="help" title="white: 흰 캔버스 + alpha 객체 (default, 객체 분리 명확). transparent: 슬라이드 기본 배경. clean: Codex 텍스트 제거 이미지 통째 (시각 충실, 객체 이동 시 같은 모양 잔존)." data-i18n-title="help.backgroundMode">?</span>
            </div>
            <select id="s_backgroundMode">
              <option value="white" data-i18n="options.bgWhite">white — 흰 캔버스 (깔끔)</option>
              <option value="transparent" data-i18n="options.bgTransparent">transparent — 슬라이드 기본</option>
              <option value="clean" data-i18n="options.bgClean">clean — 텍스트 제거 이미지 통째 (시각 충실)</option>
            </select>
          </div>
          <label class="checkline"><input id="s_useNativeShapes" type="checkbox" /> <span data-i18n="options.useNativeShapes">PPT native shape 사용 (큰 박스/카드/화살표 직접 편집 가능)</span></label>
        </div>
      </div>

      <div class="hint" style="margin-top:14px;" data-i18n-html="sidebar.sessionHint">
        ⓘ 사이드바 변경은 이 세션에만 적용됩니다. 영구 저장은 우측 상단 <strong>⚙ 설정</strong>.
      </div>
    </aside>
    <main>
      <h1 style="font-size:22px;margin-bottom:16px;" data-i18n="jobs.title">작업 상태</h1>
      <div id="jobs" class="jobs">
        <div class="empty" data-i18n="jobs.empty">아직 실행한 작업이 없습니다.</div>
      </div>
    </main>
  </div>
  <div id="modalBackdrop" class="modal-backdrop" onclick="closeModal(event)" aria-hidden="true">
    <div class="modal" role="dialog" aria-modal="true" aria-labelledby="modalTitle">
      <div class="modal-head">
        <h2 id="modalTitle">상세</h2>
        <button class="secondary" onclick="closeModal()" aria-label="모달 닫기" data-i18n="common.close" data-i18n-aria-label="modal.closeAria">닫기</button>
      </div>
      <div id="modalBody" class="modal-body"></div>
    </div>
  </div>

  <!--
    영구 설정 영역 — 평소에는 hidden 상태로 body 끝에 보관되고,
    ⚙ 설정 모달이 열리면 이 element 통째로 모달 본문으로 이동했다가
    closeModal 시 다시 body 로 복귀한다. 그 결과 기존 ID(#provider, #model,
    #baseUrl, #apiKey, #timeout, ...) 는 항상 DOM 에 단 하나만 존재해 기존
    JS 로직을 그대로 재사용한다. 사이드바는 이와 별도로 #s_* prefix 의 세션
    한정 사본을 가진다.
  -->
  <div id="advancedSettingsHost" hidden>
    <div class="setting-box settings-tab-pane" data-tab="system" style="border-top:0;margin-top:0;padding-top:0;">
      <div id="systemCheck" class="system-box">
        <div class="hint" data-i18n="system.checking">시스템 의존성을 확인하는 중입니다.</div>
      </div>
    </div>

    <div class="setting-box settings-tab-pane" data-tab="provider" style="border-top:0;margin-top:0;padding-top:0;">
      <h2>LLM Provider</h2>
      <div class="info-box" data-i18n-html="provider.info">
        LLM은 슬라이드 이미지를 layout JSON으로 해석하고 큰 이미지 영역을 판단하는 데 사용됩니다.
        기본 흐름은 슬라이드당 약 2회 호출이며, retry 설정에 따라 추가 호출될 수 있습니다.
      </div>
      <div class="provider-controls">
        <div>
          <label for="provider">Provider</label>
          <select id="provider"></select>
        </div>
        <button id="addCustom" class="secondary" type="button" data-i18n="provider.addCustom">Custom 추가</button>
      </div>
      <div id="customProviderFields" class="custom-provider-fields" hidden>
        <label for="customName">Custom Name</label>
        <input id="customName" placeholder="예: 사내 프록시" data-i18n-placeholder="provider.customNamePlaceholder" />
        <div class="actions provider-field-actions">
          <button id="deleteCustom" class="secondary" type="button" data-i18n="provider.deleteCustom">Custom 삭제</button>
        </div>
      </div>
      <label for="baseUrl">Base URL</label>
      <input id="baseUrl" placeholder="http://localhost:8300/v1" />
      <label for="model">Model</label>
      <div class="api-key-row">
        <input id="model" list="modelOptions" placeholder="gpt-5.5 (로컬은 자동 조회)" data-i18n-placeholder="model.placeholder" autocomplete="off" />
        <button id="refreshModels" class="secondary" type="button" title="사용 가능한 모델 목록 다시 가져오기" data-i18n="provider.refreshModels" data-i18n-title="model.refreshTitle">↻ 모델</button>
      </div>
      <datalist id="modelOptions"></datalist>
      <div id="modelChips" class="model-chips" hidden></div>
      <div id="modelHint" class="hint"></div>
      <label for="apiKey">API Key</label>
      <div class="api-key-row">
        <input id="apiKey" type="password" placeholder="sk-... (Ollama 등 로컬은 비워두세요)" />
        <button id="testLlm" class="secondary" type="button" data-i18n="provider.test">테스트</button>
      </div>
      <div id="testResult" class="test-result hint" hidden></div>
      <div id="providerHint" class="hint"></div>
    </div>

    <div class="setting-box settings-tab-pane" data-tab="options" style="border-top:0;margin-top:0;padding-top:0;">
      <h2 data-i18n="options.title">변환 옵션</h2>
      <div class="row">
        <div>
          <label for="timeout">Timeout</label>
          <input id="timeout" type="number" min="60" step="30" />
        </div>
        <div>
          <div class="label-row">
            <label for="retries">Retries</label>
            <span class="help" title="layout JSON 생성이 실패했을 때 같은 슬라이드를 다시 호출하는 횟수." data-i18n-title="help.retries">?</span>
          </div>
          <input id="retries" type="number" min="0" step="1" />
        </div>
      </div>
      <div class="row">
        <div>
          <div class="label-row">
            <label for="llmParallel" data-i18n="options.llmParallel">LLM 병렬 수</label>
            <span class="help" title="여러 슬라이드의 LLM 분석을 동시에 몇 개까지 실행할지." data-i18n-title="help.llmParallel">?</span>
          </div>
          <input id="llmParallel" type="number" min="1" max="10" step="1" />
        </div>
        <div>
          <div class="label-row">
            <label for="fontScale" data-i18n="options.fontScale">폰트 배율</label>
            <span class="help" title="PPTX로 렌더링할 때 텍스트 크기 배율." data-i18n-title="help.fontScale">?</span>
          </div>
          <input id="fontScale" type="number" min="0.5" max="1.5" step="0.01" />
        </div>
      </div>
      <div class="row">
        <div>
          <div class="label-row">
            <label for="hybridAllowedDelta" data-i18n="options.hybridAllowedDelta">Hybrid 허용 diff</label>
            <span class="help" title="원본 대비 픽셀 diff 가 vector 보다 이 값만큼 더 나빠도 hybrid 선택." data-i18n-title="help.hybridAllowedDelta">?</span>
          </div>
          <input id="hybridAllowedDelta" type="number" step="0.1" />
        </div>
        <div></div>
      </div>
      <div>
        <div class="label-row">
          <label for="layoutFailureMode" data-i18n="options.layoutFailureMode">Layout 실패 처리</label>
          <span class="help" title="재시도 후에도 layout JSON이 없는 슬라이드 처리 방식." data-i18n-title="help.layoutFailureMode">?</span>
        </div>
        <select id="layoutFailureMode">
          <option value="image_fallback" data-i18n="options.imageFallback">이미지 폴백으로 계속</option>
          <option value="fail" data-i18n="options.failStop">작업 실패로 중단</option>
        </select>
      </div>
      <label class="checkline"><input id="sam3" type="checkbox" /> SAM3 bbox refinement</label>
      <label class="checkline"><input id="editableEmbeddedText" type="checkbox" /> <span data-i18n="options.editableEmbeddedText">큰 이미지 내부 텍스트 편집 가능 유지</span></label>
      <label class="checkline"><input id="keepIntermediates" type="checkbox" /> <span data-i18n="options.keepIntermediates">큰 중간 산출물 보존</span></label>
      <div>
        <div class="label-row">
          <label for="watermarkMode" data-i18n="options.watermarkMode">워터마크만 제거 모드</label>
          <span class="help" title="LLM/SAM 생략. 빠름: 단순 페인트. 디테일: LaMa 인페인팅(시간 더 걸림)." data-i18n-title="help.watermarkMode">?</span>
        </div>
        <select id="watermarkMode">
          <option value="off" data-i18n="options.watermarkOff">사용 안 함 (워터마크제거포함 전체변환)</option>
          <option value="fast" data-i18n="options.watermarkFast">빠름 — 단순 페인트</option>
          <option value="detail" data-i18n="options.watermarkDetail">디테일 — LaMa 인페인팅</option>
        </select>
      </div>
      <div>
        <div class="label-row">
          <label for="reconstructionMode" data-i18n="options.reconstructionMode">재구성 모드</label>
          <span class="help" title="기본: vector + hybrid 자동. image_split: Codex Vision + image_gen + SAM2 누끼." data-i18n-title="help.reconstructionMode">?</span>
        </div>
        <select id="reconstructionMode" onchange="toggleImageSplitOptions('')">
          <option value="auto" data-i18n="options.reconstructionAuto">기본 (vector + hybrid 자동 선택)</option>
          <option value="image_split" data-i18n="options.reconstructionImageSplit">image_split — Codex + SAM2 누끼</option>
        </select>
      </div>
      <div id="imageSplitGroup" class="image-split-group">
        <div class="image-split-header" data-i18n="options.imageSplitGroup">↳ image_split 세부 옵션</div>
        <div>
          <div class="label-row">
            <label for="textEraseMode" data-i18n="options.textEraseMode">텍스트 제거 방식</label>
            <span class="help" title="codex_imagegen: Codex CLI image_gen (느림 60s, 그라데이션 보존). solid: 주변색 fill (빠름 1s)." data-i18n-title="help.textEraseMode">?</span>
          </div>
          <select id="textEraseMode">
            <option value="codex_imagegen" data-i18n="options.eraseCodex">Codex image_gen (느림, 고품질)</option>
            <option value="solid" data-i18n="options.eraseSolid">주변색 fill (빠름)</option>
          </select>
        </div>
        <div>
          <div class="label-row">
            <label for="backgroundMode" data-i18n="options.backgroundMode">슬라이드 배경</label>
            <span class="help" title="white: 흰 캔버스 + 객체 + textbox (깔끔). transparent: 기본. clean: Codex 텍스트 제거 이미지 통째." data-i18n-title="help.backgroundMode">?</span>
          </div>
          <select id="backgroundMode">
            <option value="white" data-i18n="options.bgWhite">white — 흰 캔버스</option>
            <option value="transparent" data-i18n="options.bgTransparent">transparent</option>
            <option value="clean" data-i18n="options.bgClean">clean — 통째</option>
          </select>
        </div>
        <label class="checkline"><input id="useNativeShapes" type="checkbox" /> <span data-i18n="options.useNativeShapes">PPT native shape 사용 (큰 박스/카드/화살표 직접 편집 가능)</span></label>
      </div>
    </div>
  </div>

  <script>
    const $ = (id) => document.getElementById(id);
    let presets = null;
    let selectedFile = null;
    let pollTimer = null;
    let activeProvider = "local";
    let currentPage = 1;
    const pageSize = 5;
    const jobCache = new Map();
    const sseSources = new Map();
    let lastJobIds = [];
    // 클라이언트사이드 pptx-preview 뷰어 상태. compare 모드는 두 개 인스턴스.
    // 참조: aitechdatamgr/frontend/src/components/video/PptxPreviewModal.tsx
    // compare 모드: leftKind/rightKind로 두 측의 종류 지정 (원본/변환결과, vector/hybrid 등)
    const viewerState = {
      jobId: null,
      filename: "",
      kind: "result",         // 단일 모드 종류 또는 "compare"
      leftKind: "original",   // compare 모드 좌측
      rightKind: "result",    // compare 모드 우측 (페이지 카운터 기준)
      fullscreen: false,
      instances: [],   // [{ destroy(), _host }] (단일=1개, compare=2개)
      observers: [],   // MutationObserver
      pageInfos: [],   // [{ current, total }]
    };
    let _pptxPreviewModulePromise = null;

    function loadPptxPreview() {
      if (!_pptxPreviewModulePromise) {
        // AITechHub와 동일한 메이저 버전 (^1.0.7). 첫 사용 시 ~수백 KB 로드 후 브라우저 캐시.
        _pptxPreviewModulePromise = import("https://esm.sh/pptx-preview@1.0.7").catch(err => {
          _pptxPreviewModulePromise = null;
          throw err;
        });
      }
      return _pptxPreviewModulePromise;
    }

    const optionFields = ["timeout","retries","llmParallel","fontScale","hybridAllowedDelta","layoutFailureMode","sam3","editableEmbeddedText","keepIntermediates","watermarkMode","reconstructionMode","textEraseMode","backgroundMode","useNativeShapes"];
    const settingsKey = "starSlideSettings";
    const themeKey = "starSlideTheme";
    const languageKey = "starSlideLanguage";
    const customPrefix = "custom:";
    const settingsVersion = 4;
    let currentLang = "ko";

    const I18N = {
      ko: {
        "tagline": "NotebookLM에서 생성한 이미지 슬라이드를 <strong>편집 가능한 PPTX</strong>로 변환",
        "settings.button": "⚙ 설정",
        "settings.title": "설정",
        "theme.light": "라이트 모드",
        "theme.dark": "다크 모드",
        "upload.title": "파일 업로드",
        "upload.drop": "PPTX/PDF/이미지를 드래그하거나 클릭",
        "upload.fileHint": ".pptx · .pdf · .png · .jpg (이미지는 1-슬라이드 PPTX 로 변환)",
        "upload.start": "변환 시작",
        "upload.uploading": "업로드 중",
        "upload.progress": "업로드 {percent}%",
        "upload.done": "업로드 완료, 변환 대기 중",
        "provider.info": "LLM은 슬라이드 이미지를 layout JSON으로 해석하고 큰 이미지 영역을 판단하는 데 사용됩니다.<br>기본 흐름은 슬라이드당 약 2회 호출이며, retry 설정에 따라 추가 호출될 수 있습니다.",
        "provider.sidebarHint": "Base URL · API Key · Custom Provider 관리 · 시스템 상태는 우측 상단 <strong>⚙ 설정</strong>에서.",
        "provider.registeredHint": "등록한 OpenAI 호환 provider입니다. 이름, URL, 모델명, API key를 수정한 뒤 설정 저장을 누르면 목록에 반영됩니다.",
        "provider.hint.local": "star-cliproxy 같은 로컬 OpenAI 호환 프록시",
        "provider.hint.openai": "OpenAI API key와 vision 지원 모델명을 입력",
        "provider.hint.gemini": "Gemini OpenAI-compatible endpoint 또는 로컬 프록시 모델명 사용",
        "provider.hint.ollama": "로컬 Ollama. Base URL에 /v1 prefix 필수. API key는 비워두세요.",
        "provider.groupBuiltIn": "기본 Provider",
        "provider.groupCustom": "Custom Provider",
        "provider.addCustom": "Custom 추가",
        "provider.deleteCustom": "Custom 삭제",
        "provider.customNamePlaceholder": "예: 사내 프록시",
        "provider.test": "테스트",
        "provider.testing": "확인 중...",
        "provider.refreshModels": "↻ 모델",
        "model.placeholder": "gpt-5.5 (로컬은 자동 조회)",
        "model.refreshTitle": "사용 가능한 모델 목록 다시 가져오기",
        "options.title": "변환 옵션",
        "options.llmParallel": "LLM 병렬 수",
        "options.fontScale": "폰트 배율",
        "options.hybridAllowedDelta": "Hybrid 허용 diff",
        "options.layoutFailureMode": "Layout 실패 처리",
        "options.imageFallback": "이미지 폴백으로 계속",
        "options.failStop": "작업 실패로 중단",
        "options.editableEmbeddedText": "큰 이미지 내부 텍스트 편집 가능 유지",
        "options.keepIntermediates": "큰 중간 산출물 보존",
        "options.watermarkMode": "워터마크만 제거 모드",
        "options.watermarkOff": "사용 안 함 (워터마크제거포함 전체변환)",
        "options.watermarkFast": "빠름 — 단순 페인트",
        "options.watermarkDetail": "디테일 — LaMa 인페인팅",
        "options.reconstructionMode": "재구성 모드 (단일 이미지)",
        "options.reconstructionAuto": "기본 (vector + hybrid 자동 선택)",
        "options.reconstructionImageSplit": "image_split — Codex + SAM2 누끼 (단일 이미지)",
        "options.textEraseMode": "텍스트 제거 방식 (image_split)",
        "options.eraseCodex": "Codex image_gen (느림 60s, 고품질)",
        "options.eraseSolid": "주변색 fill (빠름 1s)",
        "options.backgroundMode": "슬라이드 배경 (image_split)",
        "options.bgWhite": "white — 흰 캔버스 (깔끔)",
        "options.bgTransparent": "transparent — 슬라이드 기본",
        "options.bgClean": "clean — 텍스트 제거 이미지 통째 (시각 충실)",
        "options.useNativeShapes": "PPT native shape 사용 (큰 박스/카드/화살표 직접 편집 가능)",
        "options.imageSplitGroup": "↳ image_split 세부 옵션 (재구성 모드가 image_split 일 때만 적용)",
        "help.reconstructionMode": "단일 이미지/deck 모두 적용. image_split: Codex Vision 으로 텍스트+위치 추출 → Codex image_gen 으로 텍스트 지운 배경 → SAM2 로 객체 누끼 → editable PPTX 재조합. 슬라이드당 60-120초.",
        "help.textEraseMode": "image_split 전용. codex_imagegen: 그라데이션/복잡한 배경에 우수 (~60s). solid: 단색/카드 배경에 빠르게 적합 (~1s).",
        "help.backgroundMode": "image_split 전용. white: 흰 캔버스 위에 alpha 객체+textbox (default, 객체 분리 명확). transparent: PowerPoint 기본 슬라이드 배경. clean: Codex 텍스트 제거 이미지를 통째 깔기 (시각 100% 충실하지만 객체 이동 시 같은 모양 잔존).",
        "help.retries": "layout JSON 생성이 실패했을 때 같은 슬라이드를 다시 호출하는 횟수입니다. 기본 2회이며, 네트워크/LLM 일시 실패를 흡수합니다.",
        "help.llmParallel": "여러 슬라이드의 LLM 분석 요청을 동시에 몇 개까지 실행할지 정합니다.",
        "help.fontScale": "PPTX로 다시 렌더링할 때 모든 편집 가능 텍스트 크기에 곱하는 값입니다.",
        "help.hybridAllowedDelta": "원본 대비 픽셀 차이(diff)가 vector보다 이 값만큼 더 나빠도 hybrid를 선택합니다.",
        "help.layoutFailureMode": "재시도 후에도 layout JSON이 생성되지 않은 슬라이드 처리 방식입니다.",
        "help.watermarkMode": "LLM/SAM 변환을 생략하고 입력 슬라이드 이미지 우측 하단의 NotebookLM 워터마크만 제거합니다. 빠름은 단순 페인트, 디테일은 LaMa 인페인팅입니다.",
        "sidebar.sessionHint": "ⓘ 사이드바 변경은 이 세션에만 적용됩니다. 영구 저장은 우측 상단 <strong>⚙ 설정</strong>.",
        "jobs.title": "작업 상태",
        "jobs.empty": "아직 실행한 작업이 없습니다.",
        "jobs.selectCurrentPage": "현재 페이지 전체 선택",
        "jobs.selectedCount": "선택 {selected}건 / 총 {total}건",
        "jobs.totalCount": "총 {total}건",
        "jobs.deleteSelected": "선택 삭제",
        "jobs.deleteAll": "전체 삭제",
        "jobs.deleteTitleAll": "전체 삭제",
        "jobs.deleteTitleSelected": "선택 항목 삭제",
        "jobs.deletePrompt": "<strong>{count}건</strong>의 작업을 처리합니다. 방식을 선택하세요.",
        "jobs.deleteTrash": "_trash 폴더로 이동",
        "jobs.deletePermanent": "영구 삭제",
        "jobs.deleteConfirm": "{count}건을 디스크에서 영구 삭제합니다. 진행하시겠습니까?",
        "jobs.processing": "처리 중...",
        "jobs.deleteDone": "완료 — 삭제 {deleted}, 스킵 {skipped}",
        "jobs.deleteFailed": "삭제 실패: {error}",
        "actions.downloadPptx": "PPTX 다운로드",
        "actions.preview": "미리보기",
        "actions.original": "원본 보기",
        "actions.compare": "비교 보기",
        "actions.report": "리포트 보기",
        "actions.layout": "Layout 보기",
        "actions.cancel": "취소",
        "actions.rerun": "다시 실행",
        "actions.delete": "삭제",
        "common.close": "닫기",
        "common.cancel": "취소",
        "common.details": "상세",
        "common.saveAll": "전체 저장",
        "common.saved": "✔ 저장되었습니다.",
        "common.saveFailed": "저장 실패: {error}",
        "common.systemStatus": "시스템 상태",
        "common.required": "필수",
        "common.optional": "옵션",
        "modal.closeAria": "모달 닫기",
        "status.done": "완료",
        "status.failed": "실패",
        "status.cancelled": "취소됨",
        "status.cancelling": "취소 중",
        "status.running": "진행 중",
        "status.queued": "대기 중",
        "viewer.original": "원본",
        "viewer.result": "변환 결과",
        "viewer.compare": "비교",
        "viewer.loading": "PPTX 미리보기 라이브러리 로딩 중...",
        "viewer.failed": "미리보기 실패: {error}",
        "viewer.slideCounter": "슬라이드 {current} / {total}",
        "viewer.slideUnknown": "슬라이드 - / -",
        "viewer.prev": "이전 (←)",
        "viewer.next": "다음 (→)",
        "viewer.fullscreen": "전체화면 (F)",
        "viewer.minimize": "축소 (F)",
        "viewer.close": "닫기 (Esc)",
        "report.title": "리포트",
        "report.loading": "리포트를 불러오는 중입니다.",
        "report.loadFailed": "리포트 로드 실패: {error}",
        "report.slides": "슬라이드",
        "report.finalAvgDiff": "최종 평균 diff",
        "report.hybridChosen": "Hybrid 선택",
        "report.vectorChosen": "Vector 선택",
        "report.interpretation": "해석",
        "report.interpretationText": "diff는 원본 이미지와 변환 PPTX 렌더 결과의 평균 픽셀 차이입니다. 낮을수록 원본과 가깝습니다. Hybrid는 큰 이미지 객체를 보존한 슬라이드, Vector는 텍스트/도형 중심으로 재구성한 슬라이드입니다.",
        "report.worstSlides": "주의가 필요한 슬라이드",
        "report.decisions": "슬라이드별 선택",
        "report.vectorDownload": "Vector 다운로드",
        "report.hybridDownload": "Hybrid 다운로드",
        "report.vectorPreview": "Vector 미리보기",
        "report.hybridPreview": "Hybrid 미리보기",
        "report.outputFile": "출력 파일",
        "report.noDiff": "diff 정보가 없습니다.",
        "report.noSelection": "선택 리포트가 없습니다.",
        "table.slide": "슬라이드",
        "table.diff": "diff",
        "table.objects": "객체",
        "table.images": "이미지",
        "table.selection": "선택",
        "layout.loading": "Layout JSON을 해석하는 중입니다.",
        "layout.loadFailed": "Layout JSON 로드 실패: {error}",
        "layout.totalObjects": "전체 객체",
        "layout.textObjects": "텍스트 객체",
        "layout.imageObjects": "이미지 객체",
        "layout.help": "raster group은 원본 이미지 덩어리로 보존한 도식/일러스트 영역입니다. punched text는 이미지 안의 원본 텍스트를 지우고 PPT 텍스트로 다시 얹은 영역 수입니다.",
        "layout.title": "제목",
        "layout.text": "텍스트",
        "layout.shape": "도형",
        "layout.download": "Layout JSON 다운로드",
        "system.checking": "시스템 의존성을 확인하는 중입니다.",
        "system.failed": "시스템 의존성 확인 실패: {error}",
        "system.empty": "시스템 의존성 정보가 없습니다.",
        "alerts.selectFile": "먼저 PPTX/PDF/이미지 파일을 선택하세요.",
        "alerts.uploadFailed": "업로드 실패: {error}",
        "alerts.networkUploadFailed": "업로드 실패: 네트워크 오류",
        "alerts.cancelConfirm": "이 작업을 취소하시겠습니까? 진행 단계가 끝난 직후 중단됩니다.",
        "alerts.cancelFailed": "취소 실패: {error}",
        "alerts.rerunFailed": "다시 실행 실패: {error}",
        "phase.queued": "대기 중",
        "phase.start": "작업 시작",
        "phase.extract": "슬라이드 이미지 추출 중",
        "phase.layout": "Vision LLM layout JSON 생성 중",
        "phase.vectorQa": "vector PPTX 렌더 QA 중",
        "phase.groups": "큰 이미지 그룹 탐지 중",
        "phase.sam3": "SAM3 이미지 그룹 보정 중",
        "phase.hybrid": "hybrid layout 생성 중",
        "phase.hybridQa": "hybrid PPTX 렌더 QA 중",
        "phase.select": "최종 layout 자동 선택 중",
        "phase.final": "최종 PPTX 생성 및 미리보기 렌더링 중",
        "phase.complete": "완료",
        "phase.failed": "실패",
        "phase.cancelled": "취소됨",
        "phase.cancelling": "취소 중..."
      },
      en: {
        "tagline": "Convert NotebookLM image slides into <strong>editable PPTX</strong>",
        "settings.button": "⚙ Settings",
        "settings.title": "Settings",
        "theme.light": "Light Mode",
        "theme.dark": "Dark Mode",
        "upload.title": "Upload",
        "upload.drop": "Drag or click to upload PPTX/PDF/Image",
        "upload.fileHint": ".pptx · .pdf · .png · .jpg (single images become a 1-slide PPTX)",
        "upload.start": "Start Conversion",
        "upload.uploading": "Uploading",
        "upload.progress": "Upload {percent}%",
        "upload.done": "Upload complete, waiting for conversion",
        "provider.info": "The LLM interprets slide images into layout JSON and identifies large image regions.<br>The default flow makes about two calls per slide, plus extra calls when retries are needed.",
        "provider.sidebarHint": "Manage Base URL, API Key, Custom Providers, and system status from <strong>⚙ Settings</strong> in the top right.",
        "provider.registeredHint": "Registered OpenAI-compatible provider. Edit the name, URL, model, and API key, then save settings to update the list.",
        "provider.hint.local": "Local OpenAI-compatible proxy such as Star-CliProxy.",
        "provider.hint.openai": "Enter an OpenAI API key and a vision-capable model name.",
        "provider.hint.gemini": "Use the Gemini OpenAI-compatible endpoint or a local proxy model name.",
        "provider.hint.ollama": "Local Ollama. Base URL must include the /v1 prefix. Leave API key empty.",
        "provider.groupBuiltIn": "Built-in Providers",
        "provider.groupCustom": "Custom Providers",
        "provider.addCustom": "Add Custom",
        "provider.deleteCustom": "Delete Custom",
        "provider.customNamePlaceholder": "e.g. Internal proxy",
        "provider.test": "Test",
        "provider.testing": "Checking...",
        "provider.refreshModels": "↻ Models",
        "model.placeholder": "gpt-5.5 (local models are auto-detected)",
        "model.refreshTitle": "Refresh available model list",
        "options.title": "Conversion Options",
        "options.llmParallel": "LLM Parallelism",
        "options.fontScale": "Font Scale",
        "options.hybridAllowedDelta": "Hybrid Allowed Diff",
        "options.layoutFailureMode": "Layout Failure Handling",
        "options.imageFallback": "Continue with image fallback",
        "options.failStop": "Stop as failed",
        "options.editableEmbeddedText": "Keep text inside large images editable",
        "options.keepIntermediates": "Keep large intermediate files",
        "options.watermarkMode": "Watermark-only removal",
        "options.watermarkOff": "Off (full conversion including watermark removal)",
        "options.watermarkFast": "Fast — simple paint",
        "options.watermarkDetail": "Detail — LaMa inpaint",
        "options.reconstructionMode": "Reconstruction mode (single image)",
        "options.reconstructionAuto": "Default (vector + hybrid auto-pick)",
        "options.reconstructionImageSplit": "image_split — Codex + SAM2 cutout (single image)",
        "options.textEraseMode": "Text-erase mode (image_split)",
        "options.eraseCodex": "Codex image_gen (slow ~60s, high quality)",
        "options.eraseSolid": "Surrounding-color fill (fast ~1s)",
        "options.backgroundMode": "Slide background (image_split)",
        "options.bgWhite": "white — clean canvas (default)",
        "options.bgTransparent": "transparent — slide default",
        "options.bgClean": "clean — text-erased image full backdrop (faithful)",
        "options.useNativeShapes": "Use PPT native shapes (containers/cards/arrows are directly editable)",
        "options.imageSplitGroup": "↳ image_split sub-options (only applied when reconstruction mode is image_split)",
        "help.reconstructionMode": "Single image or deck input. image_split: Codex Vision extracts text+positions, Codex image_gen erases text into a clean background, SAM2 masks each object → editable PPTX. Takes ~60-120s per slide.",
        "help.textEraseMode": "image_split only. codex_imagegen: best for gradients/complex backgrounds (~60s). solid: faster for flat/card backgrounds (~1s).",
        "help.backgroundMode": "image_split only. white: white canvas + alpha objects + textbox (default, object separation is clearest). transparent: PowerPoint default slide background. clean: Codex's text-erased image as full backdrop (100%% visually faithful but moving an object leaves the same shape underneath).",
        "help.retries": "How many times to retry the same slide when layout JSON generation fails. Default is 2 to absorb transient network or LLM failures.",
        "help.llmParallel": "Maximum number of slide analysis requests to run in parallel.",
        "help.fontScale": "Multiplier applied to editable text sizes when rendering the PPTX.",
        "help.hybridAllowedDelta": "Choose hybrid even if its pixel diff is this much worse than vector. Larger values prefer preserving image blocks.",
        "help.layoutFailureMode": "What to do when a slide still has no layout JSON after retries.",
        "help.watermarkMode": "Skip LLM/SAM conversion and only remove the NotebookLM watermark in the bottom-right. Fast uses simple paint; Detail uses LaMa inpainting.",
        "sidebar.sessionHint": "ⓘ Sidebar changes apply only to this session. Use <strong>⚙ Settings</strong> in the top right to save permanently.",
        "jobs.title": "Jobs",
        "jobs.empty": "No jobs have been run yet.",
        "jobs.selectCurrentPage": "Select current page",
        "jobs.selectedCount": "Selected {selected} / Total {total}",
        "jobs.totalCount": "Total {total}",
        "jobs.deleteSelected": "Delete Selected",
        "jobs.deleteAll": "Delete All",
        "jobs.deleteTitleAll": "Delete All",
        "jobs.deleteTitleSelected": "Delete Selected Items",
        "jobs.deletePrompt": "Choose how to handle <strong>{count}</strong> job(s).",
        "jobs.deleteTrash": "Move to _trash",
        "jobs.deletePermanent": "Delete Permanently",
        "jobs.deleteConfirm": "Permanently delete {count} job(s) from disk?",
        "jobs.processing": "Processing...",
        "jobs.deleteDone": "Done — deleted {deleted}, skipped {skipped}",
        "jobs.deleteFailed": "Delete failed: {error}",
        "actions.downloadPptx": "Download PPTX",
        "actions.preview": "Preview",
        "actions.original": "Original",
        "actions.compare": "Compare",
        "actions.report": "Report",
        "actions.layout": "Layout",
        "actions.cancel": "Cancel",
        "actions.rerun": "Run Again",
        "actions.delete": "Delete",
        "common.close": "Close",
        "common.cancel": "Cancel",
        "common.details": "Details",
        "common.saveAll": "Save All",
        "common.saved": "✔ Saved.",
        "common.saveFailed": "Save failed: {error}",
        "common.systemStatus": "System Status",
        "common.required": "Required",
        "common.optional": "Optional",
        "modal.closeAria": "Close modal",
        "status.done": "Done",
        "status.failed": "Failed",
        "status.cancelled": "Cancelled",
        "status.cancelling": "Cancelling",
        "status.running": "Running",
        "status.queued": "Queued",
        "viewer.original": "Original",
        "viewer.result": "Converted",
        "viewer.compare": "Compare",
        "viewer.loading": "Loading PPTX preview library...",
        "viewer.failed": "Preview failed: {error}",
        "viewer.slideCounter": "Slide {current} / {total}",
        "viewer.slideUnknown": "Slide - / -",
        "viewer.prev": "Previous (←)",
        "viewer.next": "Next (→)",
        "viewer.fullscreen": "Fullscreen (F)",
        "viewer.minimize": "Exit Fullscreen (F)",
        "viewer.close": "Close (Esc)",
        "report.title": "Report",
        "report.loading": "Loading report.",
        "report.loadFailed": "Report load failed: {error}",
        "report.slides": "Slides",
        "report.finalAvgDiff": "Final Avg Diff",
        "report.hybridChosen": "Hybrid Chosen",
        "report.vectorChosen": "Vector Chosen",
        "report.interpretation": "Interpretation",
        "report.interpretationText": "diff is the average pixel difference between the original slide image and the rendered converted PPTX. Lower is closer to the original. Hybrid preserves large image objects; Vector reconstructs mostly with editable text and shapes.",
        "report.worstSlides": "Slides Needing Attention",
        "report.decisions": "Per-slide Selection",
        "report.vectorDownload": "Download Vector",
        "report.hybridDownload": "Download Hybrid",
        "report.vectorPreview": "Preview Vector",
        "report.hybridPreview": "Preview Hybrid",
        "report.outputFile": "Output file",
        "report.noDiff": "No diff information.",
        "report.noSelection": "No selection report.",
        "table.slide": "Slide",
        "table.diff": "diff",
        "table.objects": "Objects",
        "table.images": "Images",
        "table.selection": "Selection",
        "layout.loading": "Interpreting Layout JSON.",
        "layout.loadFailed": "Layout JSON load failed: {error}",
        "layout.totalObjects": "Total Objects",
        "layout.textObjects": "Text Objects",
        "layout.imageObjects": "Image Objects",
        "layout.help": "Raster groups are diagrams/illustrations preserved as original image blocks. Punched text counts text regions removed from an image and rebuilt as PPT text.",
        "layout.title": "Title",
        "layout.text": "Text",
        "layout.shape": "Shapes",
        "layout.download": "Download Layout JSON",
        "system.checking": "Checking system dependencies.",
        "system.failed": "System dependency check failed: {error}",
        "system.empty": "No system dependency information.",
        "alerts.selectFile": "Select a PPTX/PDF/image file first.",
        "alerts.uploadFailed": "Upload failed: {error}",
        "alerts.networkUploadFailed": "Upload failed: network error",
        "alerts.cancelConfirm": "Cancel this job? It will stop after the current step finishes.",
        "alerts.cancelFailed": "Cancel failed: {error}",
        "alerts.rerunFailed": "Run again failed: {error}",
        "phase.queued": "Queued",
        "phase.start": "Starting job",
        "phase.extract": "Extracting slide images",
        "phase.layout": "Generating Vision LLM layout JSON",
        "phase.vectorQa": "Rendering vector PPTX QA",
        "phase.groups": "Detecting large image groups",
        "phase.sam3": "Refining image groups with SAM3",
        "phase.hybrid": "Building hybrid layouts",
        "phase.hybridQa": "Rendering hybrid PPTX QA",
        "phase.select": "Selecting final layouts",
        "phase.final": "Generating final PPTX and preview",
        "phase.complete": "Complete",
        "phase.failed": "Failed",
        "phase.cancelled": "Cancelled",
        "phase.cancelling": "Cancelling..."
      }
    };

    function t(key, vars = {}) {
      let value = I18N[currentLang]?.[key] ?? I18N.ko[key] ?? key;
      for (const [name, replacement] of Object.entries(vars)) {
        value = value.replaceAll(`{${name}}`, String(replacement));
      }
      return value;
    }

    function applyLanguage(lang) {
      currentLang = lang === "en" ? "en" : "ko";
      localStorage.setItem(languageKey, currentLang);
      document.documentElement.lang = currentLang;
      document.querySelectorAll("[data-i18n]").forEach(el => { el.textContent = t(el.dataset.i18n); });
      document.querySelectorAll("[data-i18n-html]").forEach(el => { el.innerHTML = t(el.dataset.i18nHtml); });
      document.querySelectorAll("[data-i18n-title]").forEach(el => { el.title = t(el.dataset.i18nTitle); });
      document.querySelectorAll("[data-i18n-placeholder]").forEach(el => { el.placeholder = t(el.dataset.i18nPlaceholder); });
      document.querySelectorAll("[data-i18n-aria-label]").forEach(el => { el.setAttribute("aria-label", t(el.dataset.i18nAriaLabel)); });
      const langButton = $("languageToggle");
      if (langButton) langButton.textContent = currentLang === "ko" ? "EN" : "한글";
      applyTheme(document.body.dataset.theme || localStorage.getItem(themeKey) || "dark");
    }

    function toggleLanguage() {
      applyLanguage(currentLang === "ko" ? "en" : "ko");
      if (presets) {
        renderProviderOptions($("provider")?.value || activeProvider);
        applyProvider(false);
        loadSystemCheck();
      }
      if (selectedFile) selectFile(selectedFile);
      refreshJobs();
    }

    // === API key 암호화 keystore (WebCrypto AES-GCM + IndexedDB 비추출 키) ===
    const KEYSTORE_DB = "starSlideKeystore";
    const KEYSTORE_STORE = "keys";
    const KEYSTORE_KEY_NAME = "apiKeyMaster.v1";
    let _cryptoKeyPromise = null;
    let _settingsCache = null;
    // 진단: 한번이라도 decrypt가 실패했는가? true면 export(저장)를 막아서
    // 빈 값으로 ciphertext가 덮어써지는 데이터 손실을 차단한다.
    let _decryptFailed = false;

    function openKeystore() {
      return new Promise((resolve, reject) => {
        const req = indexedDB.open(KEYSTORE_DB, 1);
        req.onupgradeneeded = () => req.result.createObjectStore(KEYSTORE_STORE);
        req.onsuccess = () => resolve(req.result);
        req.onerror = () => reject(req.error);
      });
    }

    async function getOrCreateMasterKey() {
      const db = await openKeystore();
      const existing = await new Promise((resolve, reject) => {
        const tx = db.transaction(KEYSTORE_STORE, "readonly");
        const req = tx.objectStore(KEYSTORE_STORE).get(KEYSTORE_KEY_NAME);
        req.onsuccess = () => resolve(req.result || null);
        req.onerror = () => reject(req.error);
      });
      if (existing) { db.close(); return existing; }
      const key = await crypto.subtle.generateKey(
        { name: "AES-GCM", length: 256 },
        false,
        ["encrypt", "decrypt"]
      );
      await new Promise((resolve, reject) => {
        const tx = db.transaction(KEYSTORE_STORE, "readwrite");
        tx.objectStore(KEYSTORE_STORE).put(key, KEYSTORE_KEY_NAME);
        tx.oncomplete = () => resolve();
        tx.onerror = () => reject(tx.error);
      });
      db.close();
      return key;
    }

    function masterKey() {
      if (!_cryptoKeyPromise) _cryptoKeyPromise = getOrCreateMasterKey();
      return _cryptoKeyPromise;
    }

    function bytesToBase64(bytes) {
      let s = "";
      for (let i = 0; i < bytes.length; i++) s += String.fromCharCode(bytes[i]);
      return btoa(s);
    }

    function base64ToBytes(b64) {
      const s = atob(b64);
      const out = new Uint8Array(s.length);
      for (let i = 0; i < s.length; i++) out[i] = s.charCodeAt(i);
      return out;
    }

    async function encryptApiKey(plain) {
      if (!plain) return "";
      const key = await masterKey();
      const iv = crypto.getRandomValues(new Uint8Array(12));
      const ct = await crypto.subtle.encrypt({ name: "AES-GCM", iv }, key, new TextEncoder().encode(plain));
      const merged = new Uint8Array(iv.length + ct.byteLength);
      merged.set(iv, 0);
      merged.set(new Uint8Array(ct), iv.length);
      return bytesToBase64(merged);
    }

    async function decryptApiKey(enc) {
      if (!enc) return "";
      try {
        const key = await masterKey();
        const merged = base64ToBytes(enc);
        const iv = merged.slice(0, 12);
        const ct = merged.slice(12);
        const plain = await crypto.subtle.decrypt({ name: "AES-GCM", iv }, key, ct);
        return new TextDecoder().decode(plain);
      } catch (err) {
        console.warn("API key 복호화 실패 (키 회전 또는 저장소 손상)", err);
        _decryptFailed = true;
        return null;  // 센티넬: 호출자는 ciphertext를 그대로 보존해야 한다
      }
    }

    async function importSettingsFromStorage() {
      let raw;
      try { raw = JSON.parse(localStorage.getItem(settingsKey) || "{}"); } catch { raw = {}; }
      if (!raw || typeof raw !== "object") raw = {};
      raw.providers = raw.providers || {};
      for (const id of Object.keys(raw.providers)) {
        const p = raw.providers[id];
        if (!p || typeof p !== "object") continue;
        if (p.apiKeyEnc) {
          const plain = await decryptApiKey(p.apiKeyEnc);
          // null = decrypt 실패. apiKey는 미설정으로 두고 ciphertext는 보존한다.
          if (plain !== null) p.apiKey = plain;
        } else if (typeof p.apiKey !== "string") {
          p.apiKey = "";
        }
      }
      raw.customProviders = Array.isArray(raw.customProviders) ? raw.customProviders : [];
      for (const p of raw.customProviders) {
        if (!p || typeof p !== "object") continue;
        if (p.apiKeyEnc) {
          const plain = await decryptApiKey(p.apiKeyEnc);
          if (plain !== null) p.apiKey = plain;
        } else if (typeof p.apiKey !== "string") {
          p.apiKey = "";
        }
      }
      return raw;
    }

    async function exportSettingsToStorage(settings) {
      const out = JSON.parse(JSON.stringify(settings || {}));
      out.providers = out.providers || {};
      for (const id of Object.keys(out.providers)) {
        const p = out.providers[id];
        if (!p || typeof p !== "object") continue;
        // apiKey가 명시적 string이면 그것만 암호화. undefined(=decrypt 실패로 미설정)
        // 면 기존 apiKeyEnc를 그대로 두어 사용자의 ciphertext를 보존한다.
        if (typeof p.apiKey === "string") {
          p.apiKeyEnc = await encryptApiKey(p.apiKey);
        }
        delete p.apiKey;
      }
      out.customProviders = Array.isArray(out.customProviders) ? out.customProviders : [];
      for (const p of out.customProviders) {
        if (!p || typeof p !== "object") continue;
        if (typeof p.apiKey === "string") {
          p.apiKeyEnc = await encryptApiKey(p.apiKey);
        }
        delete p.apiKey;
      }
      out.settingsVersion = settingsVersion;
      localStorage.setItem(settingsKey, JSON.stringify(out));
    }

    function setKeystoreBanner(message) {
      const el = $("keystoreBanner");
      if (!el) return;
      if (!message) { el.hidden = true; el.textContent = ""; return; }
      el.hidden = false;
      el.textContent = message;
    }

    function hasLegacyPlaintextKey() {
      let raw;
      try { raw = JSON.parse(localStorage.getItem(settingsKey) || "{}"); } catch { return false; }
      if (!raw || typeof raw !== "object") return false;
      const providers = raw.providers || {};
      for (const id of Object.keys(providers)) {
        const p = providers[id];
        if (p && typeof p === "object" && p.apiKey && !p.apiKeyEnc) return true;
      }
      const customs = Array.isArray(raw.customProviders) ? raw.customProviders : [];
      for (const p of customs) {
        if (p && typeof p === "object" && p.apiKey && !p.apiKeyEnc) return true;
      }
      return false;
    }

    async function init() {
      currentLang = localStorage.getItem(languageKey) === "en" ? "en" : "ko";
      applyLanguage(currentLang);
      applyTheme(localStorage.getItem(themeKey) || "dark");
      presets = await (await fetch("/api/presets")).json();
      const needsMigration = hasLegacyPlaintextKey();
      _settingsCache = await importSettingsFromStorage();
      // decrypt 실패가 한 번이라도 있었다면 자동 마이그레이션을 건너뛴다 (export가
      // ciphertext를 빈값으로 덮어쓸 위험 차단). 사용자가 명시적으로 키를 다시 입력해
      // save하면 그때 정상 export 된다.
      if (needsMigration && !_decryptFailed) {
        try { await exportSettingsToStorage(_settingsCache); } catch (err) { console.warn("API key 마이그레이션 실패", err); }
      }
      if (_decryptFailed) {
        setKeystoreBanner("⚠ 저장된 API 키를 복호화할 수 없습니다 (브라우저 keystore 손상 또는 시크릿 모드). 기존 ciphertext는 보존되었으니 사용 중인 다른 브라우저에서 확인하거나, 새 키를 입력 후 저장하세요.");
      }
      loadSettings();
      $("provider").addEventListener("change", changeProviderSession);
      const sProv = document.getElementById("s_provider");
      if (sProv) sProv.addEventListener("change", onSidebarProviderChange);
      $("addCustom").addEventListener("click", addCustomProvider);
      $("deleteCustom").addEventListener("click", deleteCustomProvider);
      $("customName").addEventListener("change", renameSelectedCustomProvider);
      $("testLlm").addEventListener("click", testLlmSettings);
      $("refreshModels").addEventListener("click", () => fetchModelList(true));
      $("start").addEventListener("click", submit);
      $("themeToggle").addEventListener("click", toggleTheme);
      $("languageToggle").addEventListener("click", toggleLanguage);
      $("settingsButton").addEventListener("click", openSettingsModal);
      document.addEventListener("keydown", handleGlobalKey);
      setupDrop();
      await loadSystemCheck();
      await refreshJobs();
      pollTimer = setInterval(refreshJobs, 2500);
    }

    function handleGlobalKey(event) {
      if (!$("modalBackdrop").classList.contains("open")) return;
      if (event.key === "Escape") {
        event.preventDefault();
        closeModal();
        return;
      }
      if (viewerState.jobId == null) return;
      if (event.key === "ArrowLeft") {
        event.preventDefault();
        moveViewer(-1);
      } else if (event.key === "ArrowRight") {
        event.preventDefault();
        moveViewer(1);
      } else if (event.key === "f" || event.key === "F") {
        event.preventDefault();
        toggleViewerFullscreen();
      }
    }

    async function loadSystemCheck() {
      try {
        const payload = await (await fetch("/api/system-check")).json();
        $("systemCheck").innerHTML = renderSystemCheck(payload.items || []);
      } catch (error) {
        $("systemCheck").innerHTML = `<div class="hint" style="color:var(--danger);">${escapeHtml(t("system.failed", {error: error.message}))}</div>`;
      }
    }

    function renderSystemCheck(items) {
      if (!items.length) return `<div class="hint">${escapeHtml(t("system.empty"))}</div>`;
      return `
        <h2 style="margin-bottom:2px;">${escapeHtml(t("common.systemStatus"))}</h2>
        ${items.map(item => `
          <div class="system-row">
            <span class="system-name">${escapeHtml(item.label)}</span>
            <span class="status-pill ${item.ok ? "ok" : ""}">${item.ok ? "OK" : (item.required ? t("common.required") : t("common.optional"))}</span>
            <span class="hint">${escapeHtml(item.ok ? item.path : item.message)}</span>
          </div>
        `).join("")}
      `;
    }

    function applyTheme(theme) {
      const next = theme === "light" ? "light" : "dark";
      document.body.dataset.theme = next;
      localStorage.setItem(themeKey, next);
      const button = $("themeToggle");
      if (button) button.textContent = next === "dark" ? t("theme.light") : t("theme.dark");
    }

    function toggleTheme() {
      applyTheme(document.body.dataset.theme === "dark" ? "light" : "dark");
    }

    function loadSettings() {
      const saved = readStoredSettings();
      const selectedProvider = renderProviderOptions(saved.provider || "local");
      $("provider").value = selectedProvider;
      activeProvider = $("provider").value;
      applyProvider(false);
      const merged = {...presets.defaults, ...(saved.options || legacyOptions(saved))};
      for (const key of optionFields) {
        const el = $(key);
        if (!el) continue;
        if (el.type === "checkbox") el.checked = Boolean(merged[key]);
        else el.value = merged[key] ?? "";
      }
      // 영구 저장값을 사이드바(#s_*) 에도 채워 준다 — 새 세션의 기본값.
      applySidebarFromStorage(saved, merged, selectedProvider);
    }

    // === 사이드바(#s_*) 동기화 헬퍼 ===
    // 사이드바 input은 세션 한정이므로 영구 저장값에서 한 번만 채워주고,
    // 사용자가 사이드바에서 변경해도 localStorage 는 건드리지 않는다.
    const sidebarOptionMap = {
      timeout: "s_timeout",
      retries: "s_retries",
      llmParallel: "s_llmParallel",
      fontScale: "s_fontScale",
      hybridAllowedDelta: "s_hybridAllowedDelta",
      layoutFailureMode: "s_layoutFailureMode",
      sam3: "s_sam3",
      editableEmbeddedText: "s_editableEmbeddedText",
      keepIntermediates: "s_keepIntermediates",
      watermarkMode: "s_watermarkMode",
      reconstructionMode: "s_reconstructionMode",
      textEraseMode: "s_textEraseMode",
      backgroundMode: "s_backgroundMode",
      useNativeShapes: "s_useNativeShapes",
    };

    function applySidebarOptions(merged) {
      for (const [optKey, sidebarId] of Object.entries(sidebarOptionMap)) {
        const el = document.getElementById(sidebarId);
        if (!el) continue;
        if (el.type === "checkbox") el.checked = Boolean(merged[optKey]);
        else el.value = merged[optKey] ?? "";
      }
      // 재구성 모드에 따라 image_split sub-옵션 그룹 enable/disable
      toggleImageSplitOptions("s_");
    }

    // 재구성 모드 = image_split 일 때만 image_split 세부 옵션 그룹 활성화.
    // prefix='s_' 면 사이드바, '' 면 설정 모달.
    window.toggleImageSplitOptions = function(prefix) {
      const sel = document.getElementById(prefix + "reconstructionMode");
      const grp = document.getElementById(prefix + "imageSplitGroup");
      if (!sel || !grp) return;
      const active = sel.value === "image_split";
      grp.classList.toggle("disabled", !active);
      // 자식 입력들도 disabled 속성 토글 (form submission 영향은 없지만 시각 + 키보드)
      grp.querySelectorAll("select, input").forEach(el => {
        el.disabled = !active;
      });
    };

    function renderSidebarProviderOptions(saved, selected) {
      const builtIns = builtInProviders().map(provider =>
        `<option value="${provider.id}">${escapeHtml(provider.label)}</option>`
      ).join("");
      const customs = saved.customProviders.map(provider => {
        const label = provider.name || "Custom Provider";
        const detail = provider.baseUrl ? ` · ${provider.baseUrl}` : "";
        return `<option value="${customPrefix}${provider.id}">${escapeHtml(label + detail)}</option>`;
      }).join("");
      const sel = document.getElementById("s_provider");
      if (!sel) return;
      sel.innerHTML = `
        <optgroup label="${escapeHtml(t("provider.groupBuiltIn"))}">${builtIns}</optgroup>
        ${customs ? `<optgroup label="${escapeHtml(t("provider.groupCustom"))}">${customs}</optgroup>` : ""}
      `;
      sel.value = selected;
    }

    function applySidebarFromStorage(saved, merged, selectedProvider) {
      renderSidebarProviderOptions(saved, selectedProvider);
      const cfg = lookupProviderConfig(saved, selectedProvider);
      const sm = document.getElementById("s_model");
      if (sm) sm.value = cfg.model || "";
      applySidebarOptions(merged);
    }

    function lookupProviderConfig(saved, providerId) {
      if (isCustomProvider(providerId)) {
        const cp = findCustomProvider(saved, providerId);
        return cp ? { baseUrl: cp.baseUrl || "", apiKey: cp.apiKey || "", model: cp.model || "" }
                  : { baseUrl: "", apiKey: "", model: "" };
      }
      const stored = (saved.providers || {})[providerId] || {};
      const preset = builtInProviders().find(p => p.id === providerId) || {};
      return {
        baseUrl: stored.baseUrl ?? preset.baseUrl ?? "",
        apiKey: stored.apiKey ?? "",
        model: stored.model ?? preset.model ?? "",
      };
    }

    // 사이드바에서 #s_provider 변경 시 — 해당 provider의 영구 default를
    // 사이드바 model/변환옵션 자리에 채운다 (단, 영구 저장은 안 함).
    function onSidebarProviderChange() {
      const saved = readStoredSettings();
      const pid = document.getElementById("s_provider").value;
      const cfg = lookupProviderConfig(saved, pid);
      const sm = document.getElementById("s_model");
      if (sm) sm.value = cfg.model || "";
      // 변환옵션 default 는 provider 와 무관하므로 그대로 둠
    }

    function readStoredSettings() {
      const base = _settingsCache && typeof _settingsCache === "object"
        ? JSON.parse(JSON.stringify(_settingsCache))
        : {};
      const normalized = normalizeSettings(base);
      if (normalized.changed) {
        _settingsCache = JSON.parse(JSON.stringify(normalized.settings));
        exportSettingsToStorage(_settingsCache).catch(err => console.warn("settings export 실패", err));
      }
      return normalized.settings;
    }

    function normalizeSettings(raw) {
      const saved = raw && typeof raw === "object" ? raw : {};
      let changed = false;
        if (!saved.providers) {
          const provider = saved.provider || "local";
          saved.providers = {
            [provider]: {
              baseUrl: saved.baseUrl || "",
              model: saved.model || "",
              apiKey: saved.apiKey || "",
            },
          };
        changed = true;
      }
      if (!Array.isArray(saved.customProviders)) {
        saved.customProviders = [];
        changed = true;
      }
      if ((saved.provider === "custom" || saved.providers.custom) && !saved.customProviders.length) {
        const id = makeProviderId();
        const oldCustom = saved.providers.custom || {};
        saved.customProviders.push({
          id,
          name: "Custom Provider",
          baseUrl: oldCustom.baseUrl || saved.baseUrl || "",
          model: oldCustom.model || saved.model || "",
          apiKey: oldCustom.apiKey || saved.apiKey || "",
        });
        saved.provider = `${customPrefix}${id}`;
        delete saved.providers.custom;
        changed = true;
      }
      for (const item of saved.customProviders) {
        if (!item.id) {
          item.id = makeProviderId();
          changed = true;
        }
        if (!item.name) {
          item.name = "Custom Provider";
          changed = true;
        }
      }
      if (!saved.settingsVersion) {
        saved.options = {...(saved.options || {}), retries: 2, layoutFailureMode: "image_fallback"};
      }
      if (saved.settingsVersion !== settingsVersion) {
        saved.options = saved.options || {};
        if (saved.options.retries === undefined || Number(saved.options.retries) < 2) {
          saved.options.retries = 2;
        }
        if (!saved.options.layoutFailureMode) {
          saved.options.layoutFailureMode = "image_fallback";
        }
        saved.settingsVersion = settingsVersion;
        changed = true;
      }
      return {settings: saved, changed};
    }

    function makeProviderId() {
      if (window.crypto?.randomUUID) return window.crypto.randomUUID();
      return `custom_${Date.now()}_${Math.random().toString(16).slice(2)}`;
    }

    function writeStoredSettings(settings) {
      _settingsCache = JSON.parse(JSON.stringify(settings));
      exportSettingsToStorage(_settingsCache).catch(err => console.warn("settings export 실패", err));
    }

    function builtInProviders() {
      return presets.providers.filter(provider => provider.id !== "custom");
    }

    function providerHint(provider) {
      if (!provider) return "";
      return I18N[currentLang]?.[`provider.hint.${provider.id}`] || provider.hint || "";
    }

    function isCustomProvider(providerId) {
      return String(providerId || "").startsWith(customPrefix);
    }

    function customId(providerId) {
      return String(providerId || "").slice(customPrefix.length);
    }

    function findCustomProvider(settings, providerId) {
      if (!isCustomProvider(providerId)) return null;
      const id = customId(providerId);
      return settings.customProviders.find(item => item.id === id) || null;
    }

    function providerExists(settings, providerId) {
      if (isCustomProvider(providerId)) return Boolean(findCustomProvider(settings, providerId));
      return builtInProviders().some(provider => provider.id === providerId);
    }

    function renderProviderOptions(selectedProvider) {
      const saved = readStoredSettings();
      const selected = providerExists(saved, selectedProvider) ? selectedProvider : "local";
      const builtIns = builtInProviders().map(provider => `<option value="${provider.id}">${escapeHtml(provider.label)}</option>`).join("");
      const customs = saved.customProviders.map(provider => {
        const label = provider.name || "Custom Provider";
        const detail = provider.baseUrl ? ` · ${provider.baseUrl}` : "";
        return `<option value="${customPrefix}${provider.id}">${escapeHtml(label + detail)}</option>`;
      }).join("");
      $("provider").innerHTML = `
        <optgroup label="${escapeHtml(t("provider.groupBuiltIn"))}">${builtIns}</optgroup>
        ${customs ? `<optgroup label="${escapeHtml(t("provider.groupCustom"))}">${customs}</optgroup>` : ""}
      `;
      $("provider").value = selected;
      // 사이드바 select 도 같은 옵션·선택값으로 동기화
      renderSidebarProviderOptions(saved, selected);
      return selected;
    }

    function legacyOptions(saved) {
      const data = {};
      for (const key of optionFields) {
        if (saved[key] !== undefined) data[key] = saved[key];
      }
      return data;
    }

    function currentOptionSettings() {
      return {
        timeout: Number($("timeout").value || 600),
        retries: Number($("retries").value || 2),
        llmParallel: Number($("llmParallel").value || 5),
        fontScale: Number($("fontScale").value || 0.93),
        hybridAllowedDelta: Number($("hybridAllowedDelta").value || 0),
        layoutFailureMode: $("layoutFailureMode").value || "image_fallback",
        sam3: $("sam3").checked,
        editableEmbeddedText: $("editableEmbeddedText").checked,
        keepIntermediates: $("keepIntermediates").checked,
        watermarkMode: $("watermarkMode").value || "off",
        reconstructionMode: $("reconstructionMode").value || "auto",
        textEraseMode: $("textEraseMode").value || "codex_imagegen",
        backgroundMode: $("backgroundMode").value || "white",
        useNativeShapes: $("useNativeShapes").checked,
      };
    }

    function saveCurrentProvider(settings, providerId) {
      if (!providerId) return;
      if (isCustomProvider(providerId)) {
        const provider = findCustomProvider(settings, providerId);
        if (!provider) return;
        provider.name = $("customName").value.trim() || "Custom Provider";
        provider.baseUrl = $("baseUrl").value.trim();
        provider.model = $("model").value.trim();
        provider.apiKey = $("apiKey").value;
        return;
      }
      settings.providers = settings.providers || {};
      settings.providers[providerId] = {
        baseUrl: $("baseUrl").value.trim(),
        model: $("model").value.trim(),
        apiKey: $("apiKey").value,
      };
    }

    function saveSettings(showMessage = true) {
      const saved = readStoredSettings();
      saveCurrentProvider(saved, $("provider").value);
      saved.provider = $("provider").value;
      saved.options = currentOptionSettings();
      saved.settingsVersion = settingsVersion;
      writeStoredSettings(saved);
      renderProviderOptions(saved.provider);
      activeProvider = saved.provider;
      // 모달 [전체 저장] 시 → 사이드바도 새 default 로 sync
      const merged = {...presets.defaults, ...saved.options};
      applySidebarFromStorage(saved, merged, saved.provider);
      // 사이드바의 [설정 저장] 버튼은 ⚙ 모달로 이동했으므로 보조 토글은 제거.
      if (showMessage) {
        const btn = document.getElementById("save");
        if (btn) {
          btn.textContent = currentLang === "en" ? "Saved" : "저장됨";
          setTimeout(() => { btn.textContent = t("settings.button"); }, 900);
        }
      }
    }

    function changeProvider() {
      const saved = readStoredSettings();
      saveCurrentProvider(saved, activeProvider);
      saved.provider = $("provider").value;
      saved.options = currentOptionSettings();
      saved.settingsVersion = settingsVersion;
      writeStoredSettings(saved);
      activeProvider = $("provider").value;
      applyProvider(false);
    }

    // 세션 한정 provider 전환 — provider select 변경만 반영하고 영구 저장은 안 한다.
    // 영구 저장은 ⚙ 설정 모달의 [전체 저장] 버튼 또는 saveSettings()에서만.
    function changeProviderSession() {
      activeProvider = $("provider").value;
      applyProvider(false);
    }

    function openSettingsModal() {
      const html = `
        <div class="settings-tabs" role="tablist">
          <button type="button" class="settings-tab is-active" data-tab="system" role="tab">${escapeHtml(t("common.systemStatus"))}</button>
          <button type="button" class="settings-tab" data-tab="provider" role="tab">LLM Provider</button>
          <button type="button" class="settings-tab" data-tab="options" role="tab">${escapeHtml(t("options.title"))}</button>
        </div>
        <div id="advancedSlot"></div>
        <div class="job-actions" style="margin-top:14px;">
          <button id="settingsSave" class="primary">${escapeHtml(t("common.saveAll"))}</button>
          <button id="settingsClose" class="secondary">${escapeHtml(t("common.close"))}</button>
        </div>
        <div id="settingsSaveResult" class="hint" style="margin-top:10px;"></div>
      `;
      openModal(t("settings.button"), html);
      // settings 모달 폭 좁히기
      const modalEl = document.querySelector("#modalBackdrop .modal");
      if (modalEl) modalEl.classList.add("settings-modal");
      // 영구 설정 host 를 모달 본문 안으로 옮긴다 (DOM 이동, ID 유지)
      const host = document.getElementById("advancedSettingsHost");
      const slot = document.getElementById("advancedSlot");
      if (host && slot) {
        host.hidden = false;
        slot.appendChild(host);
        try { applyProvider(false); } catch (_) { /* noop */ }
      }
      // 모달 안의 image_split 그룹 enable/disable 동기화
      try { toggleImageSplitOptions(""); } catch (_) { /* noop */ }
      // 초기 탭: system
      switchSettingsTab("system");
      // 탭 클릭 핸들러
      document.querySelectorAll(".settings-tab").forEach(btn => {
        btn.addEventListener("click", () => switchSettingsTab(btn.dataset.tab));
      });
      const saveBtn = document.getElementById("settingsSave");
      const closeBtn = document.getElementById("settingsClose");
      const resultEl = document.getElementById("settingsSaveResult");
      if (saveBtn) {
        saveBtn.addEventListener("click", () => {
          try {
            saveSettings(false);
            if (resultEl) {
              resultEl.textContent = t("common.saved");
              resultEl.style.color = "var(--teal)";
            }
          } catch (err) {
            if (resultEl) {
              resultEl.textContent = t("common.saveFailed", {error: err && err.message ? err.message : err});
              resultEl.style.color = "var(--danger)";
            }
          }
        });
      }
      if (closeBtn) closeBtn.addEventListener("click", () => closeModal());
    }

    function switchSettingsTab(name) {
      document.querySelectorAll(".settings-tab").forEach(btn => {
        btn.classList.toggle("is-active", btn.dataset.tab === name);
      });
      document.querySelectorAll(".settings-tab-pane").forEach(pane => {
        pane.classList.toggle("is-hidden", pane.dataset.tab !== name);
      });
    }

    function applyProvider(overwrite = true) {
      // provider가 바뀌면 이전 테스트/모델 hint는 컨텍스트가 달라지므로 항상 클리어
      setTestResult(null);
      setModelHint("");
      clearModelOptions();

      const providerId = $("provider").value;
      const saved = readStoredSettings();
      const customProvider = findCustomProvider(saved, providerId);
      $("customProviderFields").hidden = !customProvider;
      if (customProvider) {
        $("customName").value = customProvider.name || "Custom Provider";
        $("baseUrl").value = customProvider.baseUrl || "";
        $("model").value = customProvider.model || "";
        $("apiKey").value = customProvider.apiKey || "";
        $("providerHint").textContent = t("provider.registeredHint");
        maybeAutoFetchModels();
        return;
      }

      const preset = builtInProviders().find(p => p.id === providerId);
      if (!preset) return;
      const providerSettings = (saved.providers || {})[providerId] || {};
      $("baseUrl").value = overwrite ? (preset.baseUrl || "") : (providerSettings.baseUrl ?? preset.baseUrl ?? "");
      $("model").value = overwrite ? (preset.model || "") : (providerSettings.model ?? preset.model ?? "");
      $("apiKey").value = overwrite ? "" : (providerSettings.apiKey ?? "");
      $("providerHint").textContent = providerHint(preset);
      maybeAutoFetchModels();
    }

    function addCustomProvider() {
      const saved = readStoredSettings();
      saveCurrentProvider(saved, activeProvider);
      const id = makeProviderId();
      const count = saved.customProviders.length + 1;
      saved.customProviders.push({
        id,
        name: `Custom Provider ${count}`,
        baseUrl: "",
        model: "",
        apiKey: "",
      });
      saved.provider = `${customPrefix}${id}`;
      saved.options = currentOptionSettings();
      saved.settingsVersion = settingsVersion;
      writeStoredSettings(saved);
      renderProviderOptions(saved.provider);
      activeProvider = saved.provider;
      applyProvider(false);
      $("customName").focus();
    }

    function deleteCustomProvider() {
      const providerId = $("provider").value;
      if (!isCustomProvider(providerId)) return;
      const saved = readStoredSettings();
      saved.customProviders = saved.customProviders.filter(item => item.id !== customId(providerId));
      saved.provider = "local";
      saved.options = currentOptionSettings();
      saved.settingsVersion = settingsVersion;
      writeStoredSettings(saved);
      renderProviderOptions(saved.provider);
      activeProvider = saved.provider;
      applyProvider(false);
    }

    function renameSelectedCustomProvider() {
      if (!isCustomProvider($("provider").value)) return;
      const saved = readStoredSettings();
      saveCurrentProvider(saved, $("provider").value);
      saved.provider = $("provider").value;
      saved.options = currentOptionSettings();
      saved.settingsVersion = settingsVersion;
      writeStoredSettings(saved);
      renderProviderOptions(saved.provider);
    }

    // 변환 시작 시 옵션 읽기 — 사이드바(#s_*) 가 권위 (세션 한정 변경 반영).
    // baseUrl/apiKey 는 사이드바에 없으므로 영구 저장값에서 lookup.
    function readOptions(includeProvider = false) {
      const $s = (id) => document.getElementById(id);
      const provider = $s("s_provider")?.value || $("provider").value || "local";
      const saved = readStoredSettings();
      const cfg = lookupProviderConfig(saved, provider);
      const v = (id, fallback) => {
        const el = $s(id);
        if (!el) return fallback;
        return el.value;
      };
      const checked = (id) => {
        const el = $s(id);
        return el ? el.checked : false;
      };
      const data = {
        baseUrl: cfg.baseUrl,
        model: ($s("s_model")?.value || cfg.model || "").trim(),
        apiKey: cfg.apiKey,
        timeout: Number(v("s_timeout", 600) || 600),
        retries: Number(v("s_retries", 2) || 2),
        llmParallel: Number(v("s_llmParallel", 5) || 5),
        fontScale: Number(v("s_fontScale", 0.93) || 0.93),
        hybridAllowedDelta: Number(v("s_hybridAllowedDelta", 0) || 0),
        layoutFailureMode: v("s_layoutFailureMode", "image_fallback") || "image_fallback",
        sam3: checked("s_sam3"),
        editableEmbeddedText: checked("s_editableEmbeddedText"),
        keepIntermediates: checked("s_keepIntermediates"),
        watermarkMode: v("s_watermarkMode", "off") || "off",
        reconstructionMode: v("s_reconstructionMode", "auto") || "auto",
        textEraseMode: v("s_textEraseMode", "codex_imagegen") || "codex_imagegen",
        backgroundMode: v("s_backgroundMode", "white") || "white",
        useNativeShapes: checked("s_useNativeShapes"),
      };
      if (includeProvider) data.provider = provider;
      return data;
    }

    function setupDrop() {
      const drop = $("drop");
      drop.addEventListener("click", () => $("file").click());
      $("file").addEventListener("change", () => selectFile($("file").files[0]));
      for (const eventName of ["dragenter", "dragover"]) {
        drop.addEventListener(eventName, (event) => { event.preventDefault(); drop.classList.add("drag"); });
      }
      for (const eventName of ["dragleave", "drop"]) {
        drop.addEventListener(eventName, (event) => { event.preventDefault(); drop.classList.remove("drag"); });
      }
      drop.addEventListener("drop", (event) => selectFile(event.dataTransfer.files[0]));
    }

    function selectFile(file) {
      selectedFile = file || null;
      $("fileLabel").textContent = selectedFile ? `${selectedFile.name} (${Math.round(selectedFile.size / 1024 / 1024 * 10) / 10} MB)` : t("upload.fileHint");
      // 새 파일을 고르면 이전 작업의 업로드 진행률/완료 메시지는 무관하므로 즉시 클리어
      setUploadProgress(false, 0, "");
    }

    let currentModelList = [];

    function renderModelChips(models) {
      currentModelList = Array.isArray(models) ? models : [];
      const root = $("modelChips");
      if (!currentModelList.length) { root.hidden = true; root.innerHTML = ""; return; }
      const current = $("model").value.trim();
      root.hidden = false;
      root.innerHTML = currentModelList.map(m => {
        const safe = escapeHtml(m);
        const active = m === current ? " active" : "";
        return `<button type="button" class="model-chip${active}" data-model="${safe}">${safe}</button>`;
      }).join("");
      root.querySelectorAll(".model-chip").forEach(btn => {
        btn.addEventListener("click", () => pickModel(btn.dataset.model));
      });
    }

    function pickModel(name) {
      $("model").value = name;
      renderModelChips(currentModelList);
    }

    function clearModelOptions() {
      $("modelOptions").innerHTML = "";
      $("modelChips").innerHTML = "";
      $("modelChips").hidden = true;
      currentModelList = [];
    }

    function isLocalHost(baseUrl) {
      if (!baseUrl) return false;
      try {
        const u = new URL(baseUrl);
        return ["localhost", "127.0.0.1", "0.0.0.0", "::1"].includes(u.hostname);
      } catch { return false; }
    }

    function setModelHint(message, state) {
      const root = $("modelHint");
      if (!message) { root.textContent = ""; root.style.color = ""; return; }
      root.textContent = message;
      root.style.color = state === "fail" ? "var(--danger)" : state === "ok" ? "var(--teal)" : "";
    }

    async function fetchModelList(manual) {
      const baseUrl = $("baseUrl").value.trim();
      const apiKey = $("apiKey").value;
      if (!baseUrl) {
        if (manual) setModelHint(currentLang === "en" ? "Enter Base URL first." : "Base URL을 먼저 입력하세요.", "fail");
        return;
      }
      const button = $("refreshModels");
      button.disabled = true;
      const original = button.textContent;
      button.textContent = currentLang === "en" ? "Loading..." : "조회 중...";
      setModelHint(currentLang === "en" ? "Loading model list..." : "모델 목록을 가져오는 중...");
      try {
        const response = await fetch("/api/list-models", {
          method: "POST",
          headers: {"Content-Type": "application/json"},
          body: JSON.stringify({baseUrl, apiKey, timeout: 10}),
        });
        const data = await response.json();
        const models = Array.isArray(data.models) ? data.models : [];
        const datalist = $("modelOptions");
        datalist.innerHTML = models.map(m => `<option value="${escapeHtml(m)}"></option>`).join("");
        renderModelChips(models);
        if (models.length === 0) {
          setModelHint(manual ? (currentLang === "en" ? "Could not load model list. The server may not support /v1/models. Enter it manually." : "모델 목록을 가져오지 못했습니다 (서버가 /v1/models를 미지원할 수 있음). 직접 입력하세요.") : "", "fail");
        } else {
          if (!$("model").value) {
            $("model").value = models[0];
            renderModelChips(models);
          }
          setModelHint(currentLang === "en" ? `${models.length} model(s) available — click a chip or type manually` : `사용 가능한 모델 ${models.length}개 — 칩을 클릭하거나 직접 입력하세요`, "ok");
        }
      } catch (error) {
        setModelHint(currentLang === "en" ? `Model list request failed: ${error.message}` : `모델 목록 요청 실패: ${error.message}`, "fail");
      } finally {
        button.disabled = false;
        button.textContent = original;
      }
    }

    function maybeAutoFetchModels() {
      // 로컬 호스트면 자동 페치, 외부 provider면 datalist/chips/hint 모두 클리어
      if (isLocalHost($("baseUrl").value.trim())) {
        fetchModelList(false);
      } else {
        clearModelOptions();
        setModelHint("");
      }
    }

    function setTestResult(state, message) {
      const root = $("testResult");
      if (!state) { root.hidden = true; root.textContent = ""; root.className = "test-result hint"; return; }
      root.hidden = false;
      root.className = `test-result hint ${state}`;
      root.textContent = message;
    }

    async function testLlmSettings() {
      const baseUrl = $("baseUrl").value.trim();
      const model = $("model").value.trim();
      const apiKey = $("apiKey").value;
      if (!baseUrl) { setTestResult("fail", currentLang === "en" ? "Enter Base URL." : "Base URL을 입력하세요."); return; }
      if (!model) { setTestResult("fail", currentLang === "en" ? "Enter a model name." : "Model 이름을 입력하세요."); return; }

      const button = $("testLlm");
      button.disabled = true;
      const original = button.textContent;
      button.textContent = t("provider.testing");
      const keyNote = apiKey ? "" : (currentLang === "en" ? " (no API key — Ollama/local proxy mode)" : " (API key 미설정 — Ollama/로컬 프록시 모드)");
      setTestResult("busy", currentLang === "en" ? `Testing ${baseUrl}...${keyNote}` : `${baseUrl} 로 테스트 호출 중...${keyNote}`);

      try {
        const response = await fetch("/api/test-llm", {
          method: "POST",
          headers: {"Content-Type": "application/json"},
          body: JSON.stringify({baseUrl, model, apiKey, timeout: 15}),
        });
        const data = await response.json();
        if (data.ok) {
          const sample = data.sample ? ` · "${data.sample}"` : "";
          const note = data.note ? ` · ${data.note}` : "";
          setTestResult("ok", currentLang === "en" ? `✓ Success · ${data.latency_ms}ms · ${data.model}${sample}${keyNote}${note}` : `✓ 성공 · ${data.latency_ms}ms · ${data.model}${sample}${keyNote}${note}`);
        } else {
          const lat = data.latency_ms !== undefined ? ` (${data.latency_ms}ms)` : "";
          setTestResult("fail", `✗ ${data.error || (currentLang === "en" ? "Unknown error" : "알 수 없는 오류")}${lat}`);
        }
      } catch (error) {
        setTestResult("fail", currentLang === "en" ? `✗ Request failed: ${error.message}` : `✗ 요청 실패: ${error.message}`);
      } finally {
        button.disabled = false;
        button.textContent = original;
      }
    }

    function setUploadProgress(visible, percent, label) {
      const root = $("uploadProgress");
      if (!visible) { root.hidden = true; return; }
      root.hidden = false;
      root.querySelector(".upload-progress-bar > div").style.width = `${Math.max(0, Math.min(100, percent))}%`;
      root.querySelector(".upload-progress-label").textContent = label || t("upload.progress", {percent: Math.round(percent)});
    }

    function submit() {
      if (!selectedFile) {
        alert(t("alerts.selectFile"));
        return;
      }
      // 변환 시작 시 자동 저장하지 않는다 — 사이드바 변경은 세션 한정.
      // 영구 저장은 우측 상단 ⚙ 설정 모달에서 명시적으로.
      const file = selectedFile;
      $("start").disabled = true;
      $("start").textContent = t("upload.uploading");
      setUploadProgress(true, 0, t("upload.progress", {percent: 0}));

      const xhr = new XMLHttpRequest();
      xhr.open("POST", `/api/jobs?filename=${encodeURIComponent(file.name)}`, true);
      xhr.setRequestHeader("x-star-slide-options", JSON.stringify(readOptions(false)));
      xhr.upload.onprogress = (event) => {
        if (event.lengthComputable) {
          const percent = (event.loaded / event.total) * 100;
          setUploadProgress(true, percent, t("upload.progress", {percent: Math.round(percent)}));
        }
      };
      xhr.onload = () => {
        $("start").disabled = false;
        $("start").textContent = t("upload.start");
        if (xhr.status >= 200 && xhr.status < 300) {
          setUploadProgress(true, 100, t("upload.done"));
          setTimeout(() => setUploadProgress(false, 0, ""), 1500);
          selectedFile = null;
          $("file").value = "";
          $("fileLabel").textContent = t("upload.fileHint");
          refreshJobs();
        } else {
          setUploadProgress(false, 0, "");
          alert(t("alerts.uploadFailed", {error: xhr.responseText || xhr.status}));
        }
      };
      xhr.onerror = () => {
        $("start").disabled = false;
        $("start").textContent = t("upload.start");
        setUploadProgress(false, 0, "");
        alert(t("alerts.networkUploadFailed"));
      };
      xhr.send(file);
    }

    // 사용자가 체크한 잡 ID 집합 (페이지 전환에도 유지)
    const selectedJobIds = new Set();

    async function refreshJobs() {
      const jobs = await (await fetch("/api/jobs")).json();
      const root = $("jobs");
      if (!jobs.length) {
        root.innerHTML = `<div class="empty">${escapeHtml(t("jobs.empty"))}</div>`;
        teardownAllSse();
        lastJobIds = [];
        selectedJobIds.clear();
        return;
      }
      jobs.forEach(job => jobCache.set(job.id, job));
      // 사라진 잡 id 는 선택 집합에서 정리
      const allIds = new Set(jobs.map(j => j.id));
      [...selectedJobIds].forEach(id => { if (!allIds.has(id)) selectedJobIds.delete(id); });

      const totalPages = Math.max(1, Math.ceil(jobs.length / pageSize));
      currentPage = Math.min(currentPage, totalPages);
      const start = (currentPage - 1) * pageSize;
      const pageJobs = jobs.slice(start, start + pageSize);
      const pageIds = pageJobs.map(job => job.id);
      const reuse = pageIds.length === lastJobIds.length && pageIds.every((id, i) => id === lastJobIds[i]);
      if (!reuse) {
        root.innerHTML =
          renderJobsToolbar(jobs, pageJobs) +
          pageJobs.map(renderJob).join("") +
          renderPager(jobs.length, totalPages);
        lastJobIds = pageIds;
        wireJobsToolbar();
      } else {
        pageJobs.forEach(updateJobCard);
        // 툴바의 카운트/체크박스 상태만 갱신
        const toolbarHtml = renderJobsToolbar(jobs, pageJobs);
        const oldToolbar = root.querySelector(".jobs-toolbar");
        if (oldToolbar) {
          oldToolbar.outerHTML = toolbarHtml;
          wireJobsToolbar();
        }
        // 페이저도 다시 그림 (총건수 변경 가능)
        const oldPager = root.querySelector(".pager");
        const pagerHtml = renderPager(jobs.length, totalPages);
        if (oldPager) {
          oldPager.outerHTML = pagerHtml;
        } else if (pagerHtml) {
          root.insertAdjacentHTML("beforeend", pagerHtml);
        }
      }
      pageJobs.forEach(maybeSubscribeSse);
      const visibleIds = new Set(pageIds);
      [...sseSources.keys()].forEach(id => {
        if (!visibleIds.has(id)) teardownSse(id);
      });
    }

    function renderJobsToolbar(allJobs, pageJobs) {
      const allChecked = pageJobs.length > 0 && pageJobs.every(j => selectedJobIds.has(j.id));
      const someChecked = pageJobs.some(j => selectedJobIds.has(j.id));
      const indet = !allChecked && someChecked ? " data-indet=\"1\"" : "";
      const selectedCount = selectedJobIds.size;
      return `
        <div class="jobs-toolbar">
          <label class="jobs-toolbar-check">
            <input type="checkbox" id="jobsSelectAll" ${allChecked ? "checked" : ""}${indet} />
            <span>${escapeHtml(t("jobs.selectCurrentPage"))}</span>
          </label>
          <span class="hint">${escapeHtml(t("jobs.selectedCount", {selected: selectedCount, total: allJobs.length}))}</span>
          <span style="flex:1;"></span>
          <button class="secondary" id="jobsDeleteSelected" ${selectedCount === 0 ? "disabled" : ""}>${escapeHtml(t("jobs.deleteSelected"))}</button>
          <button class="secondary" id="jobsDeleteAll" ${allJobs.length === 0 ? "disabled" : ""}>${escapeHtml(t("jobs.deleteAll"))}</button>
        </div>
      `;
    }

    function wireJobsToolbar() {
      const all = document.getElementById("jobsSelectAll");
      if (all) {
        // indeterminate 상태 표현 (HTML attr로 못 그리므로 JS로)
        if (all.dataset.indet === "1") all.indeterminate = true;
        all.addEventListener("change", () => {
          const pageIds = lastJobIds;
          if (all.checked) pageIds.forEach(id => selectedJobIds.add(id));
          else pageIds.forEach(id => selectedJobIds.delete(id));
          refreshJobs();
        });
      }
      const delSel = document.getElementById("jobsDeleteSelected");
      if (delSel) delSel.addEventListener("click", () => askDelete([...selectedJobIds]));
      const delAll = document.getElementById("jobsDeleteAll");
      if (delAll) delAll.addEventListener("click", () => askDelete([...jobCache.keys()], { allLabel: true }));
      // 카드 체크박스 위임 처리
      document.querySelectorAll(".job-check").forEach(cb => {
        cb.addEventListener("change", (event) => {
          const id = event.target.dataset.jobId;
          if (event.target.checked) selectedJobIds.add(id);
          else selectedJobIds.delete(id);
          refreshJobs();
        });
      });
    }

    function renderPager(total, totalPages) {
      if (totalPages <= 1) return "";
      // shadcn 류: < 1 ... 4 5 6 ... 10 >
      const pages = [];
      const range = 1;  // 현재 페이지 좌우로 보여줄 개수
      const includes = new Set([1, totalPages, currentPage]);
      for (let i = -range; i <= range; i += 1) includes.add(currentPage + i);
      const sorted = [...includes].filter(p => p >= 1 && p <= totalPages).sort((a, b) => a - b);
      let prev = 0;
      for (const p of sorted) {
        if (p - prev > 1) pages.push({ ellipsis: true });
        pages.push({ page: p });
        prev = p;
      }
      const buttons = pages.map(item => {
        if (item.ellipsis) return `<span class="pager-ellipsis">…</span>`;
        const active = item.page === currentPage ? " is-active" : "";
        return `<button class="pager-page${active}" onclick="goToPage(${item.page})">${item.page}</button>`;
      }).join("");
      return `
        <div class="pager">
          <span class="hint">${escapeHtml(t("jobs.totalCount", {total}))}</span>
          <span style="flex:1;"></span>
          <button class="pager-arrow" ${currentPage <= 1 ? "disabled" : ""} onclick="changePage(-1)">&lt;</button>
          ${buttons}
          <button class="pager-arrow" ${currentPage >= totalPages ? "disabled" : ""} onclick="changePage(1)">&gt;</button>
        </div>
      `;
    }

    function changePage(delta) {
      currentPage = Math.max(1, currentPage + delta);
      refreshJobs();
    }

    function goToPage(page) {
      currentPage = Math.max(1, page);
      refreshJobs();
    }

    async function askDelete(ids, opts = {}) {
      if (!ids.length) return;
      const html = `
        <div class="empty" style="margin-bottom:12px;">
          ${t("jobs.deletePrompt", {count: ids.length})}
        </div>
        <div class="job-actions">
          <button id="delTrash" class="secondary">${escapeHtml(t("jobs.deleteTrash"))}</button>
          <button id="delPerm" class="secondary" style="border-color:var(--danger);color:var(--danger);">${escapeHtml(t("jobs.deletePermanent"))}</button>
          <button id="delCancel" class="secondary">${escapeHtml(t("common.cancel"))}</button>
        </div>
        <div id="delResult" class="hint" style="margin-top:10px;"></div>
      `;
      openModal(opts.allLabel ? t("jobs.deleteTitleAll") : t("jobs.deleteTitleSelected"), html);
      const trashBtn = document.getElementById("delTrash");
      const permBtn = document.getElementById("delPerm");
      const cancelBtn = document.getElementById("delCancel");
      const resultEl = document.getElementById("delResult");
      const run = async (mode) => {
        if (resultEl) resultEl.textContent = t("jobs.processing");
        try {
          const response = await fetch("/api/jobs/delete", {
            method: "POST",
            headers: { "Content-Type": "application/json" },
            body: JSON.stringify({ ids, mode }),
          });
          if (!response.ok) throw new Error(await response.text());
          const data = await response.json();
          ids.forEach(id => selectedJobIds.delete(id));
          if (resultEl) resultEl.textContent = t("jobs.deleteDone", {deleted: data.deleted.length, skipped: data.skipped.length});
          await refreshJobs();
          setTimeout(() => closeModal(), 600);
        } catch (err) {
          if (resultEl) {
            resultEl.textContent = t("jobs.deleteFailed", {error: err && err.message ? err.message : err});
            resultEl.style.color = "var(--danger)";
          }
        }
      };
      if (trashBtn) trashBtn.addEventListener("click", () => run("trash"));
      if (permBtn) permBtn.addEventListener("click", () => {
        if (!confirm(t("jobs.deleteConfirm", {count: ids.length}))) return;
        run("permanent");
      });
      if (cancelBtn) cancelBtn.addEventListener("click", () => closeModal());
    }

    function statusLabel(status) {
      switch (status) {
        case "done": return t("status.done");
        case "failed": return t("status.failed");
        case "cancelled": return t("status.cancelled");
        case "cancelling": return t("status.cancelling");
        case "running": return t("status.running");
        case "queued": return t("status.queued");
        default: return status;
      }
    }

    function phaseLabel(phase) {
      const text = String(phase || "");
      const countMatch = text.match(/^(.*) \((\d+)\/(\d+)\)$/);
      const base = countMatch ? countMatch[1] : text;
      const map = {
        "대기 중": "phase.queued",
        "작업 시작": "phase.start",
        "슬라이드 이미지 추출 중": "phase.extract",
        "Vision LLM layout JSON 생성 중": "phase.layout",
        "vector PPTX 렌더 QA 중": "phase.vectorQa",
        "큰 이미지 그룹 탐지 중": "phase.groups",
        "SAM3 이미지 그룹 보정 중": "phase.sam3",
        "hybrid layout 생성 중": "phase.hybrid",
        "hybrid PPTX 렌더 QA 중": "phase.hybridQa",
        "최종 layout 자동 선택 중": "phase.select",
        "최종 PPTX 생성 및 미리보기 렌더링 중": "phase.final",
        "완료": "phase.complete",
        "실패": "phase.failed",
        "취소됨": "phase.cancelled",
        "취소 중...": "phase.cancelling",
      };
      const key = map[base];
      if (!key) return text;
      const translated = t(key);
      return countMatch ? `${translated} (${countMatch[2]}/${countMatch[3]})` : translated;
    }

    function jobActions(job) {
      const parts = [];
      const watermarkOnly = (job.options?.watermarkMode || "off") !== "off";
      if (job.status === "done") {
        // 사용자 요청 순서: PPTX 다운로드 | 미리보기 | 원본보기 | 비교 | 리포트 | Layout | 다시 실행
        parts.push(`<a class="button" href="/api/jobs/${job.id}/download">${escapeHtml(t("actions.downloadPptx"))}</a>`);
        parts.push(`<button class="secondary" onclick="openSlideViewer('${job.id}', 'result')">${escapeHtml(t("actions.preview"))}</button>`);
        parts.push(`<button class="secondary" onclick="openSlideViewer('${job.id}', 'original')">${escapeHtml(t("actions.original"))}</button>`);
        parts.push(`<button class="secondary" onclick="openCompareViewer('${job.id}')">${ICON_COMPARE_INLINE} ${escapeHtml(t("actions.compare"))}</button>`);
        // 워터마크만 제거 모드는 LLM/SAM 변환을 거치지 않아 리포트 데이터가 의미 없으므로 숨긴다.
        if (!watermarkOnly) {
          parts.push(`<button class="secondary" onclick="openReport('${job.id}')">${escapeHtml(t("actions.report"))}</button>`);
        }
        if (job.artifacts?.layout_json) {
          parts.push(`<button class="secondary" onclick="openLayoutSummary('${job.id}')">${escapeHtml(t("actions.layout"))}</button>`);
        }
      }
      if (job.status === "running" || job.status === "queued" || job.status === "cancelling") {
        parts.push(`<button class="secondary" onclick="cancelJob('${job.id}')">${escapeHtml(t("actions.cancel"))}</button>`);
      }
      if (job.status === "done" || job.status === "failed" || job.status === "cancelled") {
        parts.push(`<button class="secondary" onclick="rerunJob('${job.id}')">${escapeHtml(t("actions.rerun"))}</button>`);
        parts.push(`<button class="secondary" style="border-color:var(--danger);color:var(--danger);" onclick="askDelete(['${job.id}'])">${escapeHtml(t("actions.delete"))}</button>`);
      }
      if (!parts.length) return "";
      return `<div class="job-actions">${parts.join("")}</div>`;
    }

    function isJobActive(status) {
      return status === "running" || status === "queued" || status === "cancelling";
    }

    function renderJob(job) {
      const pct = Math.max(0, Math.min(100, job.progress || 0));
      const error = job.error ? `<div class="hint job-error" style="color:var(--danger);">${escapeHtml(job.error)}</div>` : "";
      const active = isJobActive(job.status);
      // 진행 중에만 진행률 바와 phase/percent 표시. 완료/실패/취소된 카드는 메타와 액션만 보여 군더더기를 줄인다.
      const progressBlock = active
        ? `
            <div class="hint job-phase">${escapeHtml(phaseLabel(job.phase || ""))}</div>
            <div class="bar"><div style="width:${pct}%"></div></div>
            <div class="job-meta">
              <span class="metric job-percent">${Math.round(pct)}%</span>
              <span class="metric">${new Date(job.created_at * 1000).toLocaleString()}</span>
              <span class="metric">${escapeHtml(job.options?.model || "")}</span>
            </div>`
        : `
            <div class="job-meta">
              <span class="metric">${new Date(job.created_at * 1000).toLocaleString()}</span>
              ${job.options?.model ? `<span class="metric">${escapeHtml(job.options.model)}</span>` : ""}
            </div>`;
      const thumbSrc = `/api/jobs/${job.id}/thumbnail?u=${job.updated_at || job.created_at || 0}`;
      const checked = selectedJobIds.has(job.id) ? "checked" : "";
      return `
        <section class="job" data-job-id="${job.id}" data-status="${job.status}">
          <input type="checkbox" class="job-check" data-job-id="${job.id}" ${checked} aria-label="${escapeHtml(t("table.selection"))}" />
          <img class="job-thumb" src="${thumbSrc}" alt="" loading="lazy" decoding="async" onerror="this.classList.add('is-missing')" />
          <div class="job-body">
            <div class="job-head">
              <div class="job-title">${escapeHtml(job.filename)}</div>
              <span class="badge ${job.status}">${statusLabel(job.status)}</span>
            </div>
            ${progressBlock}
            ${error}
            ${jobActions(job)}
          </div>
        </section>
      `;
    }

    function updateJobCard(job) {
      const card = document.querySelector(`section.job[data-job-id="${job.id}"]`);
      if (!card) return;
      // status 전환(active ↔ terminal)이 발생하면 progress 블록 모양 자체가 바뀌므로 통째로 다시 그린다.
      const prevStatus = card.dataset.status;
      const wasActive = isJobActive(prevStatus);
      const isActive = isJobActive(job.status);
      if (wasActive !== isActive) {
        card.outerHTML = renderJob(job);
        return;
      }
      card.dataset.status = job.status;
      const body = card.querySelector(".job-body") || card;
      const pct = Math.max(0, Math.min(100, job.progress || 0));
      const badge = card.querySelector(".badge");
      if (badge) {
        badge.className = `badge ${job.status}`;
        badge.textContent = statusLabel(job.status);
      }
      const phase = card.querySelector(".job-phase");
      if (phase) phase.textContent = phaseLabel(job.phase || "");
      const bar = card.querySelector(".bar > div");
      if (bar) bar.style.width = `${pct}%`;
      const percent = card.querySelector(".job-percent");
      if (percent) percent.textContent = `${Math.round(pct)}%`;
      const oldErr = card.querySelector(".job-error");
      if (oldErr) oldErr.remove();
      if (job.error) {
        const div = document.createElement("div");
        div.className = "hint job-error";
        div.style.color = "var(--danger)";
        div.textContent = job.error;
        const actions = body.querySelector(".job-actions");
        body.insertBefore(div, actions || null);
      }
      const oldActions = body.querySelector(".job-actions");
      const newHtml = jobActions(job);
      if (oldActions) oldActions.remove();
      if (newHtml) body.insertAdjacentHTML("beforeend", newHtml);
    }

    function maybeSubscribeSse(job) {
      if (job.status === "done" || job.status === "failed" || job.status === "cancelled") {
        teardownSse(job.id);
        return;
      }
      if (sseSources.has(job.id)) return;
      if (typeof EventSource === "undefined") return;
      const es = new EventSource(`/api/jobs/${job.id}/events`);
      sseSources.set(job.id, es);
      es.addEventListener("snapshot", (evt) => {
        try {
          const snap = JSON.parse(evt.data);
          jobCache.set(snap.id, snap);
          updateJobCard(snap);
          if (snap.status === "done" || snap.status === "failed" || snap.status === "cancelled") {
            teardownSse(snap.id);
            refreshJobs();
          }
        } catch {
          /* ignore malformed payload */
        }
      });
      es.onerror = () => teardownSse(job.id);
    }

    function teardownSse(jobId) {
      const es = sseSources.get(jobId);
      if (es) {
        es.close();
        sseSources.delete(jobId);
      }
    }

    function teardownAllSse() {
      sseSources.forEach(es => es.close());
      sseSources.clear();
    }

    async function cancelJob(jobId) {
      if (!confirm(t("alerts.cancelConfirm"))) return;
      try {
        const response = await fetch(`/api/jobs/${jobId}/cancel`, {
          method: "POST",
          headers: {"Content-Type": "application/json"},
          body: "{}",
        });
        if (!response.ok) throw new Error(await response.text());
        const snap = await response.json();
        jobCache.set(snap.id, snap);
        updateJobCard(snap);
      } catch (error) {
        alert(t("alerts.cancelFailed", {error: error.message}));
      }
    }

    async function rerunJob(jobId) {
      try {
        const response = await fetch(`/api/jobs/${jobId}/rerun`, {
          method: "POST",
          headers: {"Content-Type": "application/json"},
          body: "{}",
        });
        if (!response.ok) throw new Error(await response.text());
        currentPage = 1;
        await refreshJobs();
      } catch (error) {
        alert(t("alerts.rerunFailed", {error: error.message}));
      }
    }

    function openModal(title, html, options) {
      const modal = document.querySelector("#modalBackdrop .modal");
      const head = modal.querySelector(".modal-head");
      modal.classList.toggle("viewer-modal", Boolean(options?.viewer));
      head.style.display = options?.hideHeader ? "none" : "";
      $("modalTitle").textContent = title;
      $("modalBody").innerHTML = html;
      $("modalBackdrop").classList.add("open");
      $("modalBackdrop").setAttribute("aria-hidden", "false");
    }

    function closeModal(event) {
      if (event && event.target !== $("modalBackdrop")) return;
      const modal = document.querySelector("#modalBackdrop .modal");
      // 설정 모달의 영구 설정 host 를 다시 body 로 복귀시킨다.
      // (innerHTML="" 으로 통째로 날리기 전에 옮겨야 element 가 살아남는다)
      const host = document.getElementById("advancedSettingsHost");
      if (host && $("modalBody").contains(host)) {
        host.hidden = true;
        document.body.appendChild(host);
      }
      destroyViewerInstances();
      $("modalBackdrop").classList.remove("open");
      $("modalBackdrop").setAttribute("aria-hidden", "true");
      $("modalBody").innerHTML = "";
      modal.classList.remove("viewer-modal", "fullscreen", "settings-modal");
      modal.querySelector(".modal-head").style.display = "";
      viewerState.jobId = null;
      viewerState.kind = "result";
      viewerState.fullscreen = false;
    }

    function destroyViewerInstances() {
      for (const obs of viewerState.observers) {
        try { obs?.disconnect(); } catch { /* ignore */ }
      }
      for (const inst of viewerState.instances) {
        try { inst?.destroy?.(); } catch { /* ignore */ }
      }
      viewerState.instances = [];
      viewerState.observers = [];
      viewerState.pageInfos = [];
    }

    const ICON_COMPARE_INLINE = `<svg width="14" height="14" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round" style="vertical-align:-2px;margin-right:4px;"><rect x="3" y="5" width="7" height="14" rx="1"/><rect x="14" y="5" width="7" height="14" rx="1"/></svg>`;
    const ICON_PREV = `<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2.4" stroke-linecap="round" stroke-linejoin="round"><polyline points="15 18 9 12 15 6"/></svg>`;
    const ICON_NEXT = `<svg width="20" height="20" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2.4" stroke-linecap="round" stroke-linejoin="round"><polyline points="9 18 15 12 9 6"/></svg>`;
    const ICON_FULLSCREEN = `<svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M8 3H5a2 2 0 0 0-2 2v3"/><path d="M21 8V5a2 2 0 0 0-2-2h-3"/><path d="M3 16v3a2 2 0 0 0 2 2h3"/><path d="M16 21h3a2 2 0 0 0 2-2v-3"/></svg>`;
    const ICON_MINIMIZE = `<svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M8 3v3a2 2 0 0 1-2 2H3"/><path d="M21 8h-3a2 2 0 0 1-2-2V3"/><path d="M3 16h3a2 2 0 0 1 2 2v3"/><path d="M16 21v-3a2 2 0 0 1 2-2h3"/></svg>`;
    const ICON_DOWNLOAD = `<svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M21 15v4a2 2 0 0 1-2 2H5a2 2 0 0 1-2-2v-4"/><polyline points="7 10 12 15 17 10"/><line x1="12" y1="15" x2="12" y2="3"/></svg>`;
    const ICON_CLOSE = `<svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><line x1="18" y1="6" x2="6" y2="18"/><line x1="6" y1="6" x2="18" y2="18"/></svg>`;

    function kindLabel(kind) {
      switch (kind) {
        case "original": return t("viewer.original");
        case "result": return t("viewer.result");
        case "selected": return t("viewer.result");
        case "vector": return "Vector";
        case "hybrid": return "Hybrid";
        case "compare": return t("viewer.compare");
        default: return kind;
      }
    }

    async function openCompareViewer(jobId, leftKind, rightKind) {
      viewerState.leftKind = leftKind || "original";
      viewerState.rightKind = rightKind || "result";
      return openSlideViewer(jobId, "compare");
    }

    async function fetchPptxArrayBuffer(jobId, which) {
      const response = await fetch(`/api/jobs/${jobId}/pptx-file?which=${encodeURIComponent(which)}`);
      if (!response.ok) throw new Error(`${kindLabel(which)} PPTX: HTTP ${response.status}`);
      return response.arrayBuffer();
    }

    async function openSlideViewer(jobId, which) {
      const job = jobCache.get(jobId);
      const filename = job?.filename || "";
      destroyViewerInstances();
      viewerState.jobId = jobId;
      viewerState.filename = filename;
      viewerState.kind = which;
      viewerState.fullscreen = false;
      // 모달부터 띄우고 라이브러리 + PPTX 다운로드는 비동기.
      openModal(
        "",
        `<div class="viewer"><div class="viewer-stage"><div class="viewer-empty">${escapeHtml(t("viewer.loading"))}</div></div></div>`,
        { viewer: true, hideHeader: true }
      );
      try {
        const [{ init }, ...buffers] = await Promise.all([
          loadPptxPreview(),
          ...(which === "compare"
            ? [fetchPptxArrayBuffer(jobId, viewerState.leftKind), fetchPptxArrayBuffer(jobId, viewerState.rightKind)]
            : [fetchPptxArrayBuffer(jobId, which)]),
        ]);
        renderViewerShell();
        // shell이 DOM에 붙은 뒤 컨테이너 측정 → 라이브러리 init (AITechHub: requestAnimationFrame 패턴).
        await new Promise(r => requestAnimationFrame(() => r()));
        if (which === "compare") {
          const [origBuf, resultBuf] = buffers;
          await Promise.all([
            mountPptxInstance("compareOrigHost", origBuf, init, 0),
            mountPptxInstance("compareResultHost", resultBuf, init, 1),
          ]);
        } else {
          await mountPptxInstance("singleHost", buffers[0], init, 0);
        }
      } catch (error) {
        const msg = error?.message || String(error);
        $("modalBody").innerHTML = `<div class="viewer"><div class="viewer-stage"><div class="viewer-empty" style="color:var(--danger);">${escapeHtml(t("viewer.failed", {error: msg}))}</div></div></div>`;
      }
    }

    async function mountPptxInstance(hostId, arrayBuffer, init, slot) {
      const host = document.getElementById(hostId);
      if (!host) return;
      // hidden/측정전이면 0이 되므로 폴백 (AITechHub 동일).
      const cw = host.clientWidth || 960;
      const ch = host.clientHeight || 540;
      // 슬라이드는 16:9 — host 비율과 다르면 더 작은 쪽에 맞춰 wrapper 크기 결정.
      // 이렇게 하면 라이브러리 wrapper 자체가 16:9 가 되어 host 안의
      // letterbox 가 stage 배경(var(--surface))과 자연스럽게 섞인다.
      const target = 16 / 9;
      let w, h;
      if (cw / ch > target) {
        h = ch;
        w = Math.round(ch * target);
      } else {
        w = cw;
        h = Math.round(cw / target);
      }
      host.innerHTML = "";
      const previewer = init(host, { width: w, height: h, mode: "slide" });
      await previewer.preview(arrayBuffer);
      previewer._host = host;
      viewerState.instances[slot] = previewer;
      // 라이브러리가 그린 직후 nav/pagination 숨기고 페이지 정보 추출.
      setTimeout(() => {
        applyPptxZoom(host);
        hideNativeNavAndSync(host, slot);
      }, 50);
      // DOM 변경 시 재동기화.
      const observer = new MutationObserver(() => hideNativeNavAndSync(host, slot));
      observer.observe(host, { childList: true, subtree: true });
      viewerState.observers[slot] = observer;
    }

    function hideNativeNavAndSync(host, slot) {
      // 내장 prev/next 버튼 숨기기 (display:none !important).
      const navBtns = host.querySelectorAll(".pptx-preview-wrapper-next");
      navBtns.forEach(btn => btn.style.setProperty("display", "none", "important"));
      // 페이지네이션 텍스트 추출 후 숨김.
      const pagination = host.querySelector(".pptx-preview-wrapper-pagination");
      if (pagination) {
        const text = (pagination.innerText || "").trim();
        const m = text.match(/(\d+)\s*\/\s*(\d+)/);
        if (m) {
          viewerState.pageInfos[slot] = {
            current: parseInt(m[1], 10),
            total: parseInt(m[2], 10),
          };
          updateCounter();
        }
        pagination.style.setProperty("display", "none", "important");
      }
    }

    function clickNativeNav(host, direction) {
      // 라이브러리 내장 wrapper-next 0번=다음, 1번=이전 (AITechHub 참조 확인됨).
      const navBtns = host.querySelectorAll(".pptx-preview-wrapper-next");
      const idx = direction === "prev" ? 1 : 0;
      const btn = navBtns[idx];
      if (!btn) return;
      btn.style.setProperty("display", "block");
      btn.click();
      btn.style.setProperty("display", "none", "important");
    }

    function moveViewer(delta) {
      const dir = delta > 0 ? "next" : "prev";
      if (viewerState.kind === "compare") {
        const origHost = document.getElementById("compareOrigHost");
        const resultHost = document.getElementById("compareResultHost");
        if (origHost) clickNativeNav(origHost, dir);
        if (resultHost) clickNativeNav(resultHost, dir);
      } else {
        const host = document.getElementById("singleHost");
        if (host) clickNativeNav(host, dir);
      }
      // 페이지 변경은 MutationObserver가 자동 동기화.
      setTimeout(updateCounter, 60);
    }

    function updateCounter() {
      // 기준 슬롯: compare면 1번(변환 결과), 단일이면 0번.
      const slot = viewerState.kind === "compare" ? 1 : 0;
      const info = viewerState.pageInfos[slot];
      if (!info) return;
      const counter = document.querySelector(".viewer-title .counter");
      if (counter) counter.textContent = `${info.current} / ${info.total}`;
      const footCounter = document.querySelector(".viewer-foot span");
      if (footCounter) footCounter.textContent = t("viewer.slideCounter", {current: info.current, total: info.total});
    }

    function renderViewerShell() {
      const fullscreenIcon = viewerState.fullscreen ? ICON_MINIMIZE : ICON_FULLSCREEN;
      const fullscreenTitle = viewerState.fullscreen ? t("viewer.minimize") : t("viewer.fullscreen");
      const modal = document.querySelector("#modalBackdrop .modal");
      modal.classList.toggle("fullscreen", viewerState.fullscreen);
      const isCompare = viewerState.kind === "compare";
      const titleBlock = isCompare
        ? `<span class="kind-pill ${viewerState.leftKind === "original" ? "original" : ""}">${escapeHtml(kindLabel(viewerState.leftKind))}</span><span class="kind-pill ${viewerState.rightKind === "original" ? "original" : ""}">${escapeHtml(kindLabel(viewerState.rightKind))}</span>`
        : `<span class="kind-pill ${viewerState.kind === "original" ? "original" : ""}">${escapeHtml(kindLabel(viewerState.kind))}</span>`;
      const stage = isCompare
        ? `
            <button class="viewer-nav prev" onclick="moveViewer(-1)" title="${escapeHtml(t("viewer.prev"))}">${ICON_PREV}</button>
            <div class="compare-pane">
              <div class="pane-label">${escapeHtml(kindLabel(viewerState.leftKind))}</div>
              <div class="pane-img-wrap"><div id="compareOrigHost" class="pptx-host"></div></div>
            </div>
            <div class="compare-pane">
              <div class="pane-label result">${escapeHtml(kindLabel(viewerState.rightKind))}</div>
              <div class="pane-img-wrap"><div id="compareResultHost" class="pptx-host"></div></div>
            </div>
            <button class="viewer-nav next" onclick="moveViewer(1)" title="${escapeHtml(t("viewer.next"))}">${ICON_NEXT}</button>`
        : `
            <button class="viewer-nav prev" onclick="moveViewer(-1)" title="${escapeHtml(t("viewer.prev"))}">${ICON_PREV}</button>
            <div id="singleHost" class="pptx-host"></div>
            <button class="viewer-nav next" onclick="moveViewer(1)" title="${escapeHtml(t("viewer.next"))}">${ICON_NEXT}</button>`;
      $("modalBody").innerHTML = `
        <div class="viewer ${viewerState.fullscreen ? "fullscreen" : ""}">
          <div class="viewer-head">
            <div class="viewer-title">
              ${titleBlock}
              <span class="filename" title="${escapeHtml(viewerState.filename)}">${escapeHtml(viewerState.filename)}</span>
              <span class="counter">- / -</span>
            </div>
            <div class="viewer-actions">
              <a class="viewer-icon" href="/api/jobs/${viewerState.jobId}/pptx-file?which=${viewerState.kind === "compare" ? viewerState.rightKind : viewerState.kind}" download title="${escapeHtml(t("actions.downloadPptx"))}">${ICON_DOWNLOAD}</a>
              <button class="viewer-icon" onclick="toggleViewerFullscreen()" title="${fullscreenTitle}">${fullscreenIcon}</button>
              <button class="viewer-icon" onclick="closeModal()" title="${escapeHtml(t("viewer.close"))}">${ICON_CLOSE}</button>
            </div>
          </div>
          <div class="viewer-stage ${isCompare ? "compare" : ""}">${stage}</div>
          <div class="viewer-foot">
            <button onclick="moveViewer(-1)" title="${escapeHtml(t("viewer.prev"))}">${ICON_PREV}</button>
            <span>${escapeHtml(t("viewer.slideUnknown"))}</span>
            <button onclick="moveViewer(1)" title="${escapeHtml(t("viewer.next"))}">${ICON_NEXT}</button>
          </div>
        </div>
      `;
    }

    function toggleViewerFullscreen() {
      viewerState.fullscreen = !viewerState.fullscreen;
      const modal = document.querySelector("#modalBackdrop .modal");
      if (modal) modal.classList.toggle("fullscreen", viewerState.fullscreen);
      const viewer = document.querySelector(".viewer");
      if (viewer) viewer.classList.toggle("fullscreen", viewerState.fullscreen);
      // 컨테이너 사이즈 변경 → 라이브러리 zoom 재적용 (AITechHub 패턴).
      setTimeout(() => {
        for (const inst of viewerState.instances) {
          const host = inst?._host;
          if (host) applyPptxZoom(host);
        }
      }, 100);
    }

    function applyPptxZoom(host) {
      const wrapper = host.querySelector(".pptx-preview-wrapper");
      if (!wrapper) return;
      const baseW = parseInt(wrapper.style.width) || 960;
      const baseH = parseInt(wrapper.style.height) || 540;
      const cw = host.clientWidth;
      const ch = host.clientHeight;
      if (cw <= 0 || ch <= 0) return;
      wrapper.style.zoom = String(Math.min(cw / baseW, ch / baseH));
    }

    async function openReport(jobId) {
      openModal(t("report.title"), `<div class="empty">${escapeHtml(t("report.loading"))}</div>`);
      try {
        const report = await (await fetch(`/api/jobs/${jobId}/report-summary`)).json();
        $("modalBody").innerHTML = renderReport(report);
      } catch (error) {
        $("modalBody").innerHTML = `<div class="empty" style="color:var(--danger);">${escapeHtml(t("report.loadFailed", {error: error.message}))}</div>`;
      }
    }

    function renderReport(report) {
      const s = report.summary || {};
      const worst = report.worst_slides || [];
      const decisions = report.decisions || [];
      const files = report.files || {};
      return `
        <div class="hint">${escapeHtml(report.job?.filename || "")}</div>
        <div class="report-grid">
          <div class="report-card"><span class="hint">${escapeHtml(t("report.slides"))}</span><strong>${s.slide_count ?? "-"}</strong></div>
          <div class="report-card"><span class="hint">${escapeHtml(t("report.finalAvgDiff"))}</span><strong>${fmt(s.selected_avg_diff)}</strong></div>
          <div class="report-card"><span class="hint">${escapeHtml(t("report.hybridChosen"))}</span><strong>${s.chosen_hybrid_count ?? 0}</strong></div>
          <div class="report-card"><span class="hint">${escapeHtml(t("report.vectorChosen"))}</span><strong>${s.chosen_vector_count ?? 0}</strong></div>
        </div>
        <h2>${escapeHtml(t("report.interpretation"))}</h2>
        <p class="hint">
          ${escapeHtml(t("report.interpretationText"))}
        </p>
        <h2 style="margin-top:18px;">${escapeHtml(t("report.worstSlides"))}</h2>
        ${renderWorstTable(worst)}
        <h2 style="margin-top:18px;">${escapeHtml(t("report.decisions"))}</h2>
        ${renderDecisionTable(decisions)}
        <div class="job-actions">
          ${files.candidate_vector ? `<a class="button secondary" href="/api/jobs/${report.job.id}/artifact/vector">${escapeHtml(t("report.vectorDownload"))}</a>` : ""}
          ${files.candidate_hybrid ? `<a class="button secondary" href="/api/jobs/${report.job.id}/artifact/hybrid">${escapeHtml(t("report.hybridDownload"))}</a>` : ""}
          ${files.candidate_vector ? `<button class="secondary" onclick="openSlideViewer('${report.job.id}', 'vector')">${escapeHtml(t("report.vectorPreview"))}</button>` : ""}
          ${files.candidate_hybrid ? `<button class="secondary" onclick="openSlideViewer('${report.job.id}', 'hybrid')">${escapeHtml(t("report.hybridPreview"))}</button>` : ""}
          ${files.candidate_vector && files.candidate_hybrid ? `<button class="secondary" onclick="openCompareViewer('${report.job.id}', 'vector', 'hybrid')">${ICON_COMPARE_INLINE} ${escapeHtml(t("actions.compare"))}</button>` : ""}
        </div>
        <p class="hint">${escapeHtml(t("report.outputFile"))}: ${escapeHtml(files.output || "")} · ${formatBytes(files.output_bytes)}</p>
      `;
    }

    function renderWorstTable(items) {
      if (!items.length) return `<div class="empty">${escapeHtml(t("report.noDiff"))}</div>`;
      return `<table><thead><tr><th>${escapeHtml(t("table.slide"))}</th><th>${escapeHtml(t("table.diff"))}</th><th>${escapeHtml(t("table.objects"))}</th><th>${escapeHtml(t("table.images"))}</th></tr></thead><tbody>${
        items.map(item => `<tr><td>${item.slide_no}</td><td>${fmt(item.mean_abs_diff)}</td><td>${item.object_count ?? "-"}</td><td>${item.picture_count ?? "-"}</td></tr>`).join("")
      }</tbody></table>`;
    }

    function renderDecisionTable(items) {
      if (!items.length) return `<div class="empty">${escapeHtml(t("report.noSelection"))}</div>`;
      return `<table><thead><tr><th>${escapeHtml(t("table.slide"))}</th><th>${escapeHtml(t("table.selection"))}</th><th>Vector diff</th><th>Hybrid diff</th></tr></thead><tbody>${
        items.map(item => `<tr><td>${item.slide_no}</td><td>${escapeHtml(item.chosen || "")}</td><td>${fmt(item.vector_mean_abs_diff)}</td><td>${fmt(item.hybrid_mean_abs_diff)}</td></tr>`).join("")
      }</tbody></table>`;
    }

    async function openLayoutSummary(jobId) {
      openModal("Layout JSON", `<div class="empty">${escapeHtml(t("layout.loading"))}</div>`);
      try {
        const summary = await (await fetch(`/api/jobs/${jobId}/layout-summary`)).json();
        $("modalBody").innerHTML = renderLayoutSummary(summary);
      } catch (error) {
        $("modalBody").innerHTML = `<div class="empty" style="color:var(--danger);">${escapeHtml(t("layout.loadFailed", {error: error.message}))}</div>`;
      }
    }

    function renderLayoutSummary(summary) {
      const slides = summary.slides || [];
      const totalObjects = slides.reduce((sum, slide) => sum + Number(slide.object_count || 0), 0);
      const imageObjects = slides.reduce((sum, slide) => sum + Number(slide.type_counts?.image || 0), 0);
      const textObjects = slides.reduce((sum, slide) => sum + Number(slide.type_counts?.text || 0), 0);
      const rasterGroups = slides.reduce((sum, slide) => sum + Number(slide.raster_group_count || 0), 0);
      return `
        <div class="hint">${escapeHtml(summary.job?.filename || "")}</div>
        <div class="report-grid">
          <div class="report-card"><span class="hint">${escapeHtml(t("report.slides"))}</span><strong>${summary.slide_count ?? slides.length}</strong></div>
          <div class="report-card"><span class="hint">${escapeHtml(t("layout.totalObjects"))}</span><strong>${totalObjects}</strong></div>
          <div class="report-card"><span class="hint">${escapeHtml(t("layout.textObjects"))}</span><strong>${textObjects}</strong></div>
          <div class="report-card"><span class="hint">${escapeHtml(t("layout.imageObjects"))}</span><strong>${imageObjects}</strong></div>
        </div>
        <p class="hint">${escapeHtml(t("layout.help"))}</p>
        <table>
          <thead><tr><th>${escapeHtml(t("table.slide"))}</th><th>${escapeHtml(t("layout.title"))}</th><th>${escapeHtml(t("table.objects"))}</th><th>${escapeHtml(t("layout.text"))}</th><th>${escapeHtml(t("table.images"))}</th><th>${escapeHtml(t("layout.shape"))}</th><th>Raster group</th><th>Punched text</th></tr></thead>
          <tbody>
            ${slides.map(slide => `<tr>
              <td>${slide.slide_no}</td>
              <td>${escapeHtml(slide.title || slide.id || "")}</td>
              <td>${slide.object_count ?? 0}</td>
              <td>${slide.type_counts?.text ?? 0}</td>
              <td>${slide.type_counts?.image ?? 0}</td>
              <td>${slide.type_counts?.shape ?? 0}</td>
              <td>${slide.raster_group_count ?? 0}</td>
              <td>${slide.punched_text_regions ?? "-"}</td>
            </tr>`).join("")}
          </tbody>
        </table>
        <div class="job-actions">
          <a class="button secondary" href="/api/jobs/${summary.job.id}/artifact/layout-json">${escapeHtml(t("layout.download"))}</a>
        </div>
      `;
    }

    function fmt(value) {
      return typeof value === "number" ? value.toFixed(2) : "-";
    }

    function formatBytes(value) {
      if (typeof value !== "number") return "-";
      if (value > 1024 * 1024) return `${(value / 1024 / 1024).toFixed(1)} MB`;
      if (value > 1024) return `${(value / 1024).toFixed(1)} KB`;
      return `${value} B`;
    }

    function escapeHtml(value) {
      return String(value).replace(/[&<>"']/g, (ch) => ({
        "&": "&amp;", "<": "&lt;", ">": "&gt;", '"': "&quot;", "'": "&#039;"
      }[ch]));
    }

    init();
  </script>
</body>
</html>
"""
