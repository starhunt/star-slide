"""PPTX 슬라이드 추출 (PRD §6.1 FR-001).

NotebookLM/Gamma 같은 이미지-잠금 PPTX는 슬라이드당 단일 PICTURE shape이지만,
일반 PPTX는 텍스트박스/도형이 혼재. 본 모듈은 둘 다 지원하되 1차 타깃은 이미지-잠금.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass(frozen=True)
class PptxSlideInfo:
    """PPTX 단일 슬라이드 메타."""

    page_no: int
    width_emu: int
    height_emu: int
    has_picture: bool
    has_text: bool
    text_chars: int
    embedded_image_path: Path | None = None  # 임베드 이미지 추출 경로 (있는 경우)


def inspect_pptx(path: Path) -> tuple[int, int, list[PptxSlideInfo]]:
    """PPTX 파일 구조 분석 (이미지-잠금 여부 판단용).

    Returns:
        (slide_width_emu, slide_height_emu, [PptxSlideInfo, ...])
    """
    prs = Presentation(str(path))
    slide_w = int(prs.slide_width or 0)
    slide_h = int(prs.slide_height or 0)

    slides: list[PptxSlideInfo] = []
    for i, slide in enumerate(prs.slides, start=1):
        has_pic = False
        text_chars = 0

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_pic = True
            if shape.has_text_frame:
                text_chars += len(shape.text_frame.text)

        slides.append(
            PptxSlideInfo(
                page_no=i,
                width_emu=slide_w,
                height_emu=slide_h,
                has_picture=has_pic,
                has_text=text_chars > 0,
                text_chars=text_chars,
            )
        )

    return slide_w, slide_h, slides


def extract_embedded_images(path: Path, out_dir: Path) -> list[Path]:
    """각 슬라이드의 첫 PICTURE shape를 PNG/JPG로 추출.

    이미지-잠금 PPTX(NotebookLM 패턴)에서 가장 빠른 추출 경로.
    슬라이드 마스터/배경 합성이 필요한 경우는 rasterize 모듈 사용 권장.
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(path))

    extracted: list[Path] = []
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext or "png"
                out_path = out_dir / f"slide_{i:03d}.{ext}"
                out_path.write_bytes(image.blob)
                extracted.append(out_path)
                break  # 슬라이드당 첫 PICTURE만

    return extracted


def extract_pdf_pages(path: Path, out_dir: Path, dpi: int = 192) -> list[Path]:
    """PDF 각 페이지를 slide_001.png 형태로 렌더링한다.

    NotebookLM PDF export는 PPTX 이미지-잠금 슬라이드와 동일하게 페이지 단위
    원본 이미지로 취급한 뒤, 이후 layout JSON 재구성 파이프라인을 그대로 탄다.
    """
    try:
        from pdf2image import convert_from_path
        from pdf2image.exceptions import PDFInfoNotInstalledError, PDFPageCountError
    except ImportError as exc:  # pragma: no cover - optional runtime dependency
        raise RuntimeError("PDF 입력에는 pdf2image 의존성이 필요합니다.") from exc

    out_dir.mkdir(parents=True, exist_ok=True)
    try:
        pages = convert_from_path(str(path), dpi=dpi)
    except PDFInfoNotInstalledError as exc:  # pragma: no cover - depends on host tools
        raise RuntimeError("PDF 입력에는 poppler 설치가 필요합니다. macOS에서는 `brew install poppler`를 실행하세요.") from exc
    except PDFPageCountError as exc:  # pragma: no cover - depends on input file
        raise RuntimeError(f"PDF 페이지 수를 읽을 수 없습니다: {path}") from exc
    extracted: list[Path] = []
    for i, page in enumerate(pages, start=1):
        out_path = out_dir / f"slide_{i:03d}.png"
        page.convert("RGB").save(out_path)
        extracted.append(out_path)
    return extracted


def is_image_locked(path: Path, threshold_ratio: float = 0.9) -> bool:
    """슬라이드의 90%+ 가 이미지-잠금이면 True.

    NotebookLM/Gamma 출력 1차 타깃 패턴 감지.
    """
    _, _, slides = inspect_pptx(path)
    if not slides:
        return False
    n_locked = sum(1 for s in slides if s.has_picture and s.text_chars < 10)
    return n_locked / len(slides) >= threshold_ratio
