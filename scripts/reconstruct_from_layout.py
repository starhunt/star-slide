"""Reconstruct image-locked slides from per-slide layout JSON files.

This is the JSON-driven version of the one-off public AX reconstruction:

  image + layout JSON -> clean background -> editable PPT objects -> PPTX

The schema is intentionally small and deterministic so one slide can be fixed
without disturbing another. Unsupported slide types can still fall back to a
single picture while higher-quality layouts are added one-by-one.
"""

from __future__ import annotations

import argparse
import json
from itertools import pairwise
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

DEFAULT_FONT = "Malgun Gothic"

_DASH_ALIASES: dict[str, str] = {
    "DOT": "ROUND_DOT",
    "DOTTED": "ROUND_DOT",
    "DOTS": "ROUND_DOT",
    "ROUND": "ROUND_DOT",
    "SQUARE": "SQUARE_DOT",
    "DASHED": "DASH",
    "LONGDASH": "LONG_DASH",
    "DASHDOT": "DASH_DOT",
    "DASHDOTDOT": "DASH_DOT_DOT",
}


def _resolve_dash_style(value: Any) -> Any:
    if not isinstance(value, str):
        return None
    key = value.strip().upper().replace("-", "_").replace(" ", "_")
    if not key or key == "NONE" or key == "SOLID":
        return None
    candidate = _DASH_ALIASES.get(key, key)
    return getattr(MSO_LINE_DASH_STYLE, candidate, None)


def rgb(value: str) -> RGBColor:
    cleaned = value.strip("#")
    if len(cleaned) == 8:
        cleaned = cleaned[:6]
    return RGBColor.from_string(cleaned)


def is_no_color(value: Any) -> bool:
    if value is None:
        return True
    if not isinstance(value, str):
        return False
    cleaned = value.strip().lower()
    return cleaned in {"none", "transparent", "null"} or (
        cleaned.startswith("#") and len(cleaned) == 9 and cleaned[-2:] == "00"
    )


class LayoutRenderer:
    def __init__(
        self, layout: dict[str, Any], image_path: Path, workdir: Path, *, font_scale: float = 1.0
    ):
        self.layout = layout
        self.image_path = image_path
        self.workdir = workdir
        self.font_scale = font_scale
        canvas = layout["canvas"]
        self.w_px = int(canvas["width"])
        self.h_px = int(canvas["height"])
        slide_size = layout.get("slide_size_emu") or [16_256_000, 9_144_000]
        self.slide_w = int(slide_size[0])
        self.slide_h = int(slide_size[1])

    def px_to_emu(self, x: float, y: float) -> tuple[int, int]:
        return (
            round(x / self.w_px * self.slide_w),
            round(y / self.h_px * self.slide_h),
        )

    def box_to_emu(self, box: list[float]) -> tuple[Emu, Emu, Emu, Emu]:
        x, y, w, h = box
        x1, y1 = self.px_to_emu(x, y)
        x2, y2 = self.px_to_emu(x + w, y + h)
        return Emu(x1), Emu(y1), Emu(x2 - x1), Emu(y2 - y1)

    def make_background(self) -> Path:
        bg = self.layout.get("background", {})
        mode = bg.get("mode", "source")
        out = self.workdir / f"{self.layout['id']}_background.png"
        out.parent.mkdir(parents=True, exist_ok=True)

        if mode == "source":
            Image.open(self.image_path).convert("RGB").resize((self.w_px, self.h_px)).save(out)
            return out

        if mode == "solid":
            self._draw_decorated_background(bg, bg.get("color", "#FFFFFF")).save(out)
            return out

        if mode == "notebook_grid_cover":
            self._draw_notebook_grid_cover(bg).save(out)
            return out

        raise ValueError(f"unsupported background mode: {mode}")

    def _draw_decorated_background(self, bg: dict[str, Any], base: str) -> Image.Image:
        image = Image.new("RGB", (self.w_px, self.h_px), base)
        draw = ImageDraw.Draw(image)
        for item in bg.get("decorations", []):
            kind = item["type"]
            if kind == "rect":
                x, y, w, h = item["bbox"]
                draw.rectangle(
                    (x, y, x + w, y + h),
                    fill=None if is_no_color(item.get("fill")) else item.get("fill"),
                    outline=None if is_no_color(item.get("outline")) else item.get("outline"),
                    width=int(item.get("line_width", 1)),
                )
            elif kind == "grid":
                x, y, w, h = item["bbox"]
                step_x = int(item.get("step_x", item.get("step", 20)))
                step_y = int(item.get("step_y", item.get("step", 20)))
                color = item.get("color", "#D8DCD8")
                width = int(item.get("width", 1))
                for gx in range(round(x), round(x + w) + 1, step_x):
                    draw.line((gx, y, gx, y + h), fill=color, width=width)
                for gy in range(round(y), round(y + h) + 1, step_y):
                    draw.line((x, gy, x + w, gy), fill=color, width=width)
            elif kind == "line":
                x1, y1, x2, y2 = item["points"]
                draw.line(
                    (x1, y1, x2, y2),
                    fill=item.get("color", "#000000"),
                    width=int(item.get("width", 1)),
                )
            else:
                raise ValueError(f"unsupported background decoration: {kind}")
        return image

    def _draw_notebook_grid_cover(self, bg: dict[str, Any]) -> Image.Image:
        image = Image.new("RGB", (self.w_px, self.h_px), bg.get("base", "#F7F7F2"))
        draw = ImageDraw.Draw(image)

        grid = bg.get("grid", {})
        step = int(grid.get("step", 34))
        color = grid.get("color", "#BFC7CB")
        width = int(grid.get("width", 1))
        for x in range(0, self.w_px + 1, step):
            draw.line((x, 0, x, self.h_px), fill=color, width=width)
        for y in range(0, self.h_px + 1, step):
            draw.line((0, y, self.w_px, y), fill=color, width=width)

        content = bg.get("content_rect", {})
        if content:
            box = content["bbox"]
            fill = content.get("fill", "#F8F7F2")
            outline = content.get("outline", "#0B3558")
            draw.rectangle(
                (box[0], box[1], box[0] + box[2], box[1] + box[3]),
                fill=fill,
                outline=outline,
                width=int(content.get("line_width", 2)),
            )

        for line in bg.get("lines", []):
            x1, y1, x2, y2 = line["points"]
            draw.line(
                (x1, y1, x2, y2),
                fill=line.get("color", "#0B3558"),
                width=int(line.get("width", 2)),
            )
        return image

    def render_slide(self, prs: Presentation) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = self.make_background()
        slide.shapes.add_picture(str(bg), Emu(0), Emu(0), Emu(self.slide_w), Emu(self.slide_h))

        for obj in self.layout.get("objects", []):
            kind = obj["type"]
            if kind == "text":
                self._add_text(slide, obj)
            elif kind == "shape":
                self._add_shape(slide, obj)
            elif kind == "line":
                self._add_line(slide, obj)
            elif kind == "polyline":
                self._add_polyline(slide, obj)
            elif kind == "image":
                self._add_image(slide, obj)
            else:
                raise ValueError(f"unsupported object type: {kind}")

    def _add_text(self, slide, obj: dict[str, Any]) -> None:
        tb = slide.shapes.add_textbox(*self.box_to_emu(obj["bbox"]))
        tb.name = obj.get("name", "text")
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = bool(obj.get("word_wrap", False))
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.vertical_anchor = getattr(MSO_ANCHOR, obj.get("valign", "TOP"))

        lines = obj.get("lines")
        if lines is None:
            lines = [obj["text"]]
        for idx, line in enumerate(lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.alignment = getattr(PP_ALIGN, obj.get("align", "LEFT"))
            if "line_spacing" in obj:
                p.line_spacing = obj["line_spacing"]
            run = p.add_run()
            run.text = line
            run.font.name = obj.get("font", DEFAULT_FONT)
            run.font.size = Pt(max(1.0, float(obj.get("font_size", 18)) * self.font_scale))
            run.font.bold = bool(obj.get("bold", False))
            run.font.color.rgb = rgb(obj.get("color", "#000000"))

    def _add_shape(self, slide, obj: dict[str, Any]) -> None:
        shape_name = obj.get("shape", "RECTANGLE")
        shape_type = getattr(MSO_AUTO_SHAPE_TYPE, shape_name)
        sh = slide.shapes.add_shape(shape_type, *self.box_to_emu(obj["bbox"]))
        sh.name = obj.get("name", shape_name.lower())
        fill = obj.get("fill")
        if is_no_color(fill):
            sh.fill.background()
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = rgb(fill)
            sh.fill.transparency = int(obj.get("transparency", 0))
        stroke = obj.get("stroke")
        if is_no_color(stroke):
            sh.line.fill.background()
        else:
            sh.line.color.rgb = rgb(stroke)
            sh.line.width = Pt(float(obj.get("stroke_width", 1)))
            if "dash" in obj:
                dash_style = _resolve_dash_style(obj["dash"])
                if dash_style is not None:
                    sh.line.dash_style = dash_style

    def _add_line(self, slide, obj: dict[str, Any]) -> None:
        x1, y1, x2, y2 = obj["points"]
        ex1, ey1 = self.px_to_emu(x1, y1)
        ex2, ey2 = self.px_to_emu(x2, y2)
        line = slide.shapes.add_connector(1, Emu(ex1), Emu(ey1), Emu(ex2), Emu(ey2))
        line.name = obj.get("name", "line")
        line.line.color.rgb = rgb(obj.get("color", "#000000"))
        line.line.width = Pt(float(obj.get("width", 1)))
        if "dash" in obj:
            dash_style = _resolve_dash_style(obj["dash"])
            if dash_style is not None:
                line.line.dash_style = dash_style

    def _add_polyline(self, slide, obj: dict[str, Any]) -> None:
        points = obj["points"]
        for idx, (start, end) in enumerate(pairwise(points)):
            line_obj = {
                "name": f"{obj.get('name', 'polyline')}_{idx + 1}",
                "points": [start[0], start[1], end[0], end[1]],
                "color": obj.get("color", "#000000"),
                "width": obj.get("width", 1),
            }
            if "dash" in obj:
                line_obj["dash"] = obj["dash"]
            self._add_line(slide, line_obj)

    def _add_image(self, slide, obj: dict[str, Any]) -> None:
        path = Path(obj["path"])
        if not path.is_absolute():
            path = self.image_path.parent / path
        if "crop_bbox" in obj:
            crop = obj["crop_bbox"]
            asset = self.workdir / f"{self.layout['id']}_{obj.get('name', 'crop')}.png"
            asset.parent.mkdir(parents=True, exist_ok=True)
            with Image.open(path).convert("RGBA") as im:
                x, y, w, h = crop
                im.crop((x, y, x + w, y + h)).save(asset)
            path = asset
        slide.shapes.add_picture(str(path), *self.box_to_emu(obj["bbox"]))


def build_deck(
    layouts: list[dict[str, Any]], image_root: Path, out_path: Path, *, font_scale: float = 1.0
) -> None:
    if not layouts:
        raise ValueError("no layouts")

    prs = Presentation()
    first_size = layouts[0].get("slide_size_emu") or [16_256_000, 9_144_000]
    prs.slide_width = Emu(int(first_size[0]))
    prs.slide_height = Emu(int(first_size[1]))

    workdir = out_path.parent / f"_{out_path.stem}_assets"
    for layout in layouts:
        image_path = image_root / layout["image"]
        LayoutRenderer(layout, image_path, workdir, font_scale=font_scale).render_slide(prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def load_layouts(paths: list[Path]) -> list[dict[str, Any]]:
    layouts = []
    for path in paths:
        layouts.append(json.loads(path.read_text(encoding="utf-8")))
    return layouts


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--layout", action="append", type=Path, required=True)
    parser.add_argument("--image-root", type=Path, default=Path("."))
    parser.add_argument("-o", "--output", type=Path, required=True)
    parser.add_argument("--font-scale", type=float, default=1.0)
    args = parser.parse_args()
    layouts = load_layouts(args.layout)
    build_deck(layouts, args.image_root, args.output, font_scale=args.font_scale)
    print(args.output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
