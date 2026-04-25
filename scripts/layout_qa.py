"""Validate, build, render, and compare layout-json reconstructed decks."""

from __future__ import annotations

import argparse
import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from PIL import Image, ImageChops, ImageStat
from pptx import Presentation
from reconstruct_from_layout import build_deck, load_layouts

from star_slide.rasterize.libreoffice import render_pptx_to_pngs


@dataclass
class SlideReport:
    slide_no: int
    layout_id: str
    object_count: int
    picture_count: int
    mean_abs_diff: float | None


def _is_number_list(value: Any, length: int) -> bool:
    return (
        isinstance(value, list)
        and len(value) == length
        and all(isinstance(item, int | float) for item in value)
    )


def _validate_box(box: Any, *, path: str, errors: list[str]) -> None:
    if not _is_number_list(box, 4):
        errors.append(f"{path}: bbox must be [x, y, width, height]")
        return
    if box[2] <= 0 or box[3] <= 0:
        errors.append(f"{path}: bbox width/height must be positive")


def _validate_points(points: Any, *, path: str, errors: list[str]) -> None:
    if not _is_number_list(points, 4):
        errors.append(f"{path}: points must be [x1, y1, x2, y2]")


def validate_layout(layout: dict[str, Any], *, strict_editable: bool) -> list[str]:
    errors: list[str] = []
    layout_id = layout.get("id", "<missing-id>")

    canvas = layout.get("canvas", {})
    if not isinstance(canvas.get("width"), int) or not isinstance(canvas.get("height"), int):
        errors.append(f"{layout_id}: canvas.width/height must be integers")

    bg = layout.get("background", {})
    if strict_editable and bg.get("mode") == "source":
        errors.append(f"{layout_id}: background.mode=source is not allowed in strict mode")

    for idx, item in enumerate(bg.get("decorations", [])):
        kind = item.get("type")
        path = f"{layout_id}.background.decorations[{idx}]"
        if kind in {"rect", "grid"}:
            _validate_box(item.get("bbox"), path=path, errors=errors)
        elif kind == "line":
            _validate_points(item.get("points"), path=path, errors=errors)
        else:
            errors.append(f"{path}: unsupported decoration type {kind!r}")

    for idx, obj in enumerate(layout.get("objects", [])):
        kind = obj.get("type")
        path = f"{layout_id}.objects[{idx}]({obj.get('name', '<unnamed>')})"
        if kind in {"text", "shape", "image"}:
            _validate_box(obj.get("bbox"), path=path, errors=errors)
        elif kind == "line":
            _validate_points(obj.get("points"), path=path, errors=errors)
        elif kind == "polyline":
            points = obj.get("points")
            if (
                not isinstance(points, list)
                or len(points) < 2
                or any(not _is_number_list(point, 2) for point in points)
            ):
                errors.append(f"{path}: polyline points must contain at least two [x, y] pairs")
        else:
            errors.append(f"{path}: unsupported object type {kind!r}")

        if strict_editable and kind == "image":
            errors.append(f"{path}: image object is not allowed in strict mode")
        if strict_editable and "crop_bbox" in obj:
            errors.append(f"{path}: crop_bbox is not allowed in strict mode")

    return errors


def _make_montage(render_paths: list[Path], out_path: Path) -> None:
    if not render_paths:
        return
    thumb_w, thumb_h = 344, 194
    cols = 2
    rows = (len(render_paths) + cols - 1) // cols
    sheet = Image.new("RGB", (thumb_w * cols, thumb_h * rows), "white")
    for idx, path in enumerate(render_paths):
        image = Image.open(path).convert("RGB")
        image.thumbnail((thumb_w, thumb_h))
        sheet.paste(image, ((idx % cols) * thumb_w, (idx // cols) * thumb_h))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(out_path)


def _make_pair(original: Path, rendered: Path, out_path: Path) -> float:
    orig = Image.open(original).convert("RGB")
    rend = Image.open(rendered).convert("RGB").resize(orig.size)
    pair = Image.new("RGB", (orig.width * 2 + 20, orig.height), "white")
    pair.paste(orig, (0, 0))
    pair.paste(rend, (orig.width + 20, 0))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    pair.save(out_path)

    diff = ImageChops.difference(orig, rend)
    stat = ImageStat.Stat(diff)
    return float(sum(stat.mean) / 3)


def _inspect_pptx(pptx_path: Path, layouts: list[dict[str, Any]]) -> list[SlideReport]:
    prs = Presentation(str(pptx_path))
    reports: list[SlideReport] = []
    for idx, (slide, layout) in enumerate(zip(prs.slides, layouts, strict=True), start=1):
        shapes = list(slide.shapes)
        reports.append(
            SlideReport(
                slide_no=idx,
                layout_id=str(layout.get("id", idx)),
                object_count=len(shapes),
                picture_count=sum(1 for shape in shapes if shape.shape_type == 13),
                mean_abs_diff=None,
            )
        )
    return reports


def run_qa(
    layout_paths: list[Path],
    image_root: Path,
    output: Path,
    render_dir: Path,
    *,
    strict_editable: bool,
) -> int:
    layouts = load_layouts(layout_paths)
    errors: list[str] = []
    for layout in layouts:
        errors.extend(validate_layout(layout, strict_editable=strict_editable))
    if errors:
        for error in errors:
            print(f"ERROR {error}")
        return 2

    build_deck(layouts, image_root, output)
    render_paths = render_pptx_to_pngs(output, render_dir)
    _make_montage(render_paths, render_dir / "montage.png")

    reports = _inspect_pptx(output, layouts)
    for idx, _report in enumerate(reports):
        original = image_root / layouts[idx]["image"]
        if original.exists() and idx < len(render_paths):
            mean_abs_diff = _make_pair(
                original,
                render_paths[idx],
                render_dir / "qa_pairs" / f"slide_{idx + 1:02d}_orig_vs_render.png",
            )
            reports[idx].mean_abs_diff = round(mean_abs_diff, 2)

    report_path = render_dir / "qa_report.json"
    report_path.write_text(
        json.dumps([report.__dict__ for report in reports], ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    for report in reports:
        diff = "n/a" if report.mean_abs_diff is None else f"{report.mean_abs_diff:.2f}"
        print(
            f"slide {report.slide_no:02d} {report.layout_id}: "
            f"objects={report.object_count} pictures={report.picture_count} mean_abs_diff={diff}"
        )
    print(output)
    print(render_dir / "montage.png")
    print(report_path)
    return 0


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--layout", action="append", type=Path, required=True)
    parser.add_argument("--image-root", type=Path, required=True)
    parser.add_argument("-o", "--output", type=Path, required=True)
    parser.add_argument("--render-dir", type=Path, required=True)
    parser.add_argument("--allow-images", action="store_true")
    args = parser.parse_args()

    return run_qa(
        args.layout,
        args.image_root,
        args.output,
        args.render_dir,
        strict_editable=not args.allow_images,
    )


if __name__ == "__main__":
    raise SystemExit(main())
