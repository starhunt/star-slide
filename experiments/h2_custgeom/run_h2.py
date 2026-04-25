"""H2 검증: SVG path → custGeom XML → PowerPoint 호환 PPTX 생성.

도형 13종(simple 10 + cubic 3)을 변환해 PPTX 슬라이드에 배치.
사용자가 PowerPoint 또는 LibreOffice로 열어 "도형 편집" 진입 가능 여부 확인.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt
from rich.console import Console
from rich.table import Table

from experiments.h2_custgeom.shapes import all_shapes
from star_slide.composer.inject import replace_geometry_with_custgeom
from star_slide.composer.svg2custgeom import svg_path_to_custgeom

RESULTS = Path(__file__).parent / "results"
RESULTS.mkdir(exist_ok=True)


def main() -> None:
    console = Console()

    shapes = all_shapes()
    console.print(f"[bold]H2: SVG path → custGeom 변환 {len(shapes)}종[/]")

    # PPTX 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 그리드 4 x 4 배치 (도형당 2 x 1.5 inch)
    cols = 4
    cell_w = Inches(2)
    cell_h = Inches(1.5)
    margin = Inches(0.5)

    table = Table(title="H2 svg→custGeom 변환 결과")
    table.add_column("#", justify="right", width=3)
    table.add_column("이름", style="cyan")
    table.add_column("설명")
    table.add_column("segments", justify="right")
    table.add_column("bbox", justify="right")
    table.add_column("XML 길이", justify="right")
    table.add_column("주입", justify="center")

    results: list[dict] = []

    for i, (name, d, desc) in enumerate(shapes):
        col = i % cols
        row = i // cols
        x = margin + col * cell_w
        y = margin + row * cell_h

        try:
            cg = svg_path_to_custgeom(d, target_w=1_000_000, target_h=1_000_000)
        except Exception as e:
            console.print(f"[red]변환 실패 {name}: {e}[/]")
            results.append(
                {
                    "name": name,
                    "ok": False,
                    "error": str(e),
                }
            )
            continue

        # 셰이프 생성: 임의 preset 사용 후 custGeom으로 교체
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, cell_w - Inches(0.2), cell_h - Inches(0.2)
        )
        try:
            replace_geometry_with_custgeom(shape, cg.xml)
            inject_ok = True
        except Exception as e:
            console.print(f"[red]주입 실패 {name}: {e}[/]")
            inject_ok = False

        # 라벨 텍스트 추가 (셰이프 아래)
        tb = slide.shapes.add_textbox(
            x,
            y + cell_h - Inches(0.3),
            cell_w,
            Inches(0.3),
        )
        tf = tb.text_frame
        tf.text = name
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(9)

        bbox_str = f"({cg.bbox[0]:.0f},{cg.bbox[1]:.0f},{cg.bbox[2]:.0f},{cg.bbox[3]:.0f})"
        table.add_row(
            str(i + 1),
            name,
            desc[:30],
            str(cg.n_segments),
            bbox_str,
            f"{len(cg.xml)}자",
            "[green]✓[/]" if inject_ok else "[red]✗[/]",
        )

        results.append(
            {
                "name": name,
                "description": desc,
                "n_segments": cg.n_segments,
                "bbox": list(cg.bbox),
                "has_arc": cg.has_arc,
                "xml_length": len(cg.xml),
                "inject_ok": inject_ok,
                "ok": True,
            }
        )

    out_pptx = RESULTS / "h2_shapes.pptx"
    prs.save(out_pptx)

    out_json = RESULTS / "h2_results.json"
    out_json.write_text(
        json.dumps(
            {
                "n_shapes": len(shapes),
                "n_ok": sum(1 for r in results if r.get("ok")),
                "n_inject_ok": sum(1 for r in results if r.get("inject_ok")),
                "shapes": results,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    console.print(table)
    console.print()
    n_ok = sum(1 for r in results if r.get("inject_ok"))
    console.print(f"[bold]변환+주입 성공:[/] {n_ok}/{len(shapes)}")
    console.print(f"[bold]PPTX:[/] {out_pptx}")
    console.print(f"[bold]상세:[/] {out_json}")
    console.print()
    console.print(
        "[yellow]다음:[/] PowerPoint(또는 LibreOffice)로 PPTX 열어 도형 우클릭 "
        "→ '점 편집' 메뉴 활성화 여부 확인 (custGeom의 PowerPoint 호환성 검증)"
    )


if __name__ == "__main__":
    main()
