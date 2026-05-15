"""Phase 5 PoC — 계층적 객체 분리 (큰 컬러 박스 + 안의 작은 픽토그램 모두 보존).

v4 한계:
  - SAM2 입력 1280 으로 축소 → nearest 업샘플 → 마스크 거칠음
  - IoU dedupe 가 큰 마스크 우선 → "사용자" 박스 안의 사람 아이콘 5개가
    컬러 박스에 흡수되어 1개 crop 으로 합침
  - classify_masks 의 BG 임계가 너무 낮아 큰 컬러 박스가 BG 분류

v5 변경:
  1. SAM2 입력 1920 + points_per_side=48 (dense grid) → 97+ 마스크
  2. 마스크 업샘플 = cv2.INTER_LINEAR + threshold 0.5 (부드러움)
  3. 계층 분리: 마스크 A 가 마스크 B 안에 ≥80% 포함되면 둘 다 보존하되,
     큰 마스크에서 작은 마스크 영역을 차집합 → 시각적 중첩 방지
     ("사용자" 박스 alpha 에서 사람 5개 영역 빼고, 사람 5개는 별도 alpha)
  4. classify_masks 의 background_area_ratio=0.6 (큰 컬러 박스 보호)
  5. PPTX z-order: 큰 마스크 (배경 컬러 박스) 부터 추가 →
                   중간 마스크 (그 안의 카드들) →
                   작은 마스크 (픽토그램) → textbox 최상단
"""

from __future__ import annotations

import argparse
import json
import time
from dataclasses import dataclass, field
from pathlib import Path

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Pt

from star_slide.inpaint.lama import inpaint_with_mask
from star_slide.ocr.paddleocr_worker import OcrLine, run_ocr
from star_slide.ocr.text_attributes import (
    detect_korean,
    estimate_font_color,
    estimate_font_size_pt,
)
from star_slide.segmentation.classify import classify_masks
from star_slide.segmentation.sam2_auto import Sam2Mask, run_sam2_auto

EMU_PER_INCH = 914400
PT_PER_INCH = 72.0
DEFAULT_SLIDE_W_EMU = 12192000
DEFAULT_SLIDE_H_EMU = 6858000
DEFAULT_FONT = "Apple SD Gothic Neo"


@dataclass
class TextItem:
    text: str
    bbox: tuple[float, float, float, float]
    polygon: list[tuple[float, float]]
    is_korean: bool
    font_pt: float
    color_hex: str


@dataclass
class TextLine:
    items: list[TextItem] = field(default_factory=list)

    @property
    def text(self) -> str:
        return " ".join(it.text for it in self.items)

    @property
    def bbox(self) -> tuple[float, float, float, float]:
        xs = [it.bbox[0] for it in self.items]
        ys = [it.bbox[1] for it in self.items]
        rs = [it.bbox[0] + it.bbox[2] for it in self.items]
        bs = [it.bbox[1] + it.bbox[3] for it in self.items]
        x = min(xs); y = min(ys)
        return (x, y, max(rs) - x, max(bs) - y)

    @property
    def font_pt(self) -> float:
        return float(np.median([it.font_pt for it in self.items]))

    @property
    def color_hex(self) -> str:
        from collections import Counter
        return Counter(it.color_hex for it in self.items).most_common(1)[0][0]


@dataclass
class ObjectLayer:
    """계층 분리된 객체 레이어 (alpha PNG + 위치 + z-depth)."""

    bbox: tuple[int, int, int, int]
    z: int  # 0 = 가장 뒤, 큰 숫자일수록 앞
    area_ratio: float
    alpha_path: Path


# ---------------- OCR ----------------


def compute_polygon_height(poly: list[tuple[float, float]]) -> float:
    arr = np.asarray(poly, dtype=float)
    if len(arr) >= 4:
        side1 = np.linalg.norm(arr[1] - arr[0])
        side2 = np.linalg.norm(arr[2] - arr[1])
        return float(min(side1, side2))
    ys = arr[:, 1]
    return float(ys.max() - ys.min())


def per_polygon_mask(poly, image_size):
    w, h = image_size
    mask = np.zeros((h, w), dtype=np.uint8)
    pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
    cv2.fillPoly(mask, [pts], 1)
    return mask.astype(bool)


def lines_to_text_items(lines, image, *, effective_dpi):
    items = []
    w, h = image.size
    for line in lines:
        h_px = compute_polygon_height(line.bbox_quad)
        is_kr = detect_korean(line.text)
        font_pt = estimate_font_size_pt(h_px, is_korean=is_kr, dpi=effective_dpi)
        mask_bool = per_polygon_mask(line.bbox_quad, (w, h))
        color = estimate_font_color(image, mask_bool) or "#222222"
        items.append(TextItem(
            text=line.text, bbox=line.bbox, polygon=line.bbox_quad,
            is_korean=is_kr, font_pt=font_pt, color_hex=color,
        ))
    return items


def group_items_by_line(items):
    if not items:
        return []
    sorted_items = sorted(items, key=lambda it: (it.bbox[1] + it.bbox[3] / 2, it.bbox[0]))
    lines = []
    for it in sorted_items:
        x, y, w, h = it.bbox
        y_center = y + h / 2
        merged = False
        for line in lines:
            lx, ly, lw, lh = line.bbox
            ly_center = ly + lh / 2
            avg_h = (lh + h) / 2
            if (
                abs(ly_center - y_center) <= avg_h * 0.5
                and abs(line.font_pt - it.font_pt) / max(line.font_pt, 1) < 0.3
                and -avg_h <= (x - (lx + lw)) <= avg_h * 1.5
            ):
                line.items.append(it)
                merged = True
                break
        if not merged:
            lines.append(TextLine(items=[it]))
    return lines


# ---------------- 인페인팅 마스크 ----------------


def build_dark_pixel_inpaint_mask(image, polygons, *, luminance_threshold=160, dilate=3):
    arr = np.asarray(image.convert("RGB"))
    h, w = arr.shape[:2]
    poly_mask = np.zeros((h, w), dtype=np.uint8)
    for poly in polygons:
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        cv2.fillPoly(poly_mask, [pts], 1)
    poly_mask_bool = poly_mask.astype(bool)
    lum = (
        0.299 * arr[:, :, 0].astype(np.float32)
        + 0.587 * arr[:, :, 1].astype(np.float32)
        + 0.114 * arr[:, :, 2].astype(np.float32)
    )
    dark = lum < luminance_threshold
    out = np.zeros((h, w), dtype=np.uint8)
    for poly in polygons:
        single = np.zeros((h, w), dtype=np.uint8)
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        cv2.fillPoly(single, [pts], 1)
        sb = single.astype(bool)
        dr = float(dark[sb].mean()) if sb.any() else 0.0
        target = (~dark & sb) if dr > 0.5 else (dark & sb)
        out[target] = 255
    if dilate > 0:
        k = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (dilate, dilate))
        out = cv2.dilate(out, k, iterations=1)
    out[~poly_mask_bool] = 0
    return out


# ---------------- SAM 마스크 후처리 ----------------


def upsample_mask_smooth(seg_small, target_size):
    """nearest 대신 LINEAR + threshold 로 부드러운 마스크 업샘플."""
    seg_uint = (seg_small.astype(np.uint8)) * 255
    seg_full = cv2.resize(seg_uint, target_size, interpolation=cv2.INTER_LINEAR)
    return seg_full > 127


def mask_iou(a, b):
    if a.shape != b.shape:
        return 0.0
    inter = np.logical_and(a, b).sum()
    union = np.logical_or(a, b).sum()
    return float(inter) / float(union) if union else 0.0


def mask_containment(small, large):
    """small 마스크가 large 마스크 안에 들어간 비율 (small 면적 기준)."""
    if small.shape != large.shape:
        return 0.0
    inter = np.logical_and(small, large).sum()
    s_area = small.sum()
    return float(inter) / float(s_area) if s_area else 0.0


def hierarchical_dedupe(
    masks: list[Sam2Mask],
    *,
    iou_dup_threshold: float = 0.85,
    containment_threshold: float = 0.80,
) -> list[Sam2Mask]:
    """계층 dedupe — 거의 같은 마스크는 제거(IoU>0.85), 작은 마스크가 큰 마스크 안에
    들어가면 둘 다 보존(z-order 만 다르게).

    중복 마스크 (95%+ 겹침) 만 제거. 부모-자식 관계는 유지.
    """
    sorted_masks = sorted(masks, key=lambda m: -m.area)
    kept: list[Sam2Mask] = []
    for m in sorted_masks:
        is_dup = False
        for k in kept:
            iou = mask_iou(m.segmentation, k.segmentation)
            if iou > iou_dup_threshold:
                # 거의 같은 마스크 — 더 큰 것이 이미 kept 에 있음 → 스킵
                is_dup = True
                break
        if not is_dup:
            kept.append(m)
    return kept


def filter_by_area(masks, image_size, *, min_ratio=0.0003, max_ratio=0.55):
    w, h = image_size
    total = w * h
    return [m for m in masks if min_ratio <= m.area / total <= max_ratio]


def subtract_mask(parent: np.ndarray, children: list[np.ndarray]) -> np.ndarray:
    """parent 마스크에서 children (자식 마스크들 union) 영역 제거."""
    out = parent.copy()
    for c in children:
        if c.shape == out.shape:
            out &= ~c
    return out


# ---------------- alpha crop 추출 ----------------


def export_alpha_crop(image, mask, *, out_dir, name):
    ys, xs = np.where(mask)
    if len(ys) == 0:
        return None
    x0, x1 = int(xs.min()), int(xs.max()) + 1
    y0, y1 = int(ys.min()), int(ys.max()) + 1
    arr = np.asarray(image.convert("RGB"))
    crop_rgb = arr[y0:y1, x0:x1]
    crop_mask = mask[y0:y1, x0:x1]
    rgba = np.zeros((y1 - y0, x1 - x0, 4), dtype=np.uint8)
    rgba[:, :, :3] = crop_rgb
    rgba[:, :, 3] = (crop_mask.astype(np.uint8)) * 255
    out_dir.mkdir(parents=True, exist_ok=True)
    p = out_dir / name
    Image.fromarray(rgba, mode="RGBA").save(p)
    return (x0, y0, x1 - x0, y1 - y0), p


def build_layered_objects(
    image: Image.Image,
    masks: list[Sam2Mask],
    text_polygons: list[list[tuple[float, float]]],
    *,
    out_dir: Path,
) -> list[ObjectLayer]:
    """계층 분리: 큰 마스크에서 자식 마스크 영역 빼고 alpha crop.
    z-order = 면적 큰 순으로 0,1,2... (큰 마스크가 뒤로 깔림).
    """
    h, w = np.asarray(image).shape[:2]
    text_mask = np.zeros((h, w), dtype=bool)
    for poly in text_polygons:
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        single = np.zeros((h, w), dtype=np.uint8)
        cv2.fillPoly(single, [pts], 1)
        text_mask |= single.astype(bool)

    # 면적 큰 순으로 정렬해 부모 → 자식 처리
    sorted_masks = sorted(masks, key=lambda m: -m.area)
    layers: list[ObjectLayer] = []
    obj_dir = out_dir / "objects"

    for i, m in enumerate(sorted_masks):
        # 텍스트 픽셀 빼기
        clean = m.segmentation & ~text_mask
        # 이 마스크 안에 들어간 더 작은 마스크들 (자식) — clean 에서 빼기
        children = []
        for j, c in enumerate(sorted_masks):
            if j <= i:
                continue
            if mask_containment(c.segmentation, m.segmentation) >= 0.80:
                children.append(c.segmentation)
        if children:
            clean = subtract_mask(clean, children)

        if clean.sum() < 200:  # 자식 빼고 거의 안 남으면 (== 컨테이너) 스킵
            continue

        out = export_alpha_crop(image, clean, out_dir=obj_dir, name=f"layer_{i:03d}.png")
        if out is None:
            continue
        bbox, path = out
        layers.append(ObjectLayer(
            bbox=bbox,
            z=i,  # 큰 마스크 순서 — z 작을수록 뒤
            area_ratio=float(clean.sum()) / (h * w),
            alpha_path=path,
        ))
    return layers


# ---------------- 시각화 ----------------


def render_mask_overlay(image, masks, out_path):
    rng = np.random.default_rng(42)
    arr = np.asarray(image.convert("RGB")).copy().astype(np.float32)
    overlay = np.zeros_like(arr)
    for m in masks:
        color = rng.integers(50, 255, size=3)
        overlay[m.segmentation] = color
    blended = arr * 0.55 + overlay * 0.45
    Image.fromarray(np.clip(blended, 0, 255).astype(np.uint8)).save(out_path)


def render_layer_overlay(image, layers, out_path):
    arr = np.asarray(image.convert("RGB")).copy()
    for ly in layers:
        x, y, w, h = ly.bbox
        cv2.rectangle(arr, (x, y), (x + w, y + h), (0, 200, 0), 2)
    Image.fromarray(arr).save(out_path)


def render_text_overlay(image, lines, out_path):
    arr = np.asarray(image.convert("RGB")).copy()
    for ln in lines:
        x, y, w, h = ln.bbox
        cv2.rectangle(arr, (int(x), int(y)), (int(x + w), int(y + h)), (255, 0, 0), 2)
    Image.fromarray(arr).save(out_path)


# ---------------- PPTX ----------------


def add_textbox(slide, line, *, scale, font_name):
    x, y, w, h = line.bbox
    font_emu = int(line.font_pt / PT_PER_INCH * EMU_PER_INCH)
    left = int(x * scale) - font_emu // 2
    width = int(w * scale) + font_emu
    height = max(int(h * scale), int(font_emu * 1.6))
    y_center_emu = int((y + h / 2) * scale)
    top = y_center_emu - height // 2
    if width <= 0 or height <= 0:
        return
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, attr, Emu(0))
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = line.text
    run.font.name = font_name
    run.font.size = Pt(round(line.font_pt, 1))
    try:
        r = int(line.color_hex[1:3], 16)
        g = int(line.color_hex[3:5], 16)
        b = int(line.color_hex[5:7], 16)
        run.font.color.rgb = RGBColor(r, g, b)
    except Exception:
        pass


def build_pptx(
    *, image_size, background_path, layers, lines, out_path,
    slide_w_emu, slide_h_emu, font_name,
):
    px_w, _ = image_size
    scale = slide_w_emu / px_w
    prs = Presentation()
    prs.slide_width = Emu(slide_w_emu)
    prs.slide_height = Emu(slide_h_emu)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    if background_path is not None:
        slide.shapes.add_picture(
            str(background_path), 0, 0,
            width=Emu(slide_w_emu), height=Emu(slide_h_emu),
        )
    # z-order: 작은 z (큰 마스크) 부터 추가 → 뒤로 깔림
    for ly in sorted(layers, key=lambda l: l.z):
        x, y, w, h = ly.bbox
        slide.shapes.add_picture(
            str(ly.alpha_path),
            Emu(int(x * scale)), Emu(int(y * scale)),
            width=Emu(int(w * scale)), height=Emu(int(h * scale)),
        )
    for ln in lines:
        add_textbox(slide, ln, scale=scale, font_name=font_name)
    prs.save(str(out_path))


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", type=Path, required=True)
    ap.add_argument("--out-dir", type=Path, required=True)
    ap.add_argument("--font", default=DEFAULT_FONT)
    ap.add_argument("--slide-w-emu", type=int, default=DEFAULT_SLIDE_W_EMU)
    ap.add_argument("--slide-h-emu", type=int, default=DEFAULT_SLIDE_H_EMU)
    ap.add_argument("--sam-resize", type=int, default=1920)
    ap.add_argument("--sam-points", type=int, default=48)
    ap.add_argument("--sam-iou", type=float, default=0.80)
    ap.add_argument("--sam-stability", type=float, default=0.88)
    args = ap.parse_args()

    args.out_dir.mkdir(parents=True, exist_ok=True)
    started = time.perf_counter()

    print(f"[1/9] 입력 로드: {args.input}")
    image = Image.open(args.input).convert("RGB")
    px_w, px_h = image.size
    slide_w_inch = args.slide_w_emu / EMU_PER_INCH
    effective_dpi = px_w / slide_w_inch
    print(f"     {px_w}x{px_h}, effective_dpi={effective_dpi:.1f}")

    print("[2/9] PaddleOCR (한국어)")
    t0 = time.perf_counter()
    raw_lines = run_ocr(image, lang="korean")
    print(f"     {len(raw_lines)} 라인, {time.perf_counter() - t0:.1f}s")

    print("[3/9] 텍스트 속성 + 라인 그룹핑")
    items = lines_to_text_items(raw_lines, image, effective_dpi=effective_dpi)
    grouped = group_items_by_line(items)
    print(f"     items={len(items)} → grouped={len(grouped)}")
    render_text_overlay(image, grouped, args.out_dir / "01_text_lines.png")

    print(f"[4/9] SAM2 dense auto-mask (resize={args.sam_resize}, pps={args.sam_points})")
    sam_input = image
    if args.sam_resize and px_w > args.sam_resize:
        scale_sam = args.sam_resize / px_w
        sam_input = image.resize((args.sam_resize, int(px_h * scale_sam)))
    sam_result = run_sam2_auto(
        sam_input,
        points_per_side=args.sam_points,
        max_masks=400,
        pred_iou_thresh=args.sam_iou,
        stability_score_thresh=args.sam_stability,
    )
    print(f"     {len(sam_result.masks)} 마스크, {sam_result.elapsed_sec:.1f}s on {sam_result.device}")

    # 풀 해상도 마스크로 업샘플 (LINEAR + threshold)
    full_masks: list[Sam2Mask] = []
    for m in sam_result.masks:
        seg = upsample_mask_smooth(m.segmentation, (px_w, px_h))
        ys, xs = np.where(seg)
        if len(ys) == 0:
            continue
        bbox = (int(xs.min()), int(ys.min()),
                int(xs.max() - xs.min() + 1), int(ys.max() - ys.min() + 1))
        full_masks.append(Sam2Mask(bbox=bbox, segmentation=seg, score=m.score, area=int(seg.sum())))
    render_mask_overlay(image, full_masks, args.out_dir / "02_sam2_masks.png")

    print("[5/9] classify_masks (BG=0.6, NOISE=0.0001)")
    classification = classify_masks(
        full_masks, raw_lines, image_size=(px_w, px_h),
        background_area_ratio=0.6, noise_area_ratio=0.0001,
    )
    shapes = classification.shape_masks
    print(f"     TEXT={len(classification.text_masks)} SHAPE={len(shapes)}")

    print("[6/9] 면적 필터 + 거의-동일-마스크 제거 (IoU>0.85)")
    shapes = filter_by_area(shapes, (px_w, px_h), min_ratio=0.0003, max_ratio=0.55)
    shapes = hierarchical_dedupe(shapes, iou_dup_threshold=0.85)
    print(f"     필터 후 {len(shapes)} (계층 보존)")

    print("[7/9] 계층 alpha crop (큰 마스크에서 자식 영역 빼기)")
    layers = build_layered_objects(
        image, shapes, [it.polygon for it in items], out_dir=args.out_dir,
    )
    print(f"     {len(layers)} 레이어 alpha crop")
    render_layer_overlay(image, layers, args.out_dir / "03_layer_overlay.png")

    print("[8/9] 인페인팅 배경")
    inpaint_mask = build_dark_pixel_inpaint_mask(image, [it.polygon for it in items], dilate=3)
    Image.fromarray(inpaint_mask, mode="L").save(args.out_dir / "04_inpaint_mask.png")
    t0 = time.perf_counter()
    inpainted = inpaint_with_mask(image, Image.fromarray(inpaint_mask, mode="L"))
    inpainted_path = args.out_dir / "05_inpainted_bg.png"
    inpainted.save(inpainted_path)
    print(f"     {time.perf_counter() - t0:.1f}s")

    print("[9/9] PPTX 두 모드 생성")
    out_white = args.out_dir / "result_white_bg.pptx"
    build_pptx(
        image_size=(px_w, px_h), background_path=None,
        layers=layers, lines=grouped, out_path=out_white,
        slide_w_emu=args.slide_w_emu, slide_h_emu=args.slide_h_emu,
        font_name=args.font,
    )
    out_inpainted = args.out_dir / "result_inpainted.pptx"
    build_pptx(
        image_size=(px_w, px_h), background_path=inpainted_path,
        layers=layers, lines=grouped, out_path=out_inpainted,
        slide_w_emu=args.slide_w_emu, slide_h_emu=args.slide_h_emu,
        font_name=args.font,
    )

    report = {
        "input": str(args.input),
        "image_size": [px_w, px_h],
        "ocr_items": len(items),
        "grouped_lines": len(grouped),
        "sam_masks_total": len(full_masks),
        "sam_text_masks": len(classification.text_masks),
        "sam_shapes_after_filter": len(shapes),
        "object_layers": len(layers),
        "elapsed_sec": round(time.perf_counter() - started, 1),
    }
    (args.out_dir / "report.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    print(f"\nresult_white_bg.pptx   = 흰 배경 + 계층 객체 + textbox")
    print(f"result_inpainted.pptx  = 인페인팅 배경 + 계층 객체 + textbox")
    print(f"총 {report['elapsed_sec']}s")


if __name__ == "__main__":
    main()
