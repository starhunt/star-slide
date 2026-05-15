"""Phase 1 PoC — Text/object 완벽 분리 + 정확한 폰트 pt.

기존 6개 PoC 의 공통 함정 ("폰트 = h*72/96*0.7") 을 피하고, 이미 메인 파이프라인에
있는 검증된 모듈을 정공으로 사용한다:

  - paddleocr_worker.run_ocr  : 한국어 OCR + 정확한 polygon
  - text_attributes.estimate_font_size_pt : 한글 0.85 비율로 정확한 pt
  - text_attributes.estimate_font_color   : ink 픽셀 mode 색상
  - inpaint.lama.inpaint_with_mask        : OCR 마스크 → 텍스트 제거 깨끗한 배경

Phase 1 단순화:
  - 객체 자동 분할(SAM2/SAM3)은 미사용. 대신 텍스트만 완벽 분리하고,
    배경은 인페인팅 전 / 후 두 버전을 만들어 후속 비교용으로 보관.
  - 슬라이드 = 인페인팅된 배경 1장 + OCR 라인별 editable textbox N개.
  - 텍스트가 정확한 pt/색상으로 원본 위에 덮이는지 시각 검증이 핵심.

다음 Phase: 객체 분할(SAM2 auto + classify_masks) 추가 → 큰 아이콘/박스를
            alpha PNG crop 으로 분리해 별도 picture shape 로 삽입.

사용법:
    uv run python experiments/diagram_split/poc_v1_text_object_split.py \
        --input experiments/diagram_split/input.png \
        --out-dir experiments/diagram_split/v1_out
"""

from __future__ import annotations

import argparse
import json
import time
from dataclasses import asdict, dataclass
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

from star_slide.inpaint.lama import inpaint_with_mask
from star_slide.ocr.paddleocr_worker import OcrLine, run_ocr
from star_slide.ocr.text_attributes import (
    detect_korean,
    estimate_font_color,
    estimate_font_size_pt,
)

SLIDE_DPI = 96.0
EMU_PER_INCH = 914400


@dataclass
class TextItem:
    text: str
    confidence: float
    # 화면 픽셀 좌표
    bbox: tuple[float, float, float, float]  # x, y, w, h
    polygon: list[tuple[float, float]]
    # 추정된 시각 속성
    is_korean: bool
    font_pt: float
    color_hex: str
    mask_height_px: float


def polygons_to_text_mask(
    polygons: list[list[tuple[float, float]]],
    image_size: tuple[int, int],
    *,
    dilate: int = 3,
) -> np.ndarray:
    """OCR polygon 들을 합쳐 글자 영역 마스크 (uint8 0/255)."""
    import cv2

    w, h = image_size
    mask = np.zeros((h, w), dtype=np.uint8)
    for poly in polygons:
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        cv2.fillPoly(mask, [pts], 255)
    if dilate > 0:
        kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (dilate, dilate))
        mask = cv2.dilate(mask, kernel, iterations=1)
    return mask


def per_polygon_mask(
    poly: list[tuple[float, float]],
    image_size: tuple[int, int],
) -> np.ndarray:
    """단일 polygon 의 bool 마스크 (text_attributes 입력용)."""
    import cv2

    w, h = image_size
    mask = np.zeros((h, w), dtype=np.uint8)
    pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
    cv2.fillPoly(mask, [pts], 1)
    return mask.astype(bool)


def compute_polygon_height(poly: list[tuple[float, float]]) -> float:
    """polygon 의 axis-aligned 높이 (회전 텍스트 고려해서 quad 의 짧은 변)."""
    arr = np.asarray(poly, dtype=float)
    # 4점 polygon 이라고 가정. 짧은 변(=글자 높이) 길이.
    if len(arr) >= 4:
        side1 = np.linalg.norm(arr[1] - arr[0])
        side2 = np.linalg.norm(arr[2] - arr[1])
        return float(min(side1, side2))
    # fallback: bbox height
    ys = arr[:, 1]
    return float(ys.max() - ys.min())


def lines_to_text_items(
    lines: list[OcrLine],
    image: Image.Image,
) -> list[TextItem]:
    """OCR 결과 → TextItem (정확한 pt/색상 측정)."""
    items: list[TextItem] = []
    w, h = image.size
    for line in lines:
        h_px = compute_polygon_height(line.bbox_quad)
        is_kr = detect_korean(line.text)
        font_pt = estimate_font_size_pt(h_px, is_korean=is_kr)
        mask_bool = per_polygon_mask(line.bbox_quad, (w, h))
        color = estimate_font_color(image, mask_bool) or "#222222"
        items.append(
            TextItem(
                text=line.text,
                confidence=line.confidence,
                bbox=line.bbox,
                polygon=line.bbox_quad,
                is_korean=is_kr,
                font_pt=font_pt,
                color_hex=color,
                mask_height_px=h_px,
            )
        )
    return items


def render_debug_overlay(
    image: Image.Image,
    items: list[TextItem],
    out_path: Path,
) -> None:
    """원본 위에 OCR polygon + 폰트 pt 라벨 덮어 시각 확인."""
    import cv2

    arr = np.asarray(image.convert("RGB")).copy()
    for it in items:
        pts = np.asarray(it.polygon, dtype=np.int32).reshape(-1, 1, 2)
        cv2.polylines(arr, [pts], isClosed=True, color=(255, 0, 0), thickness=2)
        x, y, _, _ = it.bbox
        label = f"{it.font_pt:.0f}pt"
        cv2.putText(
            arr,
            label,
            (int(x), max(15, int(y) - 4)),
            cv2.FONT_HERSHEY_SIMPLEX,
            0.45,
            (0, 0, 255),
            1,
            cv2.LINE_AA,
        )
    Image.fromarray(arr).save(out_path)


def build_pptx(
    *,
    image_size: tuple[int, int],
    background_path: Path,
    items: list[TextItem],
    out_path: Path,
    font_name: str = "맑은 고딕",
) -> None:
    """배경 PNG + editable textbox 들로 1-슬라이드 PPTX 빌드.

    슬라이드 크기는 입력 이미지 픽셀 비율(96 DPI 가정)을 그대로 따라간다.
    """
    px_w, px_h = image_size
    slide_w_emu = int(px_w / SLIDE_DPI * EMU_PER_INCH)
    slide_h_emu = int(px_h / SLIDE_DPI * EMU_PER_INCH)

    prs = Presentation()
    prs.slide_width = Emu(slide_w_emu)
    prs.slide_height = Emu(slide_h_emu)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 배경 (인페인팅된 PNG) 풀스크린
    slide.shapes.add_picture(
        str(background_path), 0, 0, width=Emu(slide_w_emu), height=Emu(slide_h_emu)
    )

    # 각 OCR 라인 → editable textbox
    for it in items:
        x, y, w, h = it.bbox
        # bbox 약간 패딩 (글자 잘림 방지)
        pad_x = int(w * 0.04)
        pad_y = int(h * 0.06)
        left_emu = int((x - pad_x) / SLIDE_DPI * EMU_PER_INCH)
        top_emu = int((y - pad_y) / SLIDE_DPI * EMU_PER_INCH)
        width_emu = int((w + 2 * pad_x) / SLIDE_DPI * EMU_PER_INCH)
        height_emu = int((h + 2 * pad_y) / SLIDE_DPI * EMU_PER_INCH)

        tb = slide.shapes.add_textbox(
            Emu(left_emu), Emu(top_emu), Emu(width_emu), Emu(height_emu)
        )
        tf = tb.text_frame
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.word_wrap = False  # 1줄 OCR 결과는 wrap 비활성

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = it.text
        run.font.name = font_name
        run.font.size = Pt(it.font_pt)
        try:
            r, g, b = (
                int(it.color_hex[1:3], 16),
                int(it.color_hex[3:5], 16),
                int(it.color_hex[5:7], 16),
            )
            run.font.color.rgb = RGBColor(r, g, b)
        except Exception:
            pass

    prs.save(str(out_path))


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", type=Path, required=True)
    ap.add_argument("--out-dir", type=Path, required=True)
    ap.add_argument("--font", default="맑은 고딕")
    args = ap.parse_args()

    args.out_dir.mkdir(parents=True, exist_ok=True)
    started = time.perf_counter()

    print(f"[1/5] 입력 로드: {args.input}")
    image = Image.open(args.input).convert("RGB")
    w, h = image.size
    print(f"     크기 {w}x{h}")

    print("[2/5] PaddleOCR (한국어) 실행")
    t0 = time.perf_counter()
    lines = run_ocr(image, lang="korean")
    print(f"     {len(lines)} 라인 검출, {time.perf_counter() - t0:.1f}s")

    print("[3/5] 텍스트 속성(폰트 pt + 색상) 측정")
    items = lines_to_text_items(lines, image)
    if items:
        font_sizes = [it.font_pt for it in items]
        print(
            f"     폰트 pt: min={min(font_sizes):.1f} median={sorted(font_sizes)[len(font_sizes)//2]:.1f} "
            f"max={max(font_sizes):.1f}"
        )

    print("[4/5] LaMa 인페인팅으로 텍스트 제거")
    text_mask = polygons_to_text_mask(
        [it.polygon for it in items], (w, h), dilate=5
    )
    Image.fromarray(text_mask, mode="L").save(args.out_dir / "01_text_mask.png")
    t0 = time.perf_counter()
    inpainted = inpaint_with_mask(image, Image.fromarray(text_mask, mode="L"))
    inpainted_path = args.out_dir / "02_inpainted_bg.png"
    inpainted.save(inpainted_path)
    print(f"     인페인팅 완료, {time.perf_counter() - t0:.1f}s")

    render_debug_overlay(image, items, args.out_dir / "03_overlay.png")

    print("[5/5] PPTX 조립 (배경 + editable textbox)")
    out_pptx = args.out_dir / "result_v1.pptx"
    build_pptx(
        image_size=(w, h),
        background_path=inpainted_path,
        items=items,
        out_path=out_pptx,
        font_name=args.font,
    )

    # 메타 보고
    report = {
        "input": str(args.input),
        "image_size": [w, h],
        "ocr_lines": len(items),
        "elapsed_sec": round(time.perf_counter() - started, 1),
        "items": [asdict(it) for it in items],
    }
    (args.out_dir / "report.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    print(f"\n결과: {out_pptx}")
    print(f"디버그: {args.out_dir}/01_text_mask.png  02_inpainted_bg.png  03_overlay.png")
    print(f"총 소요: {report['elapsed_sec']}s")


if __name__ == "__main__":
    main()
