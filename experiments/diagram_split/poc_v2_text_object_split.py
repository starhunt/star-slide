"""Phase 2 PoC — v1 의 결함 수정.

v1 한계:
  - 슬라이드 크기를 입력 픽셀 비율 그대로(2560/96 → 26.67") 잡아 비표준 슬라이드.
    LibreOffice 렌더 시 한글 폰트 fallback 으로 글자 폭이 깨짐.
  - LaMa 가 큰 텍스트 영역을 흰 패치로 덮어 그라데이션 배경이 사라짐.

v2 변경:
  1. 슬라이드 = 표준 16:9 (13.33"x7.5", 12192000x6858000 EMU).
     입력이 2560x1440 (16:9) 이므로 비율 1:1 매칭.
  2. scale_factor = slide_w_emu / image_w_px 로 모든 좌표/사이즈 변환.
  3. 폰트 pt 는 estimate_font_size_pt(mask_h, dpi=effective_dpi) 로
     슬라이드 사이즈에 맞춰 자동 조정 (effective_dpi = image_w_px / slide_w_inch).
  4. 한글 폰트 = "Apple SD Gothic Neo" (macOS 기본, LibreOffice 인식).
  5. 인페인트 마스크 = polygon 내부에서 luminance < threshold 인 dark 픽셀만 →
     배경 그라데이션/색상 보존. dilate 작게.
  6. 두 가지 PPTX 동시 생성:
        result_overlay.pptx   — 원본 풀스크린 + textbox 덮음 (안전)
        result_inpainted.pptx — 인페인팅된 배경 + textbox  (편집성 우수)
"""

from __future__ import annotations

import argparse
import json
import time
from dataclasses import asdict, dataclass
from pathlib import Path

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Pt

from star_slide.inpaint.lama import inpaint_with_mask
from star_slide.ocr.paddleocr_worker import OcrLine, run_ocr
from star_slide.ocr.text_attributes import (
    detect_korean,
    estimate_font_color,
    estimate_font_size_pt,
)

EMU_PER_INCH = 914400
PT_PER_INCH = 72.0
DEFAULT_SLIDE_W_EMU = 12192000  # 13.333 inch  (16:9 표준)
DEFAULT_SLIDE_H_EMU = 6858000   #  7.5  inch
DEFAULT_FONT = "Apple SD Gothic Neo"


@dataclass
class TextItem:
    text: str
    confidence: float
    bbox: tuple[float, float, float, float]
    polygon: list[tuple[float, float]]
    is_korean: bool
    font_pt: float
    color_hex: str
    mask_height_px: float


def compute_polygon_height(poly: list[tuple[float, float]]) -> float:
    arr = np.asarray(poly, dtype=float)
    if len(arr) >= 4:
        side1 = np.linalg.norm(arr[1] - arr[0])
        side2 = np.linalg.norm(arr[2] - arr[1])
        return float(min(side1, side2))
    ys = arr[:, 1]
    return float(ys.max() - ys.min())


def per_polygon_mask(
    poly: list[tuple[float, float]],
    image_size: tuple[int, int],
) -> np.ndarray:
    w, h = image_size
    mask = np.zeros((h, w), dtype=np.uint8)
    pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
    cv2.fillPoly(mask, [pts], 1)
    return mask.astype(bool)


def lines_to_text_items(
    lines: list[OcrLine],
    image: Image.Image,
    *,
    effective_dpi: float,
) -> list[TextItem]:
    items: list[TextItem] = []
    w, h = image.size
    for line in lines:
        h_px = compute_polygon_height(line.bbox_quad)
        is_kr = detect_korean(line.text)
        font_pt = estimate_font_size_pt(h_px, is_korean=is_kr, dpi=effective_dpi)
        mask_bool = per_polygon_mask(line.bbox_quad, (w, h))
        color = estimate_font_color(image, mask_bool) or "#222222"
        items.append(
            TextItem(
                text=line.text,
                confidence=line.confidence,
                bbox=line.bbox,
                polygon=line.bbox_quad,
                is_korean=is_kr,
                font_pt=font_pt,
                color_hex=color,
                mask_height_px=h_px,
            )
        )
    return items


def build_dark_pixel_inpaint_mask(
    image: Image.Image,
    polygons: list[list[tuple[float, float]]],
    *,
    luminance_threshold: int = 160,
    dilate: int = 3,
) -> np.ndarray:
    """polygon 내부에서 dark(글자) 픽셀만 골라 인페인트 마스크 생성.

    배경 그라데이션은 보존하고 글자 stroke 만 지워서 큰 텍스트도 깔끔하게 처리.
    """
    arr = np.asarray(image.convert("RGB"))
    h, w = arr.shape[:2]
    poly_mask = np.zeros((h, w), dtype=np.uint8)
    for poly in polygons:
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        cv2.fillPoly(poly_mask, [pts], 1)
    poly_mask_bool = poly_mask.astype(bool)

    lum = (
        0.299 * arr[:, :, 0].astype(np.float32)
        + 0.587 * arr[:, :, 1].astype(np.float32)
        + 0.114 * arr[:, :, 2].astype(np.float32)
    )
    dark_pixels = lum < luminance_threshold

    # 흰 글자 케이스: 라인이 어두운 배경 위 흰 글자면 polygon 내 평균 luminance 가 낮고,
    # 이 경우 dark_pixels 가 곧 배경 → 마스크가 글자가 아닌 배경을 지움.
    # 보호 로직: polygon 내 dark 비율 > 50% 면 light pixels 를 마스크로 사용.
    inpaint_mask = np.zeros((h, w), dtype=np.uint8)
    for poly in polygons:
        single = np.zeros((h, w), dtype=np.uint8)
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        cv2.fillPoly(single, [pts], 1)
        sb = single.astype(bool)
        in_poly_dark_ratio = float(dark_pixels[sb].mean()) if sb.any() else 0.0
        if in_poly_dark_ratio > 0.5:
            # 배경이 어두운 케이스 — 흰 글자 가정
            target = ~dark_pixels & sb
        else:
            target = dark_pixels & sb
        inpaint_mask[target] = 255

    if dilate > 0:
        kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (dilate, dilate))
        inpaint_mask = cv2.dilate(inpaint_mask, kernel, iterations=1)

    # poly 외부 보호: 절대 글자 영역 밖은 건드리지 않음
    inpaint_mask[~poly_mask_bool] = 0
    return inpaint_mask


def render_debug_overlay(
    image: Image.Image,
    items: list[TextItem],
    out_path: Path,
) -> None:
    arr = np.asarray(image.convert("RGB")).copy()
    for it in items:
        pts = np.asarray(it.polygon, dtype=np.int32).reshape(-1, 1, 2)
        cv2.polylines(arr, [pts], isClosed=True, color=(255, 0, 0), thickness=2)
        x, y, _, _ = it.bbox
        label = f"{it.font_pt:.0f}pt"
        cv2.putText(
            arr, label, (int(x), max(15, int(y) - 4)),
            cv2.FONT_HERSHEY_SIMPLEX, 0.45, (0, 0, 255), 1, cv2.LINE_AA,
        )
    Image.fromarray(arr).save(out_path)


def add_textbox(
    slide,
    item: TextItem,
    *,
    scale_emu_per_px: float,
    font_name: str,
) -> None:
    x, y, w, h = item.bbox
    # textbox margin 0, 약간 패딩으로 글자 잘림 방지
    pad_x = max(2.0, w * 0.03)
    pad_y = max(2.0, h * 0.10)
    left = int((x - pad_x) * scale_emu_per_px)
    top = int((y - pad_y) * scale_emu_per_px)
    width = int((w + 2 * pad_x) * scale_emu_per_px)
    height = int((h + 2 * pad_y) * scale_emu_per_px)

    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = item.text
    run.font.name = font_name
    run.font.size = Pt(item.font_pt)
    try:
        r = int(item.color_hex[1:3], 16)
        g = int(item.color_hex[3:5], 16)
        b = int(item.color_hex[5:7], 16)
        run.font.color.rgb = RGBColor(r, g, b)
    except Exception:
        pass


def build_pptx(
    *,
    image_size: tuple[int, int],
    background_path: Path,
    items: list[TextItem],
    out_path: Path,
    slide_w_emu: int,
    slide_h_emu: int,
    font_name: str,
) -> None:
    px_w, px_h = image_size
    scale = slide_w_emu / px_w  # EMU per pixel

    prs = Presentation()
    prs.slide_width = Emu(slide_w_emu)
    prs.slide_height = Emu(slide_h_emu)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    slide.shapes.add_picture(
        str(background_path), 0, 0,
        width=Emu(slide_w_emu), height=Emu(slide_h_emu),
    )

    for it in items:
        add_textbox(slide, it, scale_emu_per_px=scale, font_name=font_name)

    prs.save(str(out_path))


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", type=Path, required=True)
    ap.add_argument("--out-dir", type=Path, required=True)
    ap.add_argument("--font", default=DEFAULT_FONT)
    ap.add_argument("--slide-w-emu", type=int, default=DEFAULT_SLIDE_W_EMU)
    ap.add_argument("--slide-h-emu", type=int, default=DEFAULT_SLIDE_H_EMU)
    args = ap.parse_args()

    args.out_dir.mkdir(parents=True, exist_ok=True)
    started = time.perf_counter()

    print(f"[1/6] 입력 로드: {args.input}")
    image = Image.open(args.input).convert("RGB")
    px_w, px_h = image.size
    print(f"     크기 {px_w}x{px_h}")

    # 슬라이드 비율 검증 — 입력이 슬라이드 비율과 다르면 크롭/패딩 필요
    img_ratio = px_w / px_h
    slide_ratio = args.slide_w_emu / args.slide_h_emu
    if abs(img_ratio - slide_ratio) > 0.02:
        print(f"     ⚠ 입력 비율 {img_ratio:.2f} vs 슬라이드 비율 {slide_ratio:.2f} 불일치")
    slide_w_inch = args.slide_w_emu / EMU_PER_INCH
    effective_dpi = px_w / slide_w_inch
    print(f"     effective_dpi = {effective_dpi:.1f} (폰트 pt 변환 기준)")

    print("[2/6] PaddleOCR (한국어) 실행")
    t0 = time.perf_counter()
    lines = run_ocr(image, lang="korean")
    print(f"     {len(lines)} 라인 검출, {time.perf_counter() - t0:.1f}s")

    print("[3/6] 텍스트 속성(폰트 pt + 색상) 측정")
    items = lines_to_text_items(lines, image, effective_dpi=effective_dpi)
    if items:
        font_sizes = sorted(it.font_pt for it in items)
        print(
            f"     폰트 pt: min={font_sizes[0]:.1f} median={font_sizes[len(font_sizes)//2]:.1f} "
            f"max={font_sizes[-1]:.1f}"
        )

    render_debug_overlay(image, items, args.out_dir / "01_overlay.png")

    print("[4/6] 텍스트 stroke 마스크(dark pixel only) 생성")
    inpaint_mask = build_dark_pixel_inpaint_mask(
        image, [it.polygon for it in items], luminance_threshold=160, dilate=3,
    )
    Image.fromarray(inpaint_mask, mode="L").save(args.out_dir / "02_inpaint_mask.png")

    print("[5/6] LaMa 인페인팅")
    t0 = time.perf_counter()
    inpainted = inpaint_with_mask(image, Image.fromarray(inpaint_mask, mode="L"))
    inpainted_path = args.out_dir / "03_inpainted_bg.png"
    inpainted.save(inpainted_path)
    orig_bg_path = args.out_dir / "03_original_bg.png"
    image.save(orig_bg_path)
    print(f"     인페인팅 완료, {time.perf_counter() - t0:.1f}s")

    print("[6/6] PPTX 두 모드 생성")
    out_overlay = args.out_dir / "result_overlay.pptx"
    build_pptx(
        image_size=(px_w, px_h),
        background_path=orig_bg_path,
        items=items,
        out_path=out_overlay,
        slide_w_emu=args.slide_w_emu,
        slide_h_emu=args.slide_h_emu,
        font_name=args.font,
    )
    out_inpainted = args.out_dir / "result_inpainted.pptx"
    build_pptx(
        image_size=(px_w, px_h),
        background_path=inpainted_path,
        items=items,
        out_path=out_inpainted,
        slide_w_emu=args.slide_w_emu,
        slide_h_emu=args.slide_h_emu,
        font_name=args.font,
    )

    report = {
        "input": str(args.input),
        "image_size": [px_w, px_h],
        "slide_size_emu": [args.slide_w_emu, args.slide_h_emu],
        "effective_dpi": effective_dpi,
        "ocr_lines": len(items),
        "elapsed_sec": round(time.perf_counter() - started, 1),
        "items": [asdict(it) for it in items],
    }
    (args.out_dir / "report.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    print(f"\nresult_overlay.pptx   = 원본 배경 + textbox (안전)")
    print(f"result_inpainted.pptx = 인페인팅 배경 + textbox (편집성)")
    print(f"총 소요: {report['elapsed_sec']}s")


if __name__ == "__main__":
    main()
