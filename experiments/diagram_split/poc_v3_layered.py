"""Phase 3 PoC — v2 결함 수정 + 객체 레이어 분리.

v2 한계:
  - OCR polygon 이 단어 단위로 잘게 쪼개져 textbox 도 단어별 → 시각 어색
  - textbox height ≈ font height → LibreOffice 가 글자 자름/축소
  - 객체(아이콘/박스/화살표) 분리 미구현

v3 변경:
  1. OCR 라인 그룹핑: 동일 baseline + 좌-우 인접 라인을 한 textbox 로 합침.
  2. textbox padding: font_pt 비례로 충분히 크게 (height = font_pt * 1.6).
  3. 객체 분리: 인페인팅된 배경에서 connected components 분석 →
     일정 면적 이상의 색상 영역(컬러 박스, 아이콘) 을 alpha PNG crop 으로 추출 →
     별도 picture shape 로 PPTX 에 삽입 (편집 가능한 객체 레이어).
  4. 두 가지 PPTX 동시 생성:
        result_safe.pptx     — 원본 PNG 풀스크린 + textbox (시각 안전)
        result_layered.pptx  — 인페인팅 배경 + 객체 alpha + textbox (full 분리)
"""

from __future__ import annotations

import argparse
import json
import time
from dataclasses import asdict, dataclass, field
from pathlib import Path

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Pt

from star_slide.inpaint.lama import inpaint_with_mask
from star_slide.ocr.paddleocr_worker import OcrLine, run_ocr
from star_slide.ocr.text_attributes import (
    detect_korean,
    estimate_font_color,
    estimate_font_size_pt,
)

EMU_PER_INCH = 914400
PT_PER_INCH = 72.0
DEFAULT_SLIDE_W_EMU = 12192000
DEFAULT_SLIDE_H_EMU = 6858000
DEFAULT_FONT = "Apple SD Gothic Neo"


@dataclass
class TextItem:
    text: str
    bbox: tuple[float, float, float, float]
    polygon: list[tuple[float, float]]
    is_korean: bool
    font_pt: float
    color_hex: str
    mask_height_px: float


@dataclass
class TextLine:
    """그룹핑된 한 줄 — N개의 OCR 결과를 좌→우로 합침."""

    items: list[TextItem] = field(default_factory=list)

    @property
    def text(self) -> str:
        return " ".join(it.text for it in self.items)

    @property
    def bbox(self) -> tuple[float, float, float, float]:
        xs = [it.bbox[0] for it in self.items]
        ys = [it.bbox[1] for it in self.items]
        rs = [it.bbox[0] + it.bbox[2] for it in self.items]
        bs = [it.bbox[1] + it.bbox[3] for it in self.items]
        x = min(xs)
        y = min(ys)
        return (x, y, max(rs) - x, max(bs) - y)

    @property
    def font_pt(self) -> float:
        return float(np.median([it.font_pt for it in self.items]))

    @property
    def color_hex(self) -> str:
        # 가장 빈도 높은 색상
        from collections import Counter
        return Counter(it.color_hex for it in self.items).most_common(1)[0][0]

    @property
    def is_korean(self) -> bool:
        return any(it.is_korean for it in self.items)


@dataclass
class ObjectCrop:
    """객체 alpha PNG crop."""

    bbox: tuple[int, int, int, int]  # x, y, w, h (px)
    alpha_path: Path


def compute_polygon_height(poly: list[tuple[float, float]]) -> float:
    arr = np.asarray(poly, dtype=float)
    if len(arr) >= 4:
        side1 = np.linalg.norm(arr[1] - arr[0])
        side2 = np.linalg.norm(arr[2] - arr[1])
        return float(min(side1, side2))
    ys = arr[:, 1]
    return float(ys.max() - ys.min())


def per_polygon_mask(
    poly: list[tuple[float, float]],
    image_size: tuple[int, int],
) -> np.ndarray:
    w, h = image_size
    mask = np.zeros((h, w), dtype=np.uint8)
    pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
    cv2.fillPoly(mask, [pts], 1)
    return mask.astype(bool)


def lines_to_text_items(
    lines: list[OcrLine],
    image: Image.Image,
    *,
    effective_dpi: float,
) -> list[TextItem]:
    items: list[TextItem] = []
    w, h = image.size
    for line in lines:
        h_px = compute_polygon_height(line.bbox_quad)
        is_kr = detect_korean(line.text)
        font_pt = estimate_font_size_pt(h_px, is_korean=is_kr, dpi=effective_dpi)
        mask_bool = per_polygon_mask(line.bbox_quad, (w, h))
        color = estimate_font_color(image, mask_bool) or "#222222"
        items.append(
            TextItem(
                text=line.text,
                bbox=line.bbox,
                polygon=line.bbox_quad,
                is_korean=is_kr,
                font_pt=font_pt,
                color_hex=color,
                mask_height_px=h_px,
            )
        )
    return items


def group_items_by_line(
    items: list[TextItem],
    *,
    y_tolerance_ratio: float = 0.5,
    x_gap_ratio: float = 1.5,
) -> list[TextLine]:
    """비슷한 baseline + 좌→우 인접 OCR 결과를 한 라인으로 묶음.

    y_tolerance: 두 item 의 y-center 차이가 평균 height * y_tolerance_ratio 이하면 동일 라인.
    x_gap: 두 item 의 좌우 gap 이 평균 height * x_gap_ratio 이하면 같은 라인 그룹.
    """
    if not items:
        return []
    # 정렬: 위→아래, 좌→우
    sorted_items = sorted(items, key=lambda it: (it.bbox[1] + it.bbox[3] / 2, it.bbox[0]))

    lines: list[TextLine] = []
    for it in sorted_items:
        x, y, w, h = it.bbox
        y_center = y + h / 2
        merged = False
        for line in lines:
            lx, ly, lw, lh = line.bbox
            ly_center = ly + lh / 2
            avg_h = (lh + h) / 2
            # 같은 baseline + 폰트 사이즈 비슷 + 가까운 거리
            if (
                abs(ly_center - y_center) <= avg_h * y_tolerance_ratio
                and abs(line.font_pt - it.font_pt) / max(line.font_pt, 1) < 0.3
                and (x - (lx + lw)) <= avg_h * x_gap_ratio
                and (x - (lx + lw)) >= -avg_h  # 너무 멀리 떨어진 경우 제외
            ):
                line.items.append(it)
                merged = True
                break
        if not merged:
            lines.append(TextLine(items=[it]))
    return lines


def build_dark_pixel_inpaint_mask(
    image: Image.Image,
    polygons: list[list[tuple[float, float]]],
    *,
    luminance_threshold: int = 160,
    dilate: int = 3,
) -> np.ndarray:
    arr = np.asarray(image.convert("RGB"))
    h, w = arr.shape[:2]
    poly_mask = np.zeros((h, w), dtype=np.uint8)
    for poly in polygons:
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        cv2.fillPoly(poly_mask, [pts], 1)
    poly_mask_bool = poly_mask.astype(bool)

    lum = (
        0.299 * arr[:, :, 0].astype(np.float32)
        + 0.587 * arr[:, :, 1].astype(np.float32)
        + 0.114 * arr[:, :, 2].astype(np.float32)
    )
    dark_pixels = lum < luminance_threshold

    inpaint_mask = np.zeros((h, w), dtype=np.uint8)
    for poly in polygons:
        single = np.zeros((h, w), dtype=np.uint8)
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        cv2.fillPoly(single, [pts], 1)
        sb = single.astype(bool)
        in_poly_dark_ratio = float(dark_pixels[sb].mean()) if sb.any() else 0.0
        if in_poly_dark_ratio > 0.5:
            target = ~dark_pixels & sb
        else:
            target = dark_pixels & sb
        inpaint_mask[target] = 255

    if dilate > 0:
        kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (dilate, dilate))
        inpaint_mask = cv2.dilate(inpaint_mask, kernel, iterations=1)

    inpaint_mask[~poly_mask_bool] = 0
    return inpaint_mask


def extract_object_crops(
    inpainted_bg: Image.Image,
    *,
    out_dir: Path,
    min_area_ratio: float = 0.005,
    max_area_ratio: float = 0.5,
) -> list[ObjectCrop]:
    """인페인팅된 배경에서 큰 컬러 영역을 alpha PNG crop 으로 추출.

    전략:
      1. 흰색에 가까운 배경 픽셀을 제외 (luminance > 235)
      2. 색상 양자화 + connected components → 컬러 영역 후보
      3. 면적 임계 [min, max] 사이의 component 만 채택
      4. 각 component → alpha PNG crop (배경 투명)
    """
    arr = np.asarray(inpainted_bg.convert("RGB"))
    h, w = arr.shape[:2]
    total_area = h * w

    lum = (
        0.299 * arr[:, :, 0].astype(np.float32)
        + 0.587 * arr[:, :, 1].astype(np.float32)
        + 0.114 * arr[:, :, 2].astype(np.float32)
    )
    # 흰 배경 제외 → 의미 있는 색상 픽셀만 남김
    fg_mask = (lum < 240).astype(np.uint8) * 255

    # 작은 점/노이즈 제거 + 약간 connectivity 위해 close
    kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (5, 5))
    fg_mask = cv2.morphologyEx(fg_mask, cv2.MORPH_CLOSE, kernel)

    n, labels, stats, _ = cv2.connectedComponentsWithStats(fg_mask, connectivity=8)
    crops: list[ObjectCrop] = []
    out_dir.mkdir(parents=True, exist_ok=True)

    for i in range(1, n):  # 0 = background
        x, y, ww, hh, area = stats[i]
        if area / total_area < min_area_ratio or area / total_area > max_area_ratio:
            continue
        # 너무 가늘거나 작은 것 제외
        if ww < 30 or hh < 30:
            continue

        # 해당 component 만 마스킹한 alpha PNG crop
        comp_mask = (labels[y : y + hh, x : x + ww] == i).astype(np.uint8) * 255
        crop_rgb = arr[y : y + hh, x : x + ww]
        # RGBA 합성
        rgba = np.zeros((hh, ww, 4), dtype=np.uint8)
        rgba[:, :, :3] = crop_rgb
        rgba[:, :, 3] = comp_mask
        crop_path = out_dir / f"obj_{i:03d}_{x}_{y}.png"
        Image.fromarray(rgba, mode="RGBA").save(crop_path)
        crops.append(ObjectCrop(bbox=(int(x), int(y), int(ww), int(hh)), alpha_path=crop_path))
    return crops


def render_object_overlay(
    image: Image.Image,
    crops: list[ObjectCrop],
    out_path: Path,
) -> None:
    arr = np.asarray(image.convert("RGB")).copy()
    for c in crops:
        x, y, w, h = c.bbox
        cv2.rectangle(arr, (x, y), (x + w, y + h), (0, 200, 0), 3)
    Image.fromarray(arr).save(out_path)


def render_text_overlay(
    image: Image.Image,
    lines: list[TextLine],
    out_path: Path,
) -> None:
    arr = np.asarray(image.convert("RGB")).copy()
    for ln in lines:
        x, y, w, h = ln.bbox
        cv2.rectangle(arr, (int(x), int(y)), (int(x + w), int(y + h)), (255, 0, 0), 2)
        label = f"{ln.font_pt:.0f}pt"
        cv2.putText(
            arr, label, (int(x), max(15, int(y) - 4)),
            cv2.FONT_HERSHEY_SIMPLEX, 0.5, (0, 0, 255), 1, cv2.LINE_AA,
        )
    Image.fromarray(arr).save(out_path)


def add_textbox(
    slide,
    line: TextLine,
    *,
    scale_emu_per_px: float,
    font_name: str,
) -> None:
    x, y, w, h = line.bbox
    # 폰트 높이 기준 padding — textbox 가 폰트보다 충분히 커야 LibreOffice/PPT 에서
    # 글자 잘림 없음.
    font_inch = line.font_pt / PT_PER_INCH
    font_emu = int(font_inch * EMU_PER_INCH)
    # 가로: 텍스트 + 양쪽 0.5em padding (한글 기준 em ~= font_pt height)
    left = int(x * scale_emu_per_px) - font_emu // 2
    width = int(w * scale_emu_per_px) + font_emu
    # 세로: 폰트 height 의 1.6배 (안전)
    height = max(int(h * scale_emu_per_px), int(font_emu * 1.6))
    # textbox 의 baseline 이 OCR bbox 중앙에 오도록
    y_center_emu = int((y + h / 2) * scale_emu_per_px)
    top = y_center_emu - height // 2

    if width <= 0 or height <= 0:
        return

    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = line.text
    run.font.name = font_name
    run.font.size = Pt(round(line.font_pt, 1))
    try:
        r = int(line.color_hex[1:3], 16)
        g = int(line.color_hex[3:5], 16)
        b = int(line.color_hex[5:7], 16)
        run.font.color.rgb = RGBColor(r, g, b)
    except Exception:
        pass


def build_safe_pptx(
    *,
    image_size: tuple[int, int],
    background_path: Path,
    lines: list[TextLine],
    out_path: Path,
    slide_w_emu: int,
    slide_h_emu: int,
    font_name: str,
) -> None:
    """원본 PNG 풀스크린 + textbox (시각 안전 모드)."""
    px_w, _ = image_size
    scale = slide_w_emu / px_w
    prs = Presentation()
    prs.slide_width = Emu(slide_w_emu)
    prs.slide_height = Emu(slide_h_emu)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(
        str(background_path), 0, 0,
        width=Emu(slide_w_emu), height=Emu(slide_h_emu),
    )
    for ln in lines:
        add_textbox(slide, ln, scale_emu_per_px=scale, font_name=font_name)
    prs.save(str(out_path))


def build_layered_pptx(
    *,
    image_size: tuple[int, int],
    inpainted_bg_path: Path,
    object_crops: list[ObjectCrop],
    lines: list[TextLine],
    out_path: Path,
    slide_w_emu: int,
    slide_h_emu: int,
    font_name: str,
) -> None:
    """인페인팅 배경 + 객체 alpha + textbox (full 분리 모드)."""
    px_w, _ = image_size
    scale = slide_w_emu / px_w
    prs = Presentation()
    prs.slide_width = Emu(slide_w_emu)
    prs.slide_height = Emu(slide_h_emu)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    slide.shapes.add_picture(
        str(inpainted_bg_path), 0, 0,
        width=Emu(slide_w_emu), height=Emu(slide_h_emu),
    )
    for c in object_crops:
        x, y, w, h = c.bbox
        slide.shapes.add_picture(
            str(c.alpha_path),
            Emu(int(x * scale)),
            Emu(int(y * scale)),
            width=Emu(int(w * scale)),
            height=Emu(int(h * scale)),
        )
    for ln in lines:
        add_textbox(slide, ln, scale_emu_per_px=scale, font_name=font_name)
    prs.save(str(out_path))


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", type=Path, required=True)
    ap.add_argument("--out-dir", type=Path, required=True)
    ap.add_argument("--font", default=DEFAULT_FONT)
    ap.add_argument("--slide-w-emu", type=int, default=DEFAULT_SLIDE_W_EMU)
    ap.add_argument("--slide-h-emu", type=int, default=DEFAULT_SLIDE_H_EMU)
    args = ap.parse_args()

    args.out_dir.mkdir(parents=True, exist_ok=True)
    started = time.perf_counter()

    print(f"[1/8] 입력 로드: {args.input}")
    image = Image.open(args.input).convert("RGB")
    px_w, px_h = image.size
    slide_w_inch = args.slide_w_emu / EMU_PER_INCH
    effective_dpi = px_w / slide_w_inch
    print(f"     크기 {px_w}x{px_h}, effective_dpi={effective_dpi:.1f}")

    print("[2/8] PaddleOCR (한국어)")
    t0 = time.perf_counter()
    raw_lines = run_ocr(image, lang="korean")
    print(f"     {len(raw_lines)} 라인, {time.perf_counter() - t0:.1f}s")

    print("[3/8] 텍스트 속성 측정")
    items = lines_to_text_items(raw_lines, image, effective_dpi=effective_dpi)
    print(f"     items={len(items)}")

    print("[4/8] 인접 OCR 결과 라인 그룹핑")
    grouped = group_items_by_line(items)
    print(f"     {len(items)} items → {len(grouped)} text lines")
    render_text_overlay(image, grouped, args.out_dir / "01_text_lines.png")

    print("[5/8] dark-pixel inpaint mask + LaMa")
    inpaint_mask = build_dark_pixel_inpaint_mask(
        image, [it.polygon for it in items], luminance_threshold=160, dilate=3,
    )
    Image.fromarray(inpaint_mask, mode="L").save(args.out_dir / "02_inpaint_mask.png")
    t0 = time.perf_counter()
    inpainted = inpaint_with_mask(image, Image.fromarray(inpaint_mask, mode="L"))
    inpainted_path = args.out_dir / "03_inpainted_bg.png"
    inpainted.save(inpainted_path)
    print(f"     {time.perf_counter() - t0:.1f}s")

    print("[6/8] 객체 alpha PNG crop 추출 (connected components)")
    crops = extract_object_crops(
        inpainted, out_dir=args.out_dir / "objects",
        min_area_ratio=0.0015, max_area_ratio=0.4,
    )
    print(f"     {len(crops)} 개 객체 crop")
    render_object_overlay(image, crops, args.out_dir / "04_object_overlay.png")

    print("[7/8] 원본 배경 (safe 모드용)")
    orig_path = args.out_dir / "00_original_bg.png"
    image.save(orig_path)

    print("[8/8] PPTX 두 모드 생성")
    out_safe = args.out_dir / "result_safe.pptx"
    build_safe_pptx(
        image_size=(px_w, px_h), background_path=orig_path,
        lines=grouped, out_path=out_safe,
        slide_w_emu=args.slide_w_emu, slide_h_emu=args.slide_h_emu,
        font_name=args.font,
    )
    out_layered = args.out_dir / "result_layered.pptx"
    build_layered_pptx(
        image_size=(px_w, px_h), inpainted_bg_path=inpainted_path,
        object_crops=crops, lines=grouped, out_path=out_layered,
        slide_w_emu=args.slide_w_emu, slide_h_emu=args.slide_h_emu,
        font_name=args.font,
    )

    report = {
        "input": str(args.input),
        "image_size": [px_w, px_h],
        "slide_size_emu": [args.slide_w_emu, args.slide_h_emu],
        "effective_dpi": effective_dpi,
        "ocr_items": len(items),
        "grouped_lines": len(grouped),
        "object_crops": len(crops),
        "elapsed_sec": round(time.perf_counter() - started, 1),
        "lines": [
            {
                "text": ln.text,
                "bbox": ln.bbox,
                "font_pt": ln.font_pt,
                "color": ln.color_hex,
                "n_items": len(ln.items),
            }
            for ln in grouped
        ],
    }
    (args.out_dir / "report.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    print(f"\nresult_safe.pptx     = 원본 PNG + textbox")
    print(f"result_layered.pptx  = 인페인팅 + alpha 객체 + textbox")
    print(f"총 소요: {report['elapsed_sec']}s")


if __name__ == "__main__":
    main()
