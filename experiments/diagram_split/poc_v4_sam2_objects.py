"""Phase 4 PoC — v3 의 객체 분리(connected components) 를 SAM2 auto-mask 로 교체.

v3 한계:
  - connected components 는 인접한 같은-색 객체를 한 덩어리로 합침.
    (예: "사용자" 영역의 5개 사람 아이콘 → 1개 crop)
  - 흰 배경 임계만으론 컬러 박스 안의 아이콘 분리 불가.

v4 접근:
  1. PaddleOCR (한국어) → 텍스트 polygon (v3 그대로)
  2. SAM2.1 auto-mask (transformers pipeline, MPS) → 슬라이드 위 모든 객체 마스크
  3. classify_masks → TEXT/SHAPE/BACKGROUND/NOISE
       - TEXT 마스크: OCR polygon 안에 ≥70% 들어가는 SAM 마스크 (글자 stroke)
       - SHAPE 마스크: 그 외 의미있는 마스크 → 객체 후보
  4. SHAPE 마스크 후처리:
       - IoU > 0.7 중복 제거 (큰 마스크 우선)
       - 면적 [min, max] 필터
       - 텍스트 polygon 영역과 차집합 (텍스트 잔재 제거)
  5. 각 SHAPE 마스크 → 정밀 alpha PNG crop (마스크 그대로, 사각 crop X)
  6. PPTX 조립:
       - 배경 1: 흰색 캔버스 (객체만으로 재구성)
       - 배경 2: 인페인팅된 원본 (그라데이션 보존)
       두 가지 PPTX 생성해서 비교.
"""

from __future__ import annotations

import argparse
import json
import time
from dataclasses import asdict, dataclass, field
from pathlib import Path

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Pt

from star_slide.inpaint.lama import inpaint_with_mask
from star_slide.ocr.paddleocr_worker import OcrLine, run_ocr
from star_slide.ocr.text_attributes import (
    detect_korean,
    estimate_font_color,
    estimate_font_size_pt,
)
from star_slide.segmentation.classify import MaskClass, classify_masks
from star_slide.segmentation.sam2_auto import Sam2Mask, run_sam2_auto

EMU_PER_INCH = 914400
PT_PER_INCH = 72.0
DEFAULT_SLIDE_W_EMU = 12192000
DEFAULT_SLIDE_H_EMU = 6858000
DEFAULT_FONT = "Apple SD Gothic Neo"


@dataclass
class TextItem:
    text: str
    bbox: tuple[float, float, float, float]
    polygon: list[tuple[float, float]]
    is_korean: bool
    font_pt: float
    color_hex: str
    mask_height_px: float


@dataclass
class TextLine:
    items: list[TextItem] = field(default_factory=list)

    @property
    def text(self) -> str:
        return " ".join(it.text for it in self.items)

    @property
    def bbox(self) -> tuple[float, float, float, float]:
        xs = [it.bbox[0] for it in self.items]
        ys = [it.bbox[1] for it in self.items]
        rs = [it.bbox[0] + it.bbox[2] for it in self.items]
        bs = [it.bbox[1] + it.bbox[3] for it in self.items]
        x = min(xs)
        y = min(ys)
        return (x, y, max(rs) - x, max(bs) - y)

    @property
    def font_pt(self) -> float:
        return float(np.median([it.font_pt for it in self.items]))

    @property
    def color_hex(self) -> str:
        from collections import Counter
        return Counter(it.color_hex for it in self.items).most_common(1)[0][0]


@dataclass
class ObjectCrop:
    """SAM2 마스크 → 정밀 alpha PNG crop."""

    bbox: tuple[int, int, int, int]
    score: float
    area_ratio: float
    alpha_path: Path


# ---------------- OCR + 텍스트 라인 그룹핑 (v3 그대로) ----------------


def compute_polygon_height(poly: list[tuple[float, float]]) -> float:
    arr = np.asarray(poly, dtype=float)
    if len(arr) >= 4:
        side1 = np.linalg.norm(arr[1] - arr[0])
        side2 = np.linalg.norm(arr[2] - arr[1])
        return float(min(side1, side2))
    ys = arr[:, 1]
    return float(ys.max() - ys.min())


def per_polygon_mask(
    poly: list[tuple[float, float]],
    image_size: tuple[int, int],
) -> np.ndarray:
    w, h = image_size
    mask = np.zeros((h, w), dtype=np.uint8)
    pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
    cv2.fillPoly(mask, [pts], 1)
    return mask.astype(bool)


def lines_to_text_items(
    lines: list[OcrLine],
    image: Image.Image,
    *,
    effective_dpi: float,
) -> list[TextItem]:
    items: list[TextItem] = []
    w, h = image.size
    for line in lines:
        h_px = compute_polygon_height(line.bbox_quad)
        is_kr = detect_korean(line.text)
        font_pt = estimate_font_size_pt(h_px, is_korean=is_kr, dpi=effective_dpi)
        mask_bool = per_polygon_mask(line.bbox_quad, (w, h))
        color = estimate_font_color(image, mask_bool) or "#222222"
        items.append(
            TextItem(
                text=line.text, bbox=line.bbox, polygon=line.bbox_quad,
                is_korean=is_kr, font_pt=font_pt, color_hex=color,
                mask_height_px=h_px,
            )
        )
    return items


def group_items_by_line(items: list[TextItem]) -> list[TextLine]:
    if not items:
        return []
    sorted_items = sorted(items, key=lambda it: (it.bbox[1] + it.bbox[3] / 2, it.bbox[0]))
    lines: list[TextLine] = []
    for it in sorted_items:
        x, y, w, h = it.bbox
        y_center = y + h / 2
        merged = False
        for line in lines:
            lx, ly, lw, lh = line.bbox
            ly_center = ly + lh / 2
            avg_h = (lh + h) / 2
            if (
                abs(ly_center - y_center) <= avg_h * 0.5
                and abs(line.font_pt - it.font_pt) / max(line.font_pt, 1) < 0.3
                and -avg_h <= (x - (lx + lw)) <= avg_h * 1.5
            ):
                line.items.append(it)
                merged = True
                break
        if not merged:
            lines.append(TextLine(items=[it]))
    return lines


# ---------------- 인페인팅 마스크 (v3 그대로 — dark pixel only) ----------------


def build_dark_pixel_inpaint_mask(
    image: Image.Image,
    polygons: list[list[tuple[float, float]]],
    *,
    luminance_threshold: int = 160,
    dilate: int = 3,
) -> np.ndarray:
    arr = np.asarray(image.convert("RGB"))
    h, w = arr.shape[:2]
    poly_mask = np.zeros((h, w), dtype=np.uint8)
    for poly in polygons:
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        cv2.fillPoly(poly_mask, [pts], 1)
    poly_mask_bool = poly_mask.astype(bool)
    lum = (
        0.299 * arr[:, :, 0].astype(np.float32)
        + 0.587 * arr[:, :, 1].astype(np.float32)
        + 0.114 * arr[:, :, 2].astype(np.float32)
    )
    dark = lum < luminance_threshold
    out = np.zeros((h, w), dtype=np.uint8)
    for poly in polygons:
        single = np.zeros((h, w), dtype=np.uint8)
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        cv2.fillPoly(single, [pts], 1)
        sb = single.astype(bool)
        dr = float(dark[sb].mean()) if sb.any() else 0.0
        target = (~dark & sb) if dr > 0.5 else (dark & sb)
        out[target] = 255
    if dilate > 0:
        k = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (dilate, dilate))
        out = cv2.dilate(out, k, iterations=1)
    out[~poly_mask_bool] = 0
    return out


# ---------------- SAM2 객체 마스크 후처리 ----------------


def mask_iou(a: np.ndarray, b: np.ndarray) -> float:
    if a.shape != b.shape:
        return 0.0
    inter = np.logical_and(a, b).sum()
    union = np.logical_or(a, b).sum()
    return float(inter) / float(union) if union else 0.0


def filter_and_dedupe_shapes(
    shapes: list[Sam2Mask],
    image_size: tuple[int, int],
    *,
    min_area_ratio: float = 0.0008,
    max_area_ratio: float = 0.45,
    iou_threshold: float = 0.7,
) -> list[Sam2Mask]:
    """SHAPE 마스크 면적 필터 + 중복 제거 (큰 마스크 우선 NMS-like).

    SAM2 auto 는 같은 객체에 여러 마스크 (small/medium/large) 를 만드는 경향이 있어
    IoU 가 큰 것끼리 묶고 가장 큰(또는 score 높은) 것만 남긴다.
    """
    w, h = image_size
    area_total = w * h
    candidates = [
        m for m in shapes
        if min_area_ratio <= m.area / area_total <= max_area_ratio
    ]
    # 면적 큰 순으로 정렬해 가장 큰 마스크가 다른 비슷한 마스크를 흡수
    candidates.sort(key=lambda m: -m.area)
    kept: list[Sam2Mask] = []
    for m in candidates:
        ok = True
        for k in kept:
            if mask_iou(m.segmentation, k.segmentation) > iou_threshold:
                ok = False
                break
        if ok:
            kept.append(m)
    return kept


def subtract_text_pixels(
    mask: np.ndarray,
    text_polygons: list[list[tuple[float, float]]],
) -> np.ndarray:
    """객체 마스크에서 텍스트 polygon 픽셀 제거 (글자 잔재 방지)."""
    out = mask.copy()
    h, w = out.shape
    text_mask = np.zeros((h, w), dtype=np.uint8)
    for poly in text_polygons:
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        cv2.fillPoly(text_mask, [pts], 1)
    out &= ~text_mask.astype(bool)
    return out


def export_alpha_crop(
    image: Image.Image,
    mask: np.ndarray,
    *,
    out_dir: Path,
    name: str,
) -> ObjectCrop | None:
    """bool 마스크 → 정밀 alpha PNG crop (사각 bbox + 마스크 외부 투명)."""
    ys, xs = np.where(mask)
    if len(ys) == 0:
        return None
    x0, x1 = int(xs.min()), int(xs.max()) + 1
    y0, y1 = int(ys.min()), int(ys.max()) + 1
    arr = np.asarray(image.convert("RGB"))
    crop_rgb = arr[y0:y1, x0:x1]
    crop_mask = mask[y0:y1, x0:x1]
    rgba = np.zeros((y1 - y0, x1 - x0, 4), dtype=np.uint8)
    rgba[:, :, :3] = crop_rgb
    rgba[:, :, 3] = (crop_mask.astype(np.uint8)) * 255
    out_dir.mkdir(parents=True, exist_ok=True)
    p = out_dir / name
    Image.fromarray(rgba, mode="RGBA").save(p)
    h, w = arr.shape[:2]
    return ObjectCrop(
        bbox=(x0, y0, x1 - x0, y1 - y0),
        score=0.0,
        area_ratio=float(mask.sum()) / (h * w),
        alpha_path=p,
    )


# ---------------- 시각화 ----------------


def render_mask_overlay(
    image: Image.Image,
    masks: list[Sam2Mask],
    out_path: Path,
) -> None:
    """SAM2 마스크 컬러 오버레이 (디버그)."""
    rng = np.random.default_rng(42)
    arr = np.asarray(image.convert("RGB")).copy().astype(np.float32)
    overlay = np.zeros_like(arr)
    for m in masks:
        color = rng.integers(50, 255, size=3)
        overlay[m.segmentation] = color
    blended = arr * 0.55 + overlay * 0.45
    Image.fromarray(np.clip(blended, 0, 255).astype(np.uint8)).save(out_path)


def render_object_overlay(
    image: Image.Image,
    crops: list[ObjectCrop],
    out_path: Path,
) -> None:
    arr = np.asarray(image.convert("RGB")).copy()
    for c in crops:
        x, y, w, h = c.bbox
        cv2.rectangle(arr, (x, y), (x + w, y + h), (0, 200, 0), 3)
    Image.fromarray(arr).save(out_path)


def render_text_overlay(
    image: Image.Image,
    lines: list[TextLine],
    out_path: Path,
) -> None:
    arr = np.asarray(image.convert("RGB")).copy()
    for ln in lines:
        x, y, w, h = ln.bbox
        cv2.rectangle(arr, (int(x), int(y)), (int(x + w), int(y + h)), (255, 0, 0), 2)
        cv2.putText(
            arr, f"{ln.font_pt:.0f}pt", (int(x), max(15, int(y) - 4)),
            cv2.FONT_HERSHEY_SIMPLEX, 0.5, (0, 0, 255), 1, cv2.LINE_AA,
        )
    Image.fromarray(arr).save(out_path)


# ---------------- PPTX 빌드 ----------------


def add_textbox(
    slide,
    line: TextLine,
    *,
    scale_emu_per_px: float,
    font_name: str,
) -> None:
    x, y, w, h = line.bbox
    font_emu = int(line.font_pt / PT_PER_INCH * EMU_PER_INCH)
    left = int(x * scale_emu_per_px) - font_emu // 2
    width = int(w * scale_emu_per_px) + font_emu
    height = max(int(h * scale_emu_per_px), int(font_emu * 1.6))
    y_center_emu = int((y + h / 2) * scale_emu_per_px)
    top = y_center_emu - height // 2
    if width <= 0 or height <= 0:
        return
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = line.text
    run.font.name = font_name
    run.font.size = Pt(round(line.font_pt, 1))
    try:
        r = int(line.color_hex[1:3], 16)
        g = int(line.color_hex[3:5], 16)
        b = int(line.color_hex[5:7], 16)
        run.font.color.rgb = RGBColor(r, g, b)
    except Exception:
        pass


def build_pptx(
    *,
    image_size: tuple[int, int],
    background_path: Path | None,
    background_color: tuple[int, int, int] | None,
    object_crops: list[ObjectCrop],
    lines: list[TextLine],
    out_path: Path,
    slide_w_emu: int,
    slide_h_emu: int,
    font_name: str,
) -> None:
    px_w, _ = image_size
    scale = slide_w_emu / px_w
    prs = Presentation()
    prs.slide_width = Emu(slide_w_emu)
    prs.slide_height = Emu(slide_h_emu)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    if background_path is not None:
        slide.shapes.add_picture(
            str(background_path), 0, 0,
            width=Emu(slide_w_emu), height=Emu(slide_h_emu),
        )
    elif background_color is not None:
        # 흰색 배경 사각형
        from pptx.enum.shapes import MSO_SHAPE
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Emu(slide_w_emu), Emu(slide_h_emu)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*background_color)
        bg.line.fill.background()

    for c in object_crops:
        x, y, w, h = c.bbox
        slide.shapes.add_picture(
            str(c.alpha_path),
            Emu(int(x * scale)), Emu(int(y * scale)),
            width=Emu(int(w * scale)), height=Emu(int(h * scale)),
        )
    for ln in lines:
        add_textbox(slide, ln, scale_emu_per_px=scale, font_name=font_name)
    prs.save(str(out_path))


# ---------------- main ----------------


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", type=Path, required=True)
    ap.add_argument("--out-dir", type=Path, required=True)
    ap.add_argument("--font", default=DEFAULT_FONT)
    ap.add_argument("--slide-w-emu", type=int, default=DEFAULT_SLIDE_W_EMU)
    ap.add_argument("--slide-h-emu", type=int, default=DEFAULT_SLIDE_H_EMU)
    ap.add_argument("--sam-points", type=int, default=32, help="SAM2 points_per_side")
    ap.add_argument("--sam-max", type=int, default=300, help="SAM2 max masks")
    ap.add_argument("--sam-resize", type=int, default=1280,
                    help="SAM2 입력 리사이즈 폭. 0=원본 사용 (정밀하지만 느림)")
    args = ap.parse_args()

    args.out_dir.mkdir(parents=True, exist_ok=True)
    started = time.perf_counter()

    print(f"[1/9] 입력 로드: {args.input}")
    image = Image.open(args.input).convert("RGB")
    px_w, px_h = image.size
    slide_w_inch = args.slide_w_emu / EMU_PER_INCH
    effective_dpi = px_w / slide_w_inch
    print(f"     크기 {px_w}x{px_h}, effective_dpi={effective_dpi:.1f}")

    print("[2/9] PaddleOCR (한국어)")
    t0 = time.perf_counter()
    raw_lines = run_ocr(image, lang="korean")
    print(f"     {len(raw_lines)} 라인, {time.perf_counter() - t0:.1f}s")

    print("[3/9] 텍스트 속성 측정 + 라인 그룹핑")
    items = lines_to_text_items(raw_lines, image, effective_dpi=effective_dpi)
    grouped = group_items_by_line(items)
    print(f"     items={len(items)} → grouped={len(grouped)}")
    render_text_overlay(image, grouped, args.out_dir / "01_text_lines.png")

    print(f"[4/9] SAM2.1 auto-mask (points_per_side={args.sam_points})")
    sam_input = image
    if args.sam_resize and args.sam_resize > 0 and px_w > args.sam_resize:
        scale_sam = args.sam_resize / px_w
        sam_input = image.resize((args.sam_resize, int(px_h * scale_sam)))
        print(f"     SAM 입력 {sam_input.size} (원본 → 리사이즈)")
    sam_result = run_sam2_auto(
        sam_input,
        points_per_side=args.sam_points,
        max_masks=args.sam_max,
    )
    print(f"     {len(sam_result.masks)} 마스크, {sam_result.elapsed_sec:.1f}s, dev={sam_result.device}")

    # SAM 마스크를 원본 해상도로 업샘플 (마스크 → 원본 크기)
    if sam_input.size != image.size:
        full_masks: list[Sam2Mask] = []
        for m in sam_result.masks:
            seg_uint = (m.segmentation.astype(np.uint8)) * 255
            seg_full = cv2.resize(seg_uint, (px_w, px_h), interpolation=cv2.INTER_NEAREST)
            seg_bool = seg_full > 127
            ys, xs = np.where(seg_bool)
            if len(ys) == 0:
                continue
            bbox = (int(xs.min()), int(ys.min()),
                    int(xs.max() - xs.min() + 1), int(ys.max() - ys.min() + 1))
            full_masks.append(
                Sam2Mask(bbox=bbox, segmentation=seg_bool, score=m.score, area=int(seg_bool.sum()))
            )
    else:
        full_masks = list(sam_result.masks)

    render_mask_overlay(image, full_masks, args.out_dir / "02_sam2_masks.png")

    print("[5/9] classify_masks → TEXT/SHAPE 분리")
    classification = classify_masks(full_masks, raw_lines, image_size=(px_w, px_h))
    shapes = classification.shape_masks
    print(
        f"     TEXT={len(classification.text_masks)} "
        f"SHAPE={len(shapes)} "
        f"(BG/NOISE 제외, {len(full_masks)} 중)"
    )

    print("[6/9] SHAPE 후처리: 면적 필터 + IoU 중복 제거")
    shapes_clean = filter_and_dedupe_shapes(
        shapes, image_size=(px_w, px_h),
        min_area_ratio=0.0008, max_area_ratio=0.45, iou_threshold=0.7,
    )
    print(f"     {len(shapes)} → {len(shapes_clean)}")

    print("[7/9] 객체 alpha PNG crop 추출")
    text_polys = [it.polygon for it in items]
    obj_dir = args.out_dir / "objects"
    crops: list[ObjectCrop] = []
    for i, m in enumerate(shapes_clean):
        clean_seg = subtract_text_pixels(m.segmentation, text_polys)
        if clean_seg.sum() < 100:  # 텍스트 빼고 거의 안 남음
            continue
        crop = export_alpha_crop(image, clean_seg, out_dir=obj_dir, name=f"obj_{i:03d}.png")
        if crop is None:
            continue
        crops.append(ObjectCrop(
            bbox=crop.bbox, score=m.score, area_ratio=crop.area_ratio,
            alpha_path=crop.alpha_path,
        ))
    print(f"     {len(crops)} 객체 alpha crop 저장")
    render_object_overlay(image, crops, args.out_dir / "03_object_overlay.png")

    print("[8/9] 인페인팅 배경 (텍스트 stroke 만 제거)")
    inpaint_mask = build_dark_pixel_inpaint_mask(image, text_polys, dilate=3)
    Image.fromarray(inpaint_mask, mode="L").save(args.out_dir / "04_inpaint_mask.png")
    t0 = time.perf_counter()
    inpainted = inpaint_with_mask(image, Image.fromarray(inpaint_mask, mode="L"))
    inpainted_path = args.out_dir / "05_inpainted_bg.png"
    inpainted.save(inpainted_path)
    print(f"     {time.perf_counter() - t0:.1f}s")

    print("[9/9] PPTX 두 모드 생성")
    out_white = args.out_dir / "result_white_bg.pptx"
    build_pptx(
        image_size=(px_w, px_h),
        background_path=None, background_color=(255, 255, 255),
        object_crops=crops, lines=grouped, out_path=out_white,
        slide_w_emu=args.slide_w_emu, slide_h_emu=args.slide_h_emu,
        font_name=args.font,
    )
    out_inpainted = args.out_dir / "result_inpainted.pptx"
    build_pptx(
        image_size=(px_w, px_h),
        background_path=inpainted_path, background_color=None,
        object_crops=crops, lines=grouped, out_path=out_inpainted,
        slide_w_emu=args.slide_w_emu, slide_h_emu=args.slide_h_emu,
        font_name=args.font,
    )

    report = {
        "input": str(args.input),
        "image_size": [px_w, px_h],
        "ocr_items": len(items),
        "grouped_lines": len(grouped),
        "sam_masks_total": len(full_masks),
        "sam_text_masks": len(classification.text_masks),
        "sam_shape_masks_raw": len(shapes),
        "sam_shape_masks_kept": len(shapes_clean),
        "object_crops": len(crops),
        "elapsed_sec": round(time.perf_counter() - started, 1),
    }
    (args.out_dir / "report.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    print(f"\nresult_white_bg.pptx   = 흰 배경 + 객체 + textbox (full 분리)")
    print(f"result_inpainted.pptx  = 인페인팅 배경 + 객체 + textbox")
    print(f"총 소요: {report['elapsed_sec']}s")


if __name__ == "__main__":
    main()
