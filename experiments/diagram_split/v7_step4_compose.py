"""v7-Step4 — 텍스트 제거 배경 + alpha 객체 + Codex textbox 재조합.

이중 노출 없음: 배경에 텍스트 없음 + 객체 alpha 에도 텍스트 없음 +
textbox 는 Codex 가 추출한 정확한 위치/사이즈/색상으로 별도 레이어.

폰트 설정:
  - font.name + East Asian (`<a:ea>`) 동시 명시 → PowerPoint/Keynote 한글 정확
  - font_size_pt = font_size_px * 72 / effective_dpi
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

EMU_PER_INCH = 914400
PT_PER_INCH = 72.0
SLIDE_W_EMU = 12192000   # 13.333"
SLIDE_H_EMU = 6858000    #  7.5"
DEFAULT_FONT = "Apple SD Gothic Neo"


def set_east_asian_font(rPr, font_name: str) -> None:
    """East Asian + Complex Script 폰트 명시 (한글 fallback 방지)."""
    for tag in ("a:ea", "a:cs"):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)
    cs = etree.SubElement(rPr, qn("a:cs"))
    cs.set("typeface", font_name)


def add_textbox(slide, item, *, scale, effective_dpi, font_name):
    x, y, w, h = item["bbox"]
    font_pt = item["font_size_px"] * PT_PER_INCH / effective_dpi
    # textbox padding — font height 의 0.5em (위/아래)
    pad_emu = int(font_pt / PT_PER_INCH * EMU_PER_INCH * 0.3)

    left = int(x * scale) - pad_emu // 2
    top = int(y * scale) - pad_emu // 2
    width = int(w * scale) + pad_emu
    # height: bbox 보다 font height 의 1.4배 보장 (글자 잘림 방지)
    min_h_emu = int(font_pt / PT_PER_INCH * EMU_PER_INCH * 1.4)
    height = max(int(h * scale) + pad_emu, min_h_emu)

    if width <= 0 or height <= 0:
        return

    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, attr, Emu(0))
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    align = item.get("alignment", "left")
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)

    run = p.add_run()
    run.text = item["text"]
    run.font.name = font_name
    run.font.size = Pt(round(font_pt, 1))
    run.font.bold = bool(item.get("is_bold", False))
    color = item.get("color", "#222222")
    try:
        r = int(color[1:3], 16)
        g = int(color[3:5], 16)
        b = int(color[5:7], 16)
        run.font.color.rgb = RGBColor(r, g, b)
    except Exception:
        pass
    set_east_asian_font(run._r.get_or_add_rPr(), font_name)


def build_pptx(
    *, image_size, background_path, layers_json, layers_dir,
    text_json, out_path, font_name, with_text, with_objects,
):
    px_w, _ = image_size
    scale = SLIDE_W_EMU / px_w
    effective_dpi = px_w / (SLIDE_W_EMU / EMU_PER_INCH)

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 배경
    if background_path is not None:
        slide.shapes.add_picture(
            str(background_path), 0, 0,
            width=Emu(SLIDE_W_EMU), height=Emu(SLIDE_H_EMU),
        )

    # 객체 alpha (z-order: 큰 마스크부터 = 뒤에 깔림)
    if with_objects and layers_json:
        layers = json.loads(layers_json.read_text(encoding="utf-8")).get("layers", [])
        for ly in sorted(layers, key=lambda l: l["z"]):
            x, y, w, h = ly["bbox"]
            slide.shapes.add_picture(
                str(layers_dir / ly["alpha_path"]),
                Emu(int(x * scale)), Emu(int(y * scale)),
                width=Emu(int(w * scale)), height=Emu(int(h * scale)),
            )

    # 텍스트 (최상단)
    if with_text and text_json:
        texts = json.loads(text_json.read_text(encoding="utf-8")).get("texts", [])
        for t in texts:
            add_textbox(slide, t, scale=scale, effective_dpi=effective_dpi, font_name=font_name)

    prs.save(str(out_path))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--clean-image", type=Path, required=True)
    ap.add_argument("--text-json", type=Path, required=True)
    ap.add_argument("--layers-json", type=Path, required=True)
    ap.add_argument("--layers-dir", type=Path, required=True,
                    help="layers_json 의 alpha_path 의 base 디렉토리")
    ap.add_argument("--out-dir", type=Path, required=True)
    ap.add_argument("--font", default=DEFAULT_FONT)
    args = ap.parse_args()

    args.out_dir.mkdir(parents=True, exist_ok=True)
    image = Image.open(args.clean_image).convert("RGB")
    px_w, px_h = image.size
    print(f"입력 {px_w}x{px_h}, slide {SLIDE_W_EMU}x{SLIDE_H_EMU} EMU")

    # 4가지 PPTX 생성 (디버그/비교용):
    #   A) 배경만 — 시각 안전성 검증
    #   B) 배경 + 객체 — 객체 분리 검증
    #   C) 배경 + 텍스트 — textbox 위치/폰트 검증
    #   D) 배경 + 객체 + 텍스트 — 최종
    out_a = args.out_dir / "v7_A_bg_only.pptx"
    build_pptx(
        image_size=(px_w, px_h), background_path=args.clean_image,
        layers_json=None, layers_dir=args.layers_dir,
        text_json=None, out_path=out_a,
        font_name=args.font, with_text=False, with_objects=False,
    )
    print(f"  A) {out_a}")

    out_b = args.out_dir / "v7_B_bg_objects.pptx"
    build_pptx(
        image_size=(px_w, px_h), background_path=args.clean_image,
        layers_json=args.layers_json, layers_dir=args.layers_dir,
        text_json=None, out_path=out_b,
        font_name=args.font, with_text=False, with_objects=True,
    )
    print(f"  B) {out_b}")

    out_c = args.out_dir / "v7_C_bg_text.pptx"
    build_pptx(
        image_size=(px_w, px_h), background_path=args.clean_image,
        layers_json=None, layers_dir=args.layers_dir,
        text_json=args.text_json, out_path=out_c,
        font_name=args.font, with_text=True, with_objects=False,
    )
    print(f"  C) {out_c}")

    out_d = args.out_dir / "v7_D_full.pptx"
    build_pptx(
        image_size=(px_w, px_h), background_path=args.clean_image,
        layers_json=args.layers_json, layers_dir=args.layers_dir,
        text_json=args.text_json, out_path=out_d,
        font_name=args.font, with_text=True, with_objects=True,
    )
    print(f"  D) {out_d}")
    print("\nv7_D_full.pptx 가 최종 — 텍스트 + 객체 + 깨끗한 배경")


if __name__ == "__main__":
    main()
