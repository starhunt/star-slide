"""Phase 6 PoC — 시각 충실도 우선 + 한글 폰트 정확화.

v5 결함 진단:
  - alpha PNG 가 "텍스트/자식 영역 다 뺀 잔해" 라 PPTX 에서 의미 없게 보임
    (layer_003 따옴표 박스 = 글자 자리에 흰 구멍, layer_009 = 텅 빈 박스)
  - 한글이 LibreOffice 에서 fallback 폰트로 깨짐 (East Asian font 메타 없음)

v6 핵심 변경:
  1. alpha 에서 자식 마스크 영역을 빼지 않음 — 큰 박스 + 안의 픽토그램 모두 보존.
     z-order 로 시각 합성: 큰 박스(뒤) → 픽토그램(중) → textbox(앞).
     텍스트 stroke 만 dark-pixel mask 로 alpha 에서 제거 → 글자 자국 사라짐.
  2. textbox 에 East Asian font 명시 (`<a:ea typeface="Apple SD Gothic Neo">`)
     → LibreOffice/PowerPoint 모두 한글 글리프에 정확한 폰트 사용.

결과:
  - 인페인팅 배경 + 알파 객체들 시각적으로 원본과 거의 동일.
  - 사용자가 객체 (사람 아이콘 등) 을 클릭하면 별도 shape 로 선택 가능.
  - 픽토그램 이동 시 큰 박스 안의 픽토그램 픽셀이 그대로 노출 (이중 픽셀 패턴).
  - textbox 는 편집 가능, 한글 정확.
"""

from __future__ import annotations

import argparse
import json
import time
from dataclasses import dataclass, field
from pathlib import Path

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from lxml import etree

from star_slide.inpaint.lama import inpaint_with_mask
from star_slide.ocr.paddleocr_worker import OcrLine, run_ocr
from star_slide.ocr.text_attributes import (
    detect_korean,
    estimate_font_color,
    estimate_font_size_pt,
)
from star_slide.segmentation.classify import classify_masks
from star_slide.segmentation.sam2_auto import Sam2Mask, run_sam2_auto

EMU_PER_INCH = 914400
PT_PER_INCH = 72.0
DEFAULT_SLIDE_W_EMU = 12192000
DEFAULT_SLIDE_H_EMU = 6858000
DEFAULT_FONT = "Apple SD Gothic Neo"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


@dataclass
class TextItem:
    text: str
    bbox: tuple[float, float, float, float]
    polygon: list[tuple[float, float]]
    is_korean: bool
    font_pt: float
    color_hex: str


@dataclass
class TextLine:
    items: list[TextItem] = field(default_factory=list)

    @property
    def text(self) -> str:
        return " ".join(it.text for it in self.items)

    @property
    def bbox(self) -> tuple[float, float, float, float]:
        xs = [it.bbox[0] for it in self.items]
        ys = [it.bbox[1] for it in self.items]
        rs = [it.bbox[0] + it.bbox[2] for it in self.items]
        bs = [it.bbox[1] + it.bbox[3] for it in self.items]
        x = min(xs); y = min(ys)
        return (x, y, max(rs) - x, max(bs) - y)

    @property
    def font_pt(self) -> float:
        return float(np.median([it.font_pt for it in self.items]))

    @property
    def color_hex(self) -> str:
        from collections import Counter
        return Counter(it.color_hex for it in self.items).most_common(1)[0][0]


@dataclass
class ObjectLayer:
    bbox: tuple[int, int, int, int]
    z: int
    area_ratio: float
    alpha_path: Path


# ---------------- OCR (v5 그대로) ----------------


def compute_polygon_height(poly):
    arr = np.asarray(poly, dtype=float)
    if len(arr) >= 4:
        side1 = np.linalg.norm(arr[1] - arr[0])
        side2 = np.linalg.norm(arr[2] - arr[1])
        return float(min(side1, side2))
    ys = arr[:, 1]
    return float(ys.max() - ys.min())


def per_polygon_mask(poly, image_size):
    w, h = image_size
    mask = np.zeros((h, w), dtype=np.uint8)
    pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
    cv2.fillPoly(mask, [pts], 1)
    return mask.astype(bool)


def lines_to_text_items(lines, image, *, effective_dpi):
    items = []
    w, h = image.size
    for line in lines:
        h_px = compute_polygon_height(line.bbox_quad)
        is_kr = detect_korean(line.text)
        font_pt = estimate_font_size_pt(h_px, is_korean=is_kr, dpi=effective_dpi)
        mask_bool = per_polygon_mask(line.bbox_quad, (w, h))
        color = estimate_font_color(image, mask_bool) or "#222222"
        items.append(TextItem(
            text=line.text, bbox=line.bbox, polygon=line.bbox_quad,
            is_korean=is_kr, font_pt=font_pt, color_hex=color,
        ))
    return items


def group_items_by_line(items):
    if not items:
        return []
    sorted_items = sorted(items, key=lambda it: (it.bbox[1] + it.bbox[3] / 2, it.bbox[0]))
    lines = []
    for it in sorted_items:
        x, y, w, h = it.bbox
        y_center = y + h / 2
        merged = False
        for line in lines:
            lx, ly, lw, lh = line.bbox
            ly_center = ly + lh / 2
            avg_h = (lh + h) / 2
            if (
                abs(ly_center - y_center) <= avg_h * 0.5
                and abs(line.font_pt - it.font_pt) / max(line.font_pt, 1) < 0.3
                and -avg_h <= (x - (lx + lw)) <= avg_h * 1.5
            ):
                line.items.append(it)
                merged = True
                break
        if not merged:
            lines.append(TextLine(items=[it]))
    return lines


# ---------------- 인페인팅 마스크 ----------------


def build_dark_pixel_text_mask(image, polygons, *, luminance_threshold=160, dilate=3):
    arr = np.asarray(image.convert("RGB"))
    h, w = arr.shape[:2]
    poly_mask = np.zeros((h, w), dtype=np.uint8)
    for poly in polygons:
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        cv2.fillPoly(poly_mask, [pts], 1)
    poly_mask_bool = poly_mask.astype(bool)
    lum = (
        0.299 * arr[:, :, 0].astype(np.float32)
        + 0.587 * arr[:, :, 1].astype(np.float32)
        + 0.114 * arr[:, :, 2].astype(np.float32)
    )
    dark = lum < luminance_threshold
    out = np.zeros((h, w), dtype=np.uint8)
    for poly in polygons:
        single = np.zeros((h, w), dtype=np.uint8)
        pts = np.asarray(poly, dtype=np.int32).reshape(-1, 1, 2)
        cv2.fillPoly(single, [pts], 1)
        sb = single.astype(bool)
        dr = float(dark[sb].mean()) if sb.any() else 0.0
        target = (~dark & sb) if dr > 0.5 else (dark & sb)
        out[target] = 255
    if dilate > 0:
        k = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (dilate, dilate))
        out = cv2.dilate(out, k, iterations=1)
    out[~poly_mask_bool] = 0
    return out


# ---------------- SAM 후처리 ----------------


def upsample_mask_smooth(seg_small, target_size):
    seg_uint = (seg_small.astype(np.uint8)) * 255
    seg_full = cv2.resize(seg_uint, target_size, interpolation=cv2.INTER_LINEAR)
    return seg_full > 127


def mask_iou(a, b):
    if a.shape != b.shape:
        return 0.0
    inter = np.logical_and(a, b).sum()
    union = np.logical_or(a, b).sum()
    return float(inter) / float(union) if union else 0.0


def filter_by_area(masks, image_size, *, min_ratio=0.0003, max_ratio=0.55):
    w, h = image_size
    total = w * h
    return [m for m in masks if min_ratio <= m.area / total <= max_ratio]


def dedupe_near_identical(masks, *, iou_threshold=0.85):
    sorted_masks = sorted(masks, key=lambda m: -m.area)
    kept = []
    for m in sorted_masks:
        if not any(mask_iou(m.segmentation, k.segmentation) > iou_threshold for k in kept):
            kept.append(m)
    return kept


# ---------------- alpha crop (v6: 자식 안 빼고 텍스트 stroke 만 제거) ----------------


def export_alpha_crop_text_only(
    image: Image.Image,
    mask: np.ndarray,
    text_stroke_mask: np.ndarray,
    *,
    out_dir: Path,
    name: str,
) -> tuple[tuple[int, int, int, int], Path] | None:
    """SAM 마스크 영역 alpha crop. 텍스트 stroke 픽셀만 alpha 에서 제거.

    자식 마스크/박스 내부는 그대로 보존 → 시각 충실도 유지.
    """
    clean = mask & ~text_stroke_mask  # 텍스트 stroke 만 alpha 에서 빼기
    ys, xs = np.where(clean)
    if len(ys) == 0:
        return None
    x0, x1 = int(xs.min()), int(xs.max()) + 1
    y0, y1 = int(ys.min()), int(ys.max()) + 1
    arr = np.asarray(image.convert("RGB"))
    crop_rgb = arr[y0:y1, x0:x1]
    crop_mask = clean[y0:y1, x0:x1]
    rgba = np.zeros((y1 - y0, x1 - x0, 4), dtype=np.uint8)
    rgba[:, :, :3] = crop_rgb
    rgba[:, :, 3] = (crop_mask.astype(np.uint8)) * 255
    out_dir.mkdir(parents=True, exist_ok=True)
    p = out_dir / name
    Image.fromarray(rgba, mode="RGBA").save(p)
    return (x0, y0, x1 - x0, y1 - y0), p


def build_layers_visual_fidelity(
    image: Image.Image,
    masks: list[Sam2Mask],
    text_polygons: list[list[tuple[float, float]]],
    *,
    out_dir: Path,
) -> list[ObjectLayer]:
    """v6: 텍스트 stroke 만 빼고 alpha crop. 자식 영역 보존."""
    text_stroke = build_dark_pixel_text_mask(image, text_polygons, dilate=2).astype(bool)
    sorted_masks = sorted(masks, key=lambda m: -m.area)
    layers: list[ObjectLayer] = []
    obj_dir = out_dir / "objects"
    h, w = np.asarray(image).shape[:2]

    for i, m in enumerate(sorted_masks):
        out = export_alpha_crop_text_only(
            image, m.segmentation, text_stroke, out_dir=obj_dir, name=f"layer_{i:03d}.png",
        )
        if out is None:
            continue
        bbox, path = out
        layers.append(ObjectLayer(
            bbox=bbox,
            z=i,  # 작은 z = 큰 마스크 = 뒤에 깔림
            area_ratio=float(m.area) / (h * w),
            alpha_path=path,
        ))
    return layers


# ---------------- 시각화 ----------------


def render_mask_overlay(image, masks, out_path):
    rng = np.random.default_rng(42)
    arr = np.asarray(image.convert("RGB")).copy().astype(np.float32)
    overlay = np.zeros_like(arr)
    for m in masks:
        color = rng.integers(50, 255, size=3)
        overlay[m.segmentation] = color
    blended = arr * 0.55 + overlay * 0.45
    Image.fromarray(np.clip(blended, 0, 255).astype(np.uint8)).save(out_path)


def render_layer_overlay(image, layers, out_path):
    arr = np.asarray(image.convert("RGB")).copy()
    for ly in layers:
        x, y, w, h = ly.bbox
        cv2.rectangle(arr, (x, y), (x + w, y + h), (0, 200, 0), 2)
    Image.fromarray(arr).save(out_path)


def render_text_overlay(image, lines, out_path):
    arr = np.asarray(image.convert("RGB")).copy()
    for ln in lines:
        x, y, w, h = ln.bbox
        cv2.rectangle(arr, (int(x), int(y)), (int(x + w), int(y + h)), (255, 0, 0), 2)
    Image.fromarray(arr).save(out_path)


# ---------------- PPTX (East Asian font 명시) ----------------


def _set_east_asian_font(rPr, font_name: str) -> None:
    """run.font.name 만 쓰면 latin font 만 설정됨. eastAsia 도 같이 설정."""
    # 기존 ea 제거
    for ea in rPr.findall(qn("a:ea")):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)
    # cs (Complex Script) 도 같이
    for cs in rPr.findall(qn("a:cs")):
        rPr.remove(cs)
    cs = etree.SubElement(rPr, qn("a:cs"))
    cs.set("typeface", font_name)


def add_textbox(slide, line, *, scale, font_name):
    x, y, w, h = line.bbox
    font_emu = int(line.font_pt / PT_PER_INCH * EMU_PER_INCH)
    left = int(x * scale) - font_emu // 2
    width = int(w * scale) + font_emu
    height = max(int(h * scale), int(font_emu * 1.6))
    y_center_emu = int((y + h / 2) * scale)
    top = y_center_emu - height // 2
    if width <= 0 or height <= 0:
        return
    tb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tb.text_frame
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, attr, Emu(0))
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = line.text
    run.font.name = font_name
    run.font.size = Pt(round(line.font_pt, 1))
    try:
        r = int(line.color_hex[1:3], 16)
        g = int(line.color_hex[3:5], 16)
        b = int(line.color_hex[5:7], 16)
        run.font.color.rgb = RGBColor(r, g, b)
    except Exception:
        pass
    # East Asian 폰트 명시 → LibreOffice/PowerPoint 한글 정확
    _set_east_asian_font(run._r.get_or_add_rPr(), font_name)


def build_pptx(
    *, image_size, background_path, layers, lines, out_path,
    slide_w_emu, slide_h_emu, font_name,
):
    px_w, _ = image_size
    scale = slide_w_emu / px_w
    prs = Presentation()
    prs.slide_width = Emu(slide_w_emu)
    prs.slide_height = Emu(slide_h_emu)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    if background_path is not None:
        slide.shapes.add_picture(
            str(background_path), 0, 0,
            width=Emu(slide_w_emu), height=Emu(slide_h_emu),
        )
    for ly in sorted(layers, key=lambda l: l.z):
        x, y, w, h = ly.bbox
        slide.shapes.add_picture(
            str(ly.alpha_path),
            Emu(int(x * scale)), Emu(int(y * scale)),
            width=Emu(int(w * scale)), height=Emu(int(h * scale)),
        )
    for ln in lines:
        add_textbox(slide, ln, scale=scale, font_name=font_name)
    prs.save(str(out_path))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", type=Path, required=True)
    ap.add_argument("--out-dir", type=Path, required=True)
    ap.add_argument("--font", default=DEFAULT_FONT)
    ap.add_argument("--slide-w-emu", type=int, default=DEFAULT_SLIDE_W_EMU)
    ap.add_argument("--slide-h-emu", type=int, default=DEFAULT_SLIDE_H_EMU)
    ap.add_argument("--sam-resize", type=int, default=1920)
    ap.add_argument("--sam-points", type=int, default=48)
    args = ap.parse_args()

    args.out_dir.mkdir(parents=True, exist_ok=True)
    started = time.perf_counter()

    print(f"[1/8] 입력: {args.input}")
    image = Image.open(args.input).convert("RGB")
    px_w, px_h = image.size
    slide_w_inch = args.slide_w_emu / EMU_PER_INCH
    effective_dpi = px_w / slide_w_inch
    print(f"     {px_w}x{px_h}, eff_dpi={effective_dpi:.1f}")

    print("[2/8] PaddleOCR")
    t0 = time.perf_counter()
    raw_lines = run_ocr(image, lang="korean")
    items = lines_to_text_items(raw_lines, image, effective_dpi=effective_dpi)
    grouped = group_items_by_line(items)
    print(f"     {len(items)} → {len(grouped)} lines, {time.perf_counter() - t0:.1f}s")
    render_text_overlay(image, grouped, args.out_dir / "01_text_lines.png")

    print(f"[3/8] SAM2 dense (resize={args.sam_resize}, pps={args.sam_points})")
    sam_input = image
    if args.sam_resize and px_w > args.sam_resize:
        scale_sam = args.sam_resize / px_w
        sam_input = image.resize((args.sam_resize, int(px_h * scale_sam)))
    sam_result = run_sam2_auto(
        sam_input, points_per_side=args.sam_points, max_masks=400,
        pred_iou_thresh=0.80, stability_score_thresh=0.88,
    )
    full_masks: list[Sam2Mask] = []
    for m in sam_result.masks:
        seg = upsample_mask_smooth(m.segmentation, (px_w, px_h))
        ys, xs = np.where(seg)
        if len(ys) == 0:
            continue
        bbox = (int(xs.min()), int(ys.min()),
                int(xs.max() - xs.min() + 1), int(ys.max() - ys.min() + 1))
        full_masks.append(Sam2Mask(bbox=bbox, segmentation=seg, score=m.score, area=int(seg.sum())))
    print(f"     {len(full_masks)} masks, {sam_result.elapsed_sec:.1f}s")
    render_mask_overlay(image, full_masks, args.out_dir / "02_sam2_masks.png")

    print("[4/8] classify + 면적 + dedupe")
    classification = classify_masks(
        full_masks, raw_lines, image_size=(px_w, px_h),
        background_area_ratio=0.6, noise_area_ratio=0.0001,
    )
    shapes = filter_by_area(classification.shape_masks, (px_w, px_h),
                            min_ratio=0.0003, max_ratio=0.55)
    shapes = dedupe_near_identical(shapes, iou_threshold=0.85)
    print(f"     SHAPE={len(shapes)}")

    print("[5/8] alpha crop (텍스트 stroke 만 제거, 자식 보존)")
    layers = build_layers_visual_fidelity(
        image, shapes, [it.polygon for it in items], out_dir=args.out_dir,
    )
    print(f"     {len(layers)} layers")
    render_layer_overlay(image, layers, args.out_dir / "03_layer_overlay.png")

    print("[6/8] 인페인팅 배경")
    inpaint_mask = build_dark_pixel_text_mask(image, [it.polygon for it in items], dilate=3)
    Image.fromarray(inpaint_mask, mode="L").save(args.out_dir / "04_inpaint_mask.png")
    t0 = time.perf_counter()
    inpainted = inpaint_with_mask(image, Image.fromarray(inpaint_mask, mode="L"))
    inpainted_path = args.out_dir / "05_inpainted_bg.png"
    inpainted.save(inpainted_path)
    print(f"     {time.perf_counter() - t0:.1f}s")

    print("[7/8] PPTX (East Asian font 명시)")
    out_inpainted = args.out_dir / "result_inpainted.pptx"
    build_pptx(
        image_size=(px_w, px_h), background_path=inpainted_path,
        layers=layers, lines=grouped, out_path=out_inpainted,
        slide_w_emu=args.slide_w_emu, slide_h_emu=args.slide_h_emu,
        font_name=args.font,
    )
    out_white = args.out_dir / "result_white_bg.pptx"
    build_pptx(
        image_size=(px_w, px_h), background_path=None,
        layers=layers, lines=grouped, out_path=out_white,
        slide_w_emu=args.slide_w_emu, slide_h_emu=args.slide_h_emu,
        font_name=args.font,
    )

    print("[8/8] 보고")
    report = {
        "input": str(args.input),
        "image_size": [px_w, px_h],
        "ocr_items": len(items),
        "grouped_lines": len(grouped),
        "sam_masks_total": len(full_masks),
        "shape_after_filter": len(shapes),
        "object_layers": len(layers),
        "elapsed_sec": round(time.perf_counter() - started, 1),
    }
    (args.out_dir / "report.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(f"\nresult_inpainted.pptx  = 인페인팅 BG + alpha 객체 + textbox")
    print(f"result_white_bg.pptx   = 흰 BG + alpha 객체 + textbox")
    print(f"총 {report['elapsed_sec']}s")


if __name__ == "__main__":
    main()
